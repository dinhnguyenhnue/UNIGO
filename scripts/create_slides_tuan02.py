# -*- coding: utf-8 -*-
"""
Script tạo Slide bài giảng Student-Facing cho Tuần 02 — v2
FIX: Không che logo, không che chân trang, tương phản tốt, z-order đúng
"""
import sys, io, os, copy, glob
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import traceback

TEMPLATE = r'D:\UNIGO\Hệ thống mẫu văn bản\Mẫu slide UNIGO.pptx'
SGK_BASE = r'D:\UNIGO\SGK'
KHBD_BASE = r'D:\UNIGO\KHBD_Tin_học'

# ─── Template master layout ───
# Picture 7 (Logo):  L=0.17 T=0.15 W=0.95 H=0.94  → ends at Y=1.09in
# Picture 9 (Footer): L=0.00 T=6.43 W=13.40 H=1.23 → starts at Y=6.43in
# SAFE ZONE: Y = 1.15in → 6.35in  (height = 5.20in)
# X safe: 0.3in → 13.0in (width = 12.7in)
SAFE_TOP = 1.15
SAFE_BOTTOM = 6.35
SAFE_LEFT = 0.3
SAFE_RIGHT = 13.0
SLIDE_W = 13.33
SLIDE_H = 7.50

# ─── 8 bộ màu (high contrast) ───
COLOR_PALETTES = [
    # Blue (Tiền TH)
    {"primary": "1B4F9B", "accent": "2D7DD2", "bg": "EBF3FE", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "1A2744", "text_on_card": "1A2744"},
    # Purple (Lớp 1)
    {"primary": "5B21B6", "accent": "7C3AED", "bg": "F3EEFF", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "2E1065", "text_on_card": "2E1065"},
    # Teal (Lớp 2)
    {"primary": "0F766E", "accent": "14B8A6", "bg": "ECFDF5", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "134E4A", "text_on_card": "134E4A"},
    # Orange (Lớp 3)
    {"primary": "C2410C", "accent": "EA580C", "bg": "FFF7ED", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "7C2D12", "text_on_card": "431407"},
    # Indigo (Lớp 4)
    {"primary": "3730A3", "accent": "4F46E5", "bg": "EEF2FF", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "1E1B4B", "text_on_card": "1E1B4B"},
    # Rose (Lớp 5)
    {"primary": "BE185D", "accent": "EC4899", "bg": "FDF2F8", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "831843", "text_on_card": "831843"},
    # Emerald (Lớp 6)
    {"primary": "047857", "accent": "10B981", "bg": "ECFDF5", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "064E3B", "text_on_card": "064E3B"},
    # Sky (Lớp 7)
    {"primary": "0369A1", "accent": "0EA5E9", "bg": "F0F9FF", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "0C4A6E", "text_on_card": "0C4A6E"},
    # Amber (Lớp 8)  
    {"primary": "B45309", "accent": "F59E0B", "bg": "FFFBEB", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "78350F", "text_on_card": "451A03"},
]

# ─── Thông tin bài học 9 lớp (same as v1) ───
LESSONS = [
    {
        "id": "Tien_TH", "folder": "Tiền_tiểu_học", "lop_label": "Tiền Tiểu học", "lop_num": 0,
        "bai": "Bài 1. Máy tính xung quanh em",
        "file_name": "Slide_Tin_hoc_Tien_TH_Bai01_May_tinh_xung_quanh_em.pptx",
        "slides_content": [
            {"type": "cover", "title": "Máy tính xung quanh em", "subtitle": "Tin học • Tiền Tiểu học • Bài 1"},
            {"type": "warmup", "title": "Các con ơi, nhìn xung quanh nào!",
             "content": "Con có thấy chiếc máy tính nào không?\nChỉ cho cô xem nhé!"},
            {"type": "learn", "title": "Đây là máy tính!",
             "content": "Máy tính ở khắp nơi:\n● Trong phòng học của chúng mình\n● Ở nhà (laptop của bố mẹ)\n● Trên tay (điện thoại cũng là máy tính nhỏ!)"},
            {"type": "learn", "title": "Máy tính có những phần nào?",
             "content": "● Màn hình — để con nhìn\n● Bàn phím — để con gõ chữ\n● Chuột — để con chỉ và nhấp\n● Thân máy — bộ não của máy tính"},
            {"type": "learn", "title": "Máy tính giúp chúng ta làm gì?",
             "content": "● Học bài vui vẻ\n● Vẽ tranh đẹp\n● Nghe nhạc hay\n● Xem hoạt hình"},
            {"type": "practice", "title": "Cùng chơi nào!",
             "content": "Con hãy chỉ vào từng bộ phận\nvà nói tên cho cô nghe:",
             "items": ["Đâu là màn hình?", "Đâu là bàn phím?", "Đâu là chuột?"]},
            {"type": "activity", "title": "Thử thách nhỏ!",
             "content": "Con hãy vẽ một chiếc máy tính\nvào giấy của mình.\nNhớ vẽ đủ các bộ phận nhé!"},
            {"type": "summary", "title": "Hôm nay con đã biết!",
             "content": "Máy tính ở xung quanh chúng ta\nMáy tính có: màn hình, bàn phím, chuột, thân máy\nMáy tính giúp con học và chơi"},
            {"type": "thanks", "title": "Các con giỏi lắm!", "content": "Hẹn gặp lại tiết sau nhé!"},
        ]
    },
    {
        "id": "Lop_1", "folder": "Lớp_1", "lop_label": "Lớp 1", "lop_num": 1,
        "bai": "Bài 1. Chiếc máy tính của em",
        "file_name": "Slide_Tin_hoc_Lop_1_Bai01_Chiec_may_tinh_cua_em.pptx",
        "slides_content": [
            {"type": "cover", "title": "Chiếc máy tính của em", "subtitle": "Tin học • Lớp 1 • Bài 1"},
            {"type": "warmup", "title": "Em đã gặp máy tính ở đâu?",
             "content": "Hãy kể cho cô nghe:\n● Em thấy máy tính ở đâu?\n● Bố mẹ dùng máy tính để làm gì?"},
            {"type": "learn", "title": "Cùng khám phá máy tính nào!",
             "content": "Máy tính có nhiều loại:\n● Máy tính để bàn (Desktop)\n● Máy tính xách tay (Laptop)\n● Máy tính bảng (Tablet)"},
            {"type": "learn", "title": "Các bộ phận của máy tính",
             "content": "● Màn hình — hiển thị hình ảnh, chữ\n● Bàn phím — em gõ chữ, số\n● Chuột — em nhấp để chọn\n● Thân máy — nơi xử lý thông tin"},
            {"type": "learn", "title": "Bật và tắt máy tính đúng cách",
             "content": "Bước 1: Nhấn nút nguồn để bật\nBước 2: Chờ máy khởi động\nBước 3: Khi muốn tắt → chọn Shut Down\n\nKhông rút dây điện đột ngột!"},
            {"type": "practice", "title": "Em hãy thực hành!",
             "content": "● Chỉ vào từng bộ phận và gọi tên\n● Thử nhấp chuột trái\n● Thử gõ tên mình bằng bàn phím"},
            {"type": "activity", "title": "Trò chơi: Nối đúng!",
             "content": "Nối tên với hình ảnh đúng:",
             "items": ["Màn hình ↔ ?", "Bàn phím ↔ ?", "Chuột ↔ ?", "Thân máy ↔ ?"]},
            {"type": "summary", "title": "Em nhớ được gì nào?",
             "content": "Máy tính có nhiều loại\n4 bộ phận chính: Màn hình, Bàn phím, Chuột, Thân máy\nBật/tắt máy đúng cách"},
            {"type": "thanks", "title": "Em giỏi lắm!", "content": "Về nhà: Hãy kể cho bố mẹ nghe\nvề chiếc máy tính ở lớp nhé!"},
        ]
    },
    {
        "id": "Lop_2", "folder": "Lớp_2", "lop_label": "Lớp 2", "lop_num": 2,
        "bai": "Bài 1. Máy tính là người bạn của em",
        "file_name": "Slide_Tin_hoc_Lop_2_Bai01_May_tinh_la_nguoi_ban_cua_em.pptx",
        "slides_content": [
            {"type": "cover", "title": "Máy tính là người bạn của em", "subtitle": "Tin học • Lớp 2 • Bài 1"},
            {"type": "warmup", "title": "Máy tính giúp em những gì?",
             "content": "Hãy kể cho cô:\n● Em dùng máy tính để làm gì?\n● Máy tính giúp bố mẹ em làm việc gì?"},
            {"type": "learn", "title": "Máy tính — Người bạn thông minh!",
             "content": "Máy tính giúp em:\n● Học bài — tra cứu kiến thức\n● Vẽ tranh — sáng tạo nghệ thuật\n● Viết văn — soạn thảo bài viết"},
            {"type": "learn", "title": "Máy tính còn giúp gì nữa?",
             "content": "● Gọi video — nói chuyện với người thân\n● Nghe nhạc, xem phim\n● Chơi game học tập\n● Tìm kiếm thông tin trên Internet"},
            {"type": "learn", "title": "Sử dụng máy tính an toàn",
             "content": "Em nhớ nhé:\n● Ngồi đúng tư thế, mắt cách màn hình 40-50cm\n● Nghỉ mắt mỗi 20-30 phút\n● Không dùng quá 1 giờ/ngày\n● Luôn hỏi bố mẹ trước khi lên mạng"},
            {"type": "practice", "title": "Em hãy thực hành!",
             "content": "● Ngồi đúng tư thế trước máy tính\n● Mở một chương trình vẽ\n● Vẽ một bức tranh đơn giản"},
            {"type": "activity", "title": "Đúng hay Sai?",
             "content": "Em hãy trả lời:",
             "items": ["Dùng máy tính cả ngày là tốt?", "Máy tính giúp em học bài?", "Em nên ngồi sát màn hình?", "Phải hỏi bố mẹ trước khi lên mạng?"]},
            {"type": "summary", "title": "Em nhớ được gì nào?",
             "content": "Máy tính là người bạn giúp em học và chơi\nCó rất nhiều việc máy tính giúp được\nDùng máy tính phải an toàn và đúng cách"},
            {"type": "thanks", "title": "Em giỏi lắm!", "content": "BTVN: Hãy kể 3 việc\nmáy tính giúp em mỗi ngày!"},
        ]
    },
    {
        "id": "Lop_3", "folder": "Lớp_3", "lop_label": "Lớp 3", "lop_num": 3,
        "bai": "Bài 1. Thông tin và quyết định",
        "file_name": "Slide_Tin_hoc_Lop_3_Bai01_Thong_tin_va_quyet_dinh.pptx",
        "slides_content": [
            {"type": "cover", "title": "Thông tin và Quyết định", "subtitle": "Tin học • Lớp 3 • Bài 1"},
            {"type": "warmup", "title": "Em hãy đoán xem!",
             "content": "Sáng nay trước khi đi học, em quyết định mặc áo gì?\nEm dựa vào đâu để quyết định?\n\nGợi ý: Xem thời tiết? Hỏi mẹ? Nhìn ra cửa sổ?"},
            {"type": "learn", "title": "Thông tin là gì?",
             "content": "Thông tin là những gì em biết được:\n● Nhìn thấy → hình ảnh, chữ viết\n● Nghe thấy → âm thanh, lời nói\n● Cảm nhận → nóng, lạnh, mềm, cứng"},
            {"type": "learn", "title": "Thông tin giúp em quyết định!",
             "content": "Ví dụ thực tế:\n● Trời mưa → em mang áo mưa\n● Đèn đỏ → em dừng lại\n● Chuông reo → em vào lớp\n\nThông tin → Suy nghĩ → Quyết định"},
            {"type": "learn", "title": "Các dạng thông tin",
             "content": "● Chữ viết (văn bản) — sách, báo, tin nhắn\n● Hình ảnh — tranh, ảnh chụp\n● Âm thanh — nhạc, tiếng nói\n● Video — phim, clip"},
            {"type": "practice", "title": "Em hãy thử nào!",
             "content": "Tình huống: Em muốn mua quà sinh nhật cho bạn.\nEm cần những thông tin gì?\n\n● Bạn thích gì?\n● Em có bao nhiêu tiền?\n● Cửa hàng ở đâu?"},
            {"type": "activity", "title": "Trò chơi: Thông tin hay không?",
             "content": "Những cái nào là thông tin?",
             "items": ["Bảng thời khóa biểu", "Tiếng chuông trường", "Một viên đá", "Tin nhắn của mẹ"]},
            {"type": "summary", "title": "Hôm nay em đã học!",
             "content": "Thông tin là những gì em biết được qua giác quan\nThông tin có nhiều dạng: chữ, hình, âm thanh, video\nThông tin giúp em đưa ra quyết định đúng"},
            {"type": "thanks", "title": "Em làm tốt lắm!", "content": "BTVN: Kể 3 thông tin em nhận được\nhôm nay và quyết định em đã đưa ra!"},
        ]
    },
    {
        "id": "Lop_4", "folder": "Lớp_4", "lop_label": "Lớp 4", "lop_num": 4,
        "bai": "Bài 1. Phần cứng và phần mềm máy tính",
        "file_name": "Slide_Tin_hoc_Lop_4_Bai01_Phan_cung_va_phan_mem.pptx",
        "slides_content": [
            {"type": "cover", "title": "Phần cứng và Phần mềm", "subtitle": "Tin học • Lớp 4 • Bài 1"},
            {"type": "warmup", "title": "Thử đoán xem!",
             "content": "Em có thể sờ được bàn phím không? → CÓ\nEm có thể sờ được trò chơi trong máy không? → KHÔNG\n\nVì sao nhỉ?"},
            {"type": "learn", "title": "Phần cứng (Hardware) là gì?",
             "content": "Phần cứng = Những thứ em NHÌN THẤY và SỜ ĐƯỢC:\n● Màn hình, bàn phím, chuột\n● Loa, tai nghe, webcam\n● Ổ cứng, RAM, CPU (bên trong thân máy)"},
            {"type": "learn", "title": "Phần mềm (Software) là gì?",
             "content": "Phần mềm = Chương trình CHẠY BÊN TRONG máy:\n● Hệ điều hành (Windows, macOS)\n● Trình duyệt web (Chrome, Firefox)\n● Trò chơi, phần mềm vẽ, phần mềm học"},
            {"type": "learn", "title": "So sánh: Phần cứng vs Phần mềm",
             "content": "Phần cứng → Như CƠ THỂ người\nPhần mềm → Như TRÍ ÓC người\n\nMáy tính cần CẢ HAI mới hoạt động được!\nKhông có phần cứng → không có máy\nKhông có phần mềm → máy không biết làm gì"},
            {"type": "practice", "title": "Em hãy phân loại!",
             "content": "Đâu là phần cứng? Đâu là phần mềm?",
             "items": ["Bàn phím → ?", "Microsoft Word → ?", "Chuột → ?", "Minecraft → ?", "Loa → ?", "Chrome → ?"]},
            {"type": "activity", "title": "Thử thách nhanh!",
             "content": "Em hãy nhìn quanh phòng máy:\n● Liệt kê 3 phần cứng em thấy\n● Kể 3 phần mềm em đã từng dùng\n\nViết vào vở trong 3 phút!"},
            {"type": "summary", "title": "Hôm nay em đã hiểu!",
             "content": "Phần cứng = sờ được, nhìn thấy (Hardware)\nPhần mềm = chương trình chạy bên trong (Software)\nMáy tính cần CẢ phần cứng VÀ phần mềm"},
            {"type": "thanks", "title": "Em giỏi lắm!", "content": "BTVN: Về nhà, liệt kê\n3 phần cứng + 3 phần mềm\ncủa máy tính ở nhà em!"},
        ]
    },
    {
        "id": "Lop_5", "folder": "Lớp_5", "lop_label": "Lớp 5", "lop_num": 5,
        "bai": "Bài 1. Em có thể làm gì với máy tính",
        "file_name": "Slide_Tin_hoc_Lop_5_Bai01_Em_co_the_lam_gi_voi_may_tinh.pptx",
        "slides_content": [
            {"type": "cover", "title": "Em có thể làm gì với máy tính?", "subtitle": "Tin học • Lớp 5 • Bài 1"},
            {"type": "warmup", "title": "Kể cho cô nghe nào!",
             "content": "Tuần trước em đã dùng máy tính để làm gì?\nKể ít nhất 2 việc nhé!\n\nGợi ý: Học bài? Xem video? Chơi game? Tìm kiếm?"},
            {"type": "learn", "title": "Máy tính — Trợ thủ đa năng!",
             "content": "Với máy tính, em có thể:\n● Soạn thảo văn bản (Word)\n● Tính toán bằng bảng tính (Excel)\n● Vẽ tranh và thiết kế (Paint, Canva)\n● Thuyết trình (PowerPoint)"},
            {"type": "learn", "title": "Kết nối và chia sẻ!",
             "content": "Nhờ Internet, em có thể:\n● Tìm kiếm thông tin (Google)\n● Gửi email cho thầy cô\n● Học trực tuyến (Zoom, Google Meet)\n● Chia sẻ bài với bạn bè"},
            {"type": "learn", "title": "Máy tính trong đời sống",
             "content": "Máy tính giúp mọi người:\n● Bác sĩ — chẩn đoán bệnh\n● Kỹ sư — thiết kế nhà\n● Giáo viên — soạn giáo án\n● Phi công — điều khiển máy bay"},
            {"type": "practice", "title": "Em hãy thực hành!",
             "content": "Bước 1: Mở phần mềm soạn thảo (Word)\nBước 2: Gõ họ tên và lớp của em\nBước 3: Viết \"Em thích dùng máy tính để...\"\nBước 4: Lưu file với tên của em"},
            {"type": "activity", "title": "Ghép đôi: Công việc ↔ Phần mềm",
             "content": "Nối công việc với phần mềm phù hợp:",
             "items": ["Viết bài văn ↔ Word", "Tính điểm TB ↔ Excel", "Làm bài trình bày ↔ PowerPoint", "Vẽ tranh ↔ Paint"]},
            {"type": "summary", "title": "Hôm nay em đã biết!",
             "content": "Máy tính làm được RẤT NHIỀU việc\nMỗi công việc có phần mềm phù hợp\nMáy tính giúp ích trong mọi nghề nghiệp"},
            {"type": "thanks", "title": "Em giỏi lắm!", "content": "BTVN: Hỏi 1 người lớn\nhọ dùng máy tính để làm gì?\nGhi lại vào vở nhé!"},
        ]
    },
    {
        "id": "Lop_6", "folder": "Lớp_6", "lop_label": "Lớp 6", "lop_num": 6,
        "bai": "Bài 1. Thông tin và dữ liệu",
        "file_name": "Slide_Tin_hoc_Lop_6_Bai01_Thong_tin_va_du_lieu.pptx",
        "slides_content": [
            {"type": "cover", "title": "Thông tin và Dữ liệu", "subtitle": "Tin học • Lớp 6 • Bài 1"},
            {"type": "warmup", "title": "Hãy thử suy nghĩ!",
             "content": "\"Ngày mai trời sẽ mưa\" — Đây là thông tin hay dữ liệu?\n\"23°C\" — Còn đây thì sao?\n\nHai khái niệm này khác nhau như thế nào?"},
            {"type": "learn", "title": "Thông tin (Information) là gì?",
             "content": "Thông tin = hiểu biết về sự vật, hiện tượng\n\nVí dụ:\n● \"Hà Nội là thủ đô Việt Nam\"\n● \"Hôm nay trời nắng\"\n● \"Lớp 6A có 35 học sinh\""},
            {"type": "learn", "title": "Dữ liệu (Data) là gì?",
             "content": "Dữ liệu = thông tin được biểu diễn\ndưới dạng máy tính xử lý được\n\nCác dạng dữ liệu:\n● Văn bản (text)\n● Hình ảnh (image)\n● Âm thanh (audio)\n● Video"},
            {"type": "learn", "title": "Thông tin → Dữ liệu: Mối quan hệ",
             "content": "Thông tin: \"Nhiệt độ hôm nay là 30 độ C\"\n         ↓ ghi lại, mã hóa\nDữ liệu: 30°C (số liệu trong máy tính)\n\nDữ liệu là CÁCH THỂ HIỆN thông tin trên máy tính"},
            {"type": "practice", "title": "Hãy phân loại!",
             "content": "Xác định dạng dữ liệu:",
             "items": ["Bài hát \"Quốc ca\" → Âm thanh", "Ảnh chụp lớp → Hình ảnh", "Bảng điểm học kỳ → Văn bản + Số", "Clip nấu ăn → Video"]},
            {"type": "activity", "title": "Thảo luận nhóm (5 phút)",
             "content": "Nhóm 4 em, thảo luận:\n\n\"Khi em đăng 1 bức ảnh lên mạng xã hội,\nem đã tạo ra dữ liệu dạng gì?\nAi có thể sử dụng dữ liệu đó?\""},
            {"type": "summary", "title": "Tóm tắt bài học",
             "content": "Thông tin = hiểu biết về thế giới xung quanh\nDữ liệu = thông tin được mã hóa cho máy tính\n4 dạng dữ liệu: văn bản, hình ảnh, âm thanh, video\nDữ liệu là cách thể hiện thông tin"},
            {"type": "thanks", "title": "Bài học kết thúc!", "content": "BTVN: Liệt kê 5 loại dữ liệu\nem tạo ra trong 1 ngày\nvà cho biết mỗi loại thuộc dạng nào."},
        ]
    },
    {
        "id": "Lop_7", "folder": "Lớp_7", "lop_label": "Lớp 7", "lop_num": 7,
        "bai": "Bài 1. Thiết bị vào - ra",
        "file_name": "Slide_Tin_hoc_Lop_7_Bai01_Thiet_bi_vao_ra.pptx",
        "slides_content": [
            {"type": "cover", "title": "Thiết bị vào - ra (I/O Devices)", "subtitle": "Tin học • Lớp 7 • Bài 1"},
            {"type": "warmup", "title": "Quan sát và trả lời!",
             "content": "Nhìn quanh phòng máy:\n\"Kể tên TẤT CẢ thiết bị em thấy\nxung quanh chiếc máy tính!\"\n\n2 phút — càng nhiều càng tốt!"},
            {"type": "learn", "title": "Thiết bị vào (Input) là gì?",
             "content": "Thiết bị vào = đưa thông tin VÀO máy tính\n\n● Bàn phím (Keyboard) — gõ chữ, số\n● Chuột (Mouse) — chỉ, nhấp, kéo thả\n● Microphone — thu âm giọng nói\n● Webcam — thu hình ảnh\n● Máy quét (Scanner) — quét tài liệu"},
            {"type": "learn", "title": "Thiết bị ra (Output) là gì?",
             "content": "Thiết bị ra = đưa thông tin RA cho người dùng\n\n● Màn hình (Monitor) — hiển thị hình ảnh\n● Loa (Speaker) — phát âm thanh\n● Máy in (Printer) — in tài liệu\n● Tai nghe (Headphone) — nghe âm thanh\n● Máy chiếu (Projector) — chiếu lên màn"},
            {"type": "learn", "title": "Input → Xử lý → Output",
             "content": "Quy trình hoạt động:\n\nBàn phím → CPU xử lý → Màn hình\n(Input)      (Processing)    (Output)\n\nMột số thiết bị vừa vào vừa ra:\nMàn hình cảm ứng, USB, Ổ cứng ngoài"},
            {"type": "practice", "title": "Hãy phân loại!",
             "content": "Sắp xếp vào đúng nhóm:",
             "items": ["Bàn phím → Input", "Màn hình → Output", "Webcam → Input", "Loa → Output", "Màn hình cảm ứng → Cả hai!"]},
            {"type": "activity", "title": "Thảo luận nhóm",
             "content": "Nhóm 4 em, trả lời:\n\n1. Nếu máy tính không có thiết bị vào,\n   em có thể làm gì với nó?\n2. Kể 3 thiết bị vào-ra mà smartphone có.\n\n5 phút thảo luận + 2 phút trình bày"},
            {"type": "summary", "title": "Tóm tắt bài học",
             "content": "Thiết bị vào (Input): đưa thông tin vào máy\nThiết bị ra (Output): đưa thông tin ra ngoài\nQuy trình: Input → Xử lý → Output\nCó thiết bị vừa vào vừa ra (I/O)"},
            {"type": "thanks", "title": "Bài học kết thúc!", "content": "BTVN: Vẽ sơ đồ Input/Output\ncủa hệ thống máy tính\ntại phòng Tin học."},
        ]
    },
    {
        "id": "Lop_8", "folder": "Lớp_8", "lop_label": "Lớp 8", "lop_num": 8,
        "bai": "Bài 1. Lược sử công cụ tính toán",
        "file_name": "Slide_Tin_hoc_Lop_8_Bai01_Luoc_su_cong_cu_tinh_toan.pptx",
        "slides_content": [
            {"type": "cover", "title": "Lược sử công cụ tính toán", "subtitle": "Tin học • Lớp 8 • Bài 1"},
            {"type": "warmup", "title": "Quan sát và so sánh!",
             "content": "Hãy nhìn 2 hình ảnh:\n● Chiếc bàn tính Abacus cổ đại\n● Siêu máy tính hiện đại\n\nCó điểm gì GIỐNG và KHÁC nhau?\n2 phút thảo luận cặp đôi"},
            {"type": "learn", "title": "Hành trình từ Abacus đến AI",
             "content": "Dòng thời gian lược sử:\n● ~2400 TCN — Bàn tính Abacus (Lưỡng Hà)\n● 1642 — Máy Pascaline (Blaise Pascal)\n● 1834 — Máy phân tích (Charles Babbage)\n● 1946 — ENIAC (máy tính điện tử đầu tiên)\n● 1971 — Microprocessor (Intel 4004)"},
            {"type": "learn", "title": "5 thế hệ máy tính",
             "content": "● Thế hệ 1 (1940-1956): Đèn chân không, rất lớn\n● Thế hệ 2 (1956-1963): Transistor, nhỏ hơn\n● Thế hệ 3 (1964-1971): Mạch tích hợp (IC)\n● Thế hệ 4 (1971-nay): Vi xử lý, PC\n● Thế hệ 5 (tương lai): Trí tuệ nhân tạo (AI)"},
            {"type": "learn", "title": "Xu hướng phát triển hiện nay",
             "content": "● Máy tính ngày càng NHỎ hơn\n● Tốc độ xử lý ngày càng NHANH hơn\n● Trí tuệ nhân tạo (AI) — máy biết \"suy nghĩ\"\n● Điện toán đám mây (Cloud Computing)\n● Internet vạn vật (IoT)"},
            {"type": "practice", "title": "Sắp xếp theo thời gian!",
             "content": "Hãy sắp xếp đúng thứ tự lịch sử:",
             "items": ["ENIAC (1946)", "Abacus (~2400 TCN)", "Microprocessor (1971)", "Pascaline (1642)", "Smartphone (2007)"]},
            {"type": "activity", "title": "Thảo luận: Tương lai?",
             "content": "Nhóm 4 em, thảo luận (5 phút):\n\n\"Máy tính thế hệ 5 (AI) sẽ thay đổi\ncuộc sống như thế nào?\nEm nghĩ 20 năm nữa máy tính sẽ ra sao?\"\n\nTrình bày 2 phút/nhóm"},
            {"type": "summary", "title": "Tóm tắt bài học",
             "content": "Công cụ tính toán phát triển hàng nghìn năm\n5 thế hệ: Đèn CK → Transistor → IC → Vi xử lý → AI\nXu hướng: nhỏ hơn, nhanh hơn, thông minh hơn"},
            {"type": "thanks", "title": "Bài học kết thúc!", "content": "BTVN: Vẽ dòng thời gian\n5 mốc lịch sử quan trọng nhất\ncủa công cụ tính toán."},
        ]
    },
]


# ═══════════════════════════════════════════════════
# HÀM HỖ TRỢ — v2 (fixed z-order, no bg cover)
# ═══════════════════════════════════════════════════

def hex_to_rgb(h):
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def set_font(run, size_pt, bold=False, color_hex="333333", font_name="Arial"):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(color_hex)
    run.font.name = font_name

def add_textbox(slide, left, top, width, height, text, size_pt=18, bold=False,
                color_hex="333333", alignment=PP_ALIGN.LEFT, font_name="Arial"):
    """Add a simple textbox — always on top of other shapes"""
    txbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size_pt, bold, color_hex, font_name)
    return txbox

def add_multiline_textbox(slide, left, top, width, height, lines, size_pt=18,
                          color_hex="333333", bold=False, font_name="Arial", alignment=PP_ALIGN.LEFT):
    """Add textbox with multiple lines as separate paragraphs"""
    txbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alignment
        run = p.add_run()
        run.text = line
        set_font(run, size_pt, bold, color_hex, font_name)
        p.space_after = Pt(6)
    return txbox

def add_shape_in_safe_zone(slide, shape_type, left, top, width, height, fill_hex,
                            border_hex=None, send_to_back=False):
    """Add a shape CLAMPED within the safe zone"""
    # Clamp to safe zone
    actual_top = max(top, SAFE_TOP)
    actual_bottom = min(top + height, SAFE_BOTTOM)
    actual_height = max(actual_bottom - actual_top, 0.1)

    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left), Inches(actual_top), Inches(width), Inches(actual_height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    if border_hex:
        shape.line.color.rgb = hex_to_rgb(border_hex)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()

    if send_to_back:
        sp = shape._element
        spTree = sp.getparent()
        spTree.remove(sp)
        spTree.insert(2, sp)
    return shape

def add_slide_transition(slide, transition_type="fade"):
    """Add XML transition effect"""
    transitions = {
        'fade': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>',
        'push': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:push dir="l"/></p:transition>',
        'wipe': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:wipe/></p:transition>',
        'cover': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:cover/></p:transition>',
    }
    xml_str = transitions.get(transition_type, transitions['fade'])
    transition_elem = etree.fromstring(xml_str)
    slide._element.append(transition_elem)


def find_best_image(lesson, slide_type):
    """Find best image for a slide"""
    lop_folder = lesson["folder"]
    lop_num = lesson["lop_num"]

    # KHBD generated images (Tien TH, Lop 1, 2)
    khbd_img_dir = os.path.join(KHBD_BASE, lop_folder, "Tuần_02", "images")
    type_map = {"cover": "cover.png", "learn": "activity.png", "practice": "practice.png",
                "activity": "activity.png", "summary": "summary.png", "warmup": "cover.png",
                "thanks": "cover.png"}
    if os.path.isdir(khbd_img_dir):
        target = type_map.get(slide_type, "activity.png")
        img_path = os.path.join(khbd_img_dir, target)
        if os.path.isfile(img_path):
            return img_path

    # SGK images (Lop 3-8)
    sgk_img_dir = os.path.join(SGK_BASE, f"Lớp_{lop_num}", "bai1_images")
    if os.path.isdir(sgk_img_dir):
        imgs = sorted([f for f in os.listdir(sgk_img_dir)
                       if f.lower().endswith(('.png','.jpg','.jpeg'))
                       and not os.path.isdir(os.path.join(sgk_img_dir, f))])
        if imgs:
            idx_map = {"cover":0,"warmup":0,"learn":1,"practice":-2,"activity":-1,"summary":0}
            idx = idx_map.get(slide_type, 0) % len(imgs)
            return os.path.join(sgk_img_dir, imgs[idx])

        full_dir = os.path.join(sgk_img_dir, "full_pages")
        if os.path.isdir(full_dir):
            pages = sorted(os.listdir(full_dir))
            if pages:
                idx = {"cover":0,"warmup":0,"learn":1,"practice":3,"activity":4,"summary":0}
                page_idx = idx.get(slide_type, 0) % len(pages)
                return os.path.join(full_dir, pages[page_idx])
    return None


def remove_template_slides(prs):
    """Remove all existing slides from template safely"""
    while len(prs.slides) > 0:
        slide = prs.slides[0]
        rId = prs.slides._sldIdLst[0].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
        )
        # Remove the relationship
        prs.part.drop_rel(rId)
        # Remove from slide ID list
        sldId = prs.slides._sldIdLst[0]
        prs.slides._sldIdLst.remove(sldId)


# ═══════════════════════════════════════════════════
# SLIDE BUILDERS — v2
# Quy tắc: KHÔNG vẽ ngoài SAFE_TOP..SAFE_BOTTOM
#           KHÔNG thêm footer (master đã có)
#           Logo master tự hiện (không che)
# ═══════════════════════════════════════════════════

def build_cover(prs, content, pal, lesson, layout):
    """Slide bìa: Banner màu NẰM TRONG safe zone + title + subtitle"""
    slide = prs.slides.add_slide(layout)

    # Colored banner INSIDE safe zone only (not full slide)
    add_shape_in_safe_zone(slide, MSO_SHAPE.RECTANGLE,
                           0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM - SAFE_TOP,
                           pal["primary"], send_to_back=True)

    # Accent diagonal stripe (decorative, within safe zone)
    add_shape_in_safe_zone(slide, MSO_SHAPE.RECTANGLE,
                           0, SAFE_TOP, SLIDE_W, 1.8,
                           pal["accent"], send_to_back=True)

    # Title (centered, large, white text on colored bg)
    add_textbox(slide, 1.5, 2.2, SLIDE_W - 3, 1.5, content["title"],
                size_pt=36, bold=True, color_hex=pal["text_on_primary"],
                alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, 1.5, 4.0, SLIDE_W - 3, 0.8, content.get("subtitle", ""),
                size_pt=20, bold=False, color_hex=pal["text_on_primary"],
                alignment=PP_ALIGN.CENTER)

    # Try image on right
    img = find_best_image(lesson, "cover")
    if img:
        try:
            # Image in right portion, within safe zone
            pic = slide.shapes.add_picture(
                img, Inches(SLIDE_W - 4.5), Inches(SAFE_TOP + 0.3),
                Inches(3.8), Inches(3.8))
        except:
            pass

    add_slide_transition(slide, "fade")
    return slide


def build_content(prs, content, pal, lesson, layout, slide_type="learn", idx=0):
    """Slide nội dung: bg nhạt + banner badge + card + optional image"""
    slide = prs.slides.add_slide(layout)

    # Light background ONLY in safe zone
    add_shape_in_safe_zone(slide, MSO_SHAPE.RECTANGLE,
                           0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM - SAFE_TOP,
                           pal["bg"], send_to_back=True)

    # Top badge banner (narrow, within safe zone)
    badge_y = SAFE_TOP + 0.05
    add_shape_in_safe_zone(slide, MSO_SHAPE.RECTANGLE,
                           0, badge_y, SLIDE_W, 0.45, pal["primary"])
    badge_text = f"  {lesson['lop_label'].upper()} • {lesson['bai'].upper()}"
    add_textbox(slide, 0.3, badge_y + 0.03, SLIDE_W - 0.6, 0.4, badge_text,
                size_pt=13, bold=True, color_hex=pal["text_on_primary"])

    # Title (below badge)
    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, SLIDE_W - 1, 0.65, content["title"],
                size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    # Content area
    content_y = title_y + 0.75
    content_h = SAFE_BOTTOM - content_y - 0.15
    text_content = content.get("content", "")
    lines = text_content.split("\n")

    img = find_best_image(lesson, slide_type)
    has_image = img is not None

    if has_image:
        # Two-column: card left (60%) + image right (35%)
        card_w = 7.5
        # White card
        add_shape_in_safe_zone(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                               0.4, content_y, card_w, content_h, pal["card"])
        # Accent bar on card left edge
        add_shape_in_safe_zone(slide, MSO_SHAPE.RECTANGLE,
                               0.4, content_y, 0.12, content_h, pal["accent"])
        # Text inside card
        add_multiline_textbox(slide, 0.8, content_y + 0.15, card_w - 0.6, content_h - 0.3,
                              lines, size_pt=18, color_hex=pal["text_on_card"])
        # Image on right
        try:
            img_x = 8.3
            img_w = SAFE_RIGHT - img_x - 0.2
            slide.shapes.add_picture(img, Inches(img_x), Inches(content_y),
                                     Inches(img_w), Inches(content_h))
        except:
            pass
    else:
        # Full-width card
        card_w = SLIDE_W - 0.8
        add_shape_in_safe_zone(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                               0.4, content_y, card_w, content_h, pal["card"])
        add_shape_in_safe_zone(slide, MSO_SHAPE.RECTANGLE,
                               0.4, content_y, 0.12, content_h, pal["accent"])
        add_multiline_textbox(slide, 0.8, content_y + 0.15, card_w - 0.6, content_h - 0.3,
                              lines, size_pt=20, color_hex=pal["text_on_card"])

    trans = ["push", "wipe", "fade", "cover"]
    add_slide_transition(slide, trans[idx % len(trans)])
    return slide


def build_practice(prs, content, pal, lesson, layout):
    """Slide luyện tập / thử thách: items as cards"""
    slide = prs.slides.add_slide(layout)

    # Light bg in safe zone
    add_shape_in_safe_zone(slide, MSO_SHAPE.RECTANGLE,
                           0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM - SAFE_TOP,
                           pal["bg"], send_to_back=True)

    # Badge
    badge_y = SAFE_TOP + 0.05
    add_shape_in_safe_zone(slide, MSO_SHAPE.RECTANGLE,
                           0, badge_y, SLIDE_W, 0.45, pal["accent"])
    add_textbox(slide, 0.3, badge_y + 0.03, SLIDE_W - 0.6, 0.4,
                f"  {lesson['lop_label'].upper()} • LUYỆN TẬP",
                size_pt=13, bold=True, color_hex=pal["text_on_primary"])

    # Title
    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, SLIDE_W - 1, 0.65, content["title"],
                size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    # Main instruction
    main_text = content.get("content", "")
    instr_y = title_y + 0.75
    if main_text:
        lines = main_text.split("\n")
        add_multiline_textbox(slide, 0.5, instr_y, SLIDE_W - 1, 0.9, lines,
                              size_pt=17, color_hex=pal["text_on_bg"])
        items_y = instr_y + 1.0
    else:
        items_y = instr_y

    # Item cards in 2 columns
    items = content.get("items", [])
    if items:
        cols = 2
        gap = 0.2
        card_w = (SLIDE_W - 1.0 - gap) / cols
        for i, item in enumerate(items[:6]):
            col = i % cols
            row = i // cols
            x = 0.5 + col * (card_w + gap)
            y = items_y + row * 0.75
            if y + 0.6 > SAFE_BOTTOM:
                break
            # Item card (white with left accent)
            add_shape_in_safe_zone(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, y, card_w, 0.6, pal["card"], pal["accent"])
            add_textbox(slide, x + 0.25, y + 0.1, card_w - 0.4, 0.4, item,
                        size_pt=15, color_hex=pal["text_on_card"])

    add_slide_transition(slide, "wipe")
    return slide


def build_summary(prs, content, pal, lesson, layout):
    """Slide tổng kết: panel trắng + items"""
    slide = prs.slides.add_slide(layout)

    # Light bg in safe zone
    add_shape_in_safe_zone(slide, MSO_SHAPE.RECTANGLE,
                           0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM - SAFE_TOP,
                           pal["bg"], send_to_back=True)

    # White panel
    panel_x, panel_y = 0.6, SAFE_TOP + 0.15
    panel_w = SLIDE_W - 1.2
    panel_h = SAFE_BOTTOM - panel_y - 0.15
    add_shape_in_safe_zone(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                           panel_x, panel_y, panel_w, panel_h, pal["card"])

    # Top accent line
    add_shape_in_safe_zone(slide, MSO_SHAPE.RECTANGLE,
                           panel_x, panel_y, panel_w, 0.06, pal["primary"])

    # Title
    add_textbox(slide, panel_x + 0.3, panel_y + 0.2, panel_w - 0.6, 0.65,
                content["title"],
                size_pt=24, bold=True, color_hex=pal["primary"], alignment=PP_ALIGN.CENTER)

    # Divider
    div_y = panel_y + 0.95
    add_shape_in_safe_zone(slide, MSO_SHAPE.RECTANGLE,
                           panel_x + 2, div_y, panel_w - 4, 0.03, pal["accent"])

    # Content items as rows with accent dot
    text_content = content.get("content", "")
    items = [l.strip() for l in text_content.split("\n") if l.strip()]
    for i, item in enumerate(items[:4]):
        y = div_y + 0.25 + i * 1.0
        if y + 0.7 > SAFE_BOTTOM:
            break
        # Accent bar
        add_shape_in_safe_zone(slide, MSO_SHAPE.RECTANGLE,
                               panel_x + 0.4, y, 0.08, 0.7, pal["accent"])
        # Item bg
        add_shape_in_safe_zone(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                               panel_x + 0.6, y, panel_w - 1.2, 0.7, pal["bg"])
        # Item text (strip leading checkmarks)
        clean = item.lstrip("✅● ").strip()
        add_textbox(slide, panel_x + 0.9, y + 0.1, panel_w - 1.8, 0.5, clean,
                    size_pt=18, color_hex=pal["text_on_bg"])

    add_slide_transition(slide, "fade")
    return slide


def build_thanks(prs, content, pal, lesson, layout):
    """Slide cảm ơn: banner màu + card BTVN"""
    slide = prs.slides.add_slide(layout)

    # Primary bg in safe zone
    add_shape_in_safe_zone(slide, MSO_SHAPE.RECTANGLE,
                           0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM - SAFE_TOP,
                           pal["primary"], send_to_back=True)

    # Big title
    add_textbox(slide, 1.0, SAFE_TOP + 0.5, SLIDE_W - 2, 1.2, content["title"],
                size_pt=36, bold=True, color_hex=pal["text_on_primary"],
                alignment=PP_ALIGN.CENTER)

    # BTVN card (white)
    hw_text = content.get("content", "")
    if hw_text:
        card_w = 8.0
        card_x = (SLIDE_W - card_w) / 2
        card_y = SAFE_TOP + 2.2
        card_h = min(2.5, SAFE_BOTTOM - card_y - 0.15)
        add_shape_in_safe_zone(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                               card_x, card_y, card_w, card_h, pal["card"])
        lines = hw_text.split("\n")
        add_multiline_textbox(slide, card_x + 0.5, card_y + 0.3, card_w - 1, card_h - 0.6,
                              lines, size_pt=18, color_hex=pal["text_on_card"],
                              alignment=PP_ALIGN.CENTER)

    add_slide_transition(slide, "fade")
    return slide


# ═══════════════════════════════════════════════════
# MAIN — Build all 9 slide decks
# ═══════════════════════════════════════════════════

def create_slide_deck(lesson, palette_idx):
    palette = COLOR_PALETTES[palette_idx % len(COLOR_PALETTES)]
    prs = Presentation(TEMPLATE)

    # Remove template slides safely
    remove_template_slides(prs)

    # Use Blank layout (index 6) — no placeholders to interfere
    layout = prs.slide_layouts[6]

    learn_idx = 0
    for sc in lesson["slides_content"]:
        stype = sc["type"]
        try:
            if stype == "cover":
                build_cover(prs, sc, palette, lesson, layout)
            elif stype in ("warmup", "learn"):
                build_content(prs, sc, palette, lesson, layout, stype, learn_idx)
                if stype == "learn":
                    learn_idx += 1
            elif stype in ("practice", "activity"):
                build_practice(prs, sc, palette, lesson, layout)
            elif stype == "summary":
                build_summary(prs, sc, palette, lesson, layout)
            elif stype == "thanks":
                build_thanks(prs, sc, palette, lesson, layout)
            else:
                build_content(prs, sc, palette, lesson, layout, stype, learn_idx)
        except Exception as e:
            print(f"  ⚠️ Error on slide '{sc.get('title','')}': {e}")

    # Save
    output_dir = os.path.join(KHBD_BASE, lesson["folder"], "Tuần_02")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, lesson["file_name"])

    try:
        prs.save(output_path)
        slide_count = len(prs.slides)
        print(f"✅ {lesson['lop_label']}: {slide_count} slides → {os.path.basename(output_path)}")
    except PermissionError:
        alt = output_path.replace(".pptx", "_v2.pptx")
        prs.save(alt)
        print(f"⚠️ {lesson['lop_label']}: File locked → saved as {os.path.basename(alt)}")
    except Exception as e:
        print(f"❌ {lesson['lop_label']}: {e}")
        traceback.print_exc()

    return output_path


if __name__ == "__main__":
    print("=" * 60)
    print("  Tạo Slide Student-Facing v2 — 9 lớp Tuần 02")
    print("  FIX: Logo/chân trang visible, contrast OK")
    print("=" * 60)

    for i, lesson in enumerate(LESSONS):
        try:
            create_slide_deck(lesson, i)
        except Exception as e:
            print(f"❌ {lesson['lop_label']}: {e}")
            traceback.print_exc()

    print("\n" + "=" * 60)
    print("  Hoàn thành! Mở file để kiểm tra.")
    print("=" * 60)
