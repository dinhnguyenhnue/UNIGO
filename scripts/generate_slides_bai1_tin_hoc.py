"""
Generate Slide Bài 1 Tin học for all grades (Tiền TH, Lớp 1-8).
Reads KHBD .docx to extract lesson content, then creates .pptx slides.

v2: Adds SGK images / AI-generated images to every activity slide.
     Minimizes empty space, maximizes visual content.
"""
import sys, os, glob, copy, re
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from docx import Document
from lxml import etree
from PIL import Image as PILImage

# ── Constants ──────────────────────────────────────────────────────
TEMPLATE = r'D:\UNIGO\Hệ thống mẫu văn bản\Mẫu slide có chân trang.pptx'
KHBD_ROOT = r'D:\UNIGO\KHBD_Tin_học'
SGK_ROOT = r'D:\UNIGO\SGK'

# Slide dimensions (matching existing Tiết 00 slides: 13.33 x 7.5 inches)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Safe zone
SAFE_TOP = Inches(1.15)
SAFE_BOTTOM = Inches(6.30)
SAFE_HEIGHT = Inches(5.15)
CONTENT_TOP = Inches(1.30)
CONTENT_LEFT = Inches(0.50)
CONTENT_RIGHT_MARGIN = Inches(0.50)

# ── 9 Color Palettes ──────────────────────────────────────────────
PALETTES = {
    'TTH': {  # Tiền tiểu học - Hồng cam vui nhộn
        'primary': RGBColor(0xFF, 0x6B, 0x6B),
        'secondary': RGBColor(0xFF, 0x8E, 0x53),
        'accent': RGBColor(0xFF, 0xD9, 0x3D),
        'bg': RGBColor(0xFF, 0xF5, 0xF5),
        'text_dark': RGBColor(0x2D, 0x13, 0x3A),
        'text_light': RGBColor(0xFF, 0xFF, 0xFF),
        'card_bg': RGBColor(0xFF, 0xFF, 0xFF),
    },
    '1': {  # Lớp 1 - Xanh lá tươi
        'primary': RGBColor(0x4E, 0xCA, 0x54),
        'secondary': RGBColor(0x2E, 0x7D, 0x32),
        'accent': RGBColor(0xFF, 0xEB, 0x3B),
        'bg': RGBColor(0xF1, 0xF8, 0xE9),
        'text_dark': RGBColor(0x1B, 0x2D, 0x1B),
        'text_light': RGBColor(0xFF, 0xFF, 0xFF),
        'card_bg': RGBColor(0xFF, 0xFF, 0xFF),
    },
    '2': {  # Lớp 2 - Tím hoa oải hương
        'primary': RGBColor(0x7C, 0x4D, 0xFF),
        'secondary': RGBColor(0x53, 0x6D, 0xFE),
        'accent': RGBColor(0xFF, 0xC1, 0x07),
        'bg': RGBColor(0xF3, 0xF0, 0xFF),
        'text_dark': RGBColor(0x1A, 0x12, 0x3E),
        'text_light': RGBColor(0xFF, 0xFF, 0xFF),
        'card_bg': RGBColor(0xFF, 0xFF, 0xFF),
    },
    '3': {  # Lớp 3 - Xanh dương tươi
        'primary': RGBColor(0x21, 0x96, 0xF3),
        'secondary': RGBColor(0x15, 0x65, 0xC0),
        'accent': RGBColor(0xFF, 0xAB, 0x00),
        'bg': RGBColor(0xE3, 0xF2, 0xFD),
        'text_dark': RGBColor(0x0D, 0x1F, 0x3C),
        'text_light': RGBColor(0xFF, 0xFF, 0xFF),
        'card_bg': RGBColor(0xFF, 0xFF, 0xFF),
    },
    '4': {  # Lớp 4 - Cam ấm
        'primary': RGBColor(0xFF, 0x98, 0x00),
        'secondary': RGBColor(0xEF, 0x6C, 0x00),
        'accent': RGBColor(0x00, 0xBC, 0xD4),
        'bg': RGBColor(0xFF, 0xF3, 0xE0),
        'text_dark': RGBColor(0x33, 0x1A, 0x00),
        'text_light': RGBColor(0xFF, 0xFF, 0xFF),
        'card_bg': RGBColor(0xFF, 0xFF, 0xFF),
    },
    '5': {  # Lớp 5 - Hồng tím
        'primary': RGBColor(0xE9, 0x1E, 0x63),
        'secondary': RGBColor(0x88, 0x0E, 0x4F),
        'accent': RGBColor(0x00, 0xE5, 0xFF),
        'bg': RGBColor(0xFC, 0xE4, 0xEC),
        'text_dark': RGBColor(0x3E, 0x0A, 0x22),
        'text_light': RGBColor(0xFF, 0xFF, 0xFF),
        'card_bg': RGBColor(0xFF, 0xFF, 0xFF),
    },
    '6': {  # Lớp 6 - Xanh teal
        'primary': RGBColor(0x00, 0x96, 0x88),
        'secondary': RGBColor(0x00, 0x4D, 0x40),
        'accent': RGBColor(0xFF, 0xD5, 0x4F),
        'bg': RGBColor(0xE0, 0xF2, 0xF1),
        'text_dark': RGBColor(0x0A, 0x2B, 0x28),
        'text_light': RGBColor(0xFF, 0xFF, 0xFF),
        'card_bg': RGBColor(0xFF, 0xFF, 0xFF),
    },
    '7': {  # Lớp 7 - Xanh đậm (Navy)
        'primary': RGBColor(0x30, 0x3F, 0x9F),
        'secondary': RGBColor(0x1A, 0x23, 0x7E),
        'accent': RGBColor(0xFF, 0xC1, 0x07),
        'bg': RGBColor(0xE8, 0xEA, 0xF6),
        'text_dark': RGBColor(0x0F, 0x14, 0x3A),
        'text_light': RGBColor(0xFF, 0xFF, 0xFF),
        'card_bg': RGBColor(0xFF, 0xFF, 0xFF),
    },
    '8': {  # Lớp 8 - Xanh ngọc bích
        'primary': RGBColor(0x00, 0x79, 0x7B),
        'secondary': RGBColor(0x00, 0x50, 0x52),
        'accent': RGBColor(0xFF, 0xB7, 0x4D),
        'bg': RGBColor(0xE0, 0xF7, 0xFA),
        'text_dark': RGBColor(0x00, 0x2B, 0x2C),
        'text_light': RGBColor(0xFF, 0xFF, 0xFF),
        'card_bg': RGBColor(0xFF, 0xFF, 0xFF),
    },
}

# ── Grade Config ──────────────────────────────────────────────────
GRADES = [
    {
        'key': 'TTH', 'folder': 'Tiền_tiểu_học', 'label': 'Tiền TH',
        'display': 'TIỀN TIỂU HỌC', 'sgk': False,
    },
    {
        'key': '1', 'folder': 'Lớp_1', 'label': 'Lớp 1',
        'display': 'LỚP 1', 'sgk': False,
    },
    {
        'key': '2', 'folder': 'Lớp_2', 'label': 'Lớp 2',
        'display': 'LỚP 2', 'sgk': False,
    },
    {
        'key': '3', 'folder': 'Lớp_3', 'label': 'Lớp 3',
        'display': 'LỚP 3', 'sgk': True, 'sgk_grade': 3,
    },
    {
        'key': '4', 'folder': 'Lớp_4', 'label': 'Lớp 4',
        'display': 'LỚP 4', 'sgk': True, 'sgk_grade': 4,
    },
    {
        'key': '5', 'folder': 'Lớp_5', 'label': 'Lớp 5',
        'display': 'LỚP 5', 'sgk': True, 'sgk_grade': 5,
    },
    {
        'key': '6', 'folder': 'Lớp_6', 'label': 'Lớp 6',
        'display': 'LỚP 6', 'sgk': True, 'sgk_grade': 6,
    },
    {
        'key': '7', 'folder': 'Lớp_7', 'label': 'Lớp 7',
        'display': 'LỚP 7', 'sgk': True, 'sgk_grade': 7,
    },
    {
        'key': '8', 'folder': 'Lớp_8', 'label': 'Lớp 8',
        'display': 'LỚP 8', 'sgk': True, 'sgk_grade': 8,
    },
]


# ══════════════════════════════════════════════════════════════════
#                     IMAGE FINDER
# ══════════════════════════════════════════════════════════════════

def find_images(grade_cfg):
    """Find available images for a grade.
    
    Returns dict with keys: 'full_pages' (sorted list), 'individual' (sorted list),
    'generated' (dict with named images like cover, activity, practice, summary).
    """
    result = {
        'full_pages': [],
        'individual': [],
        'generated': {},
    }
    
    # 1. SGK full pages
    if grade_cfg['sgk']:
        grade_num = grade_cfg['sgk_grade']
        fp_dir = os.path.join(SGK_ROOT, f'Lớp_{grade_num}', 'bai1_images', 'full_pages')
        if os.path.isdir(fp_dir):
            pages = sorted(glob.glob(os.path.join(fp_dir, '*.png')))
            result['full_pages'] = pages
        
        # Individual extracted images (skip if too many = tiled PDF)
        img_dir = os.path.join(SGK_ROOT, f'Lớp_{grade_num}', 'bai1_images')
        if os.path.isdir(img_dir):
            indiv = sorted([
                os.path.join(img_dir, f) for f in os.listdir(img_dir)
                if f.endswith(('.jpeg', '.jpg', '.png')) and os.path.isfile(os.path.join(img_dir, f))
            ])
            # If too many (tiled PDF like Lớp 8), skip individual images
            if len(indiv) <= 50:
                result['individual'] = indiv
    
    # 2. AI-generated images
    gen_dir = os.path.join(KHBD_ROOT, grade_cfg['folder'], 'Bài_01', 'images')
    if os.path.isdir(gen_dir):
        for f in os.listdir(gen_dir):
            if f.endswith(('.png', '.jpg', '.jpeg')):
                name = os.path.splitext(f)[0]  # cover, activity, practice, summary
                result['generated'][name] = os.path.join(gen_dir, f)
    
    return result


def pick_image_for_slide(images, slide_type, slide_index=0):
    """Pick the best image for a given slide type.
    
    slide_type: 'cover', 'objective', 'khởi_động', 'kiến_thức', 'luyện_tập', 'vận_dụng', 'summary', 'thankyou'
    slide_index: index within the type (e.g., 2nd knowledge slide = 1)
    
    Returns: absolute path to image, or None.
    """
    fp = images['full_pages']
    indiv = images['individual']
    gen = images['generated']
    
    # Strategy: map slide type to page numbers
    # Full pages are usually ordered: [muc_luc, bai1_page1, bai1_page2, ...]
    # We skip the first page if it's a TOC
    content_pages = fp[1:] if len(fp) > 1 else fp  # Skip potential TOC page
    
    if slide_type == 'cover':
        # Cover: first content page of SGK or generated cover
        if gen.get('cover'):
            return gen['cover']
        if content_pages:
            return content_pages[0]
    
    elif slide_type == 'objective':
        # Objective: first page showing YCCD (usually page 1 of bài)
        if content_pages:
            return content_pages[0]
        if gen.get('cover'):
            return gen['cover']
    
    elif slide_type == 'khởi_động':
        # Khởi động: first content page (has KD section)
        if content_pages:
            return content_pages[0]
        if gen.get('activity'):
            return gen['activity']
    
    elif slide_type == 'kiến_thức':
        # Knowledge slides: distribute across content pages
        # Pages 2-5 are typically knowledge content
        page_offset = 1 + slide_index  # start from 2nd content page
        if page_offset < len(content_pages):
            return content_pages[page_offset]
        # Fallback: cycle through individual images
        if indiv and slide_index < len(indiv):
            return indiv[slide_index]
        if gen.get('activity'):
            return gen['activity']
    
    elif slide_type == 'luyện_tập':
        # Practice: later pages 
        page_offset = max(3, len(content_pages) - 3) + slide_index
        if page_offset < len(content_pages):
            return content_pages[page_offset]
        if gen.get('practice'):
            return gen['practice']
        # Fallback to last content pages
        if content_pages:
            return content_pages[-1]
    
    elif slide_type == 'vận_dụng':
        # Application: second-to-last page
        if len(content_pages) >= 2:
            return content_pages[-2]
        if gen.get('practice'):
            return gen['practice']
        if content_pages:
            return content_pages[-1]
    
    elif slide_type == 'summary':
        # Summary: last page or generated summary
        if gen.get('summary'):
            return gen['summary']
        if content_pages:
            return content_pages[-1]
    
    elif slide_type == 'thankyou':
        if gen.get('cover'):
            return gen['cover']
        if content_pages:
            return content_pages[0]
    
    # Ultimate fallback: any available image
    if gen:
        return list(gen.values())[0]
    if content_pages:
        return content_pages[0]
    if indiv:
        return indiv[0]
    
    return None


# ══════════════════════════════════════════════════════════════════
#                     KHBD READER
# ══════════════════════════════════════════════════════════════════

def read_khbd(grade_cfg):
    """Read KHBD .docx and extract lesson data."""
    folder = grade_cfg['folder']
    files = glob.glob(os.path.join(KHBD_ROOT, folder, 'Bài_01', '*.docx'))
    if not files:
        raise FileNotFoundError(f"No KHBD found for {folder}")
    
    doc = Document(files[0])
    
    data = {
        'title': '',          # e.g. "BÀI 1. THÔNG TIN VÀ QUYẾT ĐỊNH"
        'topic': '',          # e.g. "Thông tin quanh ta"
        'objectives': [],     # Yêu cầu cần đạt
        'activities': [],     # List of activity dicts
    }
    
    # Extract title and topic from paragraphs
    is_thcs = grade_cfg['key'] in ('6', '7', '8')
    
    for p in doc.paragraphs:
        txt = p.text.strip()
        if 'BÀI:' in txt:
            raw = txt.replace('BÀI:', '').strip()
            raw = re.sub(r'\s*\(Tiết.*?\)', '', raw).strip()
            data['title'] = raw
        elif 'CHỦ ĐIỂM:' in txt or 'CHỦ ĐỀ:' in txt:
            data['topic'] = txt.split(':', 1)[1].strip()
        elif txt.startswith('TÊN BÀI DẠY:'):
            data['title'] = 'Bài 1. ' + txt.replace('TÊN BÀI DẠY:', '').strip()
        elif txt.startswith('Tên tiết:') and not data['topic']:
            data['topic'] = txt.replace('Tên tiết:', '').strip()
    
    # Extract objectives
    in_objectives = False
    for p in doc.paragraphs:
        txt = p.text.strip()
        if 'YÊU CẦU CẦN ĐẠT' in txt or 'Mục tiêu' == txt.strip():
            in_objectives = True
            continue
        if in_objectives:
            if txt.startswith('II.') or txt.startswith('III.') or txt.startswith('1. Kiến thức'):
                if txt.startswith('1. Kiến thức'):
                    in_objectives = True
                    continue
                break
            if txt.startswith('- ') and 'NL' not in txt and 'Năng lực' not in txt:
                obj = txt.lstrip('- ').strip()
                if len(obj) > 10:
                    data['objectives'].append(obj)
    
    # For THCS, extract objectives from "Kiến thức" section
    if is_thcs and not data['objectives']:
        for p in doc.paragraphs:
            txt = p.text.strip()
            if txt.startswith('- Sự ') or txt.startswith('- Khả năng') or txt.startswith('- Sự nhận biết'):
                data['objectives'].append(txt.lstrip('- '))
    
    # Extract activities
    if is_thcs:
        activities = _extract_thcs_activities(doc)
        data['activities'] = activities
    else:
        activities = _extract_th_activities(doc)
        data['activities'] = activities
    
    return data


def _extract_th_activities(doc):
    """Extract activities from TH-format KHBD (single table with GV|HS)."""
    activities = []
    if not doc.tables:
        return activities
    
    table = doc.tables[0]
    current_act = None
    
    for row in table.rows:
        gv_text = row.cells[0].text.strip()
        hs_text = row.cells[1].text.strip() if len(row.cells) > 1 else ''
        
        gv_lines = gv_text.split('\n')
        gv_extra = '\n'.join(gv_lines[1:]).strip() if len(gv_lines) > 1 else ''
        
        if 'MỞ ĐẦU' in gv_text or 'Hoạt động MỞ ĐẦU' in gv_text:
            current_act = {'type': 'khởi_động', 'title': 'Khởi động', 'gv': '', 'hs': ''}
            activities.append(current_act)
            if gv_extra:
                current_act['gv'] += gv_extra + '\n'
                current_act['hs'] += hs_text + '\n'
        elif 'KIẾN THỨC MỚI' in gv_text and not re.match(r'^2\.\d', gv_text):
            continue
        elif re.match(r'^2\.\d\.\s', gv_text):
            match = re.match(r'^2\.\d\.\s*(.+?)(?:\s*\(\d+ phút\))?$', gv_lines[0])
            title = match.group(1) if match else gv_lines[0][:50]
            current_act = {'type': 'kiến_thức', 'title': title, 'gv': '', 'hs': ''}
            activities.append(current_act)
        elif 'LUYỆN TẬP' in gv_text and not re.match(r'^3\.\d', gv_text):
            continue
        elif re.match(r'^3\.\d\.\s', gv_text):
            match = re.match(r'^3\.\d\.\s*(.+?)(?:\s*\(\d+ phút\))?$', gv_lines[0])
            title = match.group(1) if match else gv_lines[0][:50]
            current_act = {'type': 'luyện_tập', 'title': title, 'gv': '', 'hs': ''}
            activities.append(current_act)
        elif 'VẬN DỤNG' in gv_text or 'TRẢI NGHIỆM' in gv_text:
            current_act = {'type': 'vận_dụng', 'title': 'Vận dụng', 'gv': '', 'hs': ''}
            activities.append(current_act)
            if gv_extra:
                current_act['gv'] += gv_extra + '\n'
                current_act['hs'] += hs_text + '\n'
        elif current_act is not None:
            current_act['gv'] += gv_text + '\n'
            current_act['hs'] += hs_text + '\n'
    
    return activities


def _extract_thcs_activities(doc):
    """Extract activities from THCS-format KHBD (paragraphs + multiple tables)."""
    activities = []
    
    for p in doc.paragraphs:
        txt = p.text.strip()
        
        if re.search(r'Hoạt động\s*1', txt) and any(kw in txt for kw in ['Khởi động', 'Mở đầu', 'Xác định']):
            act = {'type': 'khởi_động', 'title': 'Khởi động', 'content': '', 'product': ''}
            activities.append(act)
        elif re.search(r'Hoạt động\s*2', txt) and any(kw in txt.lower() for kw in ['hình thành', 'kiến thức', 'giải quyết']):
            act = {'type': 'kiến_thức', 'title': 'Hình thành kiến thức mới', 'content': '', 'product': ''}
            activities.append(act)
        elif re.search(r'Hoạt động\s*3', txt) and 'Luyện tập' in txt:
            act = {'type': 'luyện_tập', 'title': 'Luyện tập', 'content': '', 'product': ''}
            activities.append(act)
        elif re.search(r'Hoạt động\s*4', txt) and any(kw in txt for kw in ['Vận dụng', 'vận dụng', 'Mở rộng', 'mở rộng']):
            act = {'type': 'vận_dụng', 'title': 'Vận dụng – Mở rộng', 'content': '', 'product': ''}
            activities.append(act)
        elif activities and txt.startswith('a) Mục tiêu:'):
            activities[-1]['objective'] = txt.replace('a) Mục tiêu:', '').strip()
        elif activities and txt.startswith('b) Nội dung:'):
            activities[-1]['content'] = txt.replace('b) Nội dung:', '').strip()
        elif activities and txt.startswith('c) Sản phẩm:'):
            activities[-1]['product'] = txt.replace('c) Sản phẩm:', '').strip()
        elif activities and txt.startswith('d) Tổ chức thực hiện:'):
            activities[-1]['organize'] = txt.replace('d) Tổ chức thực hiện:', '').strip()
    
    # Also extract GV content from tables
    act_tables = [t for t in doc.tables if len(t.columns) == 3 and len(t.rows) >= 4]
    for i, (act, tbl) in enumerate(zip(activities, act_tables)):
        gv_content = []
        for row in tbl.rows[1:]:
            gv_text = row.cells[1].text.strip() if len(row.cells) > 1 else ''
            if gv_text and len(gv_text) > 10:
                for line in gv_text.split('\n'):
                    line = line.strip()
                    if line and len(line) > 15 and not line.startswith('HS '):
                        gv_content.append(line)
                        break
        if gv_content:
            act['gv_bullets'] = gv_content
    
    return activities


# ══════════════════════════════════════════════════════════════════
#                     SLIDE BUILDER
# ══════════════════════════════════════════════════════════════════

def create_presentation(grade_cfg, lesson_data, images):
    """Create a complete slide deck for a grade's Bài 1."""
    pal = PALETTES[grade_cfg['key']]
    
    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 13.33 inches
    prs.slide_height = Emu(6858000)   # 7.5 inches
    
    blank_layout = prs.slide_layouts[6]  # Blank
    
    title_short = lesson_data['title']
    display_title = re.sub(r'^BÀI\s*\d+\.\s*', '', title_short, flags=re.IGNORECASE).strip()
    
    slides_created = []
    img_counter = {'kiến_thức': 0, 'luyện_tập': 0}
    
    # ── Slide 1: TRANG BÌA ────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    cover_img = pick_image_for_slide(images, 'cover')
    _add_cover_slide(slide, prs, grade_cfg, display_title, lesson_data['topic'], pal, cover_img)
    slides_created.append('Trang bìa')
    
    # ── Slide 2: MỤC TIÊU ────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    obj_img = pick_image_for_slide(images, 'objective')
    _add_objectives_slide(slide, prs, grade_cfg, lesson_data, pal, obj_img)
    slides_created.append('Mục tiêu')
    
    # ── Slide 3: KHỞI ĐỘNG ────────────────────────────────────────
    kd_acts = [a for a in lesson_data['activities'] if a['type'] == 'khởi_động']
    if kd_acts:
        slide = prs.slides.add_slide(blank_layout)
        kd_img = pick_image_for_slide(images, 'khởi_động')
        _add_activity_slide(slide, prs, '🎮  Khởi động', kd_acts[0], pal, 'left', kd_img)
        slides_created.append('Khởi động')
    
    # ── Slide 4-7: KIẾN THỨC MỚI ─────────────────────────────────
    kt_acts = [a for a in lesson_data['activities'] if a['type'] == 'kiến_thức']
    for i, act in enumerate(kt_acts):
        slide = prs.slides.add_slide(blank_layout)
        side = 'right' if i % 2 == 0 else 'left'
        emoji = ['📖', '📝', '💡', '🔍'][i % 4]
        title_text = f'{emoji}  {act.get("title", "Kiến thức mới")}'
        kt_img = pick_image_for_slide(images, 'kiến_thức', i)
        _add_activity_slide(slide, prs, title_text, act, pal, side, kt_img)
        slides_created.append(f'Kiến thức {i+1}')
    
    if not kt_acts:
        slide = prs.slides.add_slide(blank_layout)
        kt_img = pick_image_for_slide(images, 'kiến_thức', 0)
        _add_generic_knowledge_slide(slide, prs, display_title, pal, kt_img)
        slides_created.append('Kiến thức mới')
    
    # ── Slide 8-9: LUYỆN TẬP ─────────────────────────────────────
    lt_acts = [a for a in lesson_data['activities'] if a['type'] == 'luyện_tập']
    for i, act in enumerate(lt_acts):
        slide = prs.slides.add_slide(blank_layout)
        emoji = ['✏️', '🎯'][i % 2]
        title_text = f'{emoji}  {act.get("title", "Luyện tập")}'
        lt_img = pick_image_for_slide(images, 'luyện_tập', i)
        _add_activity_slide(slide, prs, title_text, act, pal, 'right' if i % 2 == 0 else 'left', lt_img)
        slides_created.append(f'Luyện tập {i+1}')
    
    # ── Slide 10: VẬN DỤNG ───────────────────────────────────────
    vd_acts = [a for a in lesson_data['activities'] if a['type'] == 'vận_dụng']
    if vd_acts:
        slide = prs.slides.add_slide(blank_layout)
        vd_img = pick_image_for_slide(images, 'vận_dụng')
        _add_activity_slide(slide, prs, '🚀  Vận dụng', vd_acts[0], pal, 'left', vd_img)
        slides_created.append('Vận dụng')
    
    # ── Slide 11: TỔNG KẾT ───────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    sum_img = pick_image_for_slide(images, 'summary')
    _add_summary_slide(slide, prs, display_title, lesson_data, pal, sum_img)
    slides_created.append('Tổng kết')
    
    # ── Slide 12: CẢM ƠN ─────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    ty_img = pick_image_for_slide(images, 'thankyou')
    _add_thankyou_slide(slide, prs, grade_cfg, pal, ty_img)
    slides_created.append('Cảm ơn')
    
    # ── Add footer to ALL slides ──────────────────────────────────
    for slide in prs.slides:
        _add_footer(slide, prs)
    
    return prs, slides_created


# ══════════════════════════════════════════════════════════════════
#                     SHAPE HELPERS
# ══════════════════════════════════════════════════════════════════

def _set_slide_bg(slide, color):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None):
    """Add a shape with optional fill and border."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text, font_size=Pt(18),
                 font_color=RGBColor(0, 0, 0), bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Times New Roman'):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_bullet_textbox(slide, left, top, width, height, items, pal,
                        font_size=Pt(18), bullet_char='●', line_spacing=Pt(28)):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        run_bullet = p.add_run()
        run_bullet.text = f'{bullet_char}  '
        run_bullet.font.size = Pt(14)
        run_bullet.font.color.rgb = pal['primary']
        run_bullet.font.name = 'Times New Roman'
        
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = font_size
        run_text.font.color.rgb = pal['text_dark']
        run_text.font.name = 'Times New Roman'
        
        p.space_after = Pt(8)
        p.line_spacing = line_spacing
    
    return txBox


def _add_card_with_accent(slide, left, top, width, height, pal, accent_side='left'):
    """Add a white card with colored accent bar and shadow."""
    card = _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
                      fill_color=pal['card_bg'])
    card.adjustments[0] = 0.02
    
    # Shadow
    spPr = card._element.spPr
    effectLst = etree.SubElement(spPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst')
    outerShdw = etree.SubElement(effectLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw',
                                  attrib={'blurRad': '50800', 'dist': '25400', 'dir': '5400000',
                                          'rotWithShape': '0'})
    srgbClr = etree.SubElement(outerShdw, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr',
                                attrib={'val': '000000'})
    etree.SubElement(srgbClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha',
                     attrib={'val': '25000'})
    
    # Accent bar
    bar_w = Inches(0.12)
    if accent_side == 'left':
        bar_left = left
    else:
        bar_left = left + width - bar_w
    
    bar = _add_shape(slide, MSO_SHAPE.RECTANGLE, bar_left, top, bar_w, height,
                     fill_color=pal['primary'])
    
    return card, bar


def _add_picture_fitted(slide, img_path, left, top, max_width, max_height):
    """Add a picture that fits within given bounds while maintaining aspect ratio.
    Returns the picture shape, or None if image can't be loaded.
    """
    if not img_path or not os.path.exists(img_path):
        return None
    
    try:
        with PILImage.open(img_path) as img:
            img_w, img_h = img.size
    except Exception:
        return None
    
    # Calculate scale to fit within bounds
    scale_w = max_width / img_w
    scale_h = max_height / img_h
    scale = min(scale_w, scale_h)
    
    final_w = int(img_w * scale)
    final_h = int(img_h * scale)
    
    # Center within the given area
    offset_x = (max_width - final_w) // 2
    offset_y = (max_height - final_h) // 2
    
    pic = slide.shapes.add_picture(
        img_path,
        left + offset_x,
        top + offset_y,
        final_w,
        final_h
    )
    
    return pic


def _add_picture_cropped(slide, img_path, left, top, width, height):
    """Add a picture that fills the given area (may crop edges).
    Uses fill-and-crop strategy for maximum visual impact.
    Returns the picture shape, or None.
    """
    if not img_path or not os.path.exists(img_path):
        return None
    
    try:
        with PILImage.open(img_path) as img:
            img_w, img_h = img.size
    except Exception:
        return None
    
    # Calculate scale to FILL (not fit) the area
    scale_w = width / img_w
    scale_h = height / img_h
    scale = max(scale_w, scale_h)  # Use max to fill
    
    final_w = int(img_w * scale)
    final_h = int(img_h * scale)
    
    # Center (will be cropped by PowerPoint's crop feature)
    pic = slide.shapes.add_picture(
        img_path,
        left,
        top,
        width,
        height
    )
    
    return pic


# ══════════════════════════════════════════════════════════════════
#                     SLIDE CREATORS
# ══════════════════════════════════════════════════════════════════

def _add_cover_slide(slide, prs, grade_cfg, display_title, topic, pal, img_path):
    """Create cover slide with image background."""
    _set_slide_bg(slide, pal['primary'])
    sw = prs.slide_width
    
    # Add background image (right side) if available
    if img_path:
        # Right-side large image
        img_w = Inches(5.5)
        img_h = Inches(5.0)
        img_l = sw - img_w - Inches(0.3)
        img_t = CONTENT_TOP
        pic = _add_picture_fitted(slide, img_path, img_l, img_t, img_w, img_h)
        if pic:
            # Add subtle rounded overlay
            overlay = _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                                img_l - Inches(0.1), img_t - Inches(0.1),
                                img_w + Inches(0.2), img_h + Inches(0.2),
                                fill_color=None, line_color=pal['accent'])
            if overlay:
                overlay.line.width = Pt(3)
                overlay.adjustments[0] = 0.03
    
    # White card on left
    card_w = Inches(7.5)
    card_h = Inches(4.6)
    card_l = Inches(0.4)
    card_t = Inches(1.40)
    
    card = _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, card_l, card_t, card_w, card_h,
                      fill_color=RGBColor(0xFF, 0xFF, 0xFF))
    card.adjustments[0] = 0.03
    
    # Grade badge
    badge_w = Inches(4.5)
    badge_h = Inches(0.50)
    badge_l = card_l + Inches(0.5)
    badge_t = card_t + Inches(0.3)
    
    badge = _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, badge_l, badge_t, badge_w, badge_h,
                       fill_color=pal['primary'])
    badge.adjustments[0] = 0.3
    badge_tf = badge.text_frame
    badge_tf.word_wrap = True
    badge_p = badge_tf.paragraphs[0]
    badge_p.alignment = PP_ALIGN.CENTER
    run = badge_p.add_run()
    run.text = f'TIN HỌC {grade_cfg["display"]}  •  NĂM HỌC 2026 – 2027'
    run.font.size = Pt(16)
    run.font.color.rgb = pal['text_light']
    run.font.bold = True
    run.font.name = 'Times New Roman'
    
    # Title
    title_l = card_l + Inches(0.5)
    title_t = badge_t + Inches(0.8)
    title_w = card_w - Inches(1)
    title_h = Inches(2.2)
    
    txBox = slide.shapes.add_textbox(title_l, title_t, title_w, title_h)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = 'Bài 1'
    r1.font.size = Pt(22)
    r1.font.color.rgb = pal['primary']
    r1.font.bold = True
    r1.font.name = 'Times New Roman'
    
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = display_title
    r2.font.size = Pt(28)
    r2.font.color.rgb = pal['text_dark']
    r2.font.bold = True
    r2.font.name = 'Times New Roman'
    p2.space_before = Pt(4)
    
    if topic:
        p3 = tf.add_paragraph()
        r3 = p3.add_run()
        r3.text = f'Chủ đề: {topic}'
        r3.font.size = Pt(16)
        r3.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        r3.font.italic = True
        r3.font.name = 'Times New Roman'
        p3.space_before = Pt(8)
    
    # Teacher info
    info_l = card_l + Inches(0.5)
    info_t = card_t + card_h - Inches(1.2)
    info_w = card_w - Inches(1)
    info_h = Inches(1.0)
    
    txBox2 = slide.shapes.add_textbox(info_l, info_t, info_w, info_h)
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p_school = tf2.paragraphs[0]
    r_school = p_school.add_run()
    r_school.text = 'Trường TH & THCS UNIGO'
    r_school.font.size = Pt(16)
    r_school.font.color.rgb = pal['text_dark']
    r_school.font.name = 'Times New Roman'
    
    p_gv = tf2.add_paragraph()
    r_gv = p_gv.add_run()
    r_gv.text = 'Giáo viên: Đậu Đình Nguyên'
    r_gv.font.size = Pt(16)
    r_gv.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    r_gv.font.name = 'Times New Roman'


def _add_objectives_slide(slide, prs, grade_cfg, lesson_data, pal, img_path):
    """Create objectives slide with image on right side."""
    _set_slide_bg(slide, pal['bg'])
    sw = prs.slide_width
    
    # Left card with objectives
    card_w = Inches(7.8)
    card_h = Inches(4.8)
    card_l = Inches(0.35)
    card_t = CONTENT_TOP
    
    _add_card_with_accent(slide, card_l, card_t, card_w, card_h, pal, 'left')
    
    _add_textbox(slide, card_l + Inches(0.5), card_t + Inches(0.2),
                 Inches(7), Inches(0.5),
                 '🎯  Mục tiêu bài học', Pt(24), pal['primary'], bold=True)
    
    objectives = lesson_data.get('objectives', [])
    if not objectives:
        objectives = ['Nắm được kiến thức cơ bản của bài học.',
                      'Vận dụng kiến thức vào thực hành.',
                      'Phát triển năng lực tư duy và hợp tác.']
    objectives = objectives[:4]
    
    _add_bullet_textbox(slide, card_l + Inches(0.5), card_t + Inches(0.9),
                        Inches(7), Inches(3.6), objectives, pal,
                        font_size=Pt(18))
    
    # Right image
    if img_path:
        img_l = card_l + card_w + Inches(0.25)
        img_w = sw - img_l - Inches(0.35)
        _add_picture_fitted(slide, img_path, img_l, card_t, img_w, card_h)


def _get_activity_content(act):
    """Extract content from an activity dict into bullet points."""
    bullets = []
    
    if 'gv' in act:
        gv_text = act['gv'].strip()
        for line in gv_text.split('\n'):
            line = line.strip()
            if line.startswith('*') or line.startswith('-'):
                line = line.lstrip('*- ').strip()
            if line.startswith('GV '):
                line = line[3:]
            if line.startswith('Mục tiêu:'):
                continue
            if len(line) > 15 and len(line) < 120:
                line = line[0].upper() + line[1:] if line else line
                bullets.append(line)
    
    if 'content' in act and act['content']:
        content = act['content']
        parts = re.split(r'[;]', content)
        for part in parts:
            part = part.strip().rstrip('.')
            if len(part) > 10 and len(part) < 120:
                part = part[0].upper() + part[1:] if part else part
                bullets.append(part)
    
    if 'gv_bullets' in act and act['gv_bullets']:
        for b in act['gv_bullets']:
            b = b.strip()
            if b.startswith('GV '):
                b = b[3:]
            if len(b) > 10 and len(b) < 120:
                b = b[0].upper() + b[1:] if b else b
                if b not in bullets:
                    bullets.append(b)
    
    if 'product' in act and act['product']:
        product = act['product'][:100]
        product = product[0].upper() + product[1:] if product else product
        bullets.append(f'📋 Sản phẩm: {product}')
    
    cleaned = []
    seen = set()
    for b in bullets:
        b = b.strip()
        if b and len(b) > 5 and b not in seen:
            seen.add(b)
            if len(b) > 90:
                b = b[:87] + '...'
            cleaned.append(b)
    
    return cleaned[:4]


def _add_activity_slide(slide, prs, title_text, act, pal, card_side='left', img_path=None):
    """Create an activity slide with card + image layout.
    
    The card side determines text placement, image goes on opposite side.
    """
    _set_slide_bg(slide, pal['bg'])
    sw = prs.slide_width
    
    card_w = Inches(7.0)
    card_h = Inches(4.8)
    
    if card_side == 'left':
        card_l = Inches(0.35)
        img_l = card_l + card_w + Inches(0.25)
        img_w = sw - img_l - Inches(0.35)
    else:
        img_l = Inches(0.35)
        img_w = Inches(5.5)
        card_l = img_l + img_w + Inches(0.25)
    
    card_t = CONTENT_TOP
    
    accent_side = 'left' if card_side == 'left' else 'right'
    _add_card_with_accent(slide, card_l, card_t, card_w, card_h, pal, accent_side)
    
    # Title
    title_l = card_l + Inches(0.4)
    _add_textbox(slide, title_l, card_t + Inches(0.2),
                 card_w - Inches(0.8), Inches(0.5),
                 title_text, Pt(22), pal['primary'], bold=True)
    
    # Content bullets
    bullets = _get_activity_content(act)
    if not bullets:
        act_title = act.get('title', 'Hoạt động')
        bullets = [f'Thực hiện hoạt động: {act_title}',
                   'Học sinh làm việc theo hướng dẫn của giáo viên.',
                   'Hoàn thành sản phẩm học tập.']
    
    _add_bullet_textbox(slide, title_l, card_t + Inches(0.9),
                        card_w - Inches(0.8), Inches(3.6), bullets, pal,
                        font_size=Pt(18))
    
    # IMAGE: Replace the old decorative rectangle with actual image
    if img_path and os.path.exists(img_path):
        # Add image with a subtle background card
        img_card = _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                             img_l, card_t, img_w, card_h,
                             fill_color=RGBColor(0xF8, 0xF8, 0xF8))
        img_card.adjustments[0] = 0.03
        # No border
        img_card.line.fill.background()
        
        # Add the actual image fitted inside the card (with small padding)
        pad = Inches(0.15)
        pic = _add_picture_fitted(slide, img_path,
                                  img_l + pad, card_t + pad,
                                  img_w - pad * 2, card_h - pad * 2)
        
        if pic:
            # Add a thin border around the image card
            img_card.line.color.rgb = pal['primary']
            img_card.line.width = Pt(1.5)
            
            # Add "SGK" label badge at top-right corner of image
            badge_w = Inches(1.2)
            badge_h = Inches(0.30)
            badge_l = img_l + img_w - badge_w - Inches(0.1)
            badge_t = card_t + Inches(0.1)
            badge = _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                              badge_l, badge_t, badge_w, badge_h,
                              fill_color=pal['primary'])
            badge.adjustments[0] = 0.3
            badge_tf = badge.text_frame
            badge_p = badge_tf.paragraphs[0]
            badge_p.alignment = PP_ALIGN.CENTER
            r = badge_p.add_run()
            r.text = '📚 SGK'
            r.font.size = Pt(10)
            r.font.color.rgb = pal['text_light']
            r.font.bold = True
            r.font.name = 'Times New Roman'
    else:
        # Fallback: colored area with emoji (original behavior, but improved)
        deco = _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                         img_l if card_side == 'left' else Inches(0.35),
                         card_t, img_w, card_h,
                         fill_color=pal['primary'])
        deco.adjustments[0] = 0.04
        fill_el = deco._element.spPr.solidFill
        if fill_el is not None:
            srgb = fill_el.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
            if srgb is not None:
                etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha',
                                 attrib={'val': '20000'})
        
        emoji_map = {
            'khởi_động': '🎮',
            'kiến_thức': '📖',
            'luyện_tập': '✏️',
            'vận_dụng': '🚀',
        }
        emoji = emoji_map.get(act.get('type', ''), '📌')
        _add_textbox(slide,
                     (img_l if card_side == 'left' else Inches(0.35)) + img_w // 2 - Inches(0.5),
                     card_t + card_h // 2 - Inches(0.5),
                     Inches(1), Inches(1),
                     emoji, Pt(48), pal['text_light'], alignment=PP_ALIGN.CENTER)


def _add_generic_knowledge_slide(slide, prs, display_title, pal, img_path=None):
    """Add a generic knowledge slide when no activities are extracted."""
    _set_slide_bg(slide, pal['bg'])
    sw = prs.slide_width
    
    card_w = Inches(7.0)
    card_h = Inches(4.8)
    card_l = Inches(0.35)
    card_t = CONTENT_TOP
    
    _add_card_with_accent(slide, card_l, card_t, card_w, card_h, pal, 'left')
    
    _add_textbox(slide, card_l + Inches(0.5), card_t + Inches(0.2),
                 Inches(6), Inches(0.5),
                 f'📖  {display_title}', Pt(24), pal['primary'], bold=True)
    
    bullets = [
        'Tìm hiểu nội dung bài học trong SGK.',
        'Thảo luận nhóm về các khái niệm chính.',
        'Ghi chép kiến thức trọng tâm.',
    ]
    
    _add_bullet_textbox(slide, card_l + Inches(0.5), card_t + Inches(0.9),
                        Inches(6), Inches(3.6), bullets, pal)
    
    # Right image
    if img_path:
        img_l = card_l + card_w + Inches(0.25)
        img_w = sw - img_l - Inches(0.35)
        _add_picture_fitted(slide, img_path, img_l, card_t, img_w, card_h)


def _add_summary_slide(slide, prs, display_title, lesson_data, pal, img_path=None):
    """Create summary slide with image."""
    _set_slide_bg(slide, pal['bg'])
    sw = prs.slide_width
    
    # Left card
    card_w = Inches(8.0)
    card_h = Inches(4.8)
    card_l = Inches(0.35)
    card_t = CONTENT_TOP
    
    _add_card_with_accent(slide, card_l, card_t, card_w, card_h, pal, 'left')
    
    _add_textbox(slide, card_l + Inches(0.5), card_t + Inches(0.2),
                 Inches(7), Inches(0.5),
                 f'📌  Ghi nhớ — {display_title}', Pt(24), pal['primary'], bold=True)
    
    summary = lesson_data.get('objectives', [])[:3]
    if not summary:
        summary = ['Nắm vững kiến thức cơ bản của bài học.',
                   'Hoàn thành bài tập luyện tập.',
                   'Chuẩn bị trước nội dung bài tiếp theo.']
    
    summary.append('📚 Bài tập về nhà: Xem lại nội dung bài học và chuẩn bị bài tiếp theo.')
    
    _add_bullet_textbox(slide, card_l + Inches(0.5), card_t + Inches(0.9),
                        Inches(7), Inches(3.6), summary, pal)
    
    # Right image
    if img_path:
        img_l = card_l + card_w + Inches(0.25)
        img_w = sw - img_l - Inches(0.35)
        _add_picture_fitted(slide, img_path, img_l, card_t, img_w, card_h)


def _add_thankyou_slide(slide, prs, grade_cfg, pal, img_path=None):
    """Create thank you slide with background image."""
    _set_slide_bg(slide, pal['primary'])
    sw = prs.slide_width
    sh = prs.slide_height
    
    # Background image (semi-transparent overlay)
    if img_path:
        # Place image on the right half as decoration
        img_w = Inches(5)
        img_h = Inches(4.5)
        img_l = sw - img_w - Inches(0.5)
        img_t = Inches(1.5)
        pic = _add_picture_fitted(slide, img_path, img_l, img_t, img_w, img_h)
        if pic:
            # Add semi-transparent overlay on image
            overlay = _add_shape(slide, MSO_SHAPE.RECTANGLE,
                                img_l, img_t, img_w, img_h,
                                fill_color=pal['primary'])
            fill_el = overlay._element.spPr.solidFill
            if fill_el is not None:
                srgb = fill_el.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                if srgb is not None:
                    etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha',
                                     attrib={'val': '50000'})
    
    # Center card
    card_w = Inches(8)
    card_h = Inches(3.5)
    card_l = (sw - card_w) // 2
    card_t = Inches(2.0)
    
    card = _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, card_l, card_t, card_w, card_h,
                      fill_color=RGBColor(0xFF, 0xFF, 0xFF))
    card.adjustments[0] = 0.04
    
    txBox = slide.shapes.add_textbox(card_l + Inches(0.5), card_t + Inches(0.5),
                                      card_w - Inches(1), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = 'CẢM ƠN CÁC EM! 🎉'
    r.font.size = Pt(32)
    r.font.color.rgb = pal['primary']
    r.font.bold = True
    r.font.name = 'Times New Roman'
    
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = 'Hẹn gặp lại ở tiết học sau!'
    r2.font.size = Pt(20)
    r2.font.color.rgb = pal['text_dark']
    r2.font.name = 'Times New Roman'
    p2.space_before = Pt(12)
    
    txBox2 = slide.shapes.add_textbox(card_l + Inches(0.5), card_t + Inches(2.2),
                                       card_w - Inches(1), Inches(1.0))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p3 = tf2.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = 'Trường TH & THCS UNIGO  •  GV: Đậu Đình Nguyên'
    r3.font.size = Pt(14)
    r3.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    r3.font.name = 'Times New Roman'


def _add_footer(slide, prs):
    """Add UNIGO footer bar to slide."""
    sw = prs.slide_width
    
    footer_h = Inches(0.55)
    footer_t = prs.slide_height - footer_h
    
    bar = _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, footer_t, sw, footer_h,
                     fill_color=RGBColor(0x00, 0x70, 0xC0))
    
    txBox = slide.shapes.add_textbox(Inches(0.3), footer_t + Inches(0.05),
                                      sw - Inches(0.6), footer_h - Inches(0.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = 'TRƯỜNG TIỂU HỌC VÀ THCS UNIGO'
    r.font.size = Pt(14)
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    r.font.bold = True
    r.font.name = 'Times New Roman'


# ══════════════════════════════════════════════════════════════════
#                         MAIN
# ══════════════════════════════════════════════════════════════════

def main():
    print('=' * 60)
    print('  GENERATE SLIDES BÀI 1 TIN HỌC — TẤT CẢ CÁC LỚP')
    print('  v2: Với hình ảnh SGK / AI-generated')
    print('=' * 60)
    
    results = []
    
    for grade_cfg in GRADES:
        key = grade_cfg['key']
        folder = grade_cfg['folder']
        label = grade_cfg['label']
        
        print(f'\n{"─"*50}')
        print(f'  {label}: Đang xử lý...')
        print(f'{"─"*50}')
        
        try:
            # 1. Read KHBD
            lesson_data = read_khbd(grade_cfg)
            print(f'  ✓ KHBD: {lesson_data["title"]}')
            print(f'    Topic: {lesson_data["topic"]}')
            print(f'    Objectives: {len(lesson_data["objectives"])}')
            print(f'    Activities: {len(lesson_data["activities"])}')
            
            # 2. Find images
            images = find_images(grade_cfg)
            print(f'  ✓ Images: {len(images["full_pages"])} full pages, '
                  f'{len(images["individual"])} individual, '
                  f'{len(images["generated"])} generated')
            
            # 3. Create slides
            prs, slides_created = create_presentation(grade_cfg, lesson_data, images)
            print(f'  ✓ Slides: {len(slides_created)} slides')
            
            # Count pictures in slides
            pic_count = sum(1 for s in prs.slides for sh in s.shapes if sh.shape_type == 13)
            print(f'  ✓ Pictures embedded: {pic_count}')
            
            for s in slides_created:
                print(f'    • {s}')
            
            # 4. Save
            out_dir = os.path.join(KHBD_ROOT, folder, 'Bài_01')
            os.makedirs(out_dir, exist_ok=True)
            
            title_clean = lesson_data['title']
            title_clean = re.sub(r'^BÀI\s*\d+\.\s*', '', title_clean, flags=re.IGNORECASE)
            title_clean = re.sub(r'[^\w\s]', '', title_clean)
            title_clean = title_clean.strip().replace(' ', '_')
            if key == 'TTH':
                fname = f'Slide_Tin_hoc_Tien_TH_Bai01_{title_clean}.pptx'
            else:
                fname = f'Slide_Tin_hoc_Lop_{key}_Bai01_{title_clean}.pptx'
            
            out_path = os.path.join(out_dir, fname)
            prs.save(out_path)
            print(f'  ✓ Saved: {out_path}')
            
            results.append((label, '✅', fname, len(slides_created), pic_count))
            
        except Exception as e:
            print(f'  ✗ ERROR: {e}')
            import traceback
            traceback.print_exc()
            results.append((label, '❌', str(e), 0, 0))
    
    # Summary
    print(f'\n{"="*60}')
    print('  KẾT QUẢ TỔNG HỢP')
    print(f'{"="*60}')
    for label, status, fname, count, pics in results:
        print(f'  {status} {label}: {fname} ({count} slides, {pics} images)')
    
    total_ok = sum(1 for r in results if r[1] == '✅')
    total_pics = sum(r[4] for r in results)
    print(f'\nHOÀN THÀNH! {total_ok}/{len(results)} file thành công. Tổng {total_pics} hình ảnh.')


if __name__ == '__main__':
    main()
