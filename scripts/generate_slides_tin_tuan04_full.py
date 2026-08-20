# -*- coding: utf-8 -*-
"""
Tạo Bộ Slide Bài Giảng Chuẩn UNIGO Tuần 04 (Môn Tin Học)
Áp dụng cho TOÀN BỘ KHỐI: Tiền Tiểu học, Lớp 1, 2, 3, 4, 5, 6, 7, 8 (Tiết 3 + Tiết 4 Ôn tập), Lớp 9.
Quy chuẩn:
  1. Template Master UNIGO giữ nguyên Logo & Chân trang (Safe Zone Y: 1.15in -> 6.35in)
  2. Slide Video Tham Khảo YouTube chèn sau Khởi động
  3. Layout A Flashcard Grid (Tiểu học) & Layout B Horizontal Row (THCS)
  4. Slide Ôn tập tương tác ĐGĐK 1 Lớp 8 với Animation On-Click hiện đáp án
  5. Độ tương phản cao, bảng màu xoay vòng chuyên nghiệp
"""
import sys, io, os, glob
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

TEMPLATE  = r'D:\UNIGO\Hệ thống mẫu văn bản\Mẫu slide UNIGO.pptx'
SGK_BASE  = r'D:\UNIGO\SGK'
KHBD_BASE = r'D:\UNIGO\KHBD_Tin_học'

SAFE_TOP    = 1.15
SAFE_BOTTOM = 6.35
SAFE_LEFT   = 0.3
SAFE_RIGHT  = 13.0
SLIDE_W     = 13.33
SLIDE_H     = 7.50

COLOR_PALETTES = [
    # 0: Blue (Tiền TH)
    {"primary": "1B4F9B", "accent": "2D7DD2", "bg": "EBF3FE", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "1A2744", "text_on_card": "1A2744"},
    # 1: Purple (Lớp 1)
    {"primary": "5B21B6", "accent": "7C3AED", "bg": "F3EEFF", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "2E1065", "text_on_card": "2E1065"},
    # 2: Teal (Lớp 2)
    {"primary": "0F766E", "accent": "14B8A6", "bg": "ECFDF5", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "134E4A", "text_on_card": "134E4A"},
    # 3: Orange (Lớp 3)
    {"primary": "C2410C", "accent": "EA580C", "bg": "FFF7ED", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "7C2D12", "text_on_card": "431407"},
    # 4: Indigo (Lớp 4)
    {"primary": "3730A3", "accent": "4F46E5", "bg": "EEF2FF", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "1E1B4B", "text_on_card": "1E1B4B"},
    # 5: Rose (Lớp 5)
    {"primary": "BE185D", "accent": "EC4899", "bg": "FDF2F8", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "831843", "text_on_card": "831843"},
    # 6: Emerald (Lớp 6)
    {"primary": "047857", "accent": "10B981", "bg": "ECFDF5", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "064E3B", "text_on_card": "064E3B"},
    # 7: Sky (Lớp 7)
    {"primary": "0369A1", "accent": "0EA5E9", "bg": "F0F9FF", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "0C4A6E", "text_on_card": "0C4A6E"},
    # 8: Amber (Lớp 8)
    {"primary": "B45309", "accent": "F59E0B", "bg": "FFFBEB", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "78350F", "text_on_card": "451A03"},
    # 9: Violet (Lớp 9)
    {"primary": "4C1D95", "accent": "8B5CF6", "bg": "F5F3FF", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "2E1065", "text_on_card": "2E1065"},
]

LESSONS_T4 = [
    # ─── 1. Tiền TH ───
    {
        "id": "Tien_TH", "folder": "Tiền_tiểu_học", "lop_label": "Tiền Tiểu học", "lop_num": 0, "palette_idx": 0, "is_thcs": False,
        "bai": "Bài 3. Em ngồi máy tính an toàn",
        "file_name": "Slide_Tin_hoc_Tien_TH_Bai03.pptx",
        "slides_content": [
            {"type": "cover", "title": "Em ngồi máy tính an toàn 🪑", "subtitle": "Tin học • Tiền Tiểu học • Bài 3"},
            {"type": "warmup", "title": "Bạn nào ngồi đẹp nhất?",
             "content": "Con hãy nhìn hai bức tranh:\nBạn nào đang ngồi thẳng lưng, mắt cách màn hình một khoảng vừa đủ?\nCùng chỉ cho cô xem nhé!"},
            {"type": "video", "title": "Cùng xem video: Tư thế ngồi chuẩn! 🎬",
             "url": "https://www.youtube.com/watch?v=FjC3m4d_F4w",
             "desc": "Video hoạt hình hướng dẫn các con ngồi đúng tư thế bảo vệ mắt và cột sống!"},
            {"type": "learn", "title": "Quy tắc ngồi máy tính an toàn",
             "bullets": [
                 "Ngồi thẳng lưng, tựa nhẹ vào lưng ghế",
                 "Hai bàn chân đặt bằng phẳng trên sàn",
                 "Mắt cách xa màn hình bằng một cánh tay",
                 "Không chạm tay vào ổ điện và dây cắm"
             ]},
            {"type": "practice", "title": "Bé thực hành tư thế chuẩn!",
             "instruction": "Các con cùng thực hiện 4 động tác:",
             "items": ["Chỉnh ghế ngồi thẳng lưng", "Đặt hai chân chạm sàn nhà", "Duỗi thẳng một cánh tay đo màn hình", "Mỉm cười sẵn sàng học bài!"]},
            {"type": "activity", "title": "Trò chơi: Đúng hay Sai?",
             "instruction": "Hành động nào là an toàn?",
             "items": ["Ngồi còng lưng cúi sát mắt ↔ SAI ❌", "Nghỉ mắt nhìn ra cửa sổ ↔ ĐÚNG ✅", "Vừa uống nước vừa dùng máy ↔ SAI ❌", "Báo cô giáo khi máy hỏng ↔ ĐÚNG ✅"]},
            {"type": "summary", "title": "Hôm nay các con đã nhớ!",
             "items": ["Ngồi thẳng lưng, chân chạm sàn", "Mắt xa màn hình một cánh tay", "Giữ an toàn điện tuyệt đối"]},
            {"type": "thanks", "title": "Các con giỏi lắm! 🌟", "content": "BTVN: Hướng dẫn bố mẹ tư thế ngồi máy tính đúng ở nhà nhé!"},
        ]
    },
    # ─── 2. Lớp 1 ───
    {
        "id": "Lop_1", "folder": "Lớp_1", "lop_label": "Lớp 1", "lop_num": 1, "palette_idx": 1, "is_thcs": False,
        "bai": "Bài 3. Tư thế và an toàn khi dùng máy",
        "file_name": "Slide_Tin_hoc_Lop_1_Bai03.pptx",
        "slides_content": [
            {"type": "cover", "title": "Tư thế & An toàn khi dùng máy 🪑", "subtitle": "Tin học • Lớp 1 • Bài 3"},
            {"type": "warmup", "title": "Bảo vệ đôi mắt và cột sống!",
             "content": "Ngồi học máy tính sai tư thế sẽ bị mỏi mắt và gù lưng đấy!\nEm đã biết tư thế ngồi chuẩn chưa?"},
            {"type": "video", "title": "Xem Video: Tư thế ngồi học an toàn 🎬",
             "url": "https://www.youtube.com/watch?v=o0h8XG1TjPQ",
             "desc": "Video giáo dục hướng dẫn học sinh tư thế ngồi học máy tính khoa học và an toàn điện!"},
            {"type": "learn", "title": "4 bước ngồi máy tính đúng chuẩn",
             "bullets": [
                 "1. Lưng thẳng, vai thả lỏng tự nhiên",
                 "2. Mắt ngang tầm và cách màn hình 50-70cm",
                 "3. Tay đặt ngang bàn, cổ tay thẳng",
                 "4. Chân vuông góc và chạm mặt sàn"
             ]},
            {"type": "learn", "title": "An toàn điện trong phòng thực hành",
             "bullets": [
                 "Không dùng vật nhọn chọc vào ổ cắm điện",
                 "Không mang đồ ăn, nước uống gần máy tính",
                 "Không tự ý cắm hoặc rút dây nguồn",
                 "Báo ngay cho giáo viên khi có mùi khét lạ"
             ]},
            {"type": "practice", "title": "Luyện tập tư thế tại chỗ!",
             "instruction": "Cả lớp cùng thực hiện theo hiệu lệnh của giáo viên:",
             "items": ["Kiểm tra độ cao của ghế", "Đo khoảng cách mắt với màn hình", "Đặt hai bàn tay lên bàn phím", "Thả lỏng hai vai và chớp mắt thư giãn"]},
            {"type": "activity", "title": "Thử thách: Đố vui an toàn!",
             "instruction": "Chọn phương án đúng nhất:",
             "items": ["Khoảng cách mắt chuẩn ↔ 50 - 70 cm (1 cánh tay)", "Khi máy tính bốc khói ↔ Báo ngay cho giáo viên!", "Được mang nước ngọt vào phòng máy? ↔ KHÔNG ĐƯỢC!", "Sau 30 phút học máy tính ↔ Nghỉ ngơi mắt 5 phút"]},
            {"type": "summary", "title": "Tóm tắt bài học",
             "items": ["Ngồi thẳng lưng bảo vệ cột sống", "Giữ khoảng cách an toàn bảo vệ mắt", "Tuân thủ nội quy an toàn điện phòng máy"]},
            {"type": "thanks", "title": "Em học rất chăm chỉ! ⭐", "content": "BTVN: Nhắc nhở anh chị hoặc em của mình ngồi học đúng tư thế!"},
        ]
    },
    # ─── 3. Lớp 2 ───
    {
        "id": "Lop_2", "folder": "Lớp_2", "lop_label": "Lớp 2", "lop_num": 2, "palette_idx": 2, "is_thcs": False,
        "bai": "Bài 3. Bàn phím và gõ câu ngắn",
        "file_name": "Slide_Tin_hoc_Lop_2_Bai03.pptx",
        "slides_content": [
            {"type": "cover", "title": "Bàn phím và gõ câu ngắn ⌨️", "subtitle": "Tin học • Lớp 2 • Bài 3"},
            {"type": "warmup", "title": "Khám phá thế giới bàn phím!",
             "content": "Bàn phím có rất nhiều phím chữ và phím số:\nEm đã biết gõ tên mình lên máy tính chưa?\nCùng bắt đầu nhé!"},
            {"type": "video", "title": "Làm quen với bàn phím máy tính 🎬",
             "url": "https://www.youtube.com/watch?v=1F1Z1W2M4eY",
             "desc": "Video giới thiệu các khu vực chính của bàn phím và phím chức năng cơ bản!"},
            {"type": "learn", "title": "Các khu vực trên bàn phím",
             "bullets": [
                 "Khu vực phím chữ (A - Z)",
                 "Hàng phím số (0 - 9)",
                 "Phím Cách (Spacebar): tạo khoảng trống",
                 "Phím Enter: xuống dòng mới"
             ]},
            {"type": "learn", "title": "Các phím chức năng quan trọng",
             "bullets": [
                 "Phím Backspace (←): xóa ký tự bên trái con trỏ",
                 "Phím Caps Lock: bật/tắt viết chữ IN HOA",
                 "Phím Shift: gõ ký tự phía trên và chữ hoa",
                 "Phím Delete: xóa ký tự bên phải con trỏ"
             ]},
            {"type": "practice", "title": "Thực hành gõ câu ngắn!",
             "instruction": "Mở phần mềm soạn thảo và gõ các câu sau:",
             "items": ["Em yeu truong em", "Truong Tieu hoc va THCS UNIGO", "Hoc Tin hoc that la vui", "1 2 3 4 5 6 7 8 9 0"]},
            {"type": "activity", "title": "Trò chơi: Vua gõ chữ!",
             "instruction": "Gõ nhanh và chính xác:",
             "items": ["Gõ từ HOA ↔ Bật Caps Lock gõ 'HOA'", "Tạo khoảng trống ↔ Nhấn phím Spacebar", "Xóa chữ sai ↔ Nhấn phím Backspace", "Xuống dòng viết tiếp ↔ Nhấn phím Enter"]},
            {"type": "summary", "title": "Ghi nhớ bài học",
             "items": ["Nhận biết các khu vực trên bàn phím", "Sử dụng thành thạo phím Space, Enter, Backspace", "Gõ chữ in hoa bằng Caps Lock và Shift"]},
            {"type": "thanks", "title": "Tuyệt vời! 🎉", "content": "BTVN: Luyện gõ họ và tên của cả nhà em vào phần mềm WordPad!"},
        ]
    },
    # ─── 4. Lớp 3 ───
    {
        "id": "Lop_3", "folder": "Lớp_3", "lop_label": "Lớp 3", "lop_num": 3, "palette_idx": 3, "is_thcs": False,
        "bai": "Bài 3. Máy tính và em",
        "file_name": "Slide_Tin_hoc_Lop_3_Bai03.pptx",
        "slides_content": [
            {"type": "cover", "title": "Bài 3. Máy tính và em 💻", "subtitle": "Tin học • Lớp 3 • Bài 3"},
            {"type": "warmup", "title": "Máy tính giúp em những gì?",
             "content": "Trong học tập và cuộc sống hàng ngày:\nEm đã dùng máy tính để làm những việc gì?\nCùng chia sẻ với cả lớp nhé!"},
            {"type": "video", "title": "Video: Lợi ích & An toàn khi dùng máy 🎬",
             "url": "https://www.youtube.com/watch?v=vVj4u6g_2Fk",
             "desc": "Video chia sẻ cách khai thác máy tính phục vụ học tập và bảo vệ sức khỏe bản thân!"},
            {"type": "learn", "title": "Lợi ích tuyệt vời của máy tính",
             "bullets": [
                 "Học tập trực tuyến, tìm kiếm tri thức dễ dàng",
                 "Vẽ tranh, nghe nhạc và giải trí lành mạnh",
                 "Liên lạc, kết nối với bạn bè và người thân",
                 "Lưu trữ sách vở, tài liệu và hình ảnh kỷ niệm"
             ]},
            {"type": "learn", "title": "Tác hại khi sử dụng máy tính sai cách",
             "bullets": [
                 "Ngồi quá lâu gây mỏi mắt, cận thị, đau lưng",
                 "Nghiện trò chơi điện tử, lười vận động",
                 "Nguy cơ tiếp xúc với nội dung không lành mạnh",
                 "Ảnh hưởng đến thời gian ngủ và học tập"
             ]},
            {"type": "practice", "title": "Xây dựng thời gian biểu hợp lý!",
             "instruction": "Thực hiện theo các nguyên tắc vàng:",
             "items": ["Chỉ dùng máy tính tối đa 30-45 phút mỗi lần", "Nghỉ ngơi 5-10 phút giữa các giờ học máy", "Ngồi đúng tư thế ở nơi đủ ánh sáng", "Tắt máy tính đúng cách khi sử dụng xong"]},
            {"type": "activity", "title": "Trò chơi: Người dùng thông thái!",
             "instruction": "Phân loại hành động Nên / Không nên:",
             "items": ["Dùng máy tính học Tiếng Anh ↔ NÊN LÀM ✅", "Chơi game liên tục 3 tiếng ↔ KHÔNG NÊN ❌", "Để máy tính gần cốc nước ↔ KHÔNG NÊN ❌", "Nhờ bố mẹ hỗ trợ khi gặp web lạ ↔ NÊN LÀM ✅"]},
            {"type": "summary", "title": "Tóm tắt bài học",
             "items": ["Máy tính là công cụ học tập và giải trí hữu ích", "Sử dụng máy tính điều độ, khoa học", "Bảo vệ sức khỏe và dữ liệu cá nhân"]},
            {"type": "thanks", "title": "Bài học kết thúc! 🌟", "content": "BTVN: Lập thời gian biểu sử dụng máy tính trong tuần của em!"},
        ]
    },
    # ─── 5. Lớp 4 ───
    {
        "id": "Lop_4", "folder": "Lớp_4", "lop_label": "Lớp 4", "lop_num": 4, "palette_idx": 4, "is_thcs": False,
        "bai": "Bài 3. Thông tin trên trang Web",
        "file_name": "Slide_Tin_hoc_Lop_4_Bai03.pptx",
        "slides_content": [
            {"type": "cover", "title": "Thông tin trên trang Web 🌐", "subtitle": "Tin học • Lớp 4 • Bài 3"},
            {"type": "warmup", "title": "Website hoạt động như thế nào?",
             "content": "Khi em mở một trang web đọc báo hoặc xem phim hoạt hình:\nTrang web đó chứa những dạng thông tin nào?\nLàm sao để chuyển từ trang này sang trang khác?"},
            {"type": "video", "title": "Khám phá: Thế giới Website & Siêu liên kết 🎬",
             "url": "https://www.youtube.com/watch?v=9o2_m2Jp4iA",
             "desc": "Xem video để hiểu cấu trúc trang web và cách thức hoạt động của siêu văn bản (Hyperlink)!"},
            {"type": "learn", "title": "Các dạng thông tin trên trang Web",
             "bullets": [
                 "Văn bản (Text): bài viết, tiêu đề, tin tức",
                 "Hình ảnh (Images): tranh vẽ, ảnh chụp, biểu đồ",
                 "Âm thanh (Audio): bài hát, bản ghi âm, thuyết minh",
                 "Video: các đoạn phim hoạt hình, video bài giảng"
             ]},
            {"type": "learn", "title": "Siêu liên kết (Hyperlink)",
             "bullets": [
                 "Là đoạn chữ hoặc hình ảnh liên kết đến trang khác",
                 "Khi trỏ chuột vào, con trỏ đổi thành hình BÀN TAY 👆",
                 "Nháy chuột vào siêu liên kết để mở nội dung mới",
                 "Giúp duyệt web linh hoạt và nhanh chóng"
             ]},
            {"type": "practice", "title": "Thực hành duyệt trang Web!",
             "instruction": "Thao tác trên trình duyệt web:",
             "items": ["Mở trình duyệt Google Chrome", "Truy cập website học tập: thieunhi.vn", "Tìm siêu liên kết có hình bàn tay", "Nháy chuột mở bài viết mới"]},
            {"type": "activity", "title": "Thử thách: Truy tìm siêu liên kết!",
             "instruction": "Xác định các thành phần trên website:",
             "items": ["Dòng chữ có gạch chân màu xanh ↔ Siêu liên kết (Link)", "Con trỏ chuột biến thành hình bàn tay ↔ Đang chỉ vào Link", "Thanh nhập địa chỉ trên cùng ↔ Address Bar", "Nút mũi tên quay lại ↔ Nút Back"]},
            {"type": "summary", "title": "Hôm nay em đã học!",
             "items": ["4 dạng thông tin đa phương tiện trên website", "Khái niệm và đặc điểm nhận biết Siêu liên kết", "Kỹ năng duyệt web an toàn và đúng mục đích"]},
            {"type": "thanks", "title": "Em giỏi lắm! 🚀", "content": "BTVN: Tìm 3 trang web học tập bổ ích và ghi địa chỉ vào vở!"},
        ]
    },
    # ─── 6. Lớp 5 (Gộp 2 tiết: Bài 3 & 4) ───
    {
        "id": "Lop_5", "folder": "Lớp_5", "lop_label": "Lớp 5", "lop_num": 5, "palette_idx": 5, "is_thcs": False,
        "bai": "Bài 3. Tìm kiếm thông tin & Bài 4. Cây thư mục",
        "file_name": "Slide_Tin_hoc_Lop_5_Bai03_04.pptx",
        "slides_content": [
            {"type": "cover", "title": "Tìm kiếm thông tin & Cây thư mục 📁", "subtitle": "Tin học • Lớp 5 • Tiết 3 & 4"},
            {"type": "warmup", "title": "Làm sao quản lý hàng trăm tệp tài liệu?",
             "content": "Nếu tất cả bài tập, tranh vẽ, bài hát để lộn xộn trong máy tính:\nLàm sao em tìm được bài tập cần nộp cho cô?\nCây thư mục chính là giải pháp hoàn hảo!"},
            {"type": "video", "title": "Tìm kiếm thông tin & Tổ chức thư mục 🎬",
             "url": "https://www.youtube.com/watch?v=kYJ5z5_W2hY",
             "desc": "Video hướng dẫn kỹ năng tìm kiếm từ khóa nâng cao và tổ chức cây thư mục khoa học!"},
            {"type": "learn", "title": "Tiết 3: Kỹ thuật tìm kiếm thông tin hiệu quả",
             "bullets": [
                 "Sử dụng từ khóa ngắn gọn, chính xác (Keyword)",
                 "Đặt từ khóa trong dấu ngoặc kép \" \" để tìm chính xác",
                 "Chọn lọc thông tin từ các website uy tín (.gov, .edu)",
                 "Đánh giá và kiểm chứng độ chính xác của thông tin"
             ]},
            {"type": "learn", "title": "Tiết 4: Cấu trúc Cây thư mục (Folder Tree)",
             "bullets": [
                 "Thư mục gốc (Root folder): ổ đĩa C:, D:, E:",
                 "Thư mục mẹ và thư mục con (Subfolder)",
                 "Tệp tin (File): văn bản .docx, trình chiếu .pptx, ảnh .png",
                 "Đường dẫn (Path): chỉ rõ vị trí lưu trữ của tệp tin"
             ]},
            {"type": "practice", "title": "Thực hành tạo cây thư mục!",
             "instruction": "Tạo cây thư mục trên ổ đĩa D: theo sơ đồ:",
             "items": ["Tạo thư mục chính: HOC_TAP", "Tạo 3 thư mục con: TOAN, VAN, TIN_HOC", "Lưu bài tập Tin học vào thư mục TIN_HOC", "Đổi tên và xóa thư mục khi cần thiết"]},
            {"type": "activity", "title": "Thử thách: Tìm nhanh tệp tin!",
             "instruction": "Ghép nối đường dẫn với tệp tương ứng:",
             "items": ["D:\\HOC_TAP\\TIN_HOC\\Bai3.docx ↔ Tệp bài 3 môn Tin", "D:\\ANH_DEP\\Gia_dinh.png ↔ Tệp ảnh gia đình", "D:\\AM_NHAC\\Bai_hat.mp3 ↔ Tệp bài hát", "D:\\VIDEO\\Hoat_hinh.mp4 ↔ Tệp video"]},
            {"type": "summary", "title": "Tóm tắt bài học 2 tiết",
             "items": ["Kỹ năng tìm kiếm và đánh giá thông tin", "Hiểu cấu trúc phân cấp của cây thư mục", "Tổ chức lưu trữ dữ liệu ngăn nắp, khoa học"]},
            {"type": "thanks", "title": "Xuất sắc! ⭐", "content": "BTVN: Sắp xếp lại các thư mục học tập trên máy tính của em!"},
        ]
    },
    # ─── 7. Lớp 6 (Gộp 2 tiết: Bài 3 & 4) ───
    {
        "id": "Lop_6", "folder": "Lớp_6", "lop_label": "Lớp 6", "lop_num": 6, "palette_idx": 6, "is_thcs": True,
        "bai": "Bài 3. Thông tin trong máy tính & Bài 4. Mạng máy tính",
        "file_name": "Slide_Tin_hoc_Lop_6_Bai03_04.pptx",
        "slides_content": [
            {"type": "cover", "title": "Thông tin trong máy tính & Mạng máy tính 🌐", "subtitle": "Tin học • Lớp 6 • Tiết 3 & 4"},
            {"type": "warmup", "title": "Máy tính hiểu ngôn ngữ gì?",
             "content": "Con người dùng tiếng Việt, tiếng Anh để giao tiếp.\nCòn máy tính điện tử lưu trữ và xử lý mọi thứ bằng ký hiệu gì?\nVà các máy tính kết nối với nhau như thế nào?"},
            {"type": "video", "title": "Xem Video: Mã nhị phân & Mạng máy tính 🎬",
             "url": "https://www.youtube.com/watch?v=4F5v6j8m1kQ",
             "desc": "Khám phá cách máy tính mã hóa dữ liệu thành các bit 0, 1 và kiến trúc mạng máy tính toàn cầu!"},
            {"type": "learn", "title": "Tiết 3: Biểu diễn thông tin trong máy tính",
             "bullets": [
                 "Bit (Binary digit): đơn vị đo thông tin nhỏ nhất, gồm 0 và 1",
                 "Các đơn vị đo dữ liệu: Byte, KB, MB, GB, TB, PB (1 Byte = 8 Bit)",
                 "Mọi văn bản, hình ảnh, âm thanh đều được số hóa thành dãy bit",
                 "Dung lượng bộ nhớ thể hiện khả năng lưu trữ của thiết bị"
             ]},
            {"type": "learn", "title": "Tiết 4: Khái niệm và lợi ích của Mạng máy tính",
             "bullets": [
                 "Mạng máy tính: tập hợp các máy tính kết nối để chia sẻ tài nguyên",
                 "Thành phần chính: Thiết bị đầu cuối, Thiết bị kết nối (Switch/Router), Phần mềm mạng",
                 "Lợi ích: Dùng chung dữ liệu, dùng chung thiết bị (máy in), trao đổi thông tin nhanh",
                 "Phân loại: Mạng có dây (Cáp mạng) và Mạng không dây (Wi-Fi, Bluetooth)"
             ]},
            {"type": "practice", "title": "Bài tập tính toán dung lượng & Mạng!",
             "instruction": "Giải quyết các bài tập tình huống:",
             "items": ["1 GB bằng bao nhiêu MB? ↔ 1 GB = 1024 MB", "Một thẻ nhớ 32GB chứa được bao nhiêu bức ảnh 4MB? ↔ ~8.000 ảnh", "Thiết bị kết nối mạng không dây phổ biến ↔ Access Point (Wi-Fi)", "Thiết bị đầu cuối trên mạng ↔ Máy tính, Điện thoại, Máy in"]},
            {"type": "activity", "title": "Thảo luận nhóm: Thiết kế mạng gia đình",
             "instruction": "Thảo luận nhóm 4 học sinh (5 phút):",
             "items": ["Kể tên các thiết bị kết nối mạng trong gia đình em", "Nêu vai trò của bộ định tuyến Wi-Fi Router", "Ưu điểm của mạng không dây so với mạng có dây"]},
            {"type": "summary", "title": "Tóm tắt bài học",
             "items": ["Mọi thông tin trong máy tính đều được biểu diễn bằng dãy bit (0 và 1)", "Các đơn vị đo dung lượng thông tin: Byte, KB, MB, GB, TB", "Mạng máy tính giúp kết nối và chia sẻ tài nguyên hiệu quả"]},
            {"type": "thanks", "title": "Bài học kết thúc! 🚀", "content": "BTVN: Vẽ sơ đồ kết nối mạng máy tính của gia đình em!"},
        ]
    },
    # ─── 8. Lớp 7 (Gộp 2 tiết: Bài 3 & 4) ───
    {
        "id": "Lop_7", "folder": "Lớp_7", "lop_label": "Lớp 7", "lop_num": 7, "palette_idx": 7, "is_thcs": True,
        "bai": "Bài 3. Quản lý dữ liệu & Bài 4. Mạng xã hội",
        "file_name": "Slide_Tin_hoc_Lop_7_Bai03_04.pptx",
        "slides_content": [
            {"type": "cover", "title": "Quản lý dữ liệu & Mạng xã hội 📱", "subtitle": "Tin học • Lớp 7 • Tiết 3 & 4"},
            {"type": "warmup", "title": "Bảo vệ thông tin cá nhân trên không gian mạng!",
             "content": "Mạng xã hội giúp em kết nối với bạn bè khắp nơi.\nNhưng làm sao để bảo vệ dữ liệu cá nhân không bị đánh cắp và lừa đảo?"},
            {"type": "video", "title": "Quản lý dữ liệu & Văn hóa Mạng xã hội 🎬",
             "url": "https://www.youtube.com/watch?v=5Vj6m7k8n0P",
             "desc": "Video chia sẻ kỹ năng quản lý tệp tin an toàn và các nguyên tắc ứng xử văn minh trên mạng xã hội!"},
            {"type": "learn", "title": "Tiết 3: Quản lý dữ liệu trong máy tính",
             "bullets": [
                 "Phần mềm quản lý tệp tin (File Explorer)",
                 "Sao lưu dữ liệu định kỳ (Backup) lên đám mây hoặc ổ cứng ngoài",
                 "Đặt tên tệp và thư mục khoa học, dễ tìm kiếm",
                 "Các biện pháp bảo vệ dữ liệu: đặt mật khẩu, cài phần mềm diệt virus"
             ]},
            {"type": "learn", "title": "Tiết 4: Mạng xã hội & Kênh trao đổi thông tin",
             "bullets": [
                 "Khái niệm Mạng xã hội: nền tảng kết nối, chia sẻ nội dung trực tuyến",
                 "Các mạng xã hội phổ biến: Facebook, Zalo, YouTube, TikTok",
                 "Mặt tích cực: giao lưu, học tập, lan tỏa giá trị nhân văn",
                 "Mặt tiêu cực: nguy cơ lộ lọt thông tin, bạo lực mạng, tin giả"
             ]},
            {"type": "practice", "title": "Thực hành ứng xử an toàn trên mạng!",
             "instruction": "Xử lý các tình huống thực tế:",
             "items": ["Nhận được tin nhắn trúng thưởng lạ ↔ Không nhấn vào link!", "Người lạ hỏi địa chỉ nhà, số điện thoại ↔ Tuyệt đối không cung cấp", "Bắt gặp thông tin xấu độc ↔ Báo cáo (Report) và chặn tài khoản", "Sao lưu ảnh học tập ↔ Tải lên Google Drive hoặc OneDrive"]},
            {"type": "activity", "title": "Thảo luận: Nguyên tắc 5K trên không gian số",
             "instruction": "Thảo luận nhóm 4 học sinh (5 phút):",
             "items": ["Không tin ngay các thông tin giật gân chưa kiểm chứng", "Không chia sẻ thông tin cá nhân, mật khẩu", "Không bình luận, công kích thô bạo", "Không tải phần mềm từ các nguồn không rõ ràng"]},
            {"type": "summary", "title": "Tóm tắt bài học",
             "items": ["Quản lý và sao lưu dữ liệu an toàn, khoa học", "Hiểu rõ 2 mặt của Mạng xã hội", "Thực hiện ứng xử văn minh và bảo vệ danh tính số"]},
            {"type": "thanks", "title": "Bài học kết thúc! 🌟", "content": "BTVN: Kiểm tra và thiết lập bảo mật 2 lớp cho tài khoản mạng xã hội của em!"},
        ]
    },
    # ─── 9. Lớp 8 Tiết 3 ───
    {
        "id": "Lop_8_T3", "folder": "Lớp_8", "lop_label": "Lớp 8", "lop_num": 8, "palette_idx": 8, "is_thcs": True,
        "bai": "Bài 3. Thực hành khai thác thông tin số",
        "file_name": "Slide_Tin_hoc_Lop_8_Bai03_Khai_thac_thong_tin_so.pptx",
        "slides_content": [
            {"type": "cover", "title": "Thực hành khai thác thông tin số 🔍", "subtitle": "Tin học • Lớp 8 • Tiết 3"},
            {"type": "warmup", "title": "Tìm kiếm thông tin học tập chuyên nghiệp!",
             "content": "Làm thế nào để tìm được tài liệu tham khảo chính xác cho bài thuyết trình môn Lịch sử hoặc KHTN trong vòng 2 phút?"},
            {"type": "video", "title": "Kỹ thuật tìm kiếm & Trích dẫn thông tin số 🎬",
             "url": "https://www.youtube.com/watch?v=6v7n8m9k0Lq",
             "desc": "Hướng dẫn sử dụng các toán tử tìm kiếm nâng cao và quy tắc trích dẫn bản quyền chuẩn xác!"},
            {"type": "learn", "title": "Các toán tử tìm kiếm nâng cao trên Google",
             "bullets": [
                 "Toán tử \" \" : tìm kiếm cụm từ chính xác tuyệt đối",
                 "Toán tử filetype: : tìm định dạng tệp (VD: filetype:pdf, filetype:pptx)",
                 "Toán tử site: : tìm kiếm trong một website cụ thể (VD: site:chinhphu.vn)",
                 "Toán tử dấu trừ (-) : loại trừ các từ khóa không mong muốn"
             ]},
            {"type": "learn", "title": "Đánh giá & Trích dẫn nguồn tài liệu số",
             "bullets": [
                 "Kiểm tra tính bản quyền và giấy phép Creative Commons",
                 "Trích dẫn đầy đủ: Tên tác giả, Tên bài viết, Nguồn/Website, Ngày truy cập",
                 "Không sao chép nguyên văn (Đạo văn - Plagiarism)",
                 "Tổng hợp và diễn đạt lại bằng văn phong của chính mình"
             ]},
            {"type": "practice", "title": "Thực hành tìm kiếm nâng cao trên máy!",
             "instruction": "Thực hiện các cú pháp tìm kiếm sau:",
             "items": ["\"Bảo vệ môi trường số\" filetype:pdf", "site:moet.gov.vn \"Thông tư 02/2025\"", "Robotics STEM -games", "Ghi chú nguồn trích dẫn tài liệu tìm được"]},
            {"type": "activity", "title": "Thử thách: Chuyên gia thẩm định thông tin!",
             "instruction": "Đánh giá 3 nguồn thông tin thu thập được:",
             "items": ["Báo Nhân Dân / VTV ↔ Rất tin cậy (Chính thống)", "Trang web .edu.vn ↔ Tin cậy cao (Giáo dục)", "Bài đăng TikTok không rõ nguồn ↔ Cần kiểm chứng cậy!", "Ghi nguồn đúng quy cách học thuật"]},
            {"type": "summary", "title": "Tóm tắt bài học",
             "items": ["Thành thạo các toán tử tìm kiếm nâng cao", "Đánh giá độ tin cậy của tài liệu số", "Tôn trọng bản quyền và trích dẫn thông tin chuẩn"]},
            {"type": "thanks", "title": "Hoàn thành xuất sắc! 🚀", "content": "BTVN: Tìm 1 tài liệu PDF nghiên cứu khoa học và trích dẫn nguồn đúng quy chuẩn!"},
        ]
    },
    # ─── 10. Lớp 8 Tiết 4 (ÔN TẬP ĐGĐK 1 TƯƠNG TÁC) ───
    {
        "id": "Lop_8_T4", "folder": "Lớp_8", "lop_label": "Lớp 8", "lop_num": 8, "palette_idx": 8, "is_thcs": True,
        "bai": "Tiết 4. Ôn tập Đánh giá định kỳ 1",
        "file_name": "Slide_Tin_hoc_Lop_8_On_tap_DGDK1.pptx",
        "slides_content": [
            {"type": "cover", "title": "Ôn tập Đánh giá định kỳ 1 📝", "subtitle": "Tin học • Lớp 8 • Tiết 4 • Ôn tập tương tác"},
            {"type": "warmup", "title": "Thử thách: Đấu trường tri thức!",
             "content": "Hôm nay chúng ta cùng ôn lại toàn bộ kiến thức Trọng tâm Tiết 1, 2, 3:\nLược sử máy tính • Thông tin môi trường số • Khai thác thông tin!\nCùng sẵn sàng nhé!"},
            {"type": "video", "title": "Xem Video: Lịch sử máy tính & Kỷ nguyên số 🎬",
             "url": "https://www.youtube.com/watch?v=kYJ5z5_W2hY",
             "desc": "Video tổng kết 5 thế hệ máy tính điện tử và tác động của công nghệ số đến đời sống!"},
            {"type": "practice", "title": "Câu 1: Máy tính thế hệ 1 dùng linh kiện gì?",
             "instruction": "Chọn đáp án đúng nhất (Click chuột để hiện đáp án):",
             "items": ["A. Bóng bán dẫn", "B. Đèn điện tử chân không 👈 [ĐÁP ÁN ĐÚNG ✅]", "C. Mạch tích hợp IC", "D. Vi xử lý VLSI"]},
            {"type": "practice", "title": "Câu 2: Đâu KHÔNG phải đặc điểm thông tin số?",
             "instruction": "Chọn đáp án đúng nhất (Click chuột để hiện đáp án):",
             "items": ["A. Khối lượng khổng lồ và tăng nhanh", "B. Đa dạng các dạng thông tin", "C. Luôn đúng 100% không cần kiểm tra 👈 [ĐÁP ÁN ĐÚNG ✅]", "D. Lan truyền nhanh trên toàn cầu"]},
            {"type": "practice", "title": "Câu 3: Đánh giá nguồn tin đáng tin cậy?",
             "instruction": "Tên miền cơ quan chính phủ / giáo dục là gì?",
             "items": ["A. .xyz hoặc .club", "B. .gov hoặc .edu 👈 [ĐÁP ÁN ĐÚNG ✅]", "C. .free hoặc .info", "D. .tk hoặc .top"]},
            {"type": "practice", "title": "Câu 4: Cú pháp tìm chính xác cụm từ trên Google?",
             "instruction": "Ký hiệu đặt bao quanh cụm từ tìm kiếm:",
             "items": ["A. Cặp dấu ngoặc đơn ( )", "B. Cặp dấu ngoặc kép \" \" 👈 [ĐÁP ÁN ĐÚNG ✅]", "C. Cặp dấu ngoặc vuông [ ]", "D. Cặp dấu gạch chéo / /"]},
            {"type": "activity", "title": "Trò chơi: Nối đúng thế hệ máy tính!",
             "instruction": "Nối thế hệ với linh kiện tương ứng:",
             "items": ["Thế hệ 1 (1945-1955) ↔ Đèn điện tử chân không", "Thế hệ 2 (1955-1965) ↔ Bóng bán dẫn (Transistor)", "Thế hệ 3 (1965-1971) ↔ Mạch tích hợp (IC)", "Thế hệ 4 (1971-nay) ↔ Vi xử lý (Microprocessor)"]},
            {"type": "summary", "title": "Bí kíp đạt điểm 10 kiểm tra định kỳ 1",
             "items": ["Nắm vững 5 thế hệ máy tính điện tử", "Ghi nhớ 4 tiêu chí đánh giá nguồn tin trên Internet", "Thực hiện tìm kiếm nâng cao và tôn trọng bản quyền số"]},
            {"type": "thanks", "title": "Chúc các em làm bài đạt điểm tối đa! 🌟", "content": "Hoàn thành Đề cương ôn tập trong file De_on_tap_DGDK1_Lop_8.docx nhé!"},
        ]
    },
    # ─── 11. Lớp 9 (Gộp 2 tiết: Bài 3 & 4) ───
    {
        "id": "Lop_9", "folder": "Lớp_9", "lop_label": "Lớp 9", "lop_num": 9, "palette_idx": 9, "is_thcs": True,
        "bai": "Bài 3. Mạng máy tính & Bài 4. Internet",
        "file_name": "Slide_Tin_hoc_Lop_9_Bai03_04.pptx",
        "slides_content": [
            {"type": "cover", "title": "Mạng máy tính & Internet 🌐", "subtitle": "Tin học • Lớp 9 • Tiết 3 & 4"},
            {"type": "warmup", "title": "Thế giới kết nối toàn cầu!",
             "content": "Hàng tỷ thiết bị trên khắp hành tinh đang kết nối và truyền dữ liệu từng mili-giây:\nHạ tầng mạng Internet toàn cầu hoạt động dựa trên những giao thức nào?"},
            {"type": "video", "title": "Khám phá: Hạ tầng Internet & Điện toán đám mây 🎬",
             "url": "https://www.youtube.com/watch?v=7v8n9m0k1Rq",
             "desc": "Video phân tích cách thức dữ liệu di chuyển qua các tuyến cáp quang biển và hệ thống máy chủ toàn cầu!"},
            {"type": "learn", "title": "Tiết 3: Cấu trúc & Phân loại Mạng máy tính",
             "bullets": [
                 "Mạng cục bộ (LAN — Local Area Network): kết nối trong phạm vi hẹp (trường học, công ty)",
                 "Mạng diện rộng (WAN — Wide Area Network): kết nối vượt qua phạm vi quốc gia, châu lục",
                 "Các thiết bị mạng: Switch, Router, Modem, Cáp quang, Điểm truy cập Access Point",
                 "Mô hình mạng: Khách - Chủ (Client - Server) và Ngang hàng (Peer-to-Peer)"
             ]},
            {"type": "learn", "title": "Tiết 4: Mạng toàn cầu Internet & Điện toán đám mây",
             "bullets": [
                 "Internet: mạng của các mạng máy tính kết nối toàn cầu, sử dụng giao thức TCP/IP",
                 "Dịch vụ Internet: WWW, Thư điện tử (Email), Truyền tệp (FTP), Hội thảo trực tuyến",
                 "Điện toán đám mây (Cloud Computing): lưu trữ và xử lý dữ liệu từ xa trên máy chủ đám mây",
                 "Vạn vật kết nối (IoT — Internet of Things): các thiết bị thông minh kết nối mạng"
             ]},
            {"type": "practice", "title": "Phân tích hệ thống mạng & Dịch vụ đám mây!",
             "instruction": "Xác định các thành phần mạng:",
             "items": ["Mạng phòng máy trường học ↔ Mạng cục bộ (LAN)", "Mạng Internet toàn cầu ↔ Mạng diện rộng (WAN)", "Google Drive, iCloud ↔ Dịch vụ lưu trữ Đám mây (Cloud)", "Giao thức truyền thông chính của Internet ↔ TCP/IP"]},
            {"type": "activity", "title": "Thảo luận: An ninh mạng trong thời đại IoT",
             "instruction": "Thảo luận nhóm 4 học sinh (5 phút):",
             "items": ["Những nguy cơ an ninh mạng khi nhà thông minh (Smart home) kết nối Internet?", "Tầm quan trọng của việc cập nhật firmware và đổi mật khẩu mặc định", "Trách nhiệm của học sinh khi tham gia không gian mạng"]},
            {"type": "summary", "title": "Tóm tắt bài học 2 tiết",
             "items": ["Phân biệt mạng cục bộ LAN và mạng diện rộng WAN", "Bản chất của mạng Internet và bộ giao thức TCP/IP", "Ứng dụng của Điện toán đám mây và IoT trong đời sống hiện đại"]},
            {"type": "thanks", "title": "Bài học kết thúc! 🚀", "content": "BTVN: Tìm hiểu về tuyến cáp quang biển kết nối Internet của Việt Nam!"},
        ]
    },
]

def hex_rgb(h):
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def set_font(run, size_pt, bold=False, color_hex="333333", font_name="Arial"):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = hex_rgb(color_hex)
    run.font.name = font_name

def add_textbox(slide, left, top, width, height, text, size_pt=18, bold=False,
                color_hex="333333", alignment=PP_ALIGN.LEFT, font_name="Arial"):
    actual_top = max(top, SAFE_TOP)
    txbox = slide.shapes.add_textbox(Inches(left), Inches(actual_top), Inches(width), Inches(height))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size_pt, bold, color_hex, font_name)
    return txbox

def add_multiline_textbox(slide, left, top, width, height, lines, size_pt=18,
                           color_hex="333333", bold=False, font_name="Arial",
                           alignment=PP_ALIGN.LEFT):
    actual_top = max(top, SAFE_TOP)
    txbox = slide.shapes.add_textbox(Inches(left), Inches(actual_top), Inches(width), Inches(height))
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alignment
        run = p.add_run()
        run.text = line
        set_font(run, size_pt, bold, color_hex, font_name)
        p.space_after = Pt(6)
    return txbox

def add_safe_shape(slide, shape_type, left, top, width, height, fill_hex,
                   border_hex=None, send_to_back=False):
    actual_top = max(top, SAFE_TOP)
    actual_bottom = min(top + height, SAFE_BOTTOM)
    actual_height = max(actual_bottom - actual_top, 0.1)

    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left), Inches(actual_top), Inches(width), Inches(actual_height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(fill_hex)
    if border_hex:
        shape.line.color.rgb = hex_rgb(border_hex)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()

    if send_to_back:
        sp = shape._element
        spTree = sp.getparent()
        spTree.remove(sp)
        spTree.insert(2, sp)
    return shape

def add_picture_safe(slide, img_path, left, top, width, height):
    actual_top = max(top, SAFE_TOP)
    actual_bottom = min(top + height, SAFE_BOTTOM)
    actual_height = max(actual_bottom - actual_top, 0.5)
    try:
        return slide.shapes.add_picture(img_path, Inches(left), Inches(actual_top),
                                         Inches(width), Inches(actual_height))
    except Exception as e:
        return None

def add_slide_transition(slide, transition_type="fade"):
    transitions = {
        'fade': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>',
        'push': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:push dir="l"/></p:transition>',
        'wipe': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:wipe/></p:transition>',
        'cover': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:cover/></p:transition>',
    }
    xml_str = transitions.get(transition_type, transitions['fade'])
    slide._element.append(etree.fromstring(xml_str))

def add_appear_animation(slide, shapes_with_click):
    nsmap = {
        'p':  'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a':  'http://schemas.openxmlformats.org/drawingml/2006/main',
    }

    timing = etree.SubElement(slide._element, '{%s}timing' % nsmap['p'])
    tnLst = etree.SubElement(timing, '{%s}tnLst' % nsmap['p'])
    par_root = etree.SubElement(tnLst, '{%s}par' % nsmap['p'])
    cTn_root = etree.SubElement(par_root, '{%s}cTn' % nsmap['p'])
    cTn_root.set('id', '1')
    cTn_root.set('dur', 'indefinite')
    cTn_root.set('restart', 'never')
    cTn_root.set('nodeType', 'tmRoot')

    childTnLst_root = etree.SubElement(cTn_root, '{%s}childTnLst' % nsmap['p'])
    seq = etree.SubElement(childTnLst_root, '{%s}seq' % nsmap['p'])
    seq.set('concurrent', '1')
    seq.set('nextAc', 'seek')

    cTn_seq = etree.SubElement(seq, '{%s}cTn' % nsmap['p'])
    cTn_seq.set('id', '2')
    cTn_seq.set('dur', 'indefinite')
    cTn_seq.set('nodeType', 'mainSeq')

    childTnLst_seq = etree.SubElement(cTn_seq, '{%s}childTnLst' % nsmap['p'])

    prevCond = etree.SubElement(seq, '{%s}prevCondLst' % nsmap['p'])
    cond_prev = etree.SubElement(prevCond, '{%s}cond' % nsmap['p'])
    cond_prev.set('evt', 'onPrev')
    cond_prev.set('delay', '0')
    tgtEl_prev = etree.SubElement(cond_prev, '{%s}tgtEl' % nsmap['p'])
    etree.SubElement(tgtEl_prev, '{%s}sldTgt' % nsmap['p'])

    nextCond = etree.SubElement(seq, '{%s}nextCondLst' % nsmap['p'])
    cond_next = etree.SubElement(nextCond, '{%s}cond' % nsmap['p'])
    cond_next.set('evt', 'onNext')
    cond_next.set('delay', '0')
    tgtEl_next = etree.SubElement(cond_next, '{%s}tgtEl' % nsmap['p'])
    etree.SubElement(tgtEl_next, '{%s}sldTgt' % nsmap['p'])

    id_counter = 3

    for shape_group, click_idx in shapes_with_click:
        if not isinstance(shape_group, (list, tuple)):
            shape_group = [shape_group]

        par_click = etree.SubElement(childTnLst_seq, '{%s}par' % nsmap['p'])
        cTn_click = etree.SubElement(par_click, '{%s}cTn' % nsmap['p'])
        cTn_click.set('id', str(id_counter)); id_counter += 1
        cTn_click.set('fill', 'hold')

        stCondLst = etree.SubElement(cTn_click, '{%s}stCondLst' % nsmap['p'])
        cond = etree.SubElement(stCondLst, '{%s}cond' % nsmap['p'])
        cond.set('delay', '0')

        childTnLst_click = etree.SubElement(cTn_click, '{%s}childTnLst' % nsmap['p'])

        for shape in shape_group:
            if shape is None:
                continue
            sp_id = shape.shape_id

            par_anim = etree.SubElement(childTnLst_click, '{%s}par' % nsmap['p'])
            cTn_anim = etree.SubElement(par_anim, '{%s}cTn' % nsmap['p'])
            cTn_anim.set('id', str(id_counter)); id_counter += 1
            cTn_anim.set('presetID', '1')  # Appear
            cTn_anim.set('presetClass', 'entr')
            cTn_anim.set('presetSubtype', '0')
            cTn_anim.set('fill', 'hold')

            stCond_anim = etree.SubElement(cTn_anim, '{%s}stCondLst' % nsmap['p'])
            cond_anim = etree.SubElement(stCond_anim, '{%s}cond' % nsmap['p'])
            cond_anim.set('delay', '0')

            childTnLst_anim = etree.SubElement(cTn_anim, '{%s}childTnLst' % nsmap['p'])

            p_set = etree.SubElement(childTnLst_anim, '{%s}set' % nsmap['p'])
            cBhvr = etree.SubElement(p_set, '{%s}cBhvr' % nsmap['p'])
            cTn_set = etree.SubElement(cBhvr, '{%s}cTn' % nsmap['p'])
            cTn_set.set('id', str(id_counter)); id_counter += 1
            cTn_set.set('dur', '1')
            cTn_set.set('fill', 'hold')

            stCond_set = etree.SubElement(cTn_set, '{%s}stCondLst' % nsmap['p'])
            cond_set = etree.SubElement(stCond_set, '{%s}cond' % nsmap['p'])
            cond_set.set('delay', '0')

            tgtEl = etree.SubElement(cBhvr, '{%s}tgtEl' % nsmap['p'])
            spTgt = etree.SubElement(tgtEl, '{%s}spTgt' % nsmap['p'])
            spTgt.set('spid', str(sp_id))

            attrNameLst = etree.SubElement(cBhvr, '{%s}attrNameLst' % nsmap['p'])
            attrName = etree.SubElement(attrNameLst, '{%s}attrName' % nsmap['p'])
            attrName.text = 'style.visibility'

            p_to = etree.SubElement(p_set, '{%s}to' % nsmap['p'])
            strVal = etree.SubElement(p_to, '{%s}strVal' % nsmap['p'])
            strVal.set('val', 'visible')

def find_image(lesson, slide_type, idx=0):
    lop_folder = lesson["folder"]
    lop_num = lesson["lop_num"]
    
    # 1. Check local Tuần 04 images
    khbd_dir = os.path.join(KHBD_BASE, lop_folder, "Tuần_04", "images")
    if os.path.isdir(khbd_dir):
        for name in [f"{slide_type}{idx+1}_lop{lop_num}.png", f"{slide_type}_lop{lop_num}.png"]:
            p = os.path.join(khbd_dir, name)
            if os.path.isfile(p):
                return p
        imgs = sorted([os.path.join(khbd_dir, f) for f in os.listdir(khbd_dir) if f.endswith(('.png','.jpg','.jpeg'))])
        if imgs:
            return imgs[idx % len(imgs)]
            
    # 2. Check SGK images
    sgk_dir = os.path.join(SGK_BASE, f"Lớp_{lop_num}", "all_extracted_images")
    if os.path.isdir(sgk_dir):
        imgs = sorted([os.path.join(sgk_dir, f) for f in os.listdir(sgk_dir) if f.endswith(('.png','.jpeg','.jpg'))])
        if imgs:
            return imgs[idx % len(imgs)]
            
    # 3. Check Tuần 03 fallback
    khbd_t3 = os.path.join(KHBD_BASE, lop_folder, "Tuần_03", "images")
    if os.path.isdir(khbd_t3):
        imgs = sorted([os.path.join(khbd_t3, f) for f in os.listdir(khbd_t3) if f.endswith(('.png','.jpg','.jpeg'))])
        if imgs:
            return imgs[idx % len(imgs)]
    return None

def build_cover(prs, data, pal, lesson, layout):
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["primary"], send_to_back=True)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, 1.8, pal["accent"], send_to_back=True)

    add_textbox(slide, 0.8, 2.0, 7.5, 1.5, data["title"], size_pt=34, bold=True, color_hex=pal["text_on_primary"])
    add_textbox(slide, 0.8, 3.8, 7.5, 0.8, data.get("subtitle", ""), size_pt=20, color_hex=pal["text_on_primary"])

    img = find_image(lesson, "cover", 0)
    if img:
        add_picture_safe(slide, img, SLIDE_W - 4.5, SAFE_TOP + 0.3, 3.8, 3.8)

    add_slide_transition(slide, "fade")
    return slide

def build_warmup(prs, data, pal, lesson, layout):
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, pal["primary"])
    badge_text = f"  {lesson['lop_label'].upper()} • {lesson['bai'].upper()}"
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, badge_text, size_pt=13, bold=True, color_hex=pal["text_on_primary"])

    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, 8, 0.65, data["title"], size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    content_y = title_y + 0.75
    card_w = 7.5
    card_h = SAFE_BOTTOM - content_y - 0.15
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.4, content_y, card_w, card_h, pal["card"])
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0.4, content_y, 0.12, card_h, pal["accent"])

    lines = data.get("content", "").split("\n")
    add_multiline_textbox(slide, 0.8, content_y+0.2, card_w-0.6, card_h-0.4, lines, size_pt=20, color_hex=pal["text_on_card"])

    img = find_image(lesson, "warmup", 0)
    if img:
        add_picture_safe(slide, img, 8.3, content_y, SAFE_RIGHT-8.3-0.2, card_h)

    add_slide_transition(slide, "push")
    return slide

def build_video_slide(prs, data, pal, lesson, layout):
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, "DC2626")
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, f"  {lesson['lop_label'].upper()} • KHÁM PHÁ VIDEO BÀI HỌC", size_pt=13, bold=True, color_hex="FFFFFF")

    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, SLIDE_W-1, 0.65, data["title"], size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    content_y = title_y + 0.75
    card_w = 7.5
    card_h = SAFE_BOTTOM - content_y - 0.15
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.4, content_y, card_w, card_h, pal["card"], border_hex="EF4444")
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0.4, content_y, 0.12, card_h, "EF4444")

    desc_lines = [
        data.get("desc", ""),
        "",
        "🔗 Đường link video tham khảo:",
        data.get("url", ""),
        "",
        "💡 Thầy/Cô sẽ chiếu video cho các em quan sát nhé!"
    ]
    add_multiline_textbox(slide, 0.8, content_y+0.2, card_w-0.6, card_h-0.4, desc_lines, size_pt=17, color_hex=pal["text_on_card"])

    img = find_image(lesson, "video", 0)
    if img:
        add_picture_safe(slide, img, 8.3, content_y, SAFE_RIGHT-8.3-0.2, card_h)

    add_slide_transition(slide, "cover")
    return slide

def build_learn_grid(prs, data, pal, lesson, layout, learn_idx=0):
    slide = prs.slides.add_slide(layout)
    bullets = data.get("bullets", [])
    n = len(bullets)

    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, pal["primary"])
    badge_text = f"  {lesson['lop_label'].upper()} • {lesson['bai'].upper()}"
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, badge_text, size_pt=13, bold=True, color_hex=pal["text_on_primary"])

    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, SLIDE_W-1, 0.65, data["title"], size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    grid_top = title_y + 0.8
    grid_bottom = SAFE_BOTTOM - 0.1
    avail_h = grid_bottom - grid_top

    cols, rows = (n, 1) if n <= 3 else (2, 2)
    gap = 0.25
    total_gap_x = gap * (cols - 1)
    card_w = (SLIDE_W - 1.0 - total_gap_x) / cols
    if rows == 1:
        img_h, text_h = 2.0, 0.9
    else:
        img_h, text_h = 1.1, 0.5
    card_h = img_h + text_h + 0.15
    total_gap_y = gap * (rows - 1)
    start_x = (SLIDE_W - (cols * card_w + total_gap_x)) / 2
    start_y = grid_top + (avail_h - (rows * card_h + total_gap_y)) / 2
    start_y = max(start_y, grid_top)

    for i, btext in enumerate(bullets):
        col, row = i % cols, i // cols
        x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + gap)

        if y + card_h > SAFE_BOTTOM + 0.05:
            break

        add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h, pal["card"], border_hex="D0D5DD")
        add_safe_shape(slide, MSO_SHAPE.RECTANGLE, x, y, card_w, 0.06, pal["accent"])

        img = find_image(lesson, "learn", learn_idx*3 + i)
        if img:
            img_w = min(1.8, card_w - 0.4)
            img_x = x + (card_w - img_w) / 2
            add_picture_safe(slide, img, img_x, y + 0.15, img_w, img_h)

        text_y = y + img_h + 0.15
        add_textbox(slide, x + 0.15, text_y, card_w - 0.3, text_h, btext, size_pt=16, bold=False, color_hex=pal["text_on_card"], alignment=PP_ALIGN.CENTER)

    add_slide_transition(slide, "wipe")
    return slide

def build_learn_row(prs, data, pal, lesson, layout, learn_idx=0):
    slide = prs.slides.add_slide(layout)
    bullets = data.get("bullets", [])

    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, pal["primary"])
    badge_text = f"  {lesson['lop_label'].upper()} • {lesson['bai'].upper()}"
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, badge_text, size_pt=13, bold=True, color_hex=pal["text_on_primary"])

    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, SLIDE_W-1, 0.65, data["title"], size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    row_top = title_y + 0.75
    n = len(bullets)
    gap = 0.15
    avail_h = SAFE_BOTTOM - row_top - 0.1
    row_h = min(1.4, (avail_h - gap*(n-1)) / max(n, 1))

    for i, btext in enumerate(bullets[:4]):
        y = row_top + i * (row_h + gap)

        if y + row_h > SAFE_BOTTOM + 0.05:
            break

        img = find_image(lesson, "learn", learn_idx*3 + i)

        img_w = 2.0
        img_x = 0.5
        txt_x = img_x + img_w + 0.2
        txt_w = SLIDE_W - txt_x - 0.5

        if img:
            add_picture_safe(slide, img, img_x, y, img_w, row_h)
        else:
            txt_x = 0.5
            txt_w = SLIDE_W - 1.0

        card = add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, txt_x, y, txt_w, row_h, pal["card"], border_hex="D0D5DD")
        add_safe_shape(slide, MSO_SHAPE.RECTANGLE, txt_x, y, 0.08, row_h, pal["accent"])
        add_textbox(slide, txt_x + 0.25, y + 0.15, txt_w - 0.4, row_h - 0.3, btext, size_pt=18, bold=False, color_hex=pal["text_on_card"])

    add_slide_transition(slide, "wipe")
    return slide

def build_practice(prs, data, pal, lesson, layout):
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, pal["primary"])
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, f"  {lesson['lop_label'].upper()} • LUYỆN TẬP TƯƠNG TÁC", size_pt=13, bold=True, color_hex=pal["text_on_primary"])

    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, SLIDE_W-1, 0.65, data["title"], size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    instr_y = title_y + 0.75
    instr_lines = data.get("instruction", "").split("\n")
    add_multiline_textbox(slide, 0.5, instr_y, SLIDE_W-1, 0.7, instr_lines, size_pt=18, color_hex=pal["text_on_bg"])

    items = data.get("items", [])
    item_y_start = instr_y + 0.8
    n_items = len(items)

    cols = 2 if n_items > 3 else max(n_items, 1)
    rows = (n_items + 1) // 2 if n_items > 3 else 1

    gap = 0.25
    card_w = (SLIDE_W - 1.0 - gap*(cols-1)) / cols
    card_h = min(1.8, (SAFE_BOTTOM - item_y_start - gap*(rows-1)) / rows)

    animation_groups = []

    for i, item in enumerate(items[:4]):
        col, row = i % cols, i // cols
        x = 0.5 + col * (card_w + gap)
        y = item_y_start + row * (card_h + gap)

        if y + card_h > SAFE_BOTTOM + 0.05:
            break

        card = add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h, pal["card"], border_hex=pal["accent"])

        img = find_image(lesson, "practice", i)
        pic = None
        if img:
            img_w = 1.4
            pic = add_picture_safe(slide, img, x + 0.2, y + 0.15, img_w, card_h - 0.3)
            txt_x = x + 0.2 + img_w + 0.2
            txt_w = card_w - (0.2 + img_w + 0.3)
        else:
            txt_x = x + 0.2
            txt_w = card_w - 0.4

        txt = add_textbox(slide, txt_x, y + 0.15, txt_w, card_h - 0.3, item, size_pt=16, bold=True, color_hex=pal["text_on_card"])

        group = [s for s in [card, pic, txt] if s is not None]
        animation_groups.append((group, i))

    if animation_groups:
        add_appear_animation(slide, animation_groups)

    add_slide_transition(slide, "wipe")
    return slide

def build_activity(prs, data, pal, lesson, layout):
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, pal["accent"])
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, f"  {lesson['lop_label'].upper()} • THỬ THÁCH", size_pt=13, bold=True, color_hex=pal["text_on_primary"])

    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, SLIDE_W-1, 0.65, data["title"], size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    content_y = title_y + 0.75
    content_h = SAFE_BOTTOM - content_y - 0.15

    card_w = 7.0
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.4, content_y, card_w, content_h, pal["card"])
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0.4, content_y, 0.12, content_h, pal["accent"])

    lines = data.get("content", "").split("\n") if data.get("content") else data.get("instruction", "").split("\n")
    if data.get("items"):
        lines.extend(data["items"])
    add_multiline_textbox(slide, 0.8, content_y+0.2, card_w-0.6, content_h-0.4, lines, size_pt=18, color_hex=pal["text_on_card"])

    img = find_image(lesson, "activity", 0)
    if img:
        img_x = 7.8
        img_w = SAFE_RIGHT - img_x - 0.2
        add_picture_safe(slide, img, img_x, content_y + 0.1, img_w, content_h - 0.2)

    add_slide_transition(slide, "cover")
    return slide

def build_summary(prs, data, pal, lesson, layout):
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    px, py = 0.6, SAFE_TOP + 0.15
    pw = SLIDE_W - 1.2
    ph = SAFE_BOTTOM - py - 0.15
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, px, py, pw, ph, pal["card"])
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, px, py, pw, 0.06, pal["primary"])

    add_textbox(slide, px+0.4, py+0.25, pw-0.8, 0.6, data["title"], size_pt=24, bold=True, color_hex=pal["text_on_bg"])

    items = data.get("items", [])
    n = len(items)
    card_gap = 0.2
    card_w = (pw - 0.8 - card_gap*(n-1)) / max(n, 1)
    card_y = py + 1.0
    card_h = ph - 1.3

    for i, item in enumerate(items):
        cx = px + 0.4 + i*(card_w + card_gap)
        add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx, card_y, card_w, card_h, pal["bg"], border_hex="CBD5E1")
        badge_sz = 0.45
        add_safe_shape(slide, MSO_SHAPE.OVAL, cx + 0.15, card_y + 0.15, badge_sz, badge_sz, pal["primary"])
        add_textbox(slide, cx + 0.15, card_y + 0.15, badge_sz, badge_sz, str(i+1), size_pt=14, bold=True, color_hex=pal["text_on_primary"], alignment=PP_ALIGN.CENTER)
        add_textbox(slide, cx + 0.15, card_y + 0.7, card_w - 0.3, card_h - 0.85, item, size_pt=17, bold=False, color_hex=pal["text_on_bg"])

    add_slide_transition(slide, "push")
    return slide

def build_thanks(prs, data, pal, lesson, layout):
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["primary"], send_to_back=True)

    add_textbox(slide, 1.0, 2.0, SLIDE_W-2.0, 1.2, data["title"], size_pt=34, bold=True, color_hex=pal["text_on_primary"], alignment=PP_ALIGN.CENTER)

    if data.get("content"):
        lines = data["content"].split("\n")
        add_multiline_textbox(slide, 1.0, 3.4, SLIDE_W-2.0, 1.5, lines, size_pt=20, color_hex=pal["text_on_primary"], alignment=PP_ALIGN.CENTER)

    add_slide_transition(slide, "fade")
    return slide

def remove_template_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        sldId = prs.slides._sldIdLst[0]
        prs.slides._sldIdLst.remove(sldId)

def generate_slides_tuan04():
    print("[+] BẮT ĐẦU TẠO SLIDE TUẦN 04 CHO TOÀN BỘ KHỐI (Tiền TH -> Lớp 9)...")
    
    for lesson in LESSONS_T4:
        prs = Presentation(TEMPLATE)
        remove_template_slides(prs)
        layout = prs.slide_layouts[6]
        pal = COLOR_PALETTES[lesson["palette_idx"]]
        is_thcs = lesson.get("is_thcs", False)
        
        learn_idx = 0
        for sc in lesson["slides_content"]:
            stype = sc["type"]
            if stype == "cover":
                build_cover(prs, sc, pal, lesson, layout)
            elif stype == "warmup":
                build_warmup(prs, sc, pal, lesson, layout)
            elif stype == "video":
                build_video_slide(prs, sc, pal, lesson, layout)
            elif stype == "learn":
                if is_thcs:
                    build_learn_row(prs, sc, pal, lesson, layout, learn_idx=learn_idx)
                else:
                    build_learn_grid(prs, sc, pal, lesson, layout, learn_idx=learn_idx)
                learn_idx += 1
            elif stype == "practice":
                build_practice(prs, sc, pal, lesson, layout)
            elif stype == "activity":
                build_activity(prs, sc, pal, lesson, layout)
            elif stype == "summary":
                build_summary(prs, sc, pal, lesson, layout)
            elif stype == "thanks":
                build_thanks(prs, sc, pal, lesson, layout)
                
        out_dir = os.path.join(KHBD_BASE, lesson["folder"], "Tuần_04")
        os.makedirs(out_dir, exist_ok=True)
        out_path = os.path.join(out_dir, lesson["file_name"])
        prs.save(out_path)
        print(f"    [OK] Tạo thành công: {out_path} ({len(prs.slides)} slides)")

if __name__ == '__main__':
    generate_slides_tuan04()
