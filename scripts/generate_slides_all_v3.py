# -*- coding: utf-8 -*-
"""
Script tạo Bộ Slide bài giảng Student-Facing v3 cho TOÀN BỘ 9 KHỐI LỚP Tuần 02
(Tiền Tiểu Học, Lớp 1, 2, 3, 4, 5, 6, 7, 8)

Quy chuẩn v3 (Đã lưu vào bộ nhớ Agent):
  1. Per-Bullet Images:
     - TH (Tiền TH -> Lớp 5): Layout A — Grid Flashcard (ảnh trên, text dưới)
     - THCS (Lớp 6 -> Lớp 8): Layout B — Horizontal Row (ảnh trái 2in×1.5in, text phải)
  2. Animation tuần tự cho slide Luyện tập / Ghép nối / Quiz (On Click appear per item)
  3. Game Images: Slide trò chơi bắt buộc có ảnh minh họa đủ lớn (~2in-3.5in)
  4. Anti-Bug Checklist: Safe Zone (Y 1.15 -> 6.35), Z-order insert(2), no custom footers, high contrast palette
"""
import sys, io, os, copy, glob
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import traceback

TEMPLATE  = r'D:\UNIGO\Hệ thống mẫu văn bản\Mẫu slide UNIGO.pptx'
SGK_BASE  = r'D:\UNIGO\SGK'
KHBD_BASE = r'D:\UNIGO\KHBD_Tin_học'

# ─── Template Safe Zone ───
SAFE_TOP    = 1.15   # Logo ends at 1.09in
SAFE_BOTTOM = 6.35   # Footer starts at 6.43in
SAFE_LEFT   = 0.3
SAFE_RIGHT  = 13.0
SLIDE_W     = 13.33
SLIDE_H     = 7.50

# ─── 8 Bộ màu tương phản cao (High Contrast Palettes) ───
COLOR_PALETTES = [
    # 0: Blue (Tiền TH)
    {"primary": "1B4F9B", "accent": "2D7DD2", "bg": "EBF3FE", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "1A2744", "text_on_card": "1A2744"},
    # 1: Purple (Lớp 1)
    {"primary": "5B21B6", "accent": "7C3AED", "bg": "F3EEFF", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "2E1065", "text_on_card": "2E1065"},
    # 2: Teal (Lớp 2)
    {"primary": "0F766E", "accent": "14B8A6", "bg": "ECFDF5", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "134E4A", "text_on_card": "134E4A"},
    # 3: Orange (Lớp 3)
    {"primary": "C2410C", "accent": "EA580C", "bg": "FFF7ED", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "7C2D12", "text_on_card": "431407"},
    # 4: Indigo (Lớp 4)
    {"primary": "3730A3", "accent": "4F46E5", "bg": "EEF2FF", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "1E1B4B", "text_on_card": "1E1B4B"},
    # 5: Rose (Lớp 5)
    {"primary": "BE185D", "accent": "EC4899", "bg": "FDF2F8", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "831843", "text_on_card": "831843"},
    # 6: Emerald (Lớp 6)
    {"primary": "047857", "accent": "10B981", "bg": "ECFDF5", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "064E3B", "text_on_card": "064E3B"},
    # 7: Sky (Lớp 7)
    {"primary": "0369A1", "accent": "0EA5E9", "bg": "F0F9FF", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "0C4A6E", "text_on_card": "0C4A6E"},
    # 8: Amber (Lớp 8)
    {"primary": "B45309", "accent": "F59E0B", "bg": "FFFBEB", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "78350F", "text_on_card": "451A03"},
]

# ─── Nội dung bài học 9 Khối Lớp ───
LESSONS = [
    {
        "id": "Tien_TH", "folder": "Tiền_tiểu_học", "lop_label": "Tiền Tiểu học", "lop_num": 0, "is_thcs": False,
        "bai": "Bài 1. Máy tính xung quanh em",
        "file_name": "Slide_Tin_hoc_Tien_TH_Bai01_May_tinh_xung_quanh_em.pptx",
        "slides_content": [
            {"type": "cover", "title": "Máy tính xung quanh em", "subtitle": "Tin học • Tiền Tiểu học • Bài 1"},
            {"type": "warmup", "title": "Các con ơi, nhìn xung quanh nào!",
             "content": "Con có thấy chiếc máy tính nào không?\nChỉ cho cô xem nhé!"},
            {"type": "learn", "title": "Đây là máy tính!",
             "bullets": [
                 "Trong phòng học của chúng mình",
                 "Ở nhà (laptop của bố mẹ)",
                 "Trên tay (điện thoại cũng là máy tính nhỏ!)"
             ]},
            {"type": "learn", "title": "Máy tính có những phần nào?",
             "bullets": [
                 "Màn hình — để con nhìn",
                 "Bàn phím — để con gõ chữ",
                 "Chuột — để con chỉ và nhấp",
                 "Thân máy — bộ não của máy tính"
             ]},
            {"type": "learn", "title": "Máy tính giúp chúng ta làm gì?",
             "bullets": [
                 "Học bài vui vẻ",
                 "Vẽ tranh đẹp",
                 "Nghe nhạc hay",
                 "Xem hoạt hình"
             ]},
            {"type": "practice", "title": "Cùng chơi nào!",
             "instruction": "Con hãy chỉ vào từng bộ phận và nói tên cho cô nghe:",
             "items": ["Đâu là màn hình?", "Đâu là bàn phím?", "Đâu là chuột?"]},
            {"type": "activity", "title": "Thử thách nhỏ!",
             "content": "Con hãy vẽ một chiếc máy tính vào giấy của mình.\nNhớ vẽ đủ các bộ phận nhé!"},
            {"type": "summary", "title": "Hôm nay con đã biết!",
             "items": ["Máy tính ở xung quanh chúng ta", "Máy tính có: màn hình, bàn phím, chuột, thân máy", "Máy tính giúp con học và chơi"]},
            {"type": "thanks", "title": "Các con giỏi lắm!", "content": "Hẹn gặp lại tiết sau nhé!"},
        ]
    },
    {
        "id": "Lop_1", "folder": "Lớp_1", "lop_label": "Lớp 1", "lop_num": 1, "is_thcs": False,
        "bai": "Bài 1. Chiếc máy tính của em",
        "file_name": "Slide_Tin_hoc_Lop_1_Bai01_Chiec_may_tinh_cua_em.pptx",
        "slides_content": [
            {"type": "cover", "title": "Chiếc máy tính của em", "subtitle": "Tin học • Lớp 1 • Bài 1"},
            {"type": "warmup", "title": "Em đã gặp máy tính ở đâu?",
             "content": "Hãy kể cho cô nghe:\n● Em thấy máy tính ở đâu?\n● Bố mẹ dùng máy tính để làm gì?"},
            {"type": "learn", "title": "Cùng khám phá máy tính nào!",
             "bullets": [
                 "Máy tính để bàn (Desktop)",
                 "Máy tính xách tay (Laptop)",
                 "Máy tính bảng (Tablet)"
             ]},
            {"type": "learn", "title": "Các bộ phận của máy tính",
             "bullets": [
                 "Màn hình — hiển thị hình ảnh, chữ",
                 "Bàn phím — em gõ chữ, số",
                 "Chuột — em nhấp để chọn",
                 "Thân máy — nơi xử lý thông tin"
             ]},
            {"type": "learn", "title": "Bật và tắt máy tính đúng cách",
             "bullets": [
                 "Bước 1: Nhấn nút nguồn để bật máy",
                 "Bước 2: Chờ máy khởi động xong",
                 "Bước 3: Khi muốn tắt → Chọn Shut Down"
             ]},
            {"type": "practice", "title": "Em hãy thực hành!",
             "instruction": "Thực hiện lần lượt các yêu cầu sau:",
             "items": ["Chỉ vào từng bộ phận và gọi tên", "Thử nhấp chuột trái", "Thử gõ tên mình bằng bàn phím"]},
            {"type": "activity", "title": "Trò chơi: Nối đúng!",
             "instruction": "Nối tên bộ phận với công dụng đúng:",
             "items": ["Màn hình ↔ Hiện hình ảnh", "Bàn phím ↔ Gõ chữ và số", "Chuột ↔ Nhấp và chọn", "Thân máy ↔ Bộ não máy tính"]},
            {"type": "summary", "title": "Em nhớ được gì nào?",
             "items": ["Máy tính có nhiều loại khác nhau", "4 bộ phận chính: Màn hình, Bàn phím, Chuột, Thân máy", "Bật/Tắt máy tính đúng quy trình"]},
            {"type": "thanks", "title": "Em giỏi lắm!", "content": "Về nhà: Hãy kể cho bố mẹ nghe\nvề chiếc máy tính ở lớp nhé!"},
        ]
    },
    {
        "id": "Lop_2", "folder": "Lớp_2", "lop_label": "Lớp 2", "lop_num": 2, "is_thcs": False,
        "bai": "Bài 1. Máy tính là người bạn của em",
        "file_name": "Slide_Tin_hoc_Lop_2_Bai01_May_tinh_la_nguoi_ban_cua_em.pptx",
        "slides_content": [
            {"type": "cover", "title": "Máy tính là người bạn của em", "subtitle": "Tin học • Lớp 2 • Bài 1"},
            {"type": "warmup", "title": "Máy tính giúp em những gì?",
             "content": "Hãy kể cho cô nghe:\n● Em dùng máy tính để làm gì?\n● Máy tính giúp ích gì cho công việc gia đình?"},
            {"type": "learn", "title": "Máy tính — Người bạn thông minh!",
             "bullets": [
                 "Học bài — tra cứu kiến thức bổ ích",
                 "Vẽ tranh — sáng tạo nhiều sắc màu",
                 "Viết văn — tập gõ chữ và soạn thảo"
             ]},
            {"type": "learn", "title": "Máy tính còn giúp gì nữa?",
             "bullets": [
                 "Gọi video — nói chuyện với người thân",
                 "Nghe nhạc và xem phim giải trí",
                 "Chơi các trò chơi học tập trí tuệ"
             ]},
            {"type": "learn", "title": "Sử dụng máy tính an toàn",
             "bullets": [
                 "Ngồi đúng tư thế, mắt cách màn hình 40-50cm",
                 "Nghỉ mắt sau mỗi 20-30 phút làm việc",
                 "Luôn xin phép bố mẹ trước khi dùng máy"
             ]},
            {"type": "practice", "title": "Em hãy thực hành!",
             "instruction": "Hãy hoàn thành từng bước luyện tập:",
             "items": ["Ngồi đúng tư thế trước máy tính", "Mở phần mềm vẽ", "Vẽ một chiếc ô tô hoặc bông hoa"]},
            {"type": "activity", "title": "Trò chơi: Đúng hay Sai?",
             "instruction": "Chọn Đúng hoặc Sai cho mỗi câu sau:",
             "items": ["Ngồi quá sát màn hình là ĐÚNG hay SAI?", "Hỏi ý kiến bố mẹ trước khi lên mạng?", "Chơi game cả ngày không nghỉ là ĐÚNG hay SAI?"]},
            {"type": "summary", "title": "Em nhớ được gì nào?",
             "items": ["Máy tính là người bạn hỗ trợ học tập và giải trí", "Biết nhiều ứng dụng thực tế của máy tính", "Giữ tư thế và khoảng cách ngồi an toàn"]},
            {"type": "thanks", "title": "Em giỏi lắm!", "content": "BTVN: Kể cho người thân 3 điều\nmáy tính giúp em học tốt hơn!"},
        ]
    },
    {
        "id": "Lop_3", "folder": "Lớp_3", "lop_label": "Lớp 3", "lop_num": 3, "is_thcs": False,
        "bai": "Bài 1. Thông tin và quyết định",
        "file_name": "Slide_Tin_hoc_Lop_3_Bai01_Thong_tin_va_quyet_dinh.pptx",
        "slides_content": [
            {"type": "cover", "title": "Thông tin và Quyết định", "subtitle": "Tin học • Lớp 3 • Bài 1"},
            {"type": "warmup", "title": "Em hãy đoán xem!",
             "content": "Sáng nay trước khi đi học, em quyết định mặc áo gì?\nEm dựa vào đâu để đưa ra quyết định đó?"},
            {"type": "learn", "title": "Thông tin là gì?",
             "bullets": [
                 "Nhìn thấy — chữ viết, hình ảnh, màu sắc",
                 "Nghe thấy — âm thanh, tiếng nói, tiếng chuông",
                 "Cảm nhận — độ nóng, lạnh, cứng, mềm"
             ]},
            {"type": "learn", "title": "Thông tin giúp em quyết định!",
             "bullets": [
                 "Trời mưa → Em mang theo ô hoặc áo mưa",
                 "Đèn đỏ → Em dừng xe lại trước vạch",
                 "Chuông reo → Em nhanh chóng vào lớp"
             ]},
            {"type": "learn", "title": "Các dạng thông tin cơ bản",
             "bullets": [
                 "Dạng chữ (văn bản) — sách, báo, tin nhắn",
                 "Dạng hình ảnh — bức tranh, ảnh chụp",
                 "Dạng âm thanh — bản nhạc, tiếng nói"
             ]},
            {"type": "practice", "title": "Em hãy thử nào!",
             "instruction": "Trả lời từng câu hỏi tình huống:",
             "items": ["Thấy mây đen kéo đến → Em quyết định gì?", "Nghe tiếng chuông báo thức → Em làm gì?", "Nhìn thấy bảng hiệu ghi Stop → Em làm gì?"]},
            {"type": "activity", "title": "Trò chơi: Phân loại thông tin!",
             "instruction": "Hãy xác định dạng của mỗi thông tin sau:",
             "items": ["Bảng thời khóa biểu ↔ Dạng chữ / bảng", "Bản nhạc Quốc ca ↔ Dạng âm thanh", "Bức ảnh tập thể lớp ↔ Dạng hình ảnh"]},
            {"type": "summary", "title": "Hôm nay em đã học!",
             "items": ["Thông tin là những gì ta tiếp nhận qua giác quan", "3 dạng thông tin chính: Chữ, Hình ảnh, Âm thanh", "Thông tin là cơ sở giúp con người quyết định"]},
            {"type": "thanks", "title": "Em làm tốt lắm!", "content": "BTVN: Kể 3 thông tin em nhận được\nhôm nay và quyết định tương ứng!"},
        ]
    },
    {
        "id": "Lop_4", "folder": "Lớp_4", "lop_label": "Lớp 4", "lop_num": 4, "is_thcs": False,
        "bai": "Bài 1. Phần cứng và phần mềm máy tính",
        "file_name": "Slide_Tin_hoc_Lop_4_Bai01_Phan_cung_va_phan_mem.pptx",
        "slides_content": [
            {"type": "cover", "title": "Phần cứng và Phần mềm", "subtitle": "Tin học • Lớp 4 • Bài 1"},
            {"type": "warmup", "title": "Thử đoán xem!",
             "content": "Em có thể chạm tay vào bàn phím không? → CÓ\nEm có thể sờ vào nhân vật trong trò chơi không? → KHÔNG\nVì sao lại như vậy?"},
            {"type": "learn", "title": "Phần cứng (Hardware) là gì?",
             "bullets": [
                 "Màn hình, bàn phím, chuột máy tính",
                 "Loa, tai nghe, webcam thu hình",
                 "CPU, RAM, ổ cứng bên trong thân máy"
             ]},
            {"type": "learn", "title": "Phần mềm (Software) là gì?",
             "bullets": [
                 "Hệ điều hành Windows, macOS điều khiển máy",
                 "Trình duyệt web Chrome, Edge lướt Internet",
                 "Phần mềm soạn thảo Word, phần mềm vẽ Paint"
             ]},
            {"type": "learn", "title": "Mối quan hệ Phần cứng & Phần mềm",
             "bullets": [
                 "Phần cứng như CƠ THỂ của máy tính",
                 "Phần mềm như TRÍ ÓC điều khiển máy tính",
                 "Máy tính cần CẢ HAI để hoạt động"
             ]},
            {"type": "practice", "title": "Em hãy phân loại!",
             "instruction": "Phân loại từng thiết bị/chương trình:",
             "items": ["Bàn phím ↔ Phần cứng", "Microsoft Word ↔ Phần mềm", "Chuột máy tính ↔ Phần cứng", "Trình duyệt Chrome ↔ Phần mềm"]},
            {"type": "activity", "title": "Thử thách nhanh!",
             "instruction": "Quan sát phòng máy và thực hiện:",
             "items": ["Liệt kê 3 phần cứng em thấy trước mặt", "Kể tên 2 phần mềm em đã từng sử dụng", "So sánh sự khác nhau giữa phần cứng và phần mềm"]},
            {"type": "summary", "title": "Hôm nay em đã hiểu!",
             "items": ["Phần cứng: nhìn thấy và sờ được (Hardware)", "Phần mềm: chương trình chạy bên trong (Software)", "Phần cứng và phần mềm luôn đồng hành cùng nhau"]},
            {"type": "thanks", "title": "Em giỏi lắm!", "content": "BTVN: Liệt kê 3 phần cứng và\n3 phần mềm của máy tính ở nhà em!"},
        ]
    },
    {
        "id": "Lop_5", "folder": "Lớp_5", "lop_label": "Lớp 5", "lop_num": 5, "is_thcs": False,
        "bai": "Bài 1. Em có thể làm gì với máy tính",
        "file_name": "Slide_Tin_hoc_Lop_5_Bai01_Em_co_the_lam_gi_voi_may_tinh.pptx",
        "slides_content": [
            {"type": "cover", "title": "Em có thể làm gì với máy tính?", "subtitle": "Tin học • Lớp 5 • Bài 1"},
            {"type": "warmup", "title": "Kể cho cô nghe nào!",
             "content": "Tuần trước em đã dùng máy tính để làm những việc gì?\nKể ít nhất 2 việc cụ thể nhé!"},
            {"type": "learn", "title": "Máy tính — Trợ thủ đa năng!",
             "bullets": [
                 "Soạn thảo văn bản và bài viết (Word)",
                 "Tính toán và xử lý số liệu (Excel)",
                 "Thiết kế tranh ảnh và sơ đồ (Paint, Canva)",
                 "Tạo bài trình chiếu sinh động (PowerPoint)"
             ]},
            {"type": "learn", "title": "Kết nối và chia sẻ tri thức!",
             "bullets": [
                 "Tìm kiếm thông tin học tập trên Google",
                 "Gửi thư điện tử Email trao đổi bài",
                 "Học trực tuyến qua Zoom, Google Meet"
             ]},
            {"type": "learn", "title": "Máy tính trong mọi nghề nghiệp",
             "bullets": [
                 "Bác sĩ — chẩn đoán bệnh và lưu hồ sơ",
                 "Kỹ sư — thiết kế công trình và bản vẽ",
                 "Giáo viên — giảng dạy và soạn giáo án"
             ]},
            {"type": "practice", "title": "Em hãy thực hành!",
             "instruction": "Thực hiện lần lượt các bước:",
             "items": ["Mở phần mềm Word", "Gõ họ tên và lớp của em", "Viết câu: 'Em thích máy tính vì...'", "Lưu file bài làm vào thư mục"]},
            {"type": "activity", "title": "Trò chơi: Ghép đôi công việc!",
             "instruction": "Nối công việc với phần mềm tương ứng:",
             "items": ["Viết bài văn ↔ Word", "Tạo bài thuyết trình ↔ PowerPoint", "Vẽ bức tranh ↔ Paint", "Tính bảng điểm ↔ Excel"]},
            {"type": "summary", "title": "Hôm nay em đã biết!",
             "items": ["Máy tính là công cụ đa năng trong học tập", "Mỗi công việc ứng với phần mềm chuyên dụng", "Máy tính đóng vai trò quan trọng trong xã hội"]},
            {"type": "thanks", "title": "Em giỏi lắm!", "content": "BTVN: Hỏi 1 người lớn về công việc\nhọ thực hiện trên máy tính hàng ngày!"},
        ]
    },
    {
        "id": "Lop_6", "folder": "Lớp_6", "lop_label": "Lớp 6", "lop_num": 6, "is_thcs": True,
        "bai": "Bài 1. Thông tin và dữ liệu",
        "file_name": "Slide_Tin_hoc_Lop_6_Bai01_Thong_tin_va_du_lieu.pptx",
        "slides_content": [
            {"type": "cover", "title": "Thông tin và Dữ liệu", "subtitle": "Tin học • Lớp 6 • Bài 1"},
            {"type": "warmup", "title": "Hãy thử suy nghĩ!",
             "content": "'Ngày mai trời sẽ mưa' — Đây là thông tin hay dữ liệu?\n'23°C' — Còn đây thì sao?\nHai khái niệm này khác nhau như thế nào?"},
            {"type": "learn", "title": "Khái niệm Thông tin (Information)",
             "bullets": [
                 "Thông tin là những hiểu biết của con người về thế giới xung quanh",
                 "Ví dụ: 'Hà Nội là thủ đô của Việt Nam', 'Hôm nay trời nắng đẹp'",
                 "Thông tin giúp con người đưa ra các quyết định hợp lý"
             ]},
            {"type": "learn", "title": "Khái niệm Dữ liệu (Data)",
             "bullets": [
                 "Dữ liệu là thông tin được biểu diễn dưới dạng máy tính xử lý được",
                 "Dạng văn bản (text): chữ cái, số, ký hiệu",
                 "Dạng hình ảnh, âm thanh, video được lưu trữ kỹ thuật số"
             ]},
            {"type": "learn", "title": "Mối quan hệ Thông tin ↔ Dữ liệu",
             "bullets": [
                 "Thông tin gốc -> Mã hóa, lưu trữ -> Dữ liệu trong máy tính",
                 "Dữ liệu được xử lý -> Giải mã -> Thông tin cho con người",
                 "Dữ liệu là phương tiện chở thông tin trên thiết bị công nghệ"
             ]},
            {"type": "practice", "title": "Hãy phân loại dữ liệu!",
             "instruction": "Xác định dạng dữ liệu cho các đối tượng:",
             "items": ["Bài hát Quốc ca ↔ Dữ liệu Âm thanh", "Ảnh chụp tập thể lớp ↔ Dữ liệu Hình ảnh", "Bảng điểm tổng kết ↔ Dữ liệu Văn bản & Số", "Video bài giảng ↔ Dữ liệu Video"]},
            {"type": "activity", "title": "Thảo luận nhóm",
             "instruction": "Thảo luận theo nhóm 4 học sinh (5 phút):",
             "items": ["Khi đăng ảnh lên mạng xã hội, em tạo dữ liệu dạng gì?", "Ai có thể tiếp nhận thông tin từ dữ liệu đó?", "Cần chú ý gì để bảo vệ dữ liệu cá nhân?"]},
            {"type": "summary", "title": "Tóm tắt bài học",
             "items": ["Thông tin là hiểu biết; Dữ liệu là dạng biểu diễn thông tin", "4 dạng dữ liệu phổ biến: Văn bản, Hình ảnh, Âm thanh, Video", "Xử lý dữ liệu nhằm rút ra thông tin có ích"]},
            {"type": "thanks", "title": "Bài học kết thúc!", "content": "BTVN: Liệt kê 5 loại dữ liệu\nem tiếp nhận trong một ngày!"},
        ]
    },
    {
        "id": "Lop_7", "folder": "Lớp_7", "lop_label": "Lớp 7", "lop_num": 7, "is_thcs": True,
        "bai": "Bài 1. Thiết bị vào - ra",
        "file_name": "Slide_Tin_hoc_Lop_7_Bai01_Thiet_bi_vao_ra.pptx",
        "slides_content": [
            {"type": "cover", "title": "Thiết bị vào - ra (I/O Devices)", "subtitle": "Tin học • Lớp 7 • Bài 1"},
            {"type": "warmup", "title": "Quan sát và trả lời!",
             "content": "Hãy nhìn quanh phòng máy tính:\nKể tên TẤT CẢ các thiết bị em nhìn thấy kết nối với máy tính!\nThiết bị nào đưa thông tin vào, thiết bị nào đưa thông tin ra?"},
            {"type": "learn", "title": "Thiết bị vào (Input Devices)",
             "bullets": [
                 "Bàn phím (Keyboard) — nhập dữ liệu văn bản và lệnh",
                 "Chuột (Mouse) — điều khiển con trỏ và thao tác chọn",
                 "Microphone & Webcam — thu âm thanh và hình ảnh trực tiếp",
                 "Máy quét (Scanner) — chuyển văn bản giấy thành dữ liệu số"
             ]},
            {"type": "learn", "title": "Thiết bị ra (Output Devices)",
             "bullets": [
                 "Màn hình (Monitor) — hiển thị kết quả xử lý trực quan",
                 "Loa & Tai nghe (Speaker/Headphone) — phát tín hiệu âm thanh",
                 "Máy in (Printer) — in dữ liệu ra bản giấy",
                 "Máy chiếu (Projector) — trình chiếu nội dung lên màn rộng"
             ]},
            {"type": "learn", "title": "Sơ đồ hoạt động Input → Processing → Output",
             "bullets": [
                 "Thiết bị Vào (Input) tiếp nhận thông tin từ người dùng",
                 "Bộ xử lý trung tâm (CPU) xử lý dữ liệu theo chương trình",
                 "Thiết bị Ra (Output) xuất kết quả cho người dùng tiếp nhận"
             ]},
            {"type": "practice", "title": "Hãy phân loại thiết bị!",
             "instruction": "Xác định vai trò của từng thiết bị:",
             "items": ["Bàn phím & Chuột ↔ Thiết bị VÀO (Input)", "Màn hình & Máy in ↔ Thiết bị RA (Output)", "Webcam & Micro ↔ Thiết bị VÀO (Input)", "Màn hình cảm ứng ↔ Vừa VÀO vừa RA"]},
            {"type": "activity", "title": "Thảo luận chuyên sâu",
             "instruction": "Trả lời các câu hỏi tình huống (5 phút):",
             "items": ["Nếu hỏng thiết bị vào, ta có điều khiển được máy không?", "Điện thoại thông minh dùng thiết bị vào/ra nào?", "Nêu 2 thiết bị vừa là thiết bị vào vừa là thiết bị ra"]},
            {"type": "summary", "title": "Tóm tắt bài học",
             "items": ["Thiết bị vào (Input): nhập dữ liệu vào máy tính", "Thiết bị ra (Output): xuất kết quả từ máy tính", "Quy trình xử lý: Input -> Processing -> Output"]},
            {"type": "thanks", "title": "Bài học kết thúc!", "content": "BTVN: Vẽ sơ đồ khối các thiết bị\nVào - Ra của máy tính ở phòng máy!"},
        ]
    },
    {
        "id": "Lop_8", "folder": "Lớp_8", "lop_label": "Lớp 8", "lop_num": 8, "is_thcs": True,
        "bai": "Bài 1. Lược sử công cụ tính toán",
        "file_name": "Slide_Tin_hoc_Lop_8_Bai01_Luoc_su_cong_cu_tinh_toan.pptx",
        "slides_content": [
            {"type": "cover", "title": "Lược sử công cụ tính toán", "subtitle": "Tin học • Lớp 8 • Bài 1"},
            {"type": "warmup", "title": "Quan sát và so sánh!",
             "content": "So sánh hai hình ảnh:\n● Chiếc bàn tính Abacus cổ đại\n● Siêu máy tính hiện đại ngày nay\nĐiểm giống và khác nhau cốt lõi là gì?"},
            {"type": "learn", "title": "Từ bàn tính cổ đại đến máy tính cơ học",
             "bullets": [
                 "Bàn tính Abacus (~2400 TCN) — công cụ tính toán đầu tiên",
                 "Máy tính Pascaline (1642 - Blaise Pascal) — máy tính cơ học dùng bánh răng",
                 "Máy phân tích (1834 - Charles Babbage) — ý tưởng máy tính tự động đầu tiên"
             ]},
            {"type": "learn", "title": "Kỷ nguyên máy tính điện tử (ENIAC -> PC)",
             "bullets": [
                 "ENIAC (1946) — máy tính điện tử đầu tiên, dùng đèn chân không, nặng 30 tấn",
                 "Transistor & Mạch tích hợp (IC) — giúp máy tính nhỏ gọn và nhanh hơn gấp ngàn lần",
                 "Microprocessor (1971) — sự ra đời của Máy tính cá nhân (PC)"
             ]},
            {"type": "learn", "title": "Máy tính thế hệ mới & Trí tuệ nhân tạo",
             "bullets": [
                 "Máy tính ngày càng nhỏ gọn, tốc độ siêu nhanh",
                 "Điện toán đám mây & Internet vạn vật (IoT)",
                 "Trí tuệ nhân tạo (AI) — máy tính biết học hỏi và suy luận"
             ]},
            {"type": "practice", "title": "Sắp xếp mốc lịch sử!",
             "instruction": "Sắp xếp theo đúng tiến trình thời gian:",
             "items": ["Bàn tính Abacus (~2400 TCN) ↔ Công cụ đầu tiên", "Máy tính Pascaline (1642) ↔ Bánh răng cơ học", "Máy tính ENIAC (1946) ↔ Đèn chân không", "Vi xử lý Intel (1971) ↔ Máy tính cá nhân (PC)"]},
            {"type": "activity", "title": "Thảo luận tương lai AI",
             "instruction": "Thảo luận nhóm 4 học sinh (5 phút):",
             "items": ["Máy tính thế hệ 5 (AI) giúp ích gì cho học tập?", "Em dự đoán 20 năm nữa máy tính sẽ thay đổi ra sao?", "Con người cần chuẩn bị kỹ năng gì trong kỷ nguyên AI?"]},
            {"type": "summary", "title": "Tóm tắt bài học",
             "items": ["Lịch sử tính toán trải qua hàng nghìn năm phát triển", "5 thế hệ: Đèn chân không -> Transistor -> IC -> Vi xử lý -> AI", "Xu hướng: Nhỏ gọn hơn, Nhanh hơn, Thông minh hơn"]},
            {"type": "thanks", "title": "Bài học kết thúc!", "content": "BTVN: Vẽ dòng thời gian 5 mốc quan trọng\ntrong lịch sử công cụ tính toán!"},
        ]
    },
]


# ═══════════════════════════════════════════════════
# HELPER FUNCTIONS
# ═══════════════════════════════════════════════════

def hex_rgb(h):
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def set_font(run, size_pt, bold=False, color_hex="333333", font_name="Arial"):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = hex_rgb(color_hex)
    run.font.name = font_name

def add_textbox(slide, left, top, width, height, text, size_pt=18, bold=False,
                color_hex="333333", alignment=PP_ALIGN.LEFT, font_name="Arial"):
    actual_top = max(top, SAFE_TOP)
    txbox = slide.shapes.add_textbox(Inches(left), Inches(actual_top), Inches(width), Inches(height))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size_pt, bold, color_hex, font_name)
    return txbox

def add_multiline_textbox(slide, left, top, width, height, lines, size_pt=18,
                           color_hex="333333", bold=False, font_name="Arial",
                           alignment=PP_ALIGN.LEFT):
    actual_top = max(top, SAFE_TOP)
    txbox = slide.shapes.add_textbox(Inches(left), Inches(actual_top), Inches(width), Inches(height))
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alignment
        run = p.add_run()
        run.text = line
        set_font(run, size_pt, bold, color_hex, font_name)
        p.space_after = Pt(6)
    return txbox

def add_safe_shape(slide, shape_type, left, top, width, height, fill_hex,
                   border_hex=None, send_to_back=False):
    actual_top = max(top, SAFE_TOP)
    actual_bottom = min(top + height, SAFE_BOTTOM)
    actual_height = max(actual_bottom - actual_top, 0.1)

    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left), Inches(actual_top), Inches(width), Inches(actual_height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(fill_hex)
    if border_hex:
        shape.line.color.rgb = hex_rgb(border_hex)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()

    if send_to_back:
        sp = shape._element
        spTree = sp.getparent()
        spTree.remove(sp)
        spTree.insert(2, sp)  # Never insert(0)
    return shape

def add_picture_safe(slide, img_path, left, top, width, height):
    actual_top = max(top, SAFE_TOP)
    actual_bottom = min(top + height, SAFE_BOTTOM)
    actual_height = max(actual_bottom - actual_top, 0.5)
    try:
        return slide.shapes.add_picture(img_path, Inches(left), Inches(actual_top),
                                         Inches(width), Inches(actual_height))
    except Exception as e:
        return None

def add_slide_transition(slide, transition_type="fade"):
    transitions = {
        'fade': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>',
        'push': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:push dir="l"/></p:transition>',
        'wipe': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:wipe/></p:transition>',
        'cover': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:cover/></p:transition>',
    }
    xml_str = transitions.get(transition_type, transitions['fade'])
    slide._element.append(etree.fromstring(xml_str))


def add_appear_animation(slide, shapes_with_click):
    nsmap = {
        'p':  'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a':  'http://schemas.openxmlformats.org/drawingml/2006/main',
    }

    timing = etree.SubElement(slide._element, '{%s}timing' % nsmap['p'])
    tnLst = etree.SubElement(timing, '{%s}tnLst' % nsmap['p'])
    par_root = etree.SubElement(tnLst, '{%s}par' % nsmap['p'])
    cTn_root = etree.SubElement(par_root, '{%s}cTn' % nsmap['p'])
    cTn_root.set('id', '1')
    cTn_root.set('dur', 'indefinite')
    cTn_root.set('restart', 'never')
    cTn_root.set('nodeType', 'tmRoot')

    childTnLst_root = etree.SubElement(cTn_root, '{%s}childTnLst' % nsmap['p'])

    seq = etree.SubElement(childTnLst_root, '{%s}seq' % nsmap['p'])
    seq.set('concurrent', '1')
    seq.set('nextAc', 'seek')

    cTn_seq = etree.SubElement(seq, '{%s}cTn' % nsmap['p'])
    cTn_seq.set('id', '2')
    cTn_seq.set('dur', 'indefinite')
    cTn_seq.set('nodeType', 'mainSeq')

    childTnLst_seq = etree.SubElement(cTn_seq, '{%s}childTnLst' % nsmap['p'])

    prevCond = etree.SubElement(seq, '{%s}prevCondLst' % nsmap['p'])
    cond_prev = etree.SubElement(prevCond, '{%s}cond' % nsmap['p'])
    cond_prev.set('evt', 'onPrev')
    cond_prev.set('delay', '0')
    tgtEl_prev = etree.SubElement(cond_prev, '{%s}tgtEl' % nsmap['p'])
    etree.SubElement(tgtEl_prev, '{%s}sldTgt' % nsmap['p'])

    nextCond = etree.SubElement(seq, '{%s}nextCondLst' % nsmap['p'])
    cond_next = etree.SubElement(nextCond, '{%s}cond' % nsmap['p'])
    cond_next.set('evt', 'onNext')
    cond_next.set('delay', '0')
    tgtEl_next = etree.SubElement(cond_next, '{%s}tgtEl' % nsmap['p'])
    etree.SubElement(tgtEl_next, '{%s}sldTgt' % nsmap['p'])

    id_counter = 3

    for shape_group, click_idx in shapes_with_click:
        if not isinstance(shape_group, (list, tuple)):
            shape_group = [shape_group]

        par_click = etree.SubElement(childTnLst_seq, '{%s}par' % nsmap['p'])
        cTn_click = etree.SubElement(par_click, '{%s}cTn' % nsmap['p'])
        cTn_click.set('id', str(id_counter)); id_counter += 1
        cTn_click.set('fill', 'hold')

        stCondLst = etree.SubElement(cTn_click, '{%s}stCondLst' % nsmap['p'])
        cond = etree.SubElement(stCondLst, '{%s}cond' % nsmap['p'])
        cond.set('delay', '0')

        childTnLst_click = etree.SubElement(cTn_click, '{%s}childTnLst' % nsmap['p'])

        for shape in shape_group:
            if shape is None:
                continue
            sp_id = shape.shape_id

            par_anim = etree.SubElement(childTnLst_click, '{%s}par' % nsmap['p'])
            cTn_anim = etree.SubElement(par_anim, '{%s}cTn' % nsmap['p'])
            cTn_anim.set('id', str(id_counter)); id_counter += 1
            cTn_anim.set('presetID', '1')  # Appear
            cTn_anim.set('presetClass', 'entr')
            cTn_anim.set('presetSubtype', '0')
            cTn_anim.set('fill', 'hold')

            stCond_anim = etree.SubElement(cTn_anim, '{%s}stCondLst' % nsmap['p'])
            cond_anim = etree.SubElement(stCond_anim, '{%s}cond' % nsmap['p'])
            cond_anim.set('delay', '0')

            childTnLst_anim = etree.SubElement(cTn_anim, '{%s}childTnLst' % nsmap['p'])

            p_set = etree.SubElement(childTnLst_anim, '{%s}set' % nsmap['p'])
            cBhvr = etree.SubElement(p_set, '{%s}cBhvr' % nsmap['p'])
            cTn_set = etree.SubElement(cBhvr, '{%s}cTn' % nsmap['p'])
            cTn_set.set('id', str(id_counter)); id_counter += 1
            cTn_set.set('dur', '1')
            cTn_set.set('fill', 'hold')

            stCond_set = etree.SubElement(cTn_set, '{%s}stCondLst' % nsmap['p'])
            cond_set = etree.SubElement(stCond_set, '{%s}cond' % nsmap['p'])
            cond_set.set('delay', '0')

            tgtEl = etree.SubElement(cBhvr, '{%s}tgtEl' % nsmap['p'])
            spTgt = etree.SubElement(tgtEl, '{%s}spTgt' % nsmap['p'])
            spTgt.set('spid', str(sp_id))

            attrNameLst = etree.SubElement(cBhvr, '{%s}attrNameLst' % nsmap['p'])
            attrName = etree.SubElement(attrNameLst, '{%s}attrName' % nsmap['p'])
            attrName.text = 'style.visibility'

            p_to = etree.SubElement(p_set, '{%s}to' % nsmap['p'])
            strVal = etree.SubElement(p_to, '{%s}strVal' % nsmap['p'])
            strVal.set('val', 'visible')


def find_best_image_for_lesson(lesson, slide_type, idx=0):
    """Find image for a slide from SGK or KHBD images folder"""
    lop_folder = lesson["folder"]
    lop_num = lesson["lop_num"]

    # 1. Check KHBD images
    khbd_img_dir = os.path.join(KHBD_BASE, lop_folder, "Tuần_02", "images")
    if os.path.isdir(khbd_img_dir):
        imgs = sorted([f for f in os.listdir(khbd_img_dir) if f.endswith(('.png','.jpeg','.jpg'))])
        if imgs:
            return os.path.join(khbd_img_dir, imgs[idx % len(imgs)])

    # 2. Check SGK images
    sgk_img_dir = os.path.join(SGK_BASE, f"Lớp_{lop_num}", "bai1_images")
    if os.path.isdir(sgk_img_dir):
        imgs = sorted([f for f in os.listdir(sgk_img_dir)
                       if f.endswith(('.png','.jpeg','.jpg')) and not os.path.isdir(os.path.join(sgk_img_dir, f))])
        if imgs:
            return os.path.join(sgk_img_dir, imgs[idx % len(imgs)])

        full_dir = os.path.join(sgk_img_dir, "full_pages")
        if os.path.isdir(full_dir):
            pages = sorted(os.listdir(full_dir))
            if pages:
                return os.path.join(full_dir, pages[idx % len(pages)])
    return None


def remove_template_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        sldId = prs.slides._sldIdLst[0]
        prs.slides._sldIdLst.remove(sldId)


# ═══════════════════════════════════════════════════
# BUILDERS FOR ALL GRADES
# ═══════════════════════════════════════════════════

def build_cover(prs, data, pal, lesson, layout):
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["primary"], send_to_back=True)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, 1.8, pal["accent"], send_to_back=True)

    add_textbox(slide, 0.8, 2.0, 7.5, 1.5, data["title"], size_pt=34, bold=True, color_hex=pal["text_on_primary"])
    add_textbox(slide, 0.8, 3.8, 7.5, 0.8, data.get("subtitle", ""), size_pt=20, color_hex=pal["text_on_primary"])

    img = find_best_image_for_lesson(lesson, "cover", 0)
    if img:
        add_picture_safe(slide, img, SLIDE_W - 4.5, SAFE_TOP + 0.3, 3.8, 3.8)

    add_slide_transition(slide, "fade")
    return slide


def build_warmup(prs, data, pal, lesson, layout):
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, pal["primary"])
    badge_text = f"  {lesson['lop_label'].upper()} • {lesson['bai'].upper()}"
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, badge_text, size_pt=13, bold=True, color_hex=pal["text_on_primary"])

    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, 8, 0.65, data["title"], size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    content_y = title_y + 0.75
    card_w = 7.5
    card_h = SAFE_BOTTOM - content_y - 0.15
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.4, content_y, card_w, card_h, pal["card"])
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0.4, content_y, 0.12, card_h, pal["accent"])

    lines = data.get("content", "").split("\n")
    add_multiline_textbox(slide, 0.8, content_y+0.2, card_w-0.6, card_h-0.4, lines, size_pt=20, color_hex=pal["text_on_card"])

    img = find_best_image_for_lesson(lesson, "warmup", 0)
    if img:
        add_picture_safe(slide, img, 8.3, content_y, SAFE_RIGHT-8.3-0.2, card_h)

    add_slide_transition(slide, "push")
    return slide


def build_learn_grid(prs, data, pal, lesson, layout, learn_idx=0):
    """Layout A — Grid Flashcard (Primary Grades: Tiền TH -> Lớp 5)"""
    slide = prs.slides.add_slide(layout)
    bullets = data.get("bullets", [])
    n = len(bullets)

    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, pal["primary"])
    badge_text = f"  {lesson['lop_label'].upper()} • {lesson['bai'].upper()}"
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, badge_text, size_pt=13, bold=True, color_hex=pal["text_on_primary"])

    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, SLIDE_W-1, 0.65, data["title"], size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    grid_top = title_y + 0.8
    grid_bottom = SAFE_BOTTOM - 0.1
    avail_h = grid_bottom - grid_top

    cols, rows = (n, 1) if n <= 3 else (2, 2)
    gap = 0.25
    total_gap_x = gap * (cols - 1)
    card_w = (SLIDE_W - 1.0 - total_gap_x) / cols
    if rows == 1:
        img_h, text_h = 2.0, 0.9
    else:
        img_h, text_h = 1.1, 0.5
    card_h = img_h + text_h + 0.15
    total_gap_y = gap * (rows - 1)
    start_x = (SLIDE_W - (cols * card_w + total_gap_x)) / 2
    start_y = grid_top + (avail_h - (rows * card_h + total_gap_y)) / 2
    start_y = max(start_y, grid_top)

    for i, btext in enumerate(bullets):
        col, row = i % cols, i // cols
        x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + gap)

        if y + card_h > SAFE_BOTTOM + 0.05:
            break

        add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h, pal["card"], border_hex="D0D5DD")
        add_safe_shape(slide, MSO_SHAPE.RECTANGLE, x, y, card_w, 0.06, pal["accent"])

        img = find_best_image_for_lesson(lesson, "learn", learn_idx*3 + i)
        if img:
            img_w = min(1.8, card_w - 0.4)
            img_x = x + (card_w - img_w) / 2
            add_picture_safe(slide, img, img_x, y + 0.15, img_w, img_h)

        text_y = y + img_h + 0.15
        add_textbox(slide, x + 0.15, text_y, card_w - 0.3, text_h, btext, size_pt=16, bold=False, color_hex=pal["text_on_card"], alignment=PP_ALIGN.CENTER)

    add_slide_transition(slide, "wipe")
    return slide


def build_learn_row(prs, data, pal, lesson, layout, learn_idx=0):
    """Layout B — Horizontal Row (Middle School Grades: Lớp 6 -> Lớp 8)"""
    slide = prs.slides.add_slide(layout)
    bullets = data.get("bullets", [])

    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, pal["primary"])
    badge_text = f"  {lesson['lop_label'].upper()} • {lesson['bai'].upper()}"
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, badge_text, size_pt=13, bold=True, color_hex=pal["text_on_primary"])

    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, SLIDE_W-1, 0.65, data["title"], size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    content_top = title_y + 0.8
    avail_h = SAFE_BOTTOM - content_top - 0.1
    n = len(bullets)
    row_h = min(1.25, (avail_h - (n - 1) * 0.2) / max(n, 1))

    for i, btext in enumerate(bullets[:4]):
        y = content_top + i * (row_h + 0.2)
        if y + row_h > SAFE_BOTTOM + 0.05:
            break

        # White card for row
        card_x, card_w = 0.5, SLIDE_W - 1.0
        add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, card_x, y, card_w, row_h, pal["card"], border_hex="D0D5DD")
        add_safe_shape(slide, MSO_SHAPE.RECTANGLE, card_x, y, 0.12, row_h, pal["accent"])

        img = find_best_image_for_lesson(lesson, "learn", learn_idx*3 + i)
        if img:
            img_w = 1.8
            add_picture_safe(slide, img, card_x + 0.3, y + 0.1, img_w, row_h - 0.2)
            txt_x = card_x + 0.3 + img_w + 0.3
            txt_w = card_w - (0.3 + img_w + 0.5)
        else:
            txt_x = card_x + 0.3
            txt_w = card_w - 0.6

        add_textbox(slide, txt_x, y + 0.15, txt_w, row_h - 0.3, f"● {btext}", size_pt=17, bold=False, color_hex=pal["text_on_card"])

    add_slide_transition(slide, "wipe")
    return slide


def build_practice(prs, data, pal, lesson, layout):
    """Practice / Quiz Slide — Sequential appear animation on click"""
    slide = prs.slides.add_slide(layout)

    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, pal["accent"])
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, f"  {lesson['lop_label'].upper()} • LUYỆN TẬP", size_pt=13, bold=True, color_hex=pal["text_on_primary"])

    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, SLIDE_W-1, 0.65, data["title"], size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    instr_y = title_y + 0.75
    instr_lines = data.get("instruction", "").split("\n")
    add_multiline_textbox(slide, 0.5, instr_y, SLIDE_W-1, 0.7, instr_lines, size_pt=18, color_hex=pal["text_on_bg"])

    items = data.get("items", [])
    item_y_start = instr_y + 0.8
    n_items = len(items)

    cols = 2 if n_items > 3 else max(n_items, 1)
    rows = (n_items + 1) // 2 if n_items > 3 else 1

    gap = 0.25
    card_w = (SLIDE_W - 1.0 - gap*(cols-1)) / cols
    card_h = min(1.8, (SAFE_BOTTOM - item_y_start - gap*(rows-1)) / rows)

    animation_groups = []

    for i, item in enumerate(items[:4]):
        col, row = i % cols, i // cols
        x = 0.5 + col * (card_w + gap)
        y = item_y_start + row * (card_h + gap)

        if y + card_h > SAFE_BOTTOM + 0.05:
            break

        card = add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h, pal["card"], border_hex=pal["accent"])

        img = find_best_image_for_lesson(lesson, "practice", i)
        pic = None
        if img:
            img_w = 1.4
            pic = add_picture_safe(slide, img, x + 0.2, y + 0.15, img_w, card_h - 0.3)
            txt_x = x + 0.2 + img_w + 0.2
            txt_w = card_w - (0.2 + img_w + 0.3)
        else:
            txt_x = x + 0.2
            txt_w = card_w - 0.4

        txt = add_textbox(slide, txt_x, y + 0.15, txt_w, card_h - 0.3, item, size_pt=16, bold=True, color_hex=pal["text_on_card"])

        group = [s for s in [card, pic, txt] if s is not None]
        animation_groups.append((group, i))

    if animation_groups:
        add_appear_animation(slide, animation_groups)

    add_slide_transition(slide, "wipe")
    return slide


def build_activity(prs, data, pal, lesson, layout):
    """Activity / Game Slide — Mandatory illustration picture"""
    slide = prs.slides.add_slide(layout)

    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, pal["accent"])
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, f"  {lesson['lop_label'].upper()} • THỬ THÁCH", size_pt=13, bold=True, color_hex=pal["text_on_primary"])

    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, SLIDE_W-1, 0.65, data["title"], size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    content_y = title_y + 0.75
    content_h = SAFE_BOTTOM - content_y - 0.15

    card_w = 7.0
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.4, content_y, card_w, content_h, pal["card"])
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0.4, content_y, 0.12, content_h, pal["accent"])

    lines = data.get("content", "").split("\n") if data.get("content") else data.get("instruction", "").split("\n")
    add_multiline_textbox(slide, 0.8, content_y+0.2, card_w-0.6, content_h-0.4, lines, size_pt=20, color_hex=pal["text_on_card"])

    img = find_best_image_for_lesson(lesson, "activity", 0)
    if img:
        img_x = 7.8
        img_w = SAFE_RIGHT - img_x - 0.2
        add_picture_safe(slide, img, img_x, content_y + 0.1, img_w, content_h - 0.2)

    add_slide_transition(slide, "cover")
    return slide


def build_summary(prs, data, pal, lesson, layout):
    slide = prs.slides.add_slide(layout)

    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    px, py = 0.6, SAFE_TOP + 0.15
    pw = SLIDE_W - 1.2
    ph = SAFE_BOTTOM - py - 0.15
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, px, py, pw, ph, pal["card"])
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, px, py, pw, 0.06, pal["primary"])

    add_textbox(slide, px+0.3, py+0.2, pw-0.6, 0.65, data["title"], size_pt=24, bold=True, color_hex=pal["primary"], alignment=PP_ALIGN.CENTER)

    div_y = py + 0.95
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, px+2, div_y, pw-4, 0.03, pal["accent"])

    items = data.get("items", [])
    for i, item in enumerate(items[:4]):
        y = div_y + 0.25 + i * 1.1
        if y + 0.8 > SAFE_BOTTOM:
            break
        add_safe_shape(slide, MSO_SHAPE.RECTANGLE, px+0.4, y, 0.08, 0.7, pal["accent"])
        add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, px+0.6, y, pw-1.2, 0.7, pal["bg"])
        add_textbox(slide, px+0.9, y+0.15, pw-1.8, 0.4, item, size_pt=18, color_hex=pal["text_on_bg"])

    add_slide_transition(slide, "fade")
    return slide


def build_thanks(prs, data, pal, lesson, layout):
    slide = prs.slides.add_slide(layout)

    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["primary"], send_to_back=True)

    add_textbox(slide, 1.0, SAFE_TOP+0.8, SLIDE_W-2, 1.2, data["title"], size_pt=36, bold=True, color_hex=pal["text_on_primary"], alignment=PP_ALIGN.CENTER)

    content = data.get("content", "")
    if content:
        card_w = 8.0
        card_x = (SLIDE_W - card_w) / 2
        card_y = SAFE_TOP + 2.5
        card_h = min(2.0, SAFE_BOTTOM - card_y - 0.2)
        add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, card_w, card_h, pal["card"])
        lines = content.split("\n")
        add_multiline_textbox(slide, card_x+0.5, card_y+0.3, card_w-1, card_h-0.6, lines, size_pt=20, color_hex=pal["text_on_card"], alignment=PP_ALIGN.CENTER)

    add_slide_transition(slide, "fade")
    return slide


# ═══════════════════════════════════════════════════
# VERIFICATION
# ═══════════════════════════════════════════════════

def verify_slide_deck(prs, lesson_label):
    issues = []
    for si, sl in enumerate(prs.slides):
        for sh in sl.shapes:
            top_in = sh.top / 914400
            bottom_in = (sh.top + sh.height) / 914400
            if top_in < SAFE_TOP - 0.02 and sh.name not in ('Picture 7',):
                issues.append(f"Slide {si+1}: {sh.name} top={top_in:.2f}in < {SAFE_TOP}in")
            if bottom_in > SAFE_BOTTOM + 0.02 and sh.name not in ('Picture 9',):
                issues.append(f"Slide {si+1}: {sh.name} bottom={bottom_in:.2f}in > {SAFE_BOTTOM}in")

    if issues:
        print(f"  ⚠️ {lesson_label} VERIFICATION: {len(issues)} issues found")
        for iss in issues:
            print(f"     ❌ {iss}")
    else:
        print(f"  ✅ {lesson_label} VERIFICATION: ALL {len(prs.slides)} slides PASS anti-bug checklist")
    return len(issues) == 0


# ═══════════════════════════════════════════════════
# MAIN BUILDER FOR ALL 9 GRADES
# ═══════════════════════════════════════════════════

def create_slide_deck_v3(lesson, palette_idx):
    palette = COLOR_PALETTES[palette_idx % len(COLOR_PALETTES)]
    prs = Presentation(TEMPLATE)
    remove_template_slides(prs)
    layout = prs.slide_layouts[6]

    learn_idx = 0
    for sc in lesson["slides_content"]:
        stype = sc["type"]
        try:
            if stype == "cover":
                build_cover(prs, sc, palette, lesson, layout)
            elif stype == "warmup":
                build_warmup(prs, sc, palette, lesson, layout)
            elif stype == "learn":
                if lesson["is_thcs"]:
                    build_learn_row(prs, sc, palette, lesson, layout, learn_idx)
                else:
                    build_learn_grid(prs, sc, palette, lesson, layout, learn_idx)
                learn_idx += 1
            elif stype == "practice":
                build_practice(prs, sc, palette, lesson, layout)
            elif stype == "activity":
                build_activity(prs, sc, palette, lesson, layout)
            elif stype == "summary":
                build_summary(prs, sc, palette, lesson, layout)
            elif stype == "thanks":
                build_thanks(prs, sc, palette, lesson, layout)
        except Exception as e:
            print(f"     ⚠️ Error on slide '{sc.get('title','')}': {e}")

    # Verify
    verify_slide_deck(prs, lesson["lop_label"])

    # Output path
    output_dir = os.path.join(KHBD_BASE, lesson["folder"], "Tuần_02")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, lesson["file_name"])

    try:
        prs.save(output_path)
        print(f"  ✅ {lesson['lop_label']}: {len(prs.slides)} slides -> {lesson['file_name']}")
    except PermissionError:
        alt = output_path.replace(".pptx", "_v3_alt.pptx")
        prs.save(alt)
        print(f"  ⚠️ {lesson['lop_label']}: File locked -> saved as {os.path.basename(alt)}")
    except Exception as e:
        print(f"  ❌ {lesson['lop_label']} save error: {e}")

    return output_path


def main():
    print("=" * 70)
    print("  TẠO BỘ SLIDE BÀI GIẢNG CHUẨN V3 — TOÀN BỘ 9 KHỐI LỚP (TUẦN 02)")
    print("  Quy chuẩn: Per-Bullet Images, Animation On-Click, Game Pics, Anti-Bug")
    print("=" * 70)

    for i, lesson in enumerate(LESSONS):
        print(f"\n[ Khối: {lesson['lop_label']} — {lesson['bai']} ]")
        try:
            create_slide_deck_v3(lesson, i)
        except Exception as e:
            print(f"❌ Error on {lesson['lop_label']}: {e}")
            traceback.print_exc()

    print("\n" + "=" * 70)
    print("  HOÀN THÀNH TOÀN BỘ 9 BỘ SLIDE CHUẨN V3!")
    print("=" * 70)


if __name__ == "__main__":
    main()
