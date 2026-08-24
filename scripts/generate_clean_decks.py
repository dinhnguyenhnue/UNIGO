import sys, os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

sys.path.append(r'D:\UNIGO')
from scripts.generate_slides_tuan04_lop5_lop7_expert import DECKS_CONFIG, BUILDERS, TEMPLATE

def clean_template_initial_slides():
    for config in DECKS_CONFIG:
        prs = Presentation(TEMPLATE)
        # Remove any initial slides that exist in the template
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
            
        blank_layout = prs.slide_layouts[6]
        pal = config["pal"]
        file_path = config["file_path"]
        
        print(f"Generating clean: {os.path.basename(file_path)}")
        for b_name, b_data in config["slides"]:
            if b_name in BUILDERS:
                BUILDERS[b_name](prs, b_data, pal, blank_layout)
                
        prs.save(file_path)
        print(f"  [DONE] Total slides: {len(prs.slides)}")

clean_template_initial_slides()
