"""
Tạo Slide Định hướng Tin học (Tiết 0) — Chuẩn mẫu UNIGO
V3: Font lớn + Group card+bar + Fix slide cuối không che logo.

VÙNG AN TOÀN: Y 1.15in → 6.30in
Logo master: (0.17, 0.15) 0.95×0.94in
Chân trang master: (0.00, 6.43) 13.40×1.23in
"""

import os, sys, random
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

sys.stdout.reconfigure(encoding='utf-8')

TPL_PATH = r'd:\UNIGO\Hệ thống mẫu văn bản\Mẫu slide UNIGO.pptx'
IMG_DIR = r'C:\Users\bmngu\.gemini\antigravity-ide\brain\af66fc76-d63a-4745-b893-b085a3c30f22'
BASE_OUT = r'd:\UNIGO\KHBD_Tin_học'

# ─── 8 Bộ màu hài hoà ───────────────────────────────────────────────
PALETTES = [
    { 'name': 'Ocean Tech',
        'bg': RGBColor(0xE0,0xF2,0xFE), 'card': RGBColor(0xFF,0xFF,0xFF),
        'primary': RGBColor(0x0C,0x4A,0x6E), 'second': RGBColor(0x06,0x96,0xC7),
        'accent': RGBColor(0xF9,0x73,0x16) },
    { 'name': 'Emerald Garden',
        'bg': RGBColor(0xEC,0xFD,0xF5), 'card': RGBColor(0xFF,0xFF,0xFF),
        'primary': RGBColor(0x06,0x5F,0x46), 'second': RGBColor(0x10,0xB9,0x81),
        'accent': RGBColor(0xF5,0x9E,0x0B) },
    { 'name': 'Royal Purple',
        'bg': RGBColor(0xF5,0xF3,0xFF), 'card': RGBColor(0xFF,0xFF,0xFF),
        'primary': RGBColor(0x5B,0x21,0xB6), 'second': RGBColor(0x8B,0x5C,0xF6),
        'accent': RGBColor(0x06,0xB6,0xD4) },
    { 'name': 'Sunset Coral',
        'bg': RGBColor(0xFF,0xF1,0xE6), 'card': RGBColor(0xFF,0xFF,0xFF),
        'primary': RGBColor(0x9A,0x34,0x12), 'second': RGBColor(0xEA,0x58,0x0C),
        'accent': RGBColor(0x02,0x84,0xC7) },
    { 'name': 'Teal Fresh',
        'bg': RGBColor(0xF0,0xFD,0xFA), 'card': RGBColor(0xFF,0xFF,0xFF),
        'primary': RGBColor(0x11,0x5E,0x59), 'second': RGBColor(0x14,0xB8,0xA6),
        'accent': RGBColor(0xE1,0x1D,0x48) },
    { 'name': 'Navy Elegant',
        'bg': RGBColor(0xEE,0xF2,0xFF), 'card': RGBColor(0xFF,0xFF,0xFF),
        'primary': RGBColor(0x1E,0x3A,0x8A), 'second': RGBColor(0x3B,0x82,0xF6),
        'accent': RGBColor(0xF4,0x3F,0x5E) },
    { 'name': 'Berry Academic',
        'bg': RGBColor(0xFD,0xF2,0xF8), 'card': RGBColor(0xFF,0xFF,0xFF),
        'primary': RGBColor(0x83,0x18,0x43), 'second': RGBColor(0xEC,0x48,0x99),
        'accent': RGBColor(0x0D,0x94,0x88) },
    { 'name': 'Slate Modern',
        'bg': RGBColor(0xF1,0xF5,0xF9), 'card': RGBColor(0xFF,0xFF,0xFF),
        'primary': RGBColor(0x1E,0x29,0x3B), 'second': RGBColor(0x47,0x55,0x69),
        'accent': RGBColor(0x7C,0x3A,0xED) },
]

C_WHITE = RGBColor(0xFF,0xFF,0xFF)
C_TEXT = RGBColor(0x1E,0x29,0x3B)
C_MUTED = RGBColor(0x64,0x74,0x8B)
C_BORDER = RGBColor(0xE2,0xE8,0xF0)

# ─── Font sizes chuẩn ────────────────────────────────────────────────
SZ_TITLE    = Pt(26)
SZ_TITLE_SM = Pt(24)
SZ_BODY     = Pt(18)
SZ_BODY_SM  = Pt(16)
SZ_BULLET   = Pt(14)
SZ_BADGE    = Pt(14)
SZ_INTRO_T  = Pt(28)
SZ_SUB      = Pt(20)
SZ_INFO     = Pt(16)
LINE_SP     = Pt(28)
LINE_SP_SM  = Pt(24)
SP_AFTER    = Pt(8)

IMG = {
    'classroom': os.path.join(IMG_DIR, 'classroom_computer_lab_1785861919471.png'),
    'devices':   os.path.join(IMG_DIR, 'digital_devices_around_1785861929327.png'),
    'safety':    os.path.join(IMG_DIR, 'internet_safety_kids_1785861941729.png'),
    'goals':     os.path.join(IMG_DIR, 'student_learning_goals_1785861961479.png'),
    'rules':     os.path.join(IMG_DIR, 'computer_room_rules_1785861971149.png'),
    'mindmap':   os.path.join(IMG_DIR, 'tin_hoc_mindmap_1785861980454.png'),
}

# ─── Helpers ─────────────────────────────────────────────────────────
def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def _add_rounded_card_raw(slide, left, top, width, height, fill_color=C_WHITE, border_color=C_BORDER):
    """Tạo card bo góc (trả về shape để group)."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(0.75)
    card.adjustments[0] = 0.02
    return card

def _add_accent_bar_raw(slide, left, top, height, color, width=Inches(0.15)):
    """Tạo thanh accent dọc (trả về shape để group)."""
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.adjustments[0] = 0.5
    return bar

def group_shapes(slide, shapes):
    """Nhóm (group) các shape lại với nhau trên slide."""
    spTree = slide.shapes._spTree

    # Bounding box
    min_left   = min(s.left for s in shapes)
    min_top    = min(s.top for s in shapes)
    max_right  = max(s.left + s.width for s in shapes)
    max_bottom = max(s.top + s.height for s in shapes)
    grp_w = max_right - min_left
    grp_h = max_bottom - min_top

    # Next unique ID
    used_ids = set()
    for elem in spTree.iter():
        id_val = elem.get('id')
        if id_val and id_val.isdigit():
            used_ids.add(int(id_val))
    next_id = (max(used_ids) + 1) if used_ids else 100

    grpSp_xml = (
        f'<p:grpSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f'         xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        f'         xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        f'  <p:nvGrpSpPr>'
        f'    <p:cNvPr id="{next_id}" name="Card Group {next_id}"/>'
        f'    <p:cNvGrpSpPr/>'
        f'    <p:nvPr/>'
        f'  </p:nvGrpSpPr>'
        f'  <p:grpSpPr>'
        f'    <a:xfrm>'
        f'      <a:off x="{min_left}" y="{min_top}"/>'
        f'      <a:ext cx="{grp_w}" cy="{grp_h}"/>'
        f'      <a:chOff x="{min_left}" y="{min_top}"/>'
        f'      <a:chExt cx="{grp_w}" cy="{grp_h}"/>'
        f'    </a:xfrm>'
        f'  </p:grpSpPr>'
        f'</p:grpSp>'
    )
    grpSp = parse_xml(grpSp_xml)

    # Di chuyển các shape elements vào group
    for shape in shapes:
        elem = shape._element
        spTree.remove(elem)
        grpSp.append(elem)

    spTree.append(grpSp)

def add_card_group(slide, left, top, width, height,
                   card_color=C_WHITE, bar_color=None, bar_width=Inches(0.15)):
    """Tạo card + accent bar đã group lại.
    Thanh accent bar khít hoàn toàn chiều cao card (top = card.top, height = card.height).
    """
    card = _add_rounded_card_raw(slide, left, top, width, height, card_color)
    shapes_to_group = [card]

    if bar_color is not None:
        # Bar khít đúng chiều cao card — KHÔNG có margin
        bar = _add_accent_bar_raw(slide, left, top, height, bar_color, bar_width)
        shapes_to_group.append(bar)

    group_shapes(slide, shapes_to_group)
    return card

def make_tf(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf

def add_run(para, text, font_name="Times New Roman", size=SZ_BODY,
            bold=False, italic=False, color=C_TEXT):
    r = para.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = size
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return r

def title_paragraph(tf, text, color, size=SZ_TITLE):
    p = tf.paragraphs[0]
    p.space_after = Pt(12)
    p.line_spacing = Pt(32)
    add_run(p, text, bold=True, size=size, color=color)
    return p

def bullet_items(tf, items, color=C_TEXT, bullet_color=None,
                 size=SZ_BODY, line_spacing=LINE_SP):
    for item in items:
        p = tf.add_paragraph()
        p.space_after = SP_AFTER
        p.line_spacing = line_spacing
        bc = bullet_color or color
        add_run(p, "●  ", size=SZ_BULLET, color=bc, bold=True)
        add_run(p, item, size=size, color=color)

def clear_template_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(qn('r:id'))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

def new_slide(prs, layout_idx=6):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


# ═══════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS (V3: font lớn, group card+bar, slide cuối fix logo)
# ═══════════════════════════════════════════════════════════════════════

def build_slide_intro(prs, pal, grade, subtitle):
    """Slide 1: Giới thiệu — ảnh phải, card trái."""
    s = new_slide(prs)
    set_slide_bg(s, pal['bg'])

    if os.path.exists(IMG['classroom']):
        s.shapes.add_picture(IMG['classroom'], Inches(7.0), Inches(1.3), Inches(5.6), Inches(4.6))

    add_card_group(s, Inches(0.6), Inches(1.3), Inches(6.0), Inches(4.6),
                   pal['card'], pal['accent'])

    # Badge
    badge = _add_rounded_card_raw(s, Inches(1.1), Inches(1.5), Inches(4.5), Inches(0.5),
                                  pal['primary'], pal['primary'])
    tf_b = badge.text_frame
    tf_b.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_run(tf_b.paragraphs[0], f"TIN HỌC {grade}  •  NĂM HỌC 2026 – 2027",
            size=SZ_BADGE, bold=True, color=C_WHITE)

    tf = make_tf(s, Inches(1.1), Inches(2.2), Inches(5.2), Inches(2.0))
    add_run(tf.paragraphs[0], "Tiết 0 — Định hướng", size=SZ_INTRO_T, bold=True, color=pal['primary'])
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    add_run(p2, subtitle, size=SZ_SUB, bold=True, color=pal['accent'])

    tf2 = make_tf(s, Inches(1.1), Inches(4.4), Inches(5.2), Inches(1.2))
    add_run(tf2.paragraphs[0], "Trường TH & THCS UNIGO", size=SZ_INFO, color=pal['second'], bold=True)
    p3 = tf2.add_paragraph()
    p3.space_before = Pt(4)
    add_run(p3, "Giáo viên: Đậu Đình Nguyên", size=Pt(15), color=C_MUTED, italic=True)


def build_slide_objectives(prs, pal, grade):
    """Slide 2: Mục tiêu — ảnh trái, card phải."""
    s = new_slide(prs)
    set_slide_bg(s, pal['bg'])

    if os.path.exists(IMG['goals']):
        s.shapes.add_picture(IMG['goals'], Inches(0.4), Inches(1.4), Inches(4.2), Inches(4.6))

    add_card_group(s, Inches(5.0), Inches(1.3), Inches(7.6), Inches(4.8),
                   pal['card'], pal['primary'])

    tf = make_tf(s, Inches(5.4), Inches(1.5), Inches(6.9), Inches(4.4))
    title_paragraph(tf, "🎯  Mục tiêu tiết học", pal['primary'], SZ_TITLE)
    bullet_items(tf, [
        f"Tổng quan chương trình Tin học {grade}.",
        "Nội quy an toàn phòng máy tính UNIGO.",
        "Ứng xử văn minh trên môi trường số (Digital Citizenship).",
    ], bullet_color=pal['accent'])


def build_slide_warmup(prs, pal, grade):
    """Slide 3: Khởi động — card trái, ảnh phải."""
    s = new_slide(prs)
    set_slide_bg(s, pal['bg'])

    if os.path.exists(IMG['devices']):
        s.shapes.add_picture(IMG['devices'], Inches(7.4), Inches(1.3), Inches(5.3), Inches(4.8))

    add_card_group(s, Inches(0.6), Inches(1.3), Inches(6.4), Inches(4.8),
                   pal['card'], pal['second'])

    tf = make_tf(s, Inches(1.0), Inches(1.5), Inches(5.7), Inches(4.4))
    title_paragraph(tf, "🎮  Thế giới số quanh em", pal['second'], SZ_TITLE)
    bullet_items(tf, [
        "Thiết bị nào quanh em? (máy tính, tablet...)",
        "Em đã dùng máy tính để làm gì?",
        f"Tin học {grade}: \"người dùng\" → \"nhà sáng tạo số\"!",
    ], bullet_color=pal['second'])


def build_slide_overview(prs, pal, grade):
    """Slide 4: Tổng quan — ảnh trái, card phải."""
    s = new_slide(prs)
    set_slide_bg(s, pal['bg'])

    if os.path.exists(IMG['mindmap']):
        s.shapes.add_picture(IMG['mindmap'], Inches(0.4), Inches(1.5), Inches(4.5), Inches(4.4))

    add_card_group(s, Inches(5.2), Inches(1.3), Inches(7.4), Inches(4.8),
                   pal['card'], pal['accent'])

    if grade <= 5:
        topics = [
            "Máy tính & Em (Computer & Me)",
            "Internet & Mạng (Internet Basics)",
            "Lưu trữ dữ liệu (Data Organization)",
            "Đạo đức số (Digital Ethics)",
            "Ứng dụng Tin học (IT Apps)",
            "Giải quyết vấn đề & Lập trình (Coding)",
        ]
    else:
        topics = [
            "Máy tính & Xã hội tri thức",
            "Mạng & Internet (Networking)",
            "Đạo đức & Pháp luật số",
            "Ứng dụng Tin học (IT Apps)",
            "Giải thuật & Lập trình (Algorithms)",
        ]

    tf = make_tf(s, Inches(5.6), Inches(1.5), Inches(6.7), Inches(4.4))
    title_paragraph(tf, f"📚  Chương trình Tin học {grade}", pal['accent'], SZ_TITLE_SM)
    bullet_items(tf, topics, bullet_color=pal['primary'], size=SZ_BODY_SM, line_spacing=LINE_SP_SM)


def build_slide_methods(prs, pal, grade):
    """Slide 5: 2 card — Phương pháp + Tư duy."""
    s = new_slide(prs)
    set_slide_bg(s, pal['bg'])

    # Card trái (grouped)
    add_card_group(s, Inches(0.5), Inches(1.3), Inches(5.9), Inches(4.8),
                   pal['card'], pal['primary'])
    tf1 = make_tf(s, Inches(0.9), Inches(1.5), Inches(5.2), Inches(4.4))
    title_paragraph(tf1, "🚀  Cách học hiệu quả", pal['primary'], SZ_TITLE_SM)
    bullet_items(tf1, [
        "Trải nghiệm (Hands-on)",
        "Dự án (Project-based)",
        "Tìm lỗi (Debugging)",
    ], bullet_color=pal['primary'])

    # Card phải (grouped)
    add_card_group(s, Inches(6.8), Inches(1.3), Inches(5.9), Inches(4.8),
                   pal['card'], pal['accent'])
    tf2 = make_tf(s, Inches(7.2), Inches(1.5), Inches(5.2), Inches(4.4))
    title_paragraph(tf2, "🧠  Tư duy máy tính", pal['accent'], SZ_TITLE_SM)
    bullet_items(tf2, [
        "Phân rã (Decomposition)",
        "Nhận mẫu (Pattern)",
        "Thuật toán (Algorithm)",
    ], bullet_color=pal['accent'])


def build_slide_rules(prs, pal, grade):
    """Slide 6: Nội quy — card trái, ảnh phải."""
    s = new_slide(prs)
    set_slide_bg(s, pal['bg'])

    if os.path.exists(IMG['rules']):
        s.shapes.add_picture(IMG['rules'], Inches(7.2), Inches(1.3), Inches(5.4), Inches(4.8))

    C_WARN = RGBColor(0xEA, 0x58, 0x0C)
    add_card_group(s, Inches(0.5), Inches(1.3), Inches(6.3), Inches(4.8),
                   pal['card'], C_WARN)

    tf = make_tf(s, Inches(0.9), Inches(1.5), Inches(5.6), Inches(4.4))
    title_paragraph(tf, "⚠️  Nội quy phòng máy", C_WARN, SZ_TITLE)
    bullet_items(tf, [
        "Xếp hàng, để giày dép đúng nơi.",
        "KHÔNG mang đồ ăn, nước uống vào.",
        "Bật/tắt máy đúng quy trình.",
        "Gặp sự cố → BÁO NGAY giáo viên.",
    ], bullet_color=C_WARN)


def build_slide_safety(prs, pal, grade):
    """Slide 7: An toàn số — ảnh trái, card phải."""
    s = new_slide(prs)
    set_slide_bg(s, pal['bg'])

    if os.path.exists(IMG['safety']):
        s.shapes.add_picture(IMG['safety'], Inches(0.4), Inches(1.4), Inches(4.3), Inches(4.5))

    add_card_group(s, Inches(5.0), Inches(1.3), Inches(7.6), Inches(4.8),
                   pal['card'], pal['second'])

    tf = make_tf(s, Inches(5.4), Inches(1.5), Inches(6.9), Inches(4.4))
    title_paragraph(tf, "🛡️  An toàn số (Digital Safety)", pal['second'], SZ_TITLE_SM)
    bullet_items(tf, [
        "Giữ bí mật mật khẩu (password).",
        "Không chia sẻ thông tin cá nhân.",
        "Cảnh giác link lạ, file lạ.",
        "Ứng xử văn minh, không bắt nạt mạng.",
    ], bullet_color=pal['second'])


def build_slide_assessment(prs, pal, grade):
    """Slide 8: 2 card — Đánh giá + Đồ dùng."""
    s = new_slide(prs)
    set_slide_bg(s, pal['bg'])

    add_card_group(s, Inches(0.5), Inches(1.3), Inches(5.9), Inches(4.8),
                   pal['card'], pal['primary'])
    tf1 = make_tf(s, Inches(0.9), Inches(1.5), Inches(5.2), Inches(4.4))
    title_paragraph(tf1, "📊  Cách đánh giá", pal['primary'], SZ_TITLE_SM)
    bullet_items(tf1, [
        "Đánh giá thường xuyên (ĐGTX).",
        "Đánh giá định kỳ (4 lần/năm).",
        "Sản phẩm dự án nhóm.",
    ], bullet_color=pal['primary'])

    add_card_group(s, Inches(6.8), Inches(1.3), Inches(5.9), Inches(4.8),
                   pal['card'], pal['second'])
    tf2 = make_tf(s, Inches(7.2), Inches(1.5), Inches(5.2), Inches(4.4))
    title_paragraph(tf2, "📚  Đồ dùng cần có", pal['second'], SZ_TITLE_SM)
    bullet_items(tf2, [
        f"SGK Tin học {grade}.",
        "Vở ghi & sổ tay cá nhân.",
        "USB / thẻ nhớ (THCS).",
    ], bullet_color=pal['second'])


def build_slide_homework(prs, pal, grade):
    """Slide 9: Nhiệm vụ — full width card."""
    s = new_slide(prs)
    set_slide_bg(s, pal['bg'])

    add_card_group(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(4.8),
                   pal['card'], pal['accent'])

    tf = make_tf(s, Inches(1.0), Inches(1.5), Inches(11.5), Inches(4.4))
    title_paragraph(tf, "📝  Nhiệm vụ mở rộng (Homework)", pal['accent'], SZ_TITLE)
    bullet_items(tf, [
        "Viết 2 mục tiêu em muốn đạt trong năm học.",
        "Đề xuất 1 sản phẩm số muốn tự tạo.",
        "Cam kết thực hiện nội quy phòng máy.",
        "Về nhà: Xem trước Bài 1 trong SGK.",
    ], bullet_color=pal['accent'])


def build_slide_summary(prs, pal, grade):
    """Slide 10: Tổng kết — panel màu trong VÙNG AN TOÀN, không che logo/chân."""
    s = new_slide(prs)
    set_slide_bg(s, pal['bg'])  # Nền nhạt — logo + chân trang luôn thấy rõ

    # Panel màu primary CHỈ trong vùng an toàn (dưới logo, trên chân trang)
    panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0), Inches(1.15),
                               Inches(13.33), Inches(5.15))
    panel.fill.solid()
    panel.fill.fore_color.rgb = pal['primary']
    panel.line.fill.background()

    # Card + bar (grouped)
    add_card_group(s, Inches(1.2), Inches(1.5), Inches(10.9), Inches(4.5),
                   pal['card'], pal['accent'], bar_width=Inches(0.18))

    tf = make_tf(s, Inches(1.7), Inches(1.7), Inches(10.0), Inches(4.0))
    title_paragraph(tf, f"📌  Ghi nhớ — Tin học {grade}", pal['primary'], SZ_TITLE)
    bullet_items(tf, [
        "Tin học mang đến kỹ năng số thiết thực.",
        "Tuân thủ nội quy & ứng xử văn minh.",
        "Tự tin khám phá, sáng tạo!",
    ], bullet_color=pal['accent'], size=SZ_SUB)

    p_thanks = tf.add_paragraph()
    p_thanks.space_before = Pt(20)
    p_thanks.alignment = PP_ALIGN.CENTER
    add_run(p_thanks, "Chúc các em năm học tuyệt vời! 🎉",
            size=Pt(22), bold=True, color=pal['second'])


# ─── Quality Check ───────────────────────────────────────────────────
def quality_check(filepath):
    prs = Presentation(filepath)
    issues = []
    if len(prs.slides) < 8:
        issues.append(f"WARN: Chi co {len(prs.slides)} slides")
    master = prs.slide_masters[0]
    has_logo = any(s.name == 'Picture 7' and s.shape_type == 13 for s in master.shapes)
    has_footer = any(s.name == 'Picture 9' and s.shape_type == 13 for s in master.shapes)
    if not has_logo:  issues.append("FAIL: Logo UNIGO bi mat!")
    if not has_footer: issues.append("FAIL: Chan trang bi mat!")
    # Check slide cuoi khong co bg che logo (shape at y=0)
    last = prs.slides[len(prs.slides) - 1]
    for shape in last.shapes:
        if Emu(shape.top).inches < 0.5 and Emu(shape.width).inches > 10:
            issues.append(f"WARN: Slide cuoi shape '{shape.name}' co the che logo (top={Emu(shape.top).inches:.2f}in)")
    return issues


# ─── Main ────────────────────────────────────────────────────────────
def build_deck(grade, palette_idx=None):
    prs = Presentation(TPL_PATH)
    clear_template_slides(prs)

    if palette_idx is None:
        palette_idx = (grade - 1) % len(PALETTES)
    pal = PALETTES[palette_idx]
    print(f"  Palette: {pal['name']} (#{palette_idx + 1})")

    subtitles = {
        1: "Noi quy & An toan phong may",
        2: "Nha sang tao so",
        3: "Kham pha Tin hoc 3",
        4: "Kham pha Tin hoc 4",
        5: "Ky nang so the ky 21",
        6: "Phuong phap & An toan so",
        7: "Tong quan & Ky nang so",
        8: "Tu duy may tinh",
    }

    build_slide_intro(prs, pal, grade, subtitles.get(grade, f"Kham pha Tin hoc {grade}"))
    build_slide_objectives(prs, pal, grade)
    build_slide_warmup(prs, pal, grade)
    build_slide_overview(prs, pal, grade)
    build_slide_methods(prs, pal, grade)
    build_slide_rules(prs, pal, grade)
    build_slide_safety(prs, pal, grade)
    build_slide_assessment(prs, pal, grade)
    build_slide_homework(prs, pal, grade)
    build_slide_summary(prs, pal, grade)

    return prs


def main():
    print("=== TẠO TOÀN BỘ SLIDE TIẾT 0 TIN HỌC (LỚP 1 ĐẾN LỚP 8) ===")
    total_files = 0
    all_passed = True

    for grade in range(1, 9):
        print(f"\n--- Đang xử lý Lớp {grade} ---")
        prs = build_deck(grade)

        out_folder = os.path.join(BASE_OUT, f"Lớp_{grade}", "Tiết_00")
        os.makedirs(out_folder, exist_ok=True)
        out_path = os.path.join(out_folder, f"Slide_Tin_hoc_Lop_{grade}_Tiet00_Dinh_huong.pptx")

        prs.save(out_path)
        total_files += 1
        print(f"  ✓ Đã lưu: {out_path}")
        print(f"  ✓ Số slide: {len(prs.slides)}")

        issues = quality_check(out_path)
        if issues:
            all_passed = False
            print("  ⚠ QUALITY ISSUES:")
            for iss in issues:
                print(f"    - {iss}")
        else:
            print("  ✅ QUALITY CHECK PASSED!")

    print(f"\n=== HOÀN THÀNH: Đã tạo {total_files}/8 bộ slide Tiết 0 thành công! ===")
    if all_passed:
        print("🎉 100% các bộ slide đạt chuẩn chất lượng UNIGO!")


if __name__ == '__main__':
    main()
