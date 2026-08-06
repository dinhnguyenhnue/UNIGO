"""
Tạo Slide Định hướng Robotics (Tiết 0) cho Lớp 1 - Lớp 8 — Chuẩn mẫu UNIGO
VÙNG AN TOÀN: Y 1.15in → 6.30in
Logo master: (0.17, 0.15) 0.95×0.94in
Chân trang master: (0.00, 6.43) 13.40×1.23in
"""

import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

sys.stdout.reconfigure(encoding='utf-8')

TPL_PATH = r'd:\UNIGO\Hệ thống mẫu văn bản\Mẫu slide UNIGO.pptx'
BASE_OUT = r'd:\UNIGO\KHBD_Robotics'
IMG_RULES_DIR = r'd:\UNIGO\KHBD_Robotics\images\noi_quy'

IMG_RULES = {
    'kit_ollo':     os.path.join(IMG_RULES_DIR, 'rules_kit_ollo.png'),
    'during_class': os.path.join(IMG_RULES_DIR, 'rules_during_class.png'),
    'after_class':  os.path.join(IMG_RULES_DIR, 'rules_after_class.png'),
}

# ─── 8 Bộ màu hài hoà cho Robotics ──────────────────────────────
PALETTES = [
    { 'name': 'Robo Blue',
      'bg': RGBColor(0xE0, 0xF2, 0xFE), 'card': RGBColor(0xFF, 0xFF, 0xFF),
      'primary': RGBColor(0x0C, 0x4A, 0x6E), 'second': RGBColor(0x02, 0x84, 0xC7),
      'accent': RGBColor(0xEA, 0x58, 0x0C) },
    { 'name': 'Robo Amber',
      'bg': RGBColor(0xFF, 0xFB, 0xEB), 'card': RGBColor(0xFF, 0xFF, 0xFF),
      'primary': RGBColor(0x78, 0x35, 0x0F), 'second': RGBColor(0xD9, 0x77, 0x06),
      'accent': RGBColor(0x25, 0x63, 0xEB) },
    { 'name': 'Robo Purple',
      'bg': RGBColor(0xF5, 0xF3, 0xFF), 'card': RGBColor(0xFF, 0xFF, 0xFF),
      'primary': RGBColor(0x5B, 0x21, 0xB6), 'second': RGBColor(0x7C, 0x3A, 0xED),
      'accent': RGBColor(0x05, 0x96, 0x69) },
    { 'name': 'Robo Teal',
      'bg': RGBColor(0xF0, 0xFD, 0xFA), 'card': RGBColor(0xFF, 0xFF, 0xFF),
      'primary': RGBColor(0x13, 0x4E, 0x4A), 'second': RGBColor(0x0D, 0x94, 0x88),
      'accent': RGBColor(0xDC, 0x26, 0x26) },
    { 'name': 'Robo Emerald',
      'bg': RGBColor(0xEC, 0xFD, 0xF5), 'card': RGBColor(0xFF, 0xFF, 0xFF),
      'primary': RGBColor(0x06, 0x4E, 0x3B), 'second': RGBColor(0x05, 0x96, 0x69),
      'accent': RGBColor(0xD9, 0x77, 0x06) },
    { 'name': 'Robo Crimson',
      'bg': RGBColor(0xFF, 0xF1, 0xF2), 'card': RGBColor(0xFF, 0xFF, 0xFF),
      'primary': RGBColor(0x88, 0x13, 0x37), 'second': RGBColor(0xE1, 0x1D, 0x48),
      'accent': RGBColor(0x02, 0x84, 0xC7) },
    { 'name': 'Robo Indigo',
      'bg': RGBColor(0xEE, 0xF2, 0xFF), 'card': RGBColor(0xFF, 0xFF, 0xFF),
      'primary': RGBColor(0x31, 0x2E, 0x81), 'second': RGBColor(0x4F, 0x46, 0xE5),
      'accent': RGBColor(0x05, 0x96, 0x69) },
    { 'name': 'Robo Orange',
      'bg': RGBColor(0xFF, 0xF7, 0xED), 'card': RGBColor(0xFF, 0xFF, 0xFF),
      'primary': RGBColor(0x7C, 0x2D, 0x12), 'second': RGBColor(0xEA, 0x58, 0x0C),
      'accent': RGBColor(0x02, 0x84, 0xC7) },
]

FONT_TITLE = 'Times New Roman'
FONT_BODY = 'Times New Roman'

SZ_TITLE = Pt(24)
SZ_TITLE_SM = Pt(22)
SZ_SUB = Pt(18)
SZ_BODY = Pt(18)

C_WARN = RGBColor(0xB9, 0x1C, 0x1C)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def clear_template_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]


def new_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def add_card_group(slide, left, top, width, height, card_color, bar_color, bar_width=Inches(0.12)):
    grp = slide.shapes.add_group_shape()

    card = grp.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = card_color
    card.line.fill.background()

    bar = grp.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, bar_width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = bar_color
    bar.line.fill.background()

    return grp


def make_tf(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    return tf


def add_run(p, text, font_name=FONT_BODY, size=SZ_BODY, bold=False, italic=False, color=None):
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = size
    r.bold = bold
    r.italic = italic
    if color:
        r.font.color.rgb = color
    return r


def title_paragraph(tf, text, color, size=SZ_TITLE):
    p = tf.paragraphs[0] if len(tf.paragraphs) > 0 else tf.add_paragraph()
    p.space_after = Pt(12)
    p.alignment = PP_ALIGN.LEFT
    add_run(p, text, font_name=FONT_TITLE, size=size, bold=True, color=color)


def bullet_items(tf, items, bullet_color=None, size=SZ_BODY):
    for idx, item in enumerate(items):
        p = tf.add_paragraph()
        p.space_after = Pt(8)
        p.line_spacing = Pt(28)
        p.alignment = PP_ALIGN.LEFT
        add_run(p, "● ", size=Pt(14), color=bullet_color)
        add_run(p, item, size=size, color=RGBColor(0x1E, 0x29, 0x3B))


# ─── SLIDES BUILDERS FOR ROBOTICS ──────────────────────────────

def build_slide_intro(prs, pal, grade, subtitle_text):
    s = new_slide(prs)
    set_slide_bg(s, pal['bg'])

    add_card_group(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(4.8),
                   pal['card'], pal['primary'], bar_width=Inches(0.18))

    tf = make_tf(s, Inches(1.2), Inches(1.6), Inches(11.0), Inches(4.2))

    p0 = tf.paragraphs[0]
    add_run(p0, "UNIGO COLLEGE", font_name=FONT_TITLE, size=Pt(14), bold=True, color=pal['second'])

    p1 = tf.add_paragraph()
    p1.space_before = Pt(8)
    p1.space_after = Pt(10)
    add_run(p1, f"ROBOTICS {grade}", font_name=FONT_TITLE, size=Pt(36), bold=True, color=pal['primary'])

    p2 = tf.add_paragraph()
    p2.space_after = Pt(16)
    add_run(p2, "Tiết 0: Định hướng môn học & Nội quy phòng thực hành",
            font_name=FONT_TITLE, size=Pt(22), bold=True, color=pal['second'])

    p3 = tf.add_paragraph()
    add_run(p3, f"💡 {subtitle_text}", font_name=FONT_BODY, size=Pt(18), italic=True, color=RGBColor(0x47, 0x55, 0x69))


def build_slide_objectives(prs, pal, grade):
    s = new_slide(prs)
    set_slide_bg(s, pal['bg'])

    add_card_group(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(4.8),
                   pal['card'], pal['second'])

    tf = make_tf(s, Inches(1.0), Inches(1.5), Inches(11.5), Inches(4.4))
    title_paragraph(tf, "🎯  Mục tiêu tiết học", pal['second'], SZ_TITLE)
    bullet_items(tf, [
        "Hiểu mục tiêu & vai trò của môn học Robotics.",
        "Nắm rõ phương pháp học tập trải nghiệm & làm việc nhóm.",
        "Thành thạo nội quy bảo quản bộ Kit Robotics & an toàn phòng lab.",
        "Tạo niềm hứng khởi sáng tạo công nghệ trong năm học mới.",
    ], bullet_color=pal['second'])


def build_slide_warmup(prs, pal, grade):
    s = new_slide(prs)
    set_slide_bg(s, pal['bg'])

    add_card_group(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(4.8),
                   pal['card'], pal['accent'])

    tf = make_tf(s, Inches(1.0), Inches(1.5), Inches(11.5), Inches(4.4))
    title_paragraph(tf, "🧠  Khởi động — Robot xung quanh ta!", pal['accent'], SZ_TITLE)
    bullet_items(tf, [
        "Câu hỏi: Theo em, Robot là gì? Em đã gặp robot nào trong thực tế?",
        "Chia sẻ: Tay máy công nghiệp, Robot hút bụi, Robot tự hành...",
        "Cùng khám phá: Làm thế nào để tự tay chế tạo một Robot?",
    ], bullet_color=pal['accent'])


def build_slide_overview(prs, pal, grade):
    s = new_slide(prs)
    set_slide_bg(s, pal['bg'])

    add_card_group(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(4.8),
                   pal['card'], pal['primary'])

    tf = make_tf(s, Inches(1.0), Inches(1.5), Inches(11.5), Inches(4.4))
    title_paragraph(tf, f"📚  Tổng quan chương trình Robotics Lớp {grade}", pal['primary'], SZ_TITLE)
    bullet_items(tf, [
        "Cơ cấu cơ khí: Động cơ, bánh răng, cơ chế truyền động.",
        "Cảm biến & Bộ điều khiển: Nhận biết âm thanh, khoảng cách, màu sắc.",
        "Lập trình điều khiển: Thuật toán điều khiển Robot thông minh.",
        "Dự án sáng tạo & Triển lãm Robotics UNIGO cuối năm.",
    ], bullet_color=pal['primary'])


def build_slide_methods(prs, pal, grade):
    s = new_slide(prs)
    set_slide_bg(s, pal['bg'])

    add_card_group(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(4.8),
                   pal['card'], pal['second'])

    tf = make_tf(s, Inches(1.0), Inches(1.5), Inches(11.5), Inches(4.4))
    title_paragraph(tf, "🛠️  Phương pháp học tập Robotics", pal['second'], SZ_TITLE)
    bullet_items(tf, [
        "Học qua trải nghiệm (Experiential Learning): Lắp ráp & Lập trình.",
        "Học theo nhóm (Teamwork): Phân công vai trò kỹ sư & lập trình viên.",
        "Giải quyết vấn đề (Problem Solving): Thử nghiệm & cải tiến mô hình.",
    ], bullet_color=pal['second'])


def build_slide_rules_1(prs, pal, grade):
    """Slide 6a: Quy định Kit OLLO — Quản lý linh kiện, động cơ, pin."""
    s = new_slide(prs)
    set_slide_bg(s, pal['bg'])

    if os.path.exists(IMG_RULES['kit_ollo']):
        s.shapes.add_picture(IMG_RULES['kit_ollo'], Inches(7.2), Inches(1.3), Inches(5.4), Inches(4.8))

    add_card_group(s, Inches(0.5), Inches(1.3), Inches(6.3), Inches(4.8),
                   pal['card'], C_WARN)

    tf = make_tf(s, Inches(1.0), Inches(1.5), Inches(5.6), Inches(4.4))
    title_paragraph(tf, "⚠️  II. Quy định Kit OLLO", C_WARN, Pt(22))
    bullet_items(tf, [
        "Dùng dụng cụ tháo rivet chuyên dụng; không dùng răng/móng tay cạy.",
        "Làm việc trên khay/hộp; không để linh kiện rơi xuống sàn.",
        "Cắm/rút dây cáp nhẹ nhàng, thẳng góc vào đúng cổng.",
        "Tắt nguồn khi lắp ráp; báo ngay nếu pin nóng/phồng bất thường.",
    ], bullet_color=C_WARN, size=Pt(16))


def build_slide_rules_2(prs, pal, grade):
    """Slide 6b: Nội quy trong giờ học & Lập trình."""
    s = new_slide(prs)
    set_slide_bg(s, pal['bg'])

    if os.path.exists(IMG_RULES['during_class']):
        s.shapes.add_picture(IMG_RULES['during_class'], Inches(7.2), Inches(1.3), Inches(5.4), Inches(4.8))

    add_card_group(s, Inches(0.5), Inches(1.3), Inches(6.3), Inches(4.8),
                   pal['card'], C_WARN)

    tf = make_tf(s, Inches(1.0), Inches(1.5), Inches(5.6), Inches(4.4))
    title_paragraph(tf, "⚠️  III. Trong giờ học", C_WARN, Pt(22))
    bullet_items(tf, [
        "Ngồi đúng vị trí nhóm/cá nhân được phân công.",
        "Thử nghiệm Robot trên thảm/bàn chuyên dụng, không cho chạy dưới sàn.",
        "Chỉ dùng phần mềm được chỉ định (R+ Task, ROBOTIS BLOCK...).",
    ], bullet_color=C_WARN)


def build_slide_rules_3(prs, pal, grade):
    """Slide 6c: Nội quy sau giờ thực hành."""
    s = new_slide(prs)
    set_slide_bg(s, pal['bg'])

    if os.path.exists(IMG_RULES['after_class']):
        s.shapes.add_picture(IMG_RULES['after_class'], Inches(7.2), Inches(1.3), Inches(5.4), Inches(4.8))

    add_card_group(s, Inches(0.5), Inches(1.3), Inches(6.3), Inches(4.8),
                   pal['card'], C_WARN)

    tf = make_tf(s, Inches(1.0), Inches(1.5), Inches(5.6), Inches(4.4))
    title_paragraph(tf, "⚠️  IV. Sau thực hành", C_WARN, Pt(22))
    bullet_items(tf, [
        "Tháo dỡ Robot nhẹ nhàng; phân loại linh kiện về đúng ngăn hộp.",
        "Tắt nguồn Robot, tháo dây sạc/kết nối.",
        "Kiểm tra đủ số lượng linh kiện cùng giáo viên trước khi rời.",
        "Vệ sinh mặt bàn, thu dọn rác và xếp ghế gọn gàng.",
    ], bullet_color=C_WARN, size=Pt(16))


def build_slide_assessment(prs, pal, grade):
    s = new_slide(prs)
    set_slide_bg(s, pal['bg'])

    add_card_group(s, Inches(0.5), Inches(1.3), Inches(5.9), Inches(4.8),
                   pal['card'], pal['primary'])
    tf1 = make_tf(s, Inches(0.9), Inches(1.5), Inches(5.2), Inches(4.4))
    title_paragraph(tf1, "📊  Đánh giá môn học", pal['primary'], SZ_TITLE_SM)
    bullet_items(tf1, [
        "Đánh giá kĩ năng thực hành lắp ráp.",
        "Sản phẩm Robot hoàn thiện theo bài.",
        "Dự án sáng tạo & Triển lãm cuối năm.",
    ], bullet_color=pal['primary'])

    add_card_group(s, Inches(6.8), Inches(1.3), Inches(5.9), Inches(4.8),
                   pal['card'], pal['second'])
    tf2 = make_tf(s, Inches(7.2), Inches(1.5), Inches(5.2), Inches(4.4))
    title_paragraph(tf2, "🧰  Đồ dùng cần có", pal['second'], SZ_TITLE_SM)
    bullet_items(tf2, [
        "Bộ Kit Robotics UNIGO.",
        "Vở ghi & Sổ tay ý tưởng sáng chế.",
        "Máy tính/Máy tính bảng lập trình.",
    ], bullet_color=pal['second'])


def build_slide_homework(prs, pal, grade):
    s = new_slide(prs)
    set_slide_bg(s, pal['bg'])

    add_card_group(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(4.8),
                   pal['card'], pal['accent'])

    tf = make_tf(s, Inches(1.0), Inches(1.5), Inches(11.5), Inches(4.4))
    title_paragraph(tf, "📝  Nhiệm vụ mở rộng (Homework)", pal['accent'], SZ_TITLE)
    bullet_items(tf, [
        "Phác thảo ý tưởng 1 mô hình Robot em ước mơ chế tạo.",
        "Đặt tên và mô tả tính năng chính của Robot đó.",
        "Cam kết thực hiện đúng nội quy bảo quản bộ Kit Robotics.",
        "Xem trước Bài 1 trong chương trình Robotics.",
    ], bullet_color=pal['accent'])


def build_slide_summary(prs, pal, grade):
    s = new_slide(prs)
    set_slide_bg(s, pal['bg'])

    panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0), Inches(1.15),
                               Inches(13.33), Inches(5.15))
    panel.fill.solid()
    panel.fill.fore_color.rgb = pal['primary']
    panel.line.fill.background()

    add_card_group(s, Inches(1.2), Inches(1.5), Inches(10.9), Inches(4.5),
                   pal['card'], pal['accent'], bar_width=Inches(0.18))

    tf = make_tf(s, Inches(1.7), Inches(1.7), Inches(10.0), Inches(4.0))
    title_paragraph(tf, f"📌  Ghi nhớ — Robotics Lớp {grade}", pal['primary'], SZ_TITLE)
    bullet_items(tf, [
        "Robotics giúp em tự tay biến ý tưởng thành hiện thực.",
        "Cẩn thận, tỉ mỉ & tuân thủ an toàn thiết bị.",
        "Đam mê sáng tạo — Chinh phục công nghệ!",
    ], bullet_color=pal['accent'], size=SZ_SUB)

    p_thanks = tf.add_paragraph()
    p_thanks.space_before = Pt(20)
    p_thanks.alignment = PP_ALIGN.CENTER
    add_run(p_thanks, "Chúc các em một năm học Robotics bùng nổ! 🚀🤖",
            size=Pt(22), bold=True, color=pal['second'])


# ─── QUALITY CHECK ─────────────────────────────────────────────
def quality_check(filepath):
    prs = Presentation(filepath)
    issues = []
    if len(prs.slides) < 8:
        issues.append(f"WARN: Chi co {len(prs.slides)} slides")
    master = prs.slide_masters[0]
    has_logo = any(s.name == 'Picture 7' and s.shape_type == 13 for s in master.shapes)
    has_footer = any(s.name == 'Picture 9' and s.shape_type == 13 for s in master.shapes)
    if not has_logo:  issues.append("FAIL: Logo UNIGO bi mat!")
    if not has_footer: issues.append("FAIL: Chan trang bi mat!")
    last = prs.slides[len(prs.slides) - 1]
    for shape in last.shapes:
        if Emu(shape.top).inches < 0.5 and Emu(shape.width).inches > 10:
            issues.append(f"WARN: Slide cuoi shape '{shape.name}' co the che logo (top={Emu(shape.top).inches:.2f}in)")
    return issues


# ─── MAIN GENERATOR ────────────────────────────────────────────
def build_deck(grade, palette_idx=None):
    prs = Presentation(TPL_PATH)
    clear_template_slides(prs)

    if palette_idx is None:
        palette_idx = (grade - 1) % len(PALETTES)
    pal = PALETTES[palette_idx]
    print(f"  Palette: {pal['name']} (#{palette_idx + 1})")

    subtitles = {
        1: "Kham pha the gioi Robot & Co che co ban",
        2: "Robot & Co che thong minh",
        3: "Lap trinh & Co cau chuyen dong",
        4: "Robot cam bien & Tu dong hoa",
        5: "Sang tao Robot & Tu duy ky thuat",
        6: "Dong co, Cam bien & Bo dieu khien",
        7: "He thong Robotics phuc tap",
        8: "Thiet ke & Lap trinh Robot nang cao",
    }

    build_slide_intro(prs, pal, grade, subtitles.get(grade, f"Kham pha Robotics {grade}"))
    build_slide_objectives(prs, pal, grade)
    build_slide_warmup(prs, pal, grade)
    build_slide_overview(prs, pal, grade)
    build_slide_methods(prs, pal, grade)
    build_slide_rules_1(prs, pal, grade)
    build_slide_rules_2(prs, pal, grade)
    build_slide_rules_3(prs, pal, grade)
    build_slide_assessment(prs, pal, grade)
    build_slide_homework(prs, pal, grade)
    build_slide_summary(prs, pal, grade)

    return prs


def main():
    print("=== TẠO TOÀN BỘ SLIDE TIẾT 0 ROBOTICS (LỚP 1 ĐẾN LỚP 8) ===")
    total_files = 0
    all_passed = True

    for grade in range(1, 9):
        print(f"\n--- Đang xử lý Lớp {grade} ---")
        prs = build_deck(grade)

        out_folder = os.path.join(BASE_OUT, f"Lớp_{grade}", "Tiết_00")
        os.makedirs(out_folder, exist_ok=True)
        out_path = os.path.join(out_folder, f"Slide_Robotics_Lop_{grade}_Tiet00_Dinh_huong.pptx")

        prs.save(out_path)
        total_files += 1
        print(f"  ✓ Đã lưu: {out_path}")
        print(f"  ✓ Số slide: {len(prs.slides)}")

        issues = quality_check(out_path)
        if issues:
            all_passed = False
            print("  ⚠ QUALITY ISSUES:")
            for iss in issues:
                print(f"    - {iss}")
        else:
            print("  ✅ QUALITY CHECK PASSED!")

    print(f"\n=== HOÀN THÀNH: Đã tạo {total_files}/8 bộ slide Tiết 0 Robotics thành công! ===")


if __name__ == '__main__':
    main()
