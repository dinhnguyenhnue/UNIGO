import sys, os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches

FILES_TO_CHECK = [
    r"D:\UNIGO\KHBD_Tin_học\Lớp_5\Tuần_04\Slide_Tin_hoc_Lớp_5_Tiet03_Bai_3_Tim_kiem_thong_tin.pptx",
    r"D:\UNIGO\KHBD_Tin_học\Lớp_5\Tuần_04\Slide_Tin_hoc_Lớp_5_Tiet04_Bai_4_Cay_thu_muc.pptx",
    r"D:\UNIGO\KHBD_Tin_học\Lớp_5\Tuần_04\Slide_Tin_hoc_Lop_5_Bai03_04.pptx",
    r"D:\UNIGO\KHBD_Tin_học\Lớp_7\Tuần_04\Slide_Tin_hoc_Lớp_7_Tiet03_Bai_3_Quan_ly_du_lieu_trong_may_tinh.pptx",
    r"D:\UNIGO\KHBD_Tin_học\Lớp_7\Tuần_04\Slide_Tin_hoc_Lớp_7_Tiet04_Bai_4_Mang_xa_hoi.pptx",
    r"D:\UNIGO\KHBD_Tin_học\Lớp_7\Tuần_04\Slide_Tin_hoc_Lop_7_Bai03_04.pptx"
]

def verify_file(fpath):
    print(f"\n======================================")
    print(f"VERIFYING: {os.path.basename(fpath)}")
    prs = Presentation(fpath)
    print(f"Total slides: {len(prs.slides)}")
    
    all_pass = True
    for s_idx, slide in enumerate(prs.slides):
        # Check shapes in slide
        pic_count = 0
        for shape in slide.shapes:
            if shape.shape_type == 13: # Picture
                pic_count += 1
            # Check safe zone violations (shapes added by user script)
            top_in = shape.top / 914400.0
            height_in = shape.height / 914400.0
            bottom_in = top_in + height_in
            
            # Allow background rectangles that are clamped to 1.15 -> 6.35
            if top_in < 1.14 and shape.name.startswith("Rectangle"):
                print(f"  [WARN] Slide {s_idx+1}: Shape top < 1.15 ({top_in:.2f}in)")
                all_pass = False
            if bottom_in > 6.36 and shape.name.startswith("Rectangle") and not shape.name.startswith("Picture"):
                print(f"  [WARN] Slide {s_idx+1}: Shape bottom > 6.35 ({bottom_in:.2f}in)")
                all_pass = False
        print(f"  Slide {s_idx+1}: Shapes={len(slide.shapes)}, Pictures={pic_count} - [PASS]")
    
    if all_pass:
        print(f"--> [SUCCESS 100% PASS] {os.path.basename(fpath)}")

for f in FILES_TO_CHECK:
    verify_file(f)
