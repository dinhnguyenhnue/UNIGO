# -*- coding: utf-8 -*-
"""
Tạo Bộ Slide Bài Giảng Chuẩn UNIGO Tuần 03 (Môn Tin Học)
Áp dụng cho 5 khối: Tiền Tiểu học, Lớp 1, Lớp 2, Lớp 3, Lớp 4
Quy chuẩn:
  1. Template Master UNIGO giữ nguyên Logo & Chân trang (Safe Zone Y: 1.15in -> 6.35in)
  2. Slide Video Tham Khảo YouTube chèn sau Khởi động
  3. Per-Bullet Images (Layout A Flashcard Grid)
  4. Animation tuần tự cho slide Luyện tập / Hoạt động
  5. Độ tương phản cao, bảng màu xoay vòng chuyên nghiệp
"""
import sys, io, os, glob
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

TEMPLATE  = r'D:\UNIGO\Hệ thống mẫu văn bản\Mẫu slide UNIGO.pptx'
SGK_BASE  = r'D:\UNIGO\SGK'
KHBD_BASE = r'D:\UNIGO\KHBD_Tin_học'

SAFE_TOP    = 1.15
SAFE_BOTTOM = 6.35
SAFE_LEFT   = 0.3
SAFE_RIGHT  = 13.0
SLIDE_W     = 13.33
SLIDE_H     = 7.50

COLOR_PALETTES = [
    # 0: Blue (Tiền TH)
    {"primary": "1B4F9B", "accent": "2D7DD2", "bg": "EBF3FE", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "1A2744", "text_on_card": "1A2744"},
    # 1: Purple (Lớp 1)
    {"primary": "5B21B6", "accent": "7C3AED", "bg": "F3EEFF", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "2E1065", "text_on_card": "2E1065"},
    # 2: Teal (Lớp 2)
    {"primary": "0F766E", "accent": "14B8A6", "bg": "ECFDF5", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "134E4A", "text_on_card": "134E4A"},
    # 3: Orange (Lớp 3)
    {"primary": "C2410C", "accent": "EA580C", "bg": "FFF7ED", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "7C2D12", "text_on_card": "431407"},
    # 4: Indigo (Lớp 4)
    {"primary": "3730A3", "accent": "4F46E5", "bg": "EEF2FF", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "1E1B4B", "text_on_card": "1E1B4B"},
]

LESSONS_T3 = [
    {
        "id": "Tien_TH", "folder": "Tiền_tiểu_học", "lop_label": "Tiền Tiểu học", "lop_num": 0, "palette_idx": 0,
        "bai": "Bài 2. Máy tính quanh em",
        "file_name": "Slide_Tin_hoc_Tien_TH_Bai02.pptx",
        "slides_content": [
            {"type": "cover", "title": "Máy tính quanh em 🖥️", "subtitle": "Tin học • Tiền Tiểu học • Bài 2"},
            {"type": "warmup", "title": "Ai nhìn thấy máy tính rồi?",
             "content": "Con hãy nhìn xung quanh lớp và ở nhà mình:\nMáy tính trông như thế nào?\nCùng chỉ cho cô và các bạn biết nhé!"},
            {"type": "video", "title": "Khám phá: Máy tính kỳ diệu! 🎬",
             "url": "https://www.youtube.com/watch?v=kYJ5z5_W2hY",
             "desc": "Cùng xem video khám phá thế giới máy tính thông minh nhé!"},
            {"type": "learn", "title": "4 bộ phận kỳ diệu của máy tính",
             "bullets": [
                 "Màn hình — để con nhìn tranh ảnh, video",
                 "Bàn phím — để con gõ chữ và số",
                 "Chuột — để con điều khiển mũi tên",
                 "Thân máy — bộ não lưu trữ thông tin"
             ]},
            {"type": "practice", "title": "Bé chỉ đúng bộ phận nào!",
             "instruction": "Các con cùng gọi tên 4 bộ phận nhé:",
             "items": ["Chỉ vào màn hình máy tính", "Chỉ vào bàn phím gõ chữ", "Cầm chuột máy tính trên bàn", "Chỉ vào thân máy tính"]},
            {"type": "activity", "title": "Trò chơi: Nhanh mắt nhanh tay!",
             "instruction": "Ghép đúng công dụng của từng bộ phận:",
             "items": ["Màn hình ↔ Xem phim hoạt hình", "Bàn phím ↔ Gõ tên của em", "Chuột máy ↔ Nhấp chọn tranh vẽ", "Thân máy ↔ Bộ não xử lý thông minh"]},
            {"type": "summary", "title": "Hôm nay các con đã nhớ!",
             "items": ["Máy tính có 4 bộ phận chính", "Biết gọi tên màn hình, bàn phím, chuột, thân máy", "Giữ gìn máy tính cẩn thận và sạch sẽ"]},
            {"type": "thanks", "title": "Các con học giỏi lắm! 🌟", "content": "BTVN: Kể cho bố mẹ nghe về 4 bộ phận\ncủa chiếc máy tính nhé!"},
        ]
    },
    {
        "id": "Lop_1", "folder": "Lớp_1", "lop_label": "Lớp 1", "lop_num": 1, "palette_idx": 1,
        "bai": "Bài 2. Các bộ phận của máy tính",
        "file_name": "Slide_Tin_hoc_Lop_1_Bai02.pptx",
        "slides_content": [
            {"type": "cover", "title": "Các bộ phận của máy tính 🖥️", "subtitle": "Tin học • Lớp 1 • Bài 2"},
            {"type": "warmup", "title": "Đố em đây là cái gì?",
             "content": "Một vật có màn hình sáng, có bàn phím gõ chữ, có chú chuột nhỏ xíu:\nĐó chính là máy tính để bàn đấy!"},
            {"type": "video", "title": "Xem Video: Cấu tạo máy tính 🎬",
             "url": "https://www.youtube.com/watch?v=fD09A-w-3eI",
             "desc": "Xem video để nhận biết chính xác 4 bộ phận cơ bản của máy tính!"},
            {"type": "learn", "title": "4 bộ phận chính của máy tính để bàn",
             "bullets": [
                 "Màn hình (Monitor): hiển thị chữ, hình ảnh",
                 "Thân máy (Computer Case): chứa bộ não CPU",
                 "Bàn phím (Keyboard): gõ chữ và số vào máy",
                 "Chuột (Mouse): điều khiển con trỏ trên màn hình"
             ]},
            {"type": "learn", "title": "Máy tính xách tay (Laptop)",
             "bullets": [
                 "Màn hình gắn liền với thân máy",
                 "Bàn phím và bàn di chuột tích hợp sẵn",
                 "Nhỏ gọn, dễ dàng mang đi mọi nơi",
                 "Có pin sạc dùng khi không cắm điện"
             ]},
            {"type": "practice", "title": "Nhận biết bộ phận máy tính!",
             "instruction": "Em hãy gọi tên đúng bộ phận khi thầy/cô chỉ vào:",
             "items": ["Màn hình hiển thị hình ảnh", "Thân máy tính chứa CPU", "Bàn phím gõ chữ và phím số", "Chuột máy tính có 2 nút"]},
            {"type": "activity", "title": "Trò chơi: Đúng hay Sai?",
             "instruction": "Chọn đáp án đúng nhất:",
             "items": ["Màn hình dùng để gõ chữ ↔ SAI ❌ (Bàn phím)", "Thân máy chứa bộ não CPU ↔ ĐÚNG ✅", "Chuột máy tính dùng để điều khiển ↔ ĐÚNG ✅", "Laptop có thể gấp gọn mang đi ↔ ĐÚNG ✅"]},
            {"type": "summary", "title": "Tóm tắt bài học",
             "items": ["4 bộ phận cơ bản: Màn hình, Thân máy, Bàn phím, Chuột", "Thân máy tính đóng vai trò quan trọng nhất", "Máy tính xách tay tích hợp tất cả trong một"]},
            {"type": "thanks", "title": "Em làm bài rất tốt! ⭐", "content": "BTVN: Quan sát máy tính ở nhà hoặc trường\nvà chỉ cho bạn các bộ phận nhé!"},
        ]
    },
    {
        "id": "Lop_2", "folder": "Lớp_2", "lop_label": "Lớp 2", "lop_num": 2, "palette_idx": 2,
        "bai": "Bài 2. Ôn và nâng cấp kỹ năng chuột",
        "file_name": "Slide_Tin_hoc_Lop_2_Bai02.pptx",
        "slides_content": [
            {"type": "cover", "title": "Nâng cấp kỹ năng chuột 🖱️", "subtitle": "Tin học • Lớp 2 • Bài 2"},
            {"type": "warmup", "title": "Thử thách phản xạ chuột siêu tốc!",
             "content": "Em đã thuần thục các thao tác cầm và di chuột chưa?\nHãy sẵn sàng nâng cấp kỹ năng chuột cấp độ 2 nào!"},
            {"type": "video", "title": "Kỹ năng dùng chuột máy tính 🎬",
             "url": "https://www.youtube.com/watch?v=9_Vv5v1g6Bw",
             "desc": "Video hướng dẫn các thao tác nháy đơn, nháy đôi, cuộn chuột và kéo thả!"},
            {"type": "learn", "title": "Các thao tác chuột nâng cao",
             "bullets": [
                 "Nháy đôi (Double click): mở ứng dụng và tệp tin",
                 "Nháy chuột phải (Right click): mở menu tùy chọn",
                 "Cuộn chuột (Scroll): lướt xem trang web lên xuống",
                 "Kéo thả (Drag & Drop): di chuyển biểu tượng đồ vật"
             ]},
            {"type": "practice", "title": "Thực hành trên máy tính!",
             "instruction": "Thực hiện theo các yêu cầu của giáo viên:",
             "items": ["Nháy đôi mở thư mục Học tập", "Lăn con lăn xem hết danh sách tệp", "Kéo thả biểu tượng vào thư mục", "Nháy phải chuột xem thuộc tính"]},
            {"type": "activity", "title": "Trò chơi: Bắt bong bóng chuột!",
             "instruction": "Thử thách phản xạ nhanh:",
             "items": ["Nháy trái diệt bóng thường ↔ 10 điểm", "Nháy đôi diệt bóng vàng ↔ 20 điểm", "Kéo thả bóng xanh vào giỏ ↔ 30 điểm", "Cuộn chuột né chướng ngại vật ↔ 50 điểm"]},
            {"type": "summary", "title": "Ghi nhớ bài học",
             "items": ["Thành thạo 5 thao tác chuột cơ bản và nâng cao", "Cầm chuột đúng tư thế, ngón trỏ nút trái, ngón giữa nút phải", "Kéo thả chuột chính xác và an toàn"]},
            {"type": "thanks", "title": "Xuất sắc! 🎉", "content": "BTVN: Luyện tập trò chơi luyện chuột\n10 phút trên máy tính!"},
        ]
    },
    {
        "id": "Lop_3", "folder": "Lớp_3", "lop_label": "Lớp 3", "lop_num": 3, "palette_idx": 3,
        "bai": "Bài 2. Xử lí thông tin",
        "file_name": "Slide_Tin_hoc_Lop_3_Bai02.pptx",
        "slides_content": [
            {"type": "cover", "title": "Bài 2. Xử lí thông tin 🧠", "subtitle": "Tin học • Lớp 3 • Bài 2"},
            {"type": "warmup", "title": "Đèn đỏ bật sáng — Em làm gì?",
             "content": "Khi đi đường thấy đèn giao thông chuyển sang màu đỏ:\nMắt nhìn thấy (Thông tin vào) -> Bộ não nghĩ 'Cần dừng lại' (Xử lý) -> Chân đạp phanh (Kết quả)!"},
            {"type": "video", "title": "Xem Video: Con người & Máy tính xử lý thông tin 🎬",
             "url": "https://www.youtube.com/watch?v=mQy1e8_A7C0",
             "desc": "Xem cách bộ não và máy tính cùng tiếp nhận, xử lý và đưa ra thông tin!"},
            {"type": "learn", "title": "Quy trình xử lý thông tin của con người",
             "bullets": [
                 "1. Tiếp nhận thông tin qua 5 giác quan (mắt, tai, mũi...)",
                 "2. Bộ não phân tích, suy nghĩ và xử lý thông tin",
                 "3. Đưa ra quyết định hoặc hành động phù hợp",
                 "4. Lưu trữ thông tin vào trí nhớ để dùng lần sau"
             ]},
            {"type": "learn", "title": "Quy trình xử lý thông tin của máy tính",
             "bullets": [
                 "Thu nhận thông tin vào (Input: bàn phím, chuột)",
                 "Bộ xử lý CPU tính toán, biến đổi dữ liệu",
                 "Xuất kết quả ra ngoài (Output: màn hình, loa)",
                 "Lưu trữ dữ liệu vào ổ đĩa và bộ nhớ"
             ]},
            {"type": "practice", "title": "Phân tích ví dụ thực tế!",
             "instruction": "Xác định các bước xử lý thông tin:",
             "items": ["Nghe tiếng chuông reo ↔ Thu nhận thông tin", "Biết đến giờ vào lớp ↔ Não xử lý thông tin", "Đi nhanh vào chỗ ngồi ↔ Đưa ra quyết định/hành động", "Gõ phím 2+3, màn hình hiện 5 ↔ Máy tính xử lý"]},
            {"type": "activity", "title": "Trò chơi: Thần đồng tính nhanh!",
             "instruction": "So tài giữa Não người và Máy tính:",
             "items": ["Tính 15 + 25 ↔ Cả người và máy tính đều nhanh!", "Tính 9876 x 5432 ↔ Máy tính xử lý trong chớp mắt!", "Vẽ tranh sáng tạo ↔ Con người cảm xúc và sáng tạo hơn!", "Ghi nhớ 10.000 cuốn sách ↔ Máy tính lưu trữ vượt trội!"]},
            {"type": "summary", "title": "Tóm tắt bài học",
             "items": ["Quá trình xử lý: Thu nhận -> Xử lý -> Xuất kết quả", "Bộ não là trung tâm xử lý thông tin của con người", "CPU là bộ não xử lý thông tin của máy tính"]},
            {"type": "thanks", "title": "Bài học kết thúc! 🚀", "content": "BTVN: Tìm thêm 2 ví dụ về xử lý thông tin\ntrong cuộc sống hàng ngày!"},
        ]
    },
    {
        "id": "Lop_4", "folder": "Lớp_4", "lop_label": "Lớp 4", "lop_num": 4, "palette_idx": 4,
        "bai": "Bài 2. Gõ bàn phím đúng cách",
        "file_name": "Slide_Tin_hoc_Lop_4_Bai02.pptx",
        "slides_content": [
            {"type": "cover", "title": "Gõ bàn phím đúng cách ⌨️", "subtitle": "Tin học • Lớp 4 • Bài 2"},
            {"type": "warmup", "title": "Bí mật 10 ngón tay múa trên phím!",
             "content": "Em đang gõ bàn phím bằng mấy ngón tay?\nGõ bằng 10 ngón tay sẽ giúp em gõ nhanh gấp 5 lần mà không cần nhìn bàn phím!"},
            {"type": "video", "title": "Hướng dẫn: Gõ 10 ngón tay chuẩn 🎬",
             "url": "https://www.youtube.com/watch?v=zT1gQ_w6w2s",
             "desc": "Quan sát kỹ cách đặt tay ở hàng cơ sở và phân công từng ngón tay!"},
            {"type": "learn", "title": "Hàng phím cơ sở (Home row)",
             "bullets": [
                 "Gồm các phím: A S D F G H J K L ;",
                 "Hai phím có gờ định vị: phím F và phím J",
                 "Nơi đặt 8 ngón tay xuất phát và luôn trở về",
                 "Hai ngón tay cái luôn đặt nhẹ lên phím cách (Space)"
             ]},
            {"type": "learn", "title": "Phân công 10 ngón tay gõ phím",
             "bullets": [
                 "Ngón trỏ trái phụ trách: F, V, B, R, T, 4, 5",
                 "Ngón trỏ phải phụ trách: J, N, M, U, Y, 7, 8",
                 "Các ngón còn lại gõ theo cột tương ứng",
                 "Ngón tay cái gõ phím Cách (Spacebar)"
             ]},
            {"type": "practice", "title": "Luyện đặt tay đúng vị trí!",
             "instruction": "Thực hiện theo 4 bước chuẩn:",
             "items": ["Đặt ngón trỏ trái lên phím F (tìm gờ nổi)", "Đặt ngón trỏ phải lên phím J (tìm gờ nổi)", "Thả nhẹ 6 ngón còn lại vào hàng cơ sở", "Hai ngón cái chạm nhẹ phím cách (Space)"]},
            {"type": "activity", "title": "Luyện gõ trên RapidTyping!",
             "instruction": "Thực hành chuỗi ký tự cơ sở:",
             "items": ["Gõ chuỗi: asdf jkl; asdf jkl;", "Gõ từ: fa la da ka ja ha ga", "Thi đua đạt độ chính xác trên 95%!", "Tập gõ không nhìn xuống bàn phím"]},
            {"type": "summary", "title": "Hôm nay em đã học!",
             "items": ["Hàng phím cơ sở và 2 phím có gờ F, J", "Vị trí phân công của 10 ngón tay", "Quy tắc gõ mười ngón (Touch Typing) chuẩn"]},
            {"type": "thanks", "title": "Em làm rất tuyệt vời! 🌟", "content": "BTVN: Luyện gõ 10 ngón trên phần mềm\n15 phút mỗi ngày nhé!"},
        ]
    },
]

def hex_rgb(h):
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def set_font(run, size_pt, bold=False, color_hex="333333", font_name="Arial"):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = hex_rgb(color_hex)
    run.font.name = font_name

def add_textbox(slide, left, top, width, height, text, size_pt=18, bold=False,
                color_hex="333333", alignment=PP_ALIGN.LEFT, font_name="Arial"):
    actual_top = max(top, SAFE_TOP)
    txbox = slide.shapes.add_textbox(Inches(left), Inches(actual_top), Inches(width), Inches(height))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size_pt, bold, color_hex, font_name)
    return txbox

def add_multiline_textbox(slide, left, top, width, height, lines, size_pt=18,
                           color_hex="333333", bold=False, font_name="Arial",
                           alignment=PP_ALIGN.LEFT):
    actual_top = max(top, SAFE_TOP)
    txbox = slide.shapes.add_textbox(Inches(left), Inches(actual_top), Inches(width), Inches(height))
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alignment
        run = p.add_run()
        run.text = line
        set_font(run, size_pt, bold, color_hex, font_name)
        p.space_after = Pt(6)
    return txbox

def add_safe_shape(slide, shape_type, left, top, width, height, fill_hex,
                   border_hex=None, send_to_back=False):
    actual_top = max(top, SAFE_TOP)
    actual_bottom = min(top + height, SAFE_BOTTOM)
    actual_height = max(actual_bottom - actual_top, 0.1)

    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left), Inches(actual_top), Inches(width), Inches(actual_height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(fill_hex)
    if border_hex:
        shape.line.color.rgb = hex_rgb(border_hex)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()

    if send_to_back:
        sp = shape._element
        spTree = sp.getparent()
        spTree.remove(sp)
        spTree.insert(2, sp)
    return shape

def add_picture_safe(slide, img_path, left, top, width, height):
    actual_top = max(top, SAFE_TOP)
    actual_bottom = min(top + height, SAFE_BOTTOM)
    actual_height = max(actual_bottom - actual_top, 0.5)
    try:
        return slide.shapes.add_picture(img_path, Inches(left), Inches(actual_top),
                                         Inches(width), Inches(actual_height))
    except Exception as e:
        return None

def add_slide_transition(slide, transition_type="fade"):
    transitions = {
        'fade': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>',
        'push': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:push dir="l"/></p:transition>',
        'wipe': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:wipe/></p:transition>',
        'cover': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:cover/></p:transition>',
    }
    xml_str = transitions.get(transition_type, transitions['fade'])
    slide._element.append(etree.fromstring(xml_str))

def add_appear_animation(slide, shapes_with_click):
    nsmap = {
        'p':  'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a':  'http://schemas.openxmlformats.org/drawingml/2006/main',
    }

    timing = etree.SubElement(slide._element, '{%s}timing' % nsmap['p'])
    tnLst = etree.SubElement(timing, '{%s}tnLst' % nsmap['p'])
    par_root = etree.SubElement(tnLst, '{%s}par' % nsmap['p'])
    cTn_root = etree.SubElement(par_root, '{%s}cTn' % nsmap['p'])
    cTn_root.set('id', '1')
    cTn_root.set('dur', 'indefinite')
    cTn_root.set('restart', 'never')
    cTn_root.set('nodeType', 'tmRoot')

    childTnLst_root = etree.SubElement(cTn_root, '{%s}childTnLst' % nsmap['p'])
    seq = etree.SubElement(childTnLst_root, '{%s}seq' % nsmap['p'])
    seq.set('concurrent', '1')
    seq.set('nextAc', 'seek')

    cTn_seq = etree.SubElement(seq, '{%s}cTn' % nsmap['p'])
    cTn_seq.set('id', '2')
    cTn_seq.set('dur', 'indefinite')
    cTn_seq.set('nodeType', 'mainSeq')

    childTnLst_seq = etree.SubElement(cTn_seq, '{%s}childTnLst' % nsmap['p'])

    prevCond = etree.SubElement(seq, '{%s}prevCondLst' % nsmap['p'])
    cond_prev = etree.SubElement(prevCond, '{%s}cond' % nsmap['p'])
    cond_prev.set('evt', 'onPrev')
    cond_prev.set('delay', '0')
    tgtEl_prev = etree.SubElement(cond_prev, '{%s}tgtEl' % nsmap['p'])
    etree.SubElement(tgtEl_prev, '{%s}sldTgt' % nsmap['p'])

    nextCond = etree.SubElement(seq, '{%s}nextCondLst' % nsmap['p'])
    cond_next = etree.SubElement(nextCond, '{%s}cond' % nsmap['p'])
    cond_next.set('evt', 'onNext')
    cond_next.set('delay', '0')
    tgtEl_next = etree.SubElement(cond_next, '{%s}tgtEl' % nsmap['p'])
    etree.SubElement(tgtEl_next, '{%s}sldTgt' % nsmap['p'])

    id_counter = 3

    for shape_group, click_idx in shapes_with_click:
        if not isinstance(shape_group, (list, tuple)):
            shape_group = [shape_group]

        par_click = etree.SubElement(childTnLst_seq, '{%s}par' % nsmap['p'])
        cTn_click = etree.SubElement(par_click, '{%s}cTn' % nsmap['p'])
        cTn_click.set('id', str(id_counter)); id_counter += 1
        cTn_click.set('fill', 'hold')

        stCondLst = etree.SubElement(cTn_click, '{%s}stCondLst' % nsmap['p'])
        cond = etree.SubElement(stCondLst, '{%s}cond' % nsmap['p'])
        cond.set('delay', '0')

        childTnLst_click = etree.SubElement(cTn_click, '{%s}childTnLst' % nsmap['p'])

        for shape in shape_group:
            if shape is None:
                continue
            sp_id = shape.shape_id

            par_anim = etree.SubElement(childTnLst_click, '{%s}par' % nsmap['p'])
            cTn_anim = etree.SubElement(par_anim, '{%s}cTn' % nsmap['p'])
            cTn_anim.set('id', str(id_counter)); id_counter += 1
            cTn_anim.set('presetID', '1')  # Appear
            cTn_anim.set('presetClass', 'entr')
            cTn_anim.set('presetSubtype', '0')
            cTn_anim.set('fill', 'hold')

            stCond_anim = etree.SubElement(cTn_anim, '{%s}stCondLst' % nsmap['p'])
            cond_anim = etree.SubElement(stCond_anim, '{%s}cond' % nsmap['p'])
            cond_anim.set('delay', '0')

            childTnLst_anim = etree.SubElement(cTn_anim, '{%s}childTnLst' % nsmap['p'])

            p_set = etree.SubElement(childTnLst_anim, '{%s}set' % nsmap['p'])
            cBhvr = etree.SubElement(p_set, '{%s}cBhvr' % nsmap['p'])
            cTn_set = etree.SubElement(cBhvr, '{%s}cTn' % nsmap['p'])
            cTn_set.set('id', str(id_counter)); id_counter += 1
            cTn_set.set('dur', '1')
            cTn_set.set('fill', 'hold')

            stCond_set = etree.SubElement(cTn_set, '{%s}stCondLst' % nsmap['p'])
            cond_set = etree.SubElement(stCond_set, '{%s}cond' % nsmap['p'])
            cond_set.set('delay', '0')

            tgtEl = etree.SubElement(cBhvr, '{%s}tgtEl' % nsmap['p'])
            spTgt = etree.SubElement(tgtEl, '{%s}spTgt' % nsmap['p'])
            spTgt.set('spid', str(sp_id))

            attrNameLst = etree.SubElement(cBhvr, '{%s}attrNameLst' % nsmap['p'])
            attrName = etree.SubElement(attrNameLst, '{%s}attrName' % nsmap['p'])
            attrName.text = 'style.visibility'

            p_to = etree.SubElement(p_set, '{%s}to' % nsmap['p'])
            strVal = etree.SubElement(p_to, '{%s}strVal' % nsmap['p'])
            strVal.set('val', 'visible')

def find_image(lesson, slide_type, idx=0):
    lop_folder = lesson["folder"]
    lop_num = lesson["lop_num"]
    
    # Check local Tuần 03 images
    khbd_dir = os.path.join(KHBD_BASE, lop_folder, "Tuần_03", "images")
    if os.path.isdir(khbd_dir):
        for name in [f"{slide_type}{idx+1}_lop{lop_num}.png", f"{slide_type}_lop{lop_num}.png"]:
            p = os.path.join(khbd_dir, name)
            if os.path.isfile(p):
                return p
        imgs = sorted([os.path.join(khbd_dir, f) for f in os.listdir(khbd_dir) if f.endswith(('.png','.jpg','.jpeg'))])
        if imgs:
            return imgs[idx % len(imgs)]
            
    # Check SGK images
    sgk_dir = os.path.join(SGK_BASE, f"Lớp_{lop_num}", "all_extracted_images")
    if os.path.isdir(sgk_dir):
        imgs = sorted([os.path.join(sgk_dir, f) for f in os.listdir(sgk_dir) if f.endswith(('.png','.jpeg','.jpg'))])
        if imgs:
            return imgs[idx % len(imgs)]
    return None

def build_cover(prs, data, pal, lesson, layout):
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["primary"], send_to_back=True)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, 1.8, pal["accent"], send_to_back=True)

    add_textbox(slide, 0.8, 2.0, 7.5, 1.5, data["title"], size_pt=34, bold=True, color_hex=pal["text_on_primary"])
    add_textbox(slide, 0.8, 3.8, 7.5, 0.8, data.get("subtitle", ""), size_pt=20, color_hex=pal["text_on_primary"])

    img = find_image(lesson, "cover", 0)
    if img:
        add_picture_safe(slide, img, SLIDE_W - 4.5, SAFE_TOP + 0.3, 3.8, 3.8)

    add_slide_transition(slide, "fade")
    return slide

def build_warmup(prs, data, pal, lesson, layout):
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, pal["primary"])
    badge_text = f"  {lesson['lop_label'].upper()} • {lesson['bai'].upper()}"
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, badge_text, size_pt=13, bold=True, color_hex=pal["text_on_primary"])

    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, 8, 0.65, data["title"], size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    content_y = title_y + 0.75
    card_w = 7.5
    card_h = SAFE_BOTTOM - content_y - 0.15
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.4, content_y, card_w, card_h, pal["card"])
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0.4, content_y, 0.12, card_h, pal["accent"])

    lines = data.get("content", "").split("\n")
    add_multiline_textbox(slide, 0.8, content_y+0.2, card_w-0.6, card_h-0.4, lines, size_pt=20, color_hex=pal["text_on_card"])

    img = find_image(lesson, "warmup", 0)
    if img:
        add_picture_safe(slide, img, 8.3, content_y, SAFE_RIGHT-8.3-0.2, card_h)

    add_slide_transition(slide, "push")
    return slide

def build_video_slide(prs, data, pal, lesson, layout):
    """Slide Video tham khảo YouTube chất lượng cao"""
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, "DC2626") # Youtube red accent
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, f"  {lesson['lop_label'].upper()} • KHÁM PHÁ VIDEO BÀI HỌC", size_pt=13, bold=True, color_hex="FFFFFF")

    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, SLIDE_W-1, 0.65, data["title"], size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    content_y = title_y + 0.75
    card_w = 7.5
    card_h = SAFE_BOTTOM - content_y - 0.15
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.4, content_y, card_w, card_h, pal["card"], border_hex="EF4444")
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0.4, content_y, 0.12, card_h, "EF4444")

    desc_lines = [
        data.get("desc", ""),
        "",
        "🔗 Đường link video tham khảo:",
        data.get("url", ""),
        "",
        "💡 Thầy/Cô sẽ chiếu video cho các em quan sát nhé!"
    ]
    add_multiline_textbox(slide, 0.8, content_y+0.2, card_w-0.6, card_h-0.4, desc_lines, size_pt=17, color_hex=pal["text_on_card"])

    img = find_image(lesson, "video", 0)
    if img:
        add_picture_safe(slide, img, 8.3, content_y, SAFE_RIGHT-8.3-0.2, card_h)

    add_slide_transition(slide, "cover")
    return slide

def build_learn_grid(prs, data, pal, lesson, layout, learn_idx=0):
    slide = prs.slides.add_slide(layout)
    bullets = data.get("bullets", [])
    n = len(bullets)

    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, pal["primary"])
    badge_text = f"  {lesson['lop_label'].upper()} • {lesson['bai'].upper()}"
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, badge_text, size_pt=13, bold=True, color_hex=pal["text_on_primary"])

    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, SLIDE_W-1, 0.65, data["title"], size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    grid_top = title_y + 0.8
    grid_bottom = SAFE_BOTTOM - 0.1
    avail_h = grid_bottom - grid_top

    cols, rows = (n, 1) if n <= 3 else (2, 2)
    gap = 0.25
    total_gap_x = gap * (cols - 1)
    card_w = (SLIDE_W - 1.0 - total_gap_x) / cols
    if rows == 1:
        img_h, text_h = 2.0, 0.9
    else:
        img_h, text_h = 1.1, 0.5
    card_h = img_h + text_h + 0.15
    total_gap_y = gap * (rows - 1)
    start_x = (SLIDE_W - (cols * card_w + total_gap_x)) / 2
    start_y = grid_top + (avail_h - (rows * card_h + total_gap_y)) / 2
    start_y = max(start_y, grid_top)

    for i, btext in enumerate(bullets):
        col, row = i % cols, i // cols
        x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + gap)

        if y + card_h > SAFE_BOTTOM + 0.05:
            break

        add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h, pal["card"], border_hex="D0D5DD")
        add_safe_shape(slide, MSO_SHAPE.RECTANGLE, x, y, card_w, 0.06, pal["accent"])

        img = find_image(lesson, "learn", learn_idx*3 + i)
        if img:
            img_w = min(1.8, card_w - 0.4)
            img_x = x + (card_w - img_w) / 2
            add_picture_safe(slide, img, img_x, y + 0.15, img_w, img_h)

        text_y = y + img_h + 0.15
        add_textbox(slide, x + 0.15, text_y, card_w - 0.3, text_h, btext, size_pt=16, bold=False, color_hex=pal["text_on_card"], alignment=PP_ALIGN.CENTER)

    add_slide_transition(slide, "wipe")
    return slide

def build_practice(prs, data, pal, lesson, layout):
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, pal["primary"])
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, f"  {lesson['lop_label'].upper()} • LUYỆN TẬP", size_pt=13, bold=True, color_hex=pal["text_on_primary"])

    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, SLIDE_W-1, 0.65, data["title"], size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    instr_y = title_y + 0.75
    instr_lines = data.get("instruction", "").split("\n")
    add_multiline_textbox(slide, 0.5, instr_y, SLIDE_W-1, 0.7, instr_lines, size_pt=18, color_hex=pal["text_on_bg"])

    items = data.get("items", [])
    item_y_start = instr_y + 0.8
    n_items = len(items)

    cols = 2 if n_items > 3 else max(n_items, 1)
    rows = (n_items + 1) // 2 if n_items > 3 else 1

    gap = 0.25
    card_w = (SLIDE_W - 1.0 - gap*(cols-1)) / cols
    card_h = min(1.8, (SAFE_BOTTOM - item_y_start - gap*(rows-1)) / rows)

    animation_groups = []

    for i, item in enumerate(items[:4]):
        col, row = i % cols, i // cols
        x = 0.5 + col * (card_w + gap)
        y = item_y_start + row * (card_h + gap)

        if y + card_h > SAFE_BOTTOM + 0.05:
            break

        card = add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h, pal["card"], border_hex=pal["accent"])

        img = find_image(lesson, "practice", i)
        pic = None
        if img:
            img_w = 1.4
            pic = add_picture_safe(slide, img, x + 0.2, y + 0.15, img_w, card_h - 0.3)
            txt_x = x + 0.2 + img_w + 0.2
            txt_w = card_w - (0.2 + img_w + 0.3)
        else:
            txt_x = x + 0.2
            txt_w = card_w - 0.4

        txt = add_textbox(slide, txt_x, y + 0.15, txt_w, card_h - 0.3, item, size_pt=16, bold=True, color_hex=pal["text_on_card"])

        group = [s for s in [card, pic, txt] if s is not None]
        animation_groups.append((group, i))

    if animation_groups:
        add_appear_animation(slide, animation_groups)

    add_slide_transition(slide, "wipe")
    return slide

def build_activity(prs, data, pal, lesson, layout):
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, pal["accent"])
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, f"  {lesson['lop_label'].upper()} • THỬ THÁCH", size_pt=13, bold=True, color_hex=pal["text_on_primary"])

    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, SLIDE_W-1, 0.65, data["title"], size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    content_y = title_y + 0.75
    content_h = SAFE_BOTTOM - content_y - 0.15

    card_w = 7.0
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.4, content_y, card_w, content_h, pal["card"])
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0.4, content_y, 0.12, content_h, pal["accent"])

    lines = data.get("content", "").split("\n") if data.get("content") else data.get("instruction", "").split("\n")
    if data.get("items"):
        lines.extend(data["items"])
    add_multiline_textbox(slide, 0.8, content_y+0.2, card_w-0.6, content_h-0.4, lines, size_pt=18, color_hex=pal["text_on_card"])

    img = find_image(lesson, "activity", 0)
    if img:
        img_x = 7.8
        img_w = SAFE_RIGHT - img_x - 0.2
        add_picture_safe(slide, img, img_x, content_y + 0.1, img_w, content_h - 0.2)

    add_slide_transition(slide, "cover")
    return slide

def build_summary(prs, data, pal, lesson, layout):
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    px, py = 0.6, SAFE_TOP + 0.15
    pw = SLIDE_W - 1.2
    ph = SAFE_BOTTOM - py - 0.15
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, px, py, pw, ph, pal["card"])
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, px, py, pw, 0.06, pal["primary"])

    add_textbox(slide, px+0.4, py+0.25, pw-0.8, 0.6, data["title"], size_pt=24, bold=True, color_hex=pal["text_on_bg"])

    items = data.get("items", [])
    n = len(items)
    card_gap = 0.2
    card_w = (pw - 0.8 - card_gap*(n-1)) / max(n, 1)
    card_y = py + 1.0
    card_h = ph - 1.3

    for i, item in enumerate(items):
        cx = px + 0.4 + i*(card_w + card_gap)
        add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx, card_y, card_w, card_h, pal["bg"], border_hex="CBD5E1")
        badge_sz = 0.45
        add_safe_shape(slide, MSO_SHAPE.OVAL, cx + 0.15, card_y + 0.15, badge_sz, badge_sz, pal["primary"])
        add_textbox(slide, cx + 0.15, card_y + 0.15, badge_sz, badge_sz, str(i+1), size_pt=14, bold=True, color_hex=pal["text_on_primary"], alignment=PP_ALIGN.CENTER)
        add_textbox(slide, cx + 0.15, card_y + 0.7, card_w - 0.3, card_h - 0.85, item, size_pt=17, bold=False, color_hex=pal["text_on_bg"])

    add_slide_transition(slide, "push")
    return slide

def build_thanks(prs, data, pal, lesson, layout):
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["primary"], send_to_back=True)

    add_textbox(slide, 1.0, 2.0, SLIDE_W-2.0, 1.2, data["title"], size_pt=34, bold=True, color_hex=pal["text_on_primary"], alignment=PP_ALIGN.CENTER)

    if data.get("content"):
        lines = data["content"].split("\n")
        add_multiline_textbox(slide, 1.0, 3.4, SLIDE_W-2.0, 1.5, lines, size_pt=20, color_hex=pal["text_on_primary"], alignment=PP_ALIGN.CENTER)

    add_slide_transition(slide, "fade")
    return slide

def remove_template_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        sldId = prs.slides._sldIdLst[0]
        prs.slides._sldIdLst.remove(sldId)

def generate_slides_tuan03():
    print("[+] BẮT ĐẦU TẠO SLIDE TUẦN 03 CHO 5 KHỐI (Tiền TH -> Lớp 4)...")
    
    for lesson in LESSONS_T3:
        prs = Presentation(TEMPLATE)
        remove_template_slides(prs)
        layout = prs.slide_layouts[6] # Blank
        pal = COLOR_PALETTES[lesson["palette_idx"]]
        
        learn_idx = 0
        for sc in lesson["slides_content"]:
            stype = sc["type"]
            if stype == "cover":
                build_cover(prs, sc, pal, lesson, layout)
            elif stype == "warmup":
                build_warmup(prs, sc, pal, lesson, layout)
            elif stype == "video":
                build_video_slide(prs, sc, pal, lesson, layout)
            elif stype == "learn":
                build_learn_grid(prs, sc, pal, lesson, layout, learn_idx=learn_idx)
                learn_idx += 1
            elif stype == "practice":
                build_practice(prs, sc, pal, lesson, layout)
            elif stype == "activity":
                build_activity(prs, sc, pal, lesson, layout)
            elif stype == "summary":
                build_summary(prs, sc, pal, lesson, layout)
            elif stype == "thanks":
                build_thanks(prs, sc, pal, lesson, layout)
                
        out_dir = os.path.join(KHBD_BASE, lesson["folder"], "Tuần_03")
        os.makedirs(out_dir, exist_ok=True)
        out_path = os.path.join(out_dir, lesson["file_name"])
        prs.save(out_path)
        print(f"    [OK] Tạo thành công: {out_path} ({len(prs.slides)} slides)")

if __name__ == '__main__':
    generate_slides_tuan03()
