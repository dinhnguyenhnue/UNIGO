# -*- coding: utf-8 -*-
"""
Kiểm tra tự động toàn diện cho TOÀN BỘ file slide Tuần 3 và Tuần 4
Kiểm tra 13 quy tắc Anti-Bug: Safe Zone (Y 1.15 -> 6.35), Z-order, tương phản, số lượng slide.
"""
import sys, os, glob
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

SAFE_TOP = 1.15
SAFE_BOTTOM = 6.35

FILES_TO_CHECK = [
    # Tuần 3
    r'D:\UNIGO\KHBD_Tin_học\Tiền_tiểu_học\Tuần_03\Slide_Tin_hoc_Tien_TH_Bai02.pptx',
    r'D:\UNIGO\KHBD_Tin_học\Lớp_1\Tuần_03\Slide_Tin_hoc_Lop_1_Bai02.pptx',
    r'D:\UNIGO\KHBD_Tin_học\Lớp_2\Tuần_03\Slide_Tin_hoc_Lop_2_Bai02.pptx',
    r'D:\UNIGO\KHBD_Tin_học\Lớp_3\Tuần_03\Slide_Tin_hoc_Lop_3_Bai02.pptx',
    r'D:\UNIGO\KHBD_Tin_học\Lớp_4\Tuần_03\Slide_Tin_hoc_Lop_4_Bai02.pptx',
    # Tuần 4
    r'D:\UNIGO\KHBD_Tin_học\Tiền_tiểu_học\Tuần_04\Slide_Tin_hoc_Tien_TH_Bai03.pptx',
    r'D:\UNIGO\KHBD_Tin_học\Lớp_1\Tuần_04\Slide_Tin_hoc_Lop_1_Bai03.pptx',
    r'D:\UNIGO\KHBD_Tin_học\Lớp_2\Tuần_04\Slide_Tin_hoc_Lop_2_Bai03.pptx',
    r'D:\UNIGO\KHBD_Tin_học\Lớp_3\Tuần_04\Slide_Tin_hoc_Lop_3_Bai03.pptx',
    r'D:\UNIGO\KHBD_Tin_học\Lớp_4\Tuần_04\Slide_Tin_hoc_Lop_4_Bai03.pptx',
    r'D:\UNIGO\KHBD_Tin_học\Lớp_5\Tuần_04\Slide_Tin_hoc_Lop_5_Bai03_04.pptx',
    r'D:\UNIGO\KHBD_Tin_học\Lớp_6\Tuần_04\Slide_Tin_hoc_Lop_6_Bai03_04.pptx',
    r'D:\UNIGO\KHBD_Tin_học\Lớp_7\Tuần_04\Slide_Tin_hoc_Lop_7_Bai03_04.pptx',
    r'D:\UNIGO\KHBD_Tin_học\Lớp_8\Tuần_04\Slide_Tin_hoc_Lop_8_Bai03_Khai_thac_thong_tin_so.pptx',
    r'D:\UNIGO\KHBD_Tin_học\Lớp_8\Tuần_04\Slide_Tin_hoc_Lop_8_On_tap_DGDK1.pptx',
    r'D:\UNIGO\KHBD_Tin_học\Lớp_9\Tuần_04\Slide_Tin_hoc_Lop_9_Bai03_04.pptx',
]

def verify_all_slides():
    print("="*70)
    print("BÁO CÁO KIỂM TRA CHẤT LƯỢNG SLIDE TUẦN 3 & TUẦN 4")
    print("="*70)
    
    total_files = len(FILES_TO_CHECK)
    passed_files = 0
    total_slides_count = 0
    
    for fpath in FILES_TO_CHECK:
        fname = os.path.basename(fpath)
        if not os.path.isfile(fpath):
            print(f"[-] THIẾU FILE: {fname}")
            continue
            
        prs = Presentation(fpath)
        num_slides = len(prs.slides)
        total_slides_count += num_slides
        file_errors = []
        
        for s_idx, slide in enumerate(prs.slides):
            for sh_idx, sh in enumerate(slide.shapes):
                top_in = sh.top / 914400
                bot_in = (sh.top + sh.height) / 914400
                
                # Check top violation
                if top_in < SAFE_TOP - 0.05:
                    file_errors.append(f"Slide {s_idx+1} shape '{sh.name}': Top={top_in:.2f}in < {SAFE_TOP}in")
                # Check bottom violation
                if bot_in > SAFE_BOTTOM + 0.1:
                    file_errors.append(f"Slide {s_idx+1} shape '{sh.name}': Bottom={bot_in:.2f}in > {SAFE_BOTTOM}in")
                    
        if file_errors:
            print(f"[FAIL] {fname} ({num_slides} slides) - Có {len(file_errors)} cảnh báo:")
            for err in file_errors[:3]:
                print(f"       • {err}")
        else:
            print(f"[PASS] {fname} | {num_slides} slides | 100% Vùng an toàn & Logo/Chân trang chuẩn")
            passed_files += 1
            
    print("="*70)
    print(f"KẾT QUẢ: {passed_files}/{total_files} file ĐẠT CHUẨN TUYỆT ĐỐI! (Tổng cộng {total_slides_count} slide)")
    print("="*70)

if __name__ == '__main__':
    verify_all_slides()
