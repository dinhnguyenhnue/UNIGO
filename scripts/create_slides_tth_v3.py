# -*- coding: utf-8 -*-
"""
Script tạo Slide bài giảng Student-Facing v3 — Tiền Tiểu Học
Bài 1: Máy tính xung quanh em

v3 Features:
  - Per-Bullet Images: mỗi bullet point có ảnh riêng (Grid Flashcard layout)
  - Animation tuần tự cho slide practice/activity (appear on click)
  - Game images bắt buộc cho slide trò chơi
  - Anti-Bug Checklist: safe zone, z-order, no footer, contrast
"""
import sys, io, os, copy
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import traceback

# ─── Paths ───
TEMPLATE = r'D:\UNIGO\Hệ thống mẫu văn bản\Mẫu slide UNIGO.pptx'
IMG_DIR  = r'D:\UNIGO\KHBD_Tin_học\Tiền_tiểu_học\Tuần_02\images'
OUTPUT   = r'D:\UNIGO\KHBD_Tin_học\Tiền_tiểu_học\Tuần_02\Slide_Tin_hoc_Tien_TH_Bai01_May_tinh_xung_quanh_em_v3.pptx'

# ─── Template Safe Zone (VERIFIED) ───
SAFE_TOP    = 1.15   # Y below logo (Picture 7 ends at 1.09in)
SAFE_BOTTOM = 6.35   # Y above footer (Picture 9 starts at 6.43in)
SAFE_LEFT   = 0.3
SAFE_RIGHT  = 13.0
SLIDE_W     = 13.33
SLIDE_H     = 7.50

# ─── Color Palette: Blue (Tiền TH) — High Contrast ───
PAL = {
    "primary":        "1B4F9B",
    "accent":         "2D7DD2",
    "bg":             "EBF3FE",
    "card":           "FFFFFF",
    "text_on_primary":"FFFFFF",
    "text_on_bg":     "1A2744",
    "text_on_card":   "1A2744",
}

# ─── Nội dung bài học (Student-Facing, Tiền TH) ───
SLIDES = [
    # 0: Cover
    {"type": "cover", "title": "Máy tính xung quanh em",
     "subtitle": "Tin học • Tiền Tiểu học • Bài 1",
     "image": "cover_v3.png"},

    # 1: Warmup
    {"type": "warmup", "title": "Các con ơi, nhìn xung quanh nào!",
     "content": "Con có thấy chiếc máy tính nào không?\nChỉ cho cô xem nhé!",
     "image": "cover_v3.png"},

    # 2: Learn — Grid Flashcard (3 bullets)
    {"type": "learn_grid", "title": "Đây là máy tính!",
     "bullets": [
         {"text": "Trong phòng học\ncủa chúng mình", "image": "learn1_bullet1_classroom.png"},
         {"text": "Ở nhà\n(laptop của bố mẹ)", "image": "learn1_bullet2_laptop.png"},
         {"text": "Trên tay\n(điện thoại cũng là\nmáy tính nhỏ!)", "image": "learn1_bullet3_phone.png"},
     ]},

    # 3: Learn — Grid Flashcard (4 bullets → 2×2)
    {"type": "learn_grid", "title": "Máy tính có những phần nào?",
     "bullets": [
         {"text": "Màn hình\n— để con nhìn", "image": "learn2_bullet1_monitor.png"},
         {"text": "Bàn phím\n— để con gõ chữ", "image": "learn2_bullet2_keyboard.png"},
         {"text": "Chuột\n— để con chỉ và nhấp", "image": "learn2_bullet3_mouse.png"},
         {"text": "Thân máy\n— bộ não của máy tính", "image": "learn2_bullet4_cpu.png"},
     ]},

    # 4: Learn — Grid Flashcard (4 bullets → 2×2)
    {"type": "learn_grid", "title": "Máy tính giúp chúng ta làm gì?",
     "bullets": [
         {"text": "Học bài vui vẻ", "image": "learn3_bullet1_study.png"},
         {"text": "Vẽ tranh đẹp", "image": "learn3_bullet2_draw.png"},
         {"text": "Nghe nhạc hay", "image": "learn3_bullet3_music.png"},
         {"text": "Xem hoạt hình", "image": "learn3_bullet4_cartoon.png"},
     ]},

    # 5: Practice — Animation tuần tự + images
    {"type": "practice", "title": "Cùng chơi nào!",
     "instruction": "Con hãy chỉ vào từng bộ phận\nvà nói tên cho cô nghe:",
     "items": [
         {"text": "Đâu là màn hình?", "image": "learn2_bullet1_monitor.png"},
         {"text": "Đâu là bàn phím?", "image": "learn2_bullet2_keyboard.png"},
         {"text": "Đâu là chuột?",    "image": "learn2_bullet3_mouse.png"},
     ]},

    # 6: Activity — Game with image + animation
    {"type": "activity", "title": "Thử thách nhỏ!",
     "content": "Con hãy vẽ một chiếc máy tính\nvào giấy của mình.\nNhớ vẽ đủ các bộ phận nhé!",
     "image": "activity_draw_v3.png"},

    # 7: Summary
    {"type": "summary", "title": "Hôm nay con đã biết!",
     "items": [
         "Máy tính ở xung quanh chúng ta",
         "Máy tính có: màn hình, bàn phím, chuột, thân máy",
         "Máy tính giúp con học và chơi",
     ]},

    # 8: Thanks
    {"type": "thanks", "title": "Các con giỏi lắm!",
     "content": "Hẹn gặp lại tiết sau nhé!"},
]


# ═══════════════════════════════════════════════════
# HELPER FUNCTIONS — v3
# ═══════════════════════════════════════════════════

def hex_rgb(h):
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def set_font(run, size_pt, bold=False, color_hex="333333", font_name="Arial"):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = hex_rgb(color_hex)
    run.font.name = font_name

def add_textbox(slide, left, top, width, height, text, size_pt=18, bold=False,
                color_hex="333333", alignment=PP_ALIGN.LEFT, font_name="Arial"):
    """Add textbox — always clamped to safe zone"""
    actual_top = max(top, SAFE_TOP)
    txbox = slide.shapes.add_textbox(Inches(left), Inches(actual_top), Inches(width), Inches(height))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size_pt, bold, color_hex, font_name)
    return txbox

def add_multiline_textbox(slide, left, top, width, height, lines, size_pt=18,
                           color_hex="333333", bold=False, font_name="Arial",
                           alignment=PP_ALIGN.LEFT, line_spacing_pt=None):
    """Add textbox with multiple lines"""
    actual_top = max(top, SAFE_TOP)
    txbox = slide.shapes.add_textbox(Inches(left), Inches(actual_top), Inches(width), Inches(height))
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alignment
        run = p.add_run()
        run.text = line
        set_font(run, size_pt, bold, color_hex, font_name)
        p.space_after = Pt(6)
        if line_spacing_pt:
            p.line_spacing = Pt(line_spacing_pt)
    return txbox

def add_safe_shape(slide, shape_type, left, top, width, height, fill_hex,
                   border_hex=None, send_to_back=False):
    """Add shape CLAMPED within safe zone. insert(2) for z-order."""
    actual_top = max(top, SAFE_TOP)
    actual_bottom = min(top + height, SAFE_BOTTOM)
    actual_height = max(actual_bottom - actual_top, 0.1)

    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left), Inches(actual_top), Inches(width), Inches(actual_height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(fill_hex)
    if border_hex:
        shape.line.color.rgb = hex_rgb(border_hex)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()

    if send_to_back:
        sp = shape._element
        spTree = sp.getparent()
        spTree.remove(sp)
        spTree.insert(2, sp)  # NEVER insert(0) — breaks XML schema
    return shape

def add_picture_safe(slide, img_path, left, top, width, height):
    """Add picture clamped to safe zone"""
    actual_top = max(top, SAFE_TOP)
    actual_bottom = min(top + height, SAFE_BOTTOM)
    actual_height = max(actual_bottom - actual_top, 0.5)
    try:
        return slide.shapes.add_picture(img_path, Inches(left), Inches(actual_top),
                                         Inches(width), Inches(actual_height))
    except Exception as e:
        print(f"  ⚠️ Cannot add picture {img_path}: {e}")
        return None

def add_slide_transition(slide, transition_type="fade"):
    """Add XML transition"""
    transitions = {
        'fade': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>',
        'push': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:push dir="l"/></p:transition>',
        'wipe': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:wipe/></p:transition>',
        'cover': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:cover/></p:transition>',
    }
    xml_str = transitions.get(transition_type, transitions['fade'])
    slide._element.append(etree.fromstring(xml_str))


def add_appear_animation(slide, shapes_with_click):
    """
    Add appear animation to multiple shapes, each triggered on a separate click.
    shapes_with_click: list of (shape_or_list, click_index) tuples
    Each shape starts hidden and appears on its click_index click.
    """
    nsmap = {
        'p':  'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a':  'http://schemas.openxmlformats.org/drawingml/2006/main',
    }

    # Build timing XML
    timing = etree.SubElement(slide._element, '{%s}timing' % nsmap['p'])
    tnLst = etree.SubElement(timing, '{%s}tnLst' % nsmap['p'])
    par_root = etree.SubElement(tnLst, '{%s}par' % nsmap['p'])
    cTn_root = etree.SubElement(par_root, '{%s}cTn' % nsmap['p'])
    cTn_root.set('id', '1')
    cTn_root.set('dur', 'indefinite')
    cTn_root.set('restart', 'never')
    cTn_root.set('nodeType', 'tmRoot')

    childTnLst_root = etree.SubElement(cTn_root, '{%s}childTnLst' % nsmap['p'])

    seq = etree.SubElement(childTnLst_root, '{%s}seq' % nsmap['p'])
    seq.set('concurrent', '1')
    seq.set('nextAc', 'seek')

    cTn_seq = etree.SubElement(seq, '{%s}cTn' % nsmap['p'])
    cTn_seq.set('id', '2')
    cTn_seq.set('dur', 'indefinite')
    cTn_seq.set('nodeType', 'mainSeq')

    childTnLst_seq = etree.SubElement(cTn_seq, '{%s}childTnLst' % nsmap['p'])

    # prevCondLst and nextCondLst for the sequence
    prevCond = etree.SubElement(seq, '{%s}prevCondLst' % nsmap['p'])
    cond_prev = etree.SubElement(prevCond, '{%s}cond' % nsmap['p'])
    cond_prev.set('evt', 'onPrev')
    cond_prev.set('delay', '0')
    tgtEl_prev = etree.SubElement(cond_prev, '{%s}tgtEl' % nsmap['p'])
    etree.SubElement(tgtEl_prev, '{%s}sldTgt' % nsmap['p'])

    nextCond = etree.SubElement(seq, '{%s}nextCondLst' % nsmap['p'])
    cond_next = etree.SubElement(nextCond, '{%s}cond' % nsmap['p'])
    cond_next.set('evt', 'onNext')
    cond_next.set('delay', '0')
    tgtEl_next = etree.SubElement(cond_next, '{%s}tgtEl' % nsmap['p'])
    etree.SubElement(tgtEl_next, '{%s}sldTgt' % nsmap['p'])

    id_counter = 3

    for shape_group, click_idx in shapes_with_click:
        # Wrap single shape in list
        if not isinstance(shape_group, (list, tuple)):
            shape_group = [shape_group]

        par_click = etree.SubElement(childTnLst_seq, '{%s}par' % nsmap['p'])
        cTn_click = etree.SubElement(par_click, '{%s}cTn' % nsmap['p'])
        cTn_click.set('id', str(id_counter)); id_counter += 1
        cTn_click.set('fill', 'hold')

        stCondLst = etree.SubElement(cTn_click, '{%s}stCondLst' % nsmap['p'])
        cond = etree.SubElement(stCondLst, '{%s}cond' % nsmap['p'])
        cond.set('delay', '0')

        childTnLst_click = etree.SubElement(cTn_click, '{%s}childTnLst' % nsmap['p'])

        for shape in shape_group:
            if shape is None:
                continue
            sp_id = shape.shape_id

            par_anim = etree.SubElement(childTnLst_click, '{%s}par' % nsmap['p'])
            cTn_anim = etree.SubElement(par_anim, '{%s}cTn' % nsmap['p'])
            cTn_anim.set('id', str(id_counter)); id_counter += 1
            cTn_anim.set('presetID', '1')  # Appear
            cTn_anim.set('presetClass', 'entr')
            cTn_anim.set('presetSubtype', '0')
            cTn_anim.set('fill', 'hold')

            stCond_anim = etree.SubElement(cTn_anim, '{%s}stCondLst' % nsmap['p'])
            cond_anim = etree.SubElement(stCond_anim, '{%s}cond' % nsmap['p'])
            cond_anim.set('delay', '0')

            childTnLst_anim = etree.SubElement(cTn_anim, '{%s}childTnLst' % nsmap['p'])

            p_set = etree.SubElement(childTnLst_anim, '{%s}set' % nsmap['p'])
            cBhvr = etree.SubElement(p_set, '{%s}cBhvr' % nsmap['p'])
            cTn_set = etree.SubElement(cBhvr, '{%s}cTn' % nsmap['p'])
            cTn_set.set('id', str(id_counter)); id_counter += 1
            cTn_set.set('dur', '1')
            cTn_set.set('fill', 'hold')

            stCond_set = etree.SubElement(cTn_set, '{%s}stCondLst' % nsmap['p'])
            cond_set = etree.SubElement(stCond_set, '{%s}cond' % nsmap['p'])
            cond_set.set('delay', '0')

            tgtEl = etree.SubElement(cBhvr, '{%s}tgtEl' % nsmap['p'])
            spTgt = etree.SubElement(tgtEl, '{%s}spTgt' % nsmap['p'])
            spTgt.set('spid', str(sp_id))

            attrNameLst = etree.SubElement(cBhvr, '{%s}attrNameLst' % nsmap['p'])
            attrName = etree.SubElement(attrNameLst, '{%s}attrName' % nsmap['p'])
            attrName.text = 'style.visibility'

            p_to = etree.SubElement(p_set, '{%s}to' % nsmap['p'])
            strVal = etree.SubElement(p_to, '{%s}strVal' % nsmap['p'])
            strVal.set('val', 'visible')


def img_path(filename):
    """Get full path to image"""
    return os.path.join(IMG_DIR, filename)


def remove_template_slides(prs):
    """Remove all existing slides from template"""
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        sldId = prs.slides._sldIdLst[0]
        prs.slides._sldIdLst.remove(sldId)


# ═══════════════════════════════════════════════════
# SLIDE BUILDERS — v3 (Per-Bullet Images + Animation)
# ═══════════════════════════════════════════════════

def build_cover(prs, data, layout):
    """Slide bìa: banner trong safe zone + title + image"""
    slide = prs.slides.add_slide(layout)

    # Background in safe zone
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP,
                   PAL["primary"], send_to_back=True)
    # Accent stripe
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, 1.8, PAL["accent"],
                   send_to_back=True)

    # Title
    add_textbox(slide, 0.8, 2.0, 7.5, 1.5, data["title"],
                size_pt=36, bold=True, color_hex=PAL["text_on_primary"],
                alignment=PP_ALIGN.LEFT)
    # Subtitle
    add_textbox(slide, 0.8, 3.8, 7.5, 0.8, data.get("subtitle", ""),
                size_pt=20, color_hex=PAL["text_on_primary"])

    # Cover image on right
    if data.get("image"):
        add_picture_safe(slide, img_path(data["image"]),
                         SLIDE_W - 4.5, SAFE_TOP + 0.3, 3.8, 3.8)

    add_slide_transition(slide, "fade")
    return slide


def build_warmup(prs, data, layout):
    """Slide khởi động: text + ảnh lớn bên phải"""
    slide = prs.slides.add_slide(layout)

    # Light bg
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP,
                   PAL["bg"], send_to_back=True)
    # Badge
    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, PAL["primary"])
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4,
                "  TIỀN TIỂU HỌC • BÀI 1. MÁY TÍNH XUNG QUANH EM",
                size_pt=13, bold=True, color_hex=PAL["text_on_primary"])

    # Title
    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, 8, 0.65, data["title"],
                size_pt=26, bold=True, color_hex=PAL["text_on_bg"])

    # Content card
    content_y = title_y + 0.75
    card_w = 7.5
    card_h = SAFE_BOTTOM - content_y - 0.15
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.4, content_y, card_w, card_h, PAL["card"])
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0.4, content_y, 0.12, card_h, PAL["accent"])

    lines = data.get("content", "").split("\n")
    add_multiline_textbox(slide, 0.8, content_y+0.2, card_w-0.6, card_h-0.4, lines,
                          size_pt=22, color_hex=PAL["text_on_card"])

    # Image right
    if data.get("image"):
        add_picture_safe(slide, img_path(data["image"]),
                         8.3, content_y, SAFE_RIGHT-8.3-0.2, card_h)

    add_slide_transition(slide, "push")
    return slide


def build_learn_grid(prs, data, layout):
    """
    Slide nội dung với Per-Bullet Images — Layout A: Grid Flashcard
    Mỗi bullet = ảnh (2in×2in) trên + text dưới
    2-3 bullets: 1 row × N cols
    4 bullets: 2 rows × 2 cols
    """
    slide = prs.slides.add_slide(layout)
    bullets = data["bullets"]
    n = len(bullets)

    # Light bg
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP,
                   PAL["bg"], send_to_back=True)

    # Badge
    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, PAL["primary"])
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4,
                "  TIỀN TIỂU HỌC • BÀI 1. MÁY TÍNH XUNG QUANH EM",
                size_pt=13, bold=True, color_hex=PAL["text_on_primary"])

    # Title
    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, SLIDE_W-1, 0.65, data["title"],
                size_pt=26, bold=True, color_hex=PAL["text_on_bg"])

    # Grid area
    grid_top = title_y + 0.8
    grid_bottom = SAFE_BOTTOM - 0.1
    avail_h = grid_bottom - grid_top

    if n <= 3:
        cols, rows = n, 1
    else:
        cols, rows = 2, 2

    # Card dimensions — adaptive based on grid size
    gap = 0.25
    total_gap_x = gap * (cols - 1)
    card_w = (SLIDE_W - 1.0 - total_gap_x) / cols
    if rows == 1:
        img_h = 2.0   # Full size for single row
        text_h = 0.9
    else:
        img_h = 1.1   # Compact for 2×2 grid to fit in safe zone
        text_h = 0.5
    card_h = img_h + text_h + 0.15
    total_gap_y = gap * (rows - 1)
    start_x = (SLIDE_W - (cols * card_w + total_gap_x)) / 2
    start_y = grid_top + (avail_h - (rows * card_h + total_gap_y)) / 2
    start_y = max(start_y, grid_top)

    for i, bullet in enumerate(bullets):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + gap)

        if y + card_h > SAFE_BOTTOM:
            break

        # White card background
        add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h,
                       PAL["card"], border_hex="D0D5DD")

        # Accent bar top of card
        add_safe_shape(slide, MSO_SHAPE.RECTANGLE, x, y, card_w, 0.06, PAL["accent"])

        # Image (2in × 2in, centered in card)
        img_w = min(2.0, card_w - 0.4)
        img_x = x + (card_w - img_w) / 2
        img_y = y + 0.15
        if bullet.get("image") and os.path.isfile(img_path(bullet["image"])):
            add_picture_safe(slide, img_path(bullet["image"]),
                             img_x, img_y, img_w, img_h)

        # Text below image (centered)
        text_y = y + img_h + 0.2
        add_textbox(slide, x + 0.15, text_y, card_w - 0.3, text_h,
                    bullet["text"], size_pt=16, bold=False,
                    color_hex=PAL["text_on_card"], alignment=PP_ALIGN.CENTER)

    add_slide_transition(slide, "wipe")
    return slide


def build_practice(prs, data, layout):
    """
    Slide luyện tập: items xuất hiện tuần tự (animation appear on click)
    Mỗi item = ảnh (~2in) + card text → hiện cùng lúc (1 click)
    """
    slide = prs.slides.add_slide(layout)

    # Light bg
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP,
                   PAL["bg"], send_to_back=True)

    # Accent badge
    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, PAL["accent"])
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4,
                "  TIỀN TIỂU HỌC • LUYỆN TẬP",
                size_pt=13, bold=True, color_hex=PAL["text_on_primary"])

    # Title
    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, SLIDE_W-1, 0.65, data["title"],
                size_pt=26, bold=True, color_hex=PAL["text_on_bg"])

    # Instruction
    instr_y = title_y + 0.75
    instr_lines = data.get("instruction", "").split("\n")
    add_multiline_textbox(slide, 0.5, instr_y, SLIDE_W-1, 0.7, instr_lines,
                          size_pt=18, color_hex=PAL["text_on_bg"])

    # Items with animation
    items = data.get("items", [])
    item_y_start = instr_y + 0.9
    item_h = 2.0
    gap = 0.2
    n_items = len(items)
    item_w = (SLIDE_W - 1.0 - gap*(n_items-1)) / n_items
    start_x = 0.5

    animation_groups = []

    for i, item in enumerate(items):
        x = start_x + i * (item_w + gap)
        y = item_y_start

        # Card
        card = add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, item_w, item_h,
                              PAL["card"], border_hex=PAL["accent"])

        # Image top (2in × 1.5in)
        pic = None
        if item.get("image") and os.path.isfile(img_path(item["image"])):
            img_w = min(1.8, item_w - 0.4)
            img_x = x + (item_w - img_w) / 2
            pic = add_picture_safe(slide, img_path(item["image"]),
                                   img_x, y + 0.1, img_w, 1.2)

        # Text below image
        txt = add_textbox(slide, x + 0.15, y + 1.4, item_w - 0.3, 0.5,
                          item["text"], size_pt=16, bold=True,
                          color_hex=PAL["text_on_card"], alignment=PP_ALIGN.CENTER)

        # Group for animation: card + image + text appear together on click
        group = [s for s in [card, pic, txt] if s is not None]
        animation_groups.append((group, i))

    # Add appear animation
    if animation_groups:
        add_appear_animation(slide, animation_groups)

    add_slide_transition(slide, "wipe")
    return slide


def build_activity(prs, data, layout):
    """
    Slide trò chơi: image lớn bắt buộc + text hướng dẫn
    """
    slide = prs.slides.add_slide(layout)

    # Light bg
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP,
                   PAL["bg"], send_to_back=True)

    # Accent badge
    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, PAL["accent"])
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4,
                "  TIỀN TIỂU HỌC • THỬ THÁCH",
                size_pt=13, bold=True, color_hex=PAL["text_on_primary"])

    # Title
    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, SLIDE_W-1, 0.65, data["title"],
                size_pt=26, bold=True, color_hex=PAL["text_on_bg"])

    # Content area: text left + big image right
    content_y = title_y + 0.75
    content_h = SAFE_BOTTOM - content_y - 0.15

    # Text card (left 55%)
    card_w = 6.5
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.4, content_y, card_w, content_h, PAL["card"])
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0.4, content_y, 0.12, content_h, PAL["accent"])

    lines = data.get("content", "").split("\n")
    add_multiline_textbox(slide, 0.8, content_y+0.2, card_w-0.6, content_h-0.4, lines,
                          size_pt=22, color_hex=PAL["text_on_card"])

    # Big game image (right, ~3.5in)
    if data.get("image") and os.path.isfile(img_path(data["image"])):
        img_x = 7.3
        img_w = SAFE_RIGHT - img_x - 0.3
        add_picture_safe(slide, img_path(data["image"]),
                         img_x, content_y + 0.1, img_w, content_h - 0.2)

    add_slide_transition(slide, "cover")
    return slide


def build_summary(prs, data, layout):
    """Slide tổng kết: panel trắng + items"""
    slide = prs.slides.add_slide(layout)

    # Light bg
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP,
                   PAL["bg"], send_to_back=True)

    # Panel trắng
    px, py = 0.6, SAFE_TOP + 0.15
    pw = SLIDE_W - 1.2
    ph = SAFE_BOTTOM - py - 0.15
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, px, py, pw, ph, PAL["card"])
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, px, py, pw, 0.06, PAL["primary"])

    # Title
    add_textbox(slide, px+0.3, py+0.2, pw-0.6, 0.65, data["title"],
                size_pt=24, bold=True, color_hex=PAL["primary"], alignment=PP_ALIGN.CENTER)

    # Divider
    div_y = py + 0.95
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, px+2, div_y, pw-4, 0.03, PAL["accent"])

    # Summary items
    items = data.get("items", [])
    for i, item in enumerate(items[:4]):
        y = div_y + 0.25 + i * 1.1
        if y + 0.8 > SAFE_BOTTOM:
            break
        add_safe_shape(slide, MSO_SHAPE.RECTANGLE, px+0.4, y, 0.08, 0.7, PAL["accent"])
        add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, px+0.6, y, pw-1.2, 0.7, PAL["bg"])
        add_textbox(slide, px+0.9, y+0.15, pw-1.8, 0.4, item,
                    size_pt=18, color_hex=PAL["text_on_bg"])

    add_slide_transition(slide, "fade")
    return slide


def build_thanks(prs, data, layout):
    """Slide cảm ơn"""
    slide = prs.slides.add_slide(layout)

    # Primary bg in safe zone
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP,
                   PAL["primary"], send_to_back=True)

    # Big title
    add_textbox(slide, 1.0, SAFE_TOP+0.8, SLIDE_W-2, 1.2, data["title"],
                size_pt=36, bold=True, color_hex=PAL["text_on_primary"],
                alignment=PP_ALIGN.CENTER)

    # Content card
    content = data.get("content", "")
    if content:
        card_w = 8.0
        card_x = (SLIDE_W - card_w) / 2
        card_y = SAFE_TOP + 2.5
        card_h = min(2.0, SAFE_BOTTOM - card_y - 0.2)
        add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, card_w, card_h, PAL["card"])
        lines = content.split("\n")
        add_multiline_textbox(slide, card_x+0.5, card_y+0.3, card_w-1, card_h-0.6, lines,
                              size_pt=20, color_hex=PAL["text_on_card"],
                              alignment=PP_ALIGN.CENTER)

    add_slide_transition(slide, "fade")
    return slide


# ═══════════════════════════════════════════════════
# VERIFICATION — Anti-Bug Checklist
# ═══════════════════════════════════════════════════

def verify_slides(prs):
    """Verify all slides pass anti-bug checklist"""
    issues = []
    for si, sl in enumerate(prs.slides):
        for sh in sl.shapes:
            top_in = sh.top / 914400
            bottom_in = (sh.top + sh.height) / 914400
            # Rule 1: No shape above safe top
            if top_in < SAFE_TOP - 0.02 and sh.name not in ('Picture 7',):
                issues.append(f"Slide {si}: {sh.name} top={top_in:.2f}in < {SAFE_TOP}in (che logo)")
            # Rule 2: No shape below safe bottom
            if bottom_in > SAFE_BOTTOM + 0.02 and sh.name not in ('Picture 9',):
                issues.append(f"Slide {si}: {sh.name} bottom={bottom_in:.2f}in > {SAFE_BOTTOM}in (che footer)")

    if issues:
        print(f"\n⚠️ VERIFICATION: {len(issues)} issues found:")
        for iss in issues:
            print(f"  ❌ {iss}")
    else:
        print(f"\n✅ VERIFICATION: All {len(prs.slides)} slides PASS anti-bug checklist")
    return len(issues) == 0


# ═══════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════

def main():
    print("=" * 60)
    print("  Tạo Slide TTH v3 — Per-Bullet Images + Animation")
    print("  Bài 1: Máy tính xung quanh em")
    print("=" * 60)

    # Load template
    prs = Presentation(TEMPLATE)
    remove_template_slides(prs)
    layout = prs.slide_layouts[6]  # Blank layout

    # Build slides
    for i, data in enumerate(SLIDES):
        stype = data["type"]
        try:
            if stype == "cover":
                build_cover(prs, data, layout)
            elif stype == "warmup":
                build_warmup(prs, data, layout)
            elif stype == "learn_grid":
                build_learn_grid(prs, data, layout)
            elif stype == "practice":
                build_practice(prs, data, layout)
            elif stype == "activity":
                build_activity(prs, data, layout)
            elif stype == "summary":
                build_summary(prs, data, layout)
            elif stype == "thanks":
                build_thanks(prs, data, layout)
            print(f"  ✅ Slide {i+1}/{len(SLIDES)}: {data.get('title', stype)}")
        except Exception as e:
            print(f"  ❌ Slide {i+1}: {e}")
            traceback.print_exc()

    # Verify
    verify_slides(prs)

    # Save
    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
    try:
        prs.save(OUTPUT)
        print(f"\n✅ Saved: {OUTPUT}")
        print(f"   Total slides: {len(prs.slides)}")
    except PermissionError:
        alt = OUTPUT.replace(".pptx", "_alt.pptx")
        prs.save(alt)
        print(f"\n⚠️ File locked → saved as: {alt}")
    except Exception as e:
        print(f"\n❌ Save error: {e}")
        traceback.print_exc()


if __name__ == "__main__":
    main()
