# -*- coding: utf-8 -*-
"""
Hệ thống Tạo Slide Bài Giảng Nghệ Thuật Cao Cấp - Tin học Lớp 5 & Lớp 7 (Tuần 04)
Đạt chuẩn 100% UNIGO:
- Sử dụng Hình ảnh AI 3D/Cartoon sống động, bắt mắt, không dùng ảnh scan thô ráp.
- Thiết kế bố cục chuẩn Visual-First, màu sắc hiện đại, tương phản sắc nét.
- Tuân thủ Vùng An Toàn (Y: 1.15in -> 6.35in), bảo tồn Logo & Chân trang UNIGO.
"""
import sys, os
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

TEMPLATE    = r'D:\UNIGO\Hệ thống mẫu văn bản\Mẫu slide UNIGO.pptx'
SAFE_TOP    = 1.15
SAFE_BOTTOM = 6.35
SAFE_LEFT   = 0.4
SLIDE_W     = 13.33
SLIDE_H     = 7.50

# Palettes
PAL_L5_T3 = {"primary": "BE185D", "accent": "EC4899", "bg": "FDF2F8", "card": "FFFFFF", "text_on_card": "0F172A", "text_on_bg": "831843"}
PAL_L5_T4 = {"primary": "9F1239", "accent": "F43F5E", "bg": "FFF1F2", "card": "FFFFFF", "text_on_card": "0F172A", "text_on_bg": "881337"}
PAL_L7_T3 = {"primary": "0369A1", "accent": "0EA5E9", "bg": "F0F9FF", "card": "FFFFFF", "text_on_card": "0F172A", "text_on_bg": "0C4A6E"}
PAL_L7_T4 = {"primary": "0E7490", "accent": "06B6D4", "bg": "ECFEFF", "card": "FFFFFF", "text_on_card": "0F172A", "text_on_bg": "164E63"}

IMG5_DIR = r"D:\UNIGO\KHBD_Tin_học\Lớp_5\Tuần_04\images"
IMG7_DIR = r"D:\UNIGO\KHBD_Tin_học\Lớp_7\Tuần_04\images"

def hex_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_safe_shape(slide, shape_type, left, top, width, height, fill_hex, border_hex=None, send_to_back=False):
    actual_top = max(top, SAFE_TOP)
    actual_bottom = min(top + height, SAFE_BOTTOM)
    actual_height = max(actual_bottom - actual_top, 0.1)

    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(actual_top), Inches(width), Inches(actual_height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(fill_hex)
    if border_hex:
        shape.line.color.rgb = hex_rgb(border_hex)
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()

    if send_to_back:
        spTree = shape._element.getparent()
        spTree.remove(shape._element)
        spTree.insert(2, shape._element)
    return shape

def add_textbox(slide, left, top, width, height, text, size_pt=18, bold=False, color_hex="0F172A", alignment=PP_ALIGN.LEFT, font_name="Arial"):
    actual_top = max(top, SAFE_TOP)
    actual_bottom = min(top + height, SAFE_BOTTOM)
    actual_height = max(actual_bottom - actual_top, 0.1)

    txBox = slide.shapes.add_textbox(Inches(left), Inches(actual_top), Inches(width), Inches(actual_height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)

    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(size_pt)
    p.font.bold = bold
    p.font.color.rgb = hex_rgb(color_hex)
    p.alignment = alignment
    return txBox

def add_picture_safe(slide, img_path, left, top, width, height):
    if not os.path.exists(img_path):
        return None
    actual_top = max(top, SAFE_TOP)
    actual_bottom = min(top + height, SAFE_BOTTOM)
    actual_height = max(actual_bottom - actual_top, 0.1)
    try:
        pic = slide.shapes.add_picture(img_path, Inches(left), Inches(actual_top), Inches(width), Inches(actual_height))
        return pic
    except Exception as e:
        print(f"Error picture {img_path}: {e}")
        return None

def add_slide_transition(slide, transition_type="fade"):
    slide_xml = slide._element
    existing_trans = slide_xml.xpath('./p:transition')
    if existing_trans:
        for t in existing_trans:
            slide_xml.remove(t)
    trans = etree.SubElement(slide_xml, '{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
    trans.set('spd', 'med')
    etree.SubElement(trans, f'{{http://schemas.openxmlformats.org/presentationml/2006/main}}{transition_type}')

# ─── BUILDERS ───

def build_cover(prs, data, pal, layout):
    slide = prs.slides.add_slide(layout)
    add_slide_transition(slide, "wipe")
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM - SAFE_TOP, pal["primary"], send_to_back=True)
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.2, 1.6, 10.93, 4.0, "FFFFFF", border_hex=pal["accent"])
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.2, 1.6, 10.93, 0.25, pal["accent"])
    
    add_textbox(slide, 1.5, 2.0, 10.33, 0.5, "TRƯỜNG TIỂU HỌC & THCS UNIGO", size_pt=15, bold=True, color_hex=pal["primary"], alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1.5, 2.6, 10.33, 1.4, data["title"], size_pt=28, bold=True, color_hex=pal["text_on_card"], alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1.5, 4.1, 10.33, 0.6, data["subtitle"], size_pt=18, bold=False, color_hex=pal["primary"], alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1.5, 4.8, 10.33, 0.4, "GV: Đậu Đình Nguyên • Bộ môn Tin học", size_pt=14, bold=False, color_hex="64748B", alignment=PP_ALIGN.CENTER)

def build_section_banner(prs, data, pal, layout):
    slide = prs.slides.add_slide(layout)
    add_slide_transition(slide, "push")
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM - SAFE_TOP, pal["bg"], send_to_back=True)
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.5, 1.8, 10.33, 3.8, pal["primary"])
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.5, 1.8, 0.4, 3.8, pal["accent"])
    
    add_textbox(slide, 2.3, 2.2, 8.8, 0.6, data.get("badge", "KHÁM PHÁ KIẾN THỨC MỚI"), size_pt=16, bold=True, color_hex="FEF08A", alignment=PP_ALIGN.LEFT)
    add_textbox(slide, 2.3, 2.9, 8.8, 1.5, data["title"], size_pt=26, bold=True, color_hex="FFFFFF", alignment=PP_ALIGN.LEFT)
    if "desc" in data:
        add_textbox(slide, 2.3, 4.4, 8.8, 0.8, data["desc"], size_pt=16, bold=False, color_hex="F1F5F9", alignment=PP_ALIGN.LEFT)

def build_hero_example(prs, data, pal, layout):
    """Layout 5 & 15: Cột trái (38% Text) | Cột phải (62% Hero Image)"""
    slide = prs.slides.add_slide(layout)
    add_slide_transition(slide, "fade")
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM - SAFE_TOP, pal["bg"], send_to_back=True)
    
    # Header
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, SAFE_LEFT, 1.25, 3.2, 0.38, pal["primary"])
    add_textbox(slide, SAFE_LEFT, 1.25, 3.2, 0.38, data.get("badge", "HOẠT ĐỘNG KHÁM PHÁ"), size_pt=13, bold=True, color_hex="FFFFFF", alignment=PP_ALIGN.CENTER)
    add_textbox(slide, SAFE_LEFT + 3.4, 1.22, 8.8, 0.45, data["title"], size_pt=21, bold=True, color_hex=pal["text_on_bg"], alignment=PP_ALIGN.LEFT)
    
    # Left Column: Card Text
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, SAFE_LEFT, 1.8, 4.4, 4.35, "FFFFFF", border_hex=pal["accent"])
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, SAFE_LEFT, 1.8, 0.18, 4.35, pal["accent"])
    
    text_content = data["text"]
    lines = [l.strip() for l in text_content.split('\n') if l.strip()]
    top_pos = 1.95
    for l in lines:
        is_heading = l.startswith('•') or (l[0].isdigit() and l[1] in ['.', ')'])
        add_textbox(slide, SAFE_LEFT + 0.3, top_pos, 3.95, 0.7, l, size_pt=15 if is_heading else 14, bold=is_heading, color_hex=pal["text_on_card"])
        top_pos += 0.75
    
    # Right Column: Hero Image
    img_path = data.get("img")
    if img_path and os.path.exists(img_path):
        add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, SAFE_LEFT + 4.7, 1.8, 7.7, 4.35, "FFFFFF", border_hex="CBD5E1")
        add_picture_safe(slide, img_path, SAFE_LEFT + 4.8, 1.9, 7.5, 4.15)

def build_full_hero_slide(prs, data, pal, layout):
    """Layout Hero Image to tràn slide kèm hộp thoại câu hỏi dẫn dắt"""
    slide = prs.slides.add_slide(layout)
    add_slide_transition(slide, "fade")
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM - SAFE_TOP, pal["bg"], send_to_back=True)
    
    # Header
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, SAFE_LEFT, 1.25, 3.2, 0.38, pal["primary"])
    add_textbox(slide, SAFE_LEFT, 1.25, 3.2, 0.38, data.get("badge", "KHỞI ĐỘNG & TÌNH HUỐNG"), size_pt=13, bold=True, color_hex="FFFFFF", alignment=PP_ALIGN.CENTER)
    add_textbox(slide, SAFE_LEFT + 3.4, 1.22, 8.8, 0.45, data["title"], size_pt=21, bold=True, color_hex=pal["text_on_bg"], alignment=PP_ALIGN.LEFT)
    
    # Hero Image in Top/Middle (12.3in width x 3.2in height)
    img_path = data.get("img")
    if img_path and os.path.exists(img_path):
        add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, SAFE_LEFT, 1.8, 12.3, 3.25, "FFFFFF", border_hex=pal["accent"])
        add_picture_safe(slide, img_path, SAFE_LEFT + 0.1, 1.85, 12.1, 3.15)
        
    # Question Card below image (12.3in width x 1.0in height)
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, SAFE_LEFT, 5.15, 12.3, 1.05, "FFFFFF", border_hex=pal["primary"])
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, SAFE_LEFT, 5.15, 0.2, 1.05, pal["primary"])
    add_textbox(slide, SAFE_LEFT + 0.35, 5.25, 11.6, 0.85, data["question"], size_pt=15, bold=True, color_hex=pal["text_on_card"])

def build_arrow_badges_diagram(prs, data, pal, layout):
    """Layout 19: Cột trái (Badges quy trình) | Cột phải (Ảnh Hero)"""
    slide = prs.slides.add_slide(layout)
    add_slide_transition(slide, "fade")
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM - SAFE_TOP, pal["bg"], send_to_back=True)
    
    # Header
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, SAFE_LEFT, 1.25, 3.2, 0.38, pal["primary"])
    add_textbox(slide, SAFE_LEFT, 1.25, 3.2, 0.38, data.get("badge", "QUY TRÌNH & THAO TÁC"), size_pt=13, bold=True, color_hex="FFFFFF", alignment=PP_ALIGN.CENTER)
    add_textbox(slide, SAFE_LEFT + 3.4, 1.22, 8.8, 0.45, data["title"], size_pt=21, bold=True, color_hex=pal["text_on_bg"], alignment=PP_ALIGN.LEFT)
    
    badges = data.get("badges", [])
    badge_colors = [pal["primary"], pal["accent"], "0D9488", "D97706", "7C3AED"]
    
    # Left column: Badges
    cur_top = 1.8
    card_h = 0.95 if len(badges) <= 4 else 0.75
    for i, b in enumerate(badges):
        b_col = badge_colors[i % len(badge_colors)]
        add_safe_shape(slide, MSO_SHAPE.CHEVRON, SAFE_LEFT, cur_top, 5.2, card_h, b_col)
        add_textbox(slide, SAFE_LEFT + 0.3, cur_top + 0.08, 4.6, card_h - 0.15, b, size_pt=14, bold=True, color_hex="FFFFFF", alignment=PP_ALIGN.LEFT)
        cur_top += (card_h + 0.12)
        
    # Right column: Image
    img_path = data.get("img")
    if img_path and os.path.exists(img_path):
        add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, SAFE_LEFT + 5.5, 1.8, 6.9, 4.35, "FFFFFF", border_hex="CBD5E1")
        add_picture_safe(slide, img_path, SAFE_LEFT + 5.6, 1.9, 6.7, 4.15)

def build_three_cards_diagram(prs, data, pal, layout):
    """Layout 3 Thẻ Card trực quan đối sánh / phân loại ngang"""
    slide = prs.slides.add_slide(layout)
    add_slide_transition(slide, "fade")
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM - SAFE_TOP, pal["bg"], send_to_back=True)
    
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, SAFE_LEFT, 1.25, 3.2, 0.38, pal["primary"])
    add_textbox(slide, SAFE_LEFT, 1.25, 3.2, 0.38, data.get("badge", "KIẾN THỨC CỐT LÕI"), size_pt=13, bold=True, color_hex="FFFFFF", alignment=PP_ALIGN.CENTER)
    add_textbox(slide, SAFE_LEFT + 3.4, 1.22, 8.8, 0.45, data["title"], size_pt=21, bold=True, color_hex=pal["text_on_bg"], alignment=PP_ALIGN.LEFT)
    
    cards = data.get("cards", [])
    card_w = (12.3 - (len(cards)-1)*0.3) / len(cards)
    colors = [pal["primary"], "0D9488", "D97706", "7C3AED"]
    
    for i, cd in enumerate(cards):
        c_left = SAFE_LEFT + i*(card_w + 0.3)
        add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, c_left, 1.8, card_w, 4.35, "FFFFFF", border_hex=colors[i % len(colors)])
        add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, c_left, 1.8, card_w, 0.5, colors[i % len(colors)])
        add_textbox(slide, c_left, 1.85, card_w, 0.45, cd["header"], size_pt=15, bold=True, color_hex="FFFFFF", alignment=PP_ALIGN.CENTER)
        
        # Bullets / Desc inside card
        lines = cd["text"].split('\n')
        cur_y = 2.45
        for l in lines:
            if l.strip():
                add_textbox(slide, c_left + 0.15, cur_y, card_w - 0.3, 0.65, l.strip(), size_pt=14, bold=False, color_hex=pal["text_on_card"])
                cur_y += 0.7

def build_conclusion_orange(prs, data, pal, layout):
    """Layout 14: Bóng thoại câu hỏi + Gợi ý + HỘP KẾT LUẬN CAM BO GÓC"""
    slide = prs.slides.add_slide(layout)
    add_slide_transition(slide, "fade")
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM - SAFE_TOP, pal["bg"], send_to_back=True)
    
    # Question Bubble
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, SAFE_LEFT, 1.3, 12.3, 0.9, pal["primary"])
    add_textbox(slide, SAFE_LEFT + 0.3, 1.35, 11.7, 0.8, "❓ " + data["question"], size_pt=18, bold=True, color_hex="FFFFFF", alignment=PP_ALIGN.LEFT)
    
    # Suggestions & Image
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, SAFE_LEFT, 2.3, 7.0, 2.4, "FFFFFF", border_hex=pal["accent"])
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, SAFE_LEFT, 2.3, 0.18, 2.4, pal["accent"])
    
    sugg_top = 2.45
    for sg in data["suggestions"]:
        add_textbox(slide, SAFE_LEFT + 0.3, sugg_top, 6.5, 0.9, sg, size_pt=15, bold=False, color_hex=pal["text_on_card"])
        sugg_top += 1.05
        
    img_path = data.get("img")
    if img_path and os.path.exists(img_path):
        add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, SAFE_LEFT + 7.2, 2.3, 5.1, 2.4, "FFFFFF", border_hex="CBD5E1")
        add_picture_safe(slide, img_path, SAFE_LEFT + 7.3, 2.35, 4.9, 2.3)
    
    # Orange Conclusion Box
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, SAFE_LEFT, 4.85, 12.3, 1.3, "EA580C")
    add_textbox(slide, SAFE_LEFT + 0.3, 4.9, 11.7, 0.35, "📌 GHI NHỚ TRỌNG TÂM (SGK)", size_pt=14, bold=True, color_hex="FEF08A", alignment=PP_ALIGN.LEFT)
    add_textbox(slide, SAFE_LEFT + 0.3, 5.25, 11.7, 0.85, data["conclusion"], size_pt=16, bold=True, color_hex="FFFFFF", alignment=PP_ALIGN.LEFT)

def build_quiz_mascot(prs, data, pal, layout):
    """Layout 10: Trắc nghiệm 4 đáp án A-B-C-D"""
    slide = prs.slides.add_slide(layout)
    add_slide_transition(slide, "fade")
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM - SAFE_TOP, pal["bg"], send_to_back=True)
    
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, SAFE_LEFT, 1.25, 3.2, 0.38, "EA580C")
    add_textbox(slide, SAFE_LEFT, 1.25, 3.2, 0.38, "LUYỆN TẬP NHANH", size_pt=13, bold=True, color_hex="FFFFFF", alignment=PP_ALIGN.CENTER)
    
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, SAFE_LEFT, 1.75, 12.3, 1.0, "FFFFFF", border_hex=pal["primary"])
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, SAFE_LEFT, 1.75, 0.2, 1.0, pal["primary"])
    add_textbox(slide, SAFE_LEFT + 0.3, 1.85, 11.7, 0.8, "❓ " + data["question"], size_pt=17, bold=True, color_hex=pal["text_on_card"])
    
    opts = data["options"]
    card_w = 5.95
    card_h = 1.45
    positions = [(SAFE_LEFT, 2.9), (SAFE_LEFT + 6.35, 2.9), (SAFE_LEFT, 4.5), (SAFE_LEFT + 6.35, 4.5)]
    colors = [pal["primary"], "0D9488", "D97706", "7C3AED"]
    for idx, (pos_x, pos_y) in enumerate(positions):
        if idx < len(opts):
            opt_text = opts[idx]
            add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, pos_x, pos_y, card_w, card_h, "FFFFFF", border_hex="CBD5E1")
            add_safe_shape(slide, MSO_SHAPE.OVAL, pos_x + 0.2, pos_y + 0.25, 0.9, 0.9, colors[idx])
            letter = opt_text[:2].strip()
            rest_text = opt_text[2:].strip()
            add_textbox(slide, pos_x + 0.2, pos_y + 0.45, 0.9, 0.5, letter, size_pt=18, bold=True, color_hex="FFFFFF", alignment=PP_ALIGN.CENTER)
            add_textbox(slide, pos_x + 1.25, pos_y + 0.25, card_w - 1.4, card_h - 0.4, rest_text, size_pt=15, bold=False, color_hex=pal["text_on_card"])

def build_thanks(prs, data, pal, layout):
    slide = prs.slides.add_slide(layout)
    add_slide_transition(slide, "cover")
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM - SAFE_TOP, pal["primary"], send_to_back=True)
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 2.0, 1.8, 9.33, 3.8, "FFFFFF", border_hex=pal["accent"])
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 2.0, 1.8, 9.33, 0.2, pal["accent"])
    
    add_textbox(slide, 2.2, 2.2, 8.93, 0.8, data.get("title", "CẢM ƠN CÁC EM ĐÃ CHÚ Ý LẮNG NGHE! 🌟"), size_pt=24, bold=True, color_hex=pal["primary"], alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 2.2, 3.1, 8.93, 1.2, data["content"], size_pt=17, bold=False, color_hex=pal["text_on_card"], alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 2.2, 4.5, 8.93, 0.6, "Hẹn gặp lại các em ở bài học tuần sau! 🚀", size_pt=16, bold=True, color_hex=pal["accent"], alignment=PP_ALIGN.CENTER)

# ─── MASTER DECK CONFIGURATIONS ───

DECKS_MASTER = [
    # ── 1. LỚP 5 - TIẾT 3 (BÀI 3) ──
    {
        "file_path": r"D:\UNIGO\KHBD_Tin_học\Lớp_5\Tuần_04\Slide_Tin_hoc_Lớp_5_Tiet03_Bai_3_Tim_kiem_thong_tin.pptx",
        "pal": PAL_L5_T3,
        "slides": [
            ("cover", {"title": "BÀI 3. TÌM KIẾM THÔNG TIN\nTRONG GIẢI QUYẾT VẤN ĐỀ", "subtitle": "Tin học • Lớp 5 • Tiết 3 (Chủ đề 3)"}),
            ("full_hero", {
                "title": "Tình huống: Chuẩn bị hành lý đi nghỉ hè",
                "img": os.path.join(IMG5_DIR, "ai_lop5_bai3_travel.jpg"),
                "question": "👉 An đi biển Nha Trang (nắng nóng), Khoa đi Đà Lạt (se lạnh), Minh đi Úc (mùa đông tuyết).\nVì sao mỗi bạn phải chuẩn bị hành lý hoàn toàn khác nhau?"
            }),
            ("three_cards", {
                "title": "Sự cần thiết của việc thu thập thông tin",
                "cards": [
                    {"header": "1. Hiểu rõ hoàn cảnh", "text": "• Nắm bắt thời tiết thực tế\n• Biết địa điểm tham quan\n• Chuẩn bị trang phục đúng"},
                    {"header": "2. Tránh rủi ro", "text": "• Không bị cảm lạnh bất ngờ\n• Tránh lãng phí hành lý\n• Luôn chủ động mọi tình huống"},
                    {"header": "3. Giải quyết hiệu quả", "text": "• Đưa ra quyết định sáng suốt\n• Chuyến đi an toàn, vui vẻ\n• Tiết kiệm thời gian & công sức"}
                ]
            }),
            ("hero_example", {
                "title": "Kỹ thuật sử dụng Từ khóa tìm kiếm thần kỳ",
                "text": "• Chọn từ khóa ngắn gọn, chính xác (Keyword)\n• Đặt từ khóa trong dấu ngoặc kép \" \"\n• Ví dụ: \"dự báo thời tiết Đà Lạt\"\n• Máy tìm kiếm sẽ lọc chính xác cụm từ em cần!",
                "img": os.path.join(IMG5_DIR, "ai_lop5_bai3_search.jpg")
            }),
            ("arrow_badges", {
                "title": "Quy trình tìm kiếm và kiểm chứng thông tin",
                "badges": [
                    "Bước 1: Xác định rõ thông tin cần tìm",
                    "Bước 2: Chọn từ khóa phù hợp đặt trong \" \"",
                    "Bước 3: Ưu tiên website uy tín (.gov.vn, .edu.vn)",
                    "Bước 4: Tổng hợp & áp dụng giải quyết vấn đề"
                ],
                "img": os.path.join(IMG5_DIR, "ai_lop5_bai3_search.jpg")
            }),
            ("conclusion_orange", {
                "question": "Thu thập và tìm kiếm thông tin giúp ích gì cho em?",
                "suggestions": [
                    "👉 Giúp nắm bắt đầy đủ thông tin thực tế trước khi hành động.",
                    "👉 Đưa ra quyết định thông minh, giải quyết vấn đề khoa học."
                ],
                "img": os.path.join(IMG5_DIR, "ai_lop5_bai3_travel.jpg"),
                "conclusion": "Thu thập và tìm kiếm thông tin là bước đầu tiên cực kỳ quan trọng giúp em giải quyết mọi vấn đề một cách tự tin và hiệu quả!"
            }),
            ("quiz_mascot", {
                "question": "Để tìm chính xác cụm từ 'dự báo thời tiết Đà Lạt', em nên gõ như thế nào vào Google?",
                "options": [
                    "A. \"dự báo thời tiết Đà Lạt\"",
                    "B. thời tiết và mọi thứ ở Đà Lạt",
                    "C. tìm cho tôi thời tiết",
                    "D. da lat hom nay the nao"
                ]
            }),
            ("thanks", {
                "title": "HOÀN THÀNH TIẾT 3! 🌟",
                "content": "BTVN: Cùng bố mẹ tìm hiểu thông tin thời tiết địa điểm du lịch cuối tuần này nhé!"
            })
        ]
    },

    # ── 2. LỚP 5 - TIẾT 4 (BÀI 4) ──
    {
        "file_path": r"D:\UNIGO\KHBD_Tin_học\Lớp_5\Tuần_04\Slide_Tin_hoc_Lớp_5_Tiet04_Bai_4_Cay_thu_muc.pptx",
        "pal": PAL_L5_T4,
        "slides": [
            ("cover", {"title": "BÀI 4. CÂY THƯ MỤC\n(FOLDER TREE)", "subtitle": "Tin học • Lớp 5 • Tiết 4 (Chủ đề 3)"}),
            ("full_hero", {
                "title": "Cây tri thức — Sơ đồ Cây thư mục kỳ diệu",
                "img": os.path.join(IMG5_DIR, "ai_lop5_bai4_tree.jpg"),
                "question": "👉 Gốc cây là ổ đĩa C:, D:. Cành lớn là Thư mục mẹ, cành nhỏ là Thư mục con, quả là các Tệp tin!\nEm thấy cây thư mục giúp máy tính gọn gàng như thế nào?"
            }),
            ("three_cards", {
                "title": "Các thành phần trong Cây thư mục",
                "cards": [
                    {"header": "1. Ổ đĩa gốc (Root)", "text": "• Ổ đĩa C:, D:, E:\n• Nơi chứa toàn bộ cây dữ liệu\n• Điểm xuất phát của mọi đường dẫn"},
                    {"header": "2. Thư mục (Folder)", "text": "• Thư mục mẹ chứa thư mục con\n• Dùng để phân loại theo chủ đề\n• Ví dụ: TOAN, VAN, TIN_HOC"},
                    {"header": "3. Tệp tin (File)", "text": "• Bài văn .docx, Trình chiếu .pptx\n• Tranh vẽ .png, Bài hát .mp3\n• Được lưu gọn trong thư mục"}
                ]
            }),
            ("arrow_badges", {
                "title": "4 Thao tác quản lý thư mục cơ bản",
                "badges": [
                    "1. Tạo thư mục mới (New Folder)",
                    "2. Đổi tên thư mục (Rename)",
                    "3. Sao chép / Di chuyển (Copy / Cut & Paste)",
                    "4. Xóa thư mục không dùng (Delete)"
                ],
                "img": os.path.join(IMG5_DIR, "ai_lop5_bai4_tree.jpg")
            }),
            ("hero_example", {
                "title": "Thực hành xây dựng Cây thư mục học tập",
                "text": "• Tạo thư mục mẹ HOC_TAP trên ổ đĩa D:\n• Tạo 3 thư mục con: TOAN, TIENG_VIET, TIN_HOC\n• Lưu bài tập Tin học vào đúng thư mục TIN_HOC\n• Giúp dữ liệu học tập luôn ngăn nắp, dễ tìm!",
                "img": os.path.join(IMG5_DIR, "ai_lop5_bai4_tree.jpg")
            }),
            ("conclusion_orange", {
                "question": "Tổ chức tệp tin theo Cây thư mục mang lại lợi ích gì?",
                "suggestions": [
                    "👉 Sắp xếp dữ liệu ngăn nắp, khoa học theo từng phân cấp.",
                    "👉 Tìm kiếm bài tập nhanh chóng, không bị thất lạc hay xóa nhầm."
                ],
                "img": os.path.join(IMG5_DIR, "ai_lop5_bai4_tree.jpg"),
                "conclusion": "Cây thư mục là phương pháp quản lý dữ liệu hình cây phân cấp chuẩn mực giúp máy tính của em luôn ngăn nắp và khoa học!"
            }),
            ("quiz_mascot", {
                "question": "Trong đường dẫn D:\\HOC_TAP\\TIN_HOC\\BaiTap.docx, đâu là thư mục con?",
                "options": [
                    "A. Ổ đĩa D:",
                    "B. Thư mục TIN_HOC",
                    "C. Tệp BaiTap.docx",
                    "D. Thư mục gốc"
                ]
            }),
            ("thanks", {
                "title": "XUẤT SẮC HOÀN THÀNH BÀI HỌC! 🚀",
                "content": "BTVN: Tạo cây thư mục học tập trên máy tính ở nhà của em nhé!"
            })
        ]
    },

    # ── 3. LỚP 5 - TỔNG HỢP TUẦN 04 (BÀI 3 & 4) ──
    {
        "file_path": r"D:\UNIGO\KHBD_Tin_học\Lớp_5\Tuần_04\Slide_Tin_hoc_Lop_5_Bai03_04.pptx",
        "pal": PAL_L5_T3,
        "slides": [
            ("cover", {"title": "TÌM KIẾM THÔNG TIN &\nCẤU TRÚC CÂY THƯ MỤC", "subtitle": "Tin học • Lớp 5 • Tuần 04 (Tiết 3 & Tiết 4)"}),
            ("section_banner", {"title": "TIẾT 3: TÌM KIẾM THÔNG TIN TRONG GIẢI QUYẾT VẤN ĐỀ", "desc": "Khám phá kỹ thuật tìm kiếm chính xác và chọn lọc thông tin đáng tin cậy"}),
            ("full_hero", {
                "title": "Tình huống: Chuẩn bị hành lý đi nghỉ hè",
                "img": os.path.join(IMG5_DIR, "ai_lop5_bai3_travel.jpg"),
                "question": "👉 An đi biển Nha Trang (nắng nóng), Khoa đi Đà Lạt (se lạnh), Minh đi Úc (mùa đông tuyết).\nVì sao mỗi bạn phải chuẩn bị hành lý hoàn toàn khác nhau?"
            }),
            ("hero_example", {
                "title": "Kỹ thuật sử dụng Từ khóa tìm kiếm thần kỳ",
                "text": "• Chọn từ khóa ngắn gọn, chính xác\n• Đặt từ khóa trong dấu ngoặc kép \" \"\n• Ưu tiên nguồn tin chính thống, đáng tin cậy\n• Giải quyết vấn đề nhanh chóng và hiệu quả",
                "img": os.path.join(IMG5_DIR, "ai_lop5_bai3_search.jpg")
            }),
            ("section_banner", {"title": "TIẾT 4: BÀI 4. CÂY THƯ MỤC TRONG MÁY TÍNH", "desc": "Tổ chức lưu trữ dữ liệu khoa học theo cấu trúc hình cây phân cấp"}),
            ("full_hero", {
                "title": "Cây tri thức — Sơ đồ Cây thư mục kỳ diệu",
                "img": os.path.join(IMG5_DIR, "ai_lop5_bai4_tree.jpg"),
                "question": "👉 Gốc là ổ đĩa (C:, D:), Cành lớn là Thư mục mẹ, Cành nhỏ là Thư mục con, Quả là các Tệp tin!"
            }),
            ("three_cards", {
                "title": "Các thành phần trong Cây thư mục",
                "cards": [
                    {"header": "1. Ổ đĩa gốc (Root)", "text": "• Ổ đĩa C:, D:, E:\n• Nơi chứa toàn bộ cây dữ liệu"},
                    {"header": "2. Thư mục (Folder)", "text": "• Thư mục mẹ chứa thư mục con\n• Ví dụ: TOAN, VAN, TIN_HOC"},
                    {"header": "3. Tệp tin (File)", "text": "• Bài văn .docx, Bài giảng .pptx\n• Tranh vẽ .png, Âm thanh .mp3"}
                ]
            }),
            ("conclusion_orange", {
                "question": "Tổng kết kiến thức trọng tâm Tuần 04 (Lớp 5)",
                "suggestions": [
                    "👉 Tiết 3: Thu thập và tìm kiếm thông tin chính xác bằng từ khóa \" \".",
                    "👉 Tiết 4: Sắp xếp dữ liệu ngăn nắp theo cấu trúc Cây thư mục."
                ],
                "img": os.path.join(IMG5_DIR, "ai_lop5_bai4_tree.jpg"),
                "conclusion": "Làm chủ kỹ năng tìm kiếm và tổ chức dữ liệu giúp em trở thành người sử dụng máy tính thông minh và hiệu quả!"
            }),
            ("quiz_mascot", {
                "question": "Trong đường dẫn D:\\HOC_TAP\\TIN_HOC\\BaiTap.docx, đâu là thư mục con?",
                "options": [
                    "A. Ổ đĩa D:",
                    "B. Thư mục TIN_HOC",
                    "C. Tệp BaiTap.docx",
                    "D. Thư mục gốc"
                ]
            }),
            ("thanks", {
                "title": "BÀI HỌC KẾT THÚC! 🌟",
                "content": "BTVN: Hoàn thành bài tập thực hành trên máy tính và chuẩn bị cho bài tuần tới nhé!"
            })
        ]
    },

    # ── 4. LỚP 7 - TIẾT 3 (BÀI 3) ──
    {
        "file_path": r"D:\UNIGO\KHBD_Tin_học\Lớp_7\Tuần_04\Slide_Tin_hoc_Lớp_7_Tiet03_Bai_3_Quan_ly_du_lieu_trong_may_tinh.pptx",
        "pal": PAL_L7_T3,
        "slides": [
            ("cover", {"title": "BÀI 3. QUẢN LÝ DỮ LIỆU\nTRONG MÁY TÍNH", "subtitle": "Tin học • Lớp 7 • Tiết 3 (Chủ đề 2)"}),
            ("full_hero", {
                "title": "Không gian Dữ liệu số & Bảo mật công nghệ",
                "img": os.path.join(IMG7_DIR, "ai_lop7_bai3_data_hub.jpg"),
                "question": "👉 Thế giới số chứa hàng triệu tệp tin đa dạng (.docx, .xlsx, .pptx, .jpg, .mp4).\nLàm sao để vừa sắp xếp khoa học vừa bảo vệ dữ liệu an toàn trước virus?"
            }),
            ("three_cards", {
                "title": "Phân loại tệp theo Phần mở rộng (File Extension)",
                "cards": [
                    {"header": "1. Tài liệu văn phòng", "text": "• .docx: Văn bản Word\n• .xlsx: Bảng tính Excel\n• .pptx: Trình chiếu Slide\n• .pdf: Tài liệu chuẩn cố định"},
                    {"header": "2. Đa phương tiện", "text": "• .jpg / .png: Hình ảnh đồ họa\n• .mp3: Âm thanh, bài hát\n• .mp4: Video, phim bài giảng\n• .gif: Ảnh động sinh động"},
                    {"header": "3. Tệp thực thi & Cảnh báo", "text": "• .exe / .bat: Chương trình cài đặt\n• Cảnh báo: Không bấm vào tệp .exe lạ\n• Nguy cơ chứa virus, mã độc\n• Luôn quét virus trước khi mở"}
                ]
            }),
            ("arrow_badges", {
                "title": "Chiến lược Sao lưu dữ liệu & Phòng chống Virus",
                "badges": [
                    "1. Sao lưu định kỳ (Backup) lên Cloud / Ổ cứng ngoài",
                    "2. Đặt mật khẩu mạnh và không chia sẻ cho người khác",
                    "3. Bật phần mềm diệt virus (Windows Defender)",
                    "4. Tuyệt đối không mở liên kết lạ, tệp đính kèm đáng ngờ"
                ],
                "img": os.path.join(IMG7_DIR, "ai_lop7_bai3_data_hub.jpg")
            }),
            ("hero_example", {
                "title": "Bộ phím tắt thần tốc quản lý tệp tin",
                "text": "• F2: Đổi tên tệp / thư mục (Rename)\n• Ctrl + Shift + N: Tạo nhanh thư mục mới\n• Ctrl + C / Ctrl + X: Sao chép / Cắt tệp\n• Ctrl + V: Dán tệp vào thư mục đích\n• Shift + Delete: Xóa vĩnh viễn tệp rác",
                "img": os.path.join(IMG7_DIR, "ai_lop7_bai3_data_hub.jpg")
            }),
            ("conclusion_orange", {
                "question": "Biện pháp an toàn nhất để không bị mất dữ liệu quan trọng là gì?",
                "suggestions": [
                    "👉 Sao lưu dữ liệu dự phòng (Backup) thường xuyên lên đám mây hoặc thiết bị ngoài.",
                    "👉 Đặt mật khẩu bảo vệ tài khoản và quét virus định kỳ."
                ],
                "img": os.path.join(IMG7_DIR, "ai_lop7_bai3_data_hub.jpg"),
                "conclusion": "Quản lý dữ liệu khoa học kết hợp với sao lưu dự phòng thường xuyên là kỹ năng sống còn của mọi công dân số!"
            }),
            ("quiz_mascot", {
                "question": "Đâu là phần mềm diệt virus và bảo vệ máy tính được tích hợp sẵn trong Windows?",
                "options": [
                    "A. Microsoft Word",
                    "B. Windows Defender",
                    "C. File Explorer",
                    "D. Google Chrome"
                ]
            }),
            ("thanks", {
                "title": "HOÀN THÀNH TIẾT 3! 🌟",
                "content": "BTVN: Kiểm tra và sao lưu các tài liệu học tập của em lên Google Drive hoặc USB!"
            })
        ]
    },

    # ── 5. LỚP 7 - TIẾT 4 (BÀI 4) ──
    {
        "file_path": r"D:\UNIGO\KHBD_Tin_học\Lớp_7\Tuần_04\Slide_Tin_hoc_Lớp_7_Tiet04_Bai_4_Mang_xa_hoi.pptx",
        "pal": PAL_L7_T4,
        "slides": [
            ("cover", {"title": "BÀI 4. MẠNG XÃ HỘI &\nTRAO ĐỔI THÔNG TIN INTERNET", "subtitle": "Tin học • Lớp 7 • Tiết 4 (Chủ đề 2)"}),
            ("full_hero", {
                "title": "Thế giới Mạng xã hội & Tấm khiên bảo vệ số",
                "img": os.path.join(IMG7_DIR, "ai_lop7_bai4_social_networks.jpg"),
                "question": "👉 Mạng xã hội kết nối hàng triệu bạn bè khắp năm châu.\nLàm thế nào để tận dụng tiện ích giao lưu học tập mà vẫn an toàn trước cạm bẫy số?"
            }),
            ("three_cards", {
                "title": "Các kênh trao đổi thông tin phổ biến",
                "cards": [
                    {"header": "1. Thư điện tử (E-mail)", "text": "• Trao đổi thư từ học tập, công việc\n• Đính kèm tệp văn bản, hình ảnh\n• Trang trọng, lưu vết lâu dài"},
                    {"header": "2. Diễn đàn (Forum)", "text": "• Thảo luận chuyên sâu theo chủ đề\n• Hỏi đáp kiến thức học thuật\n• Chia sẻ kinh nghiệm cộng đồng"},
                    {"header": "3. Mạng xã hội (Social Media)", "text": "• Kết nối bạn bè, chia sẻ khoảnh khắc\n• Nhắn tin tức thời, gọi video\n• Lan tỏa thông điệp tích cực"}
                ]
            }),
            ("three_cards", {
                "title": "Nhận diện 2 mặt của Mạng xã hội",
                "cards": [
                    {"header": "🌟 Mặt tích cực", "text": "• Giao lưu bạn bè bốn phương\n• Học tập trực tuyến bổ ích\n• Cập nhật tin tức đời sống nhanh"},
                    {"header": "⚠️ Cạm bẫy & Rủi ro", "text": "• Nguy cơ lộ thông tin cá nhân\n• Tiếp xúc tin giả (Fake News)\n• Nghiện mạng, giảm thị lực"},
                    {"header": "🛡️ Phòng ngừa", "text": "• Không tin thông tin giật gân\n• Không chia sẻ mật khẩu cá nhân\n• Báo xấu nội dung độc hại"}
                ]
            }),
            ("arrow_badges", {
                "title": "Nguyên tắc 5K trở thành Công dân số thông thái",
                "badges": [
                    "1. Không chia sẻ mật khẩu & thông tin nhạy cảm",
                    "2. Không đăng ảnh người khác khi chưa được phép",
                    "3. Kiểm chứng độ tin cậy trước khi chia sẻ tin tức",
                    "4. Giữ thái độ văn minh, tôn trọng, không bạo lực mạng",
                    "5. Kịp thời báo cho bố mẹ, thầy cô khi gặp nguy hiểm"
                ],
                "img": os.path.join(IMG7_DIR, "ai_lop7_bai4_social_networks.jpg")
            }),
            ("conclusion_orange", {
                "question": "Em cần làm gì để sử dụng Mạng xã hội an toàn và văn minh?",
                "suggestions": [
                    "👉 Bảo vệ danh tính số: Đặt mật khẩu mạnh, bảo mật 2 lớp.",
                    "👉 Sử dụng điều độ, ứng xử văn minh và luôn tôn trọng pháp luật."
                ],
                "img": os.path.join(IMG7_DIR, "ai_lop7_bai4_social_networks.jpg"),
                "conclusion": "Mạng xã hội là công cụ kết nối tuyệt vời. Hãy là người dùng thông thái, văn minh và luôn biết bảo vệ bản thân trên không gian mạng!"
            }),
            ("quiz_mascot", {
                "question": "Hành vi nào sau đây là KHÔNG NÊN LÀM trên mạng xã hội?",
                "options": [
                    "A. Chia sẻ bài viết học tập hay cho bạn bè",
                    "B. Đăng số CCCD và mật khẩu lên trang cá nhân",
                    "C. Báo cáo bài viết lừa đảo cho quản trị viên",
                    "D. Chúc mừng sinh nhật bạn bè lịch sự"
                ]
            }),
            ("thanks", {
                "title": "HOÀN THÀNH BÀI HỌC! 🚀",
                "content": "BTVN: Kiểm tra và thiết lập chế độ bảo mật riêng tư cho tài khoản mạng xã hội của em!"
            })
        ]
    },

    # ── 6. LỚP 7 - TỔNG HỢP TUẦN 04 (BÀI 3 & 4) ──
    {
        "file_path": r"D:\UNIGO\KHBD_Tin_học\Lớp_7\Tuần_04\Slide_Tin_hoc_Lop_7_Bai03_04.pptx",
        "pal": PAL_L7_T3,
        "slides": [
            ("cover", {"title": "QUẢN LÝ DỮ LIỆU &\nMẠNG XÃ HỘI INTERNET", "subtitle": "Tin học • Lớp 7 • Tuần 04 (Tiết 3 & Tiết 4)"}),
            ("section_banner", {"title": "TIẾT 3: BÀI 3. QUẢN LÝ DỮ LIỆU TRONG MÁY TÍNH", "desc": "Nắm vững phần mở rộng tệp, thao tác File Explorer và kỹ thuật sao lưu dự phòng"}),
            ("full_hero", {
                "title": "Không gian Dữ liệu số & Bảo mật công nghệ",
                "img": os.path.join(IMG7_DIR, "ai_lop7_bai3_data_hub.jpg"),
                "question": "👉 Nhận diện phần mở rộng (.docx, .xlsx, .pptx, .jpg, .mp4) & Sao lưu dự phòng an toàn!"
            }),
            ("arrow_badges", {
                "title": "Chiến lược Sao lưu dữ liệu & Phòng chống Virus",
                "badges": [
                    "1. Sao lưu định kỳ (Backup) lên Cloud / Ổ cứng ngoài",
                    "2. Đặt mật khẩu mạnh và không chia sẻ cho người khác",
                    "3. Bật phần mềm diệt virus (Windows Defender)",
                    "4. Tuyệt đối không mở liên kết lạ, tệp đính kèm đáng ngờ"
                ],
                "img": os.path.join(IMG7_DIR, "ai_lop7_bai3_data_hub.jpg")
            }),
            ("section_banner", {"title": "TIẾT 4: BÀI 4. MẠNG XÃ HỘI & KÊNH TRAO ĐỔI THÔNG TIN", "desc": "Hiểu rõ 2 mặt của Mạng xã hội, kỹ năng giao tiếp văn minh và bảo mật thông tin"}),
            ("full_hero", {
                "title": "Thế giới Mạng xã hội & Tấm khiên bảo vệ số",
                "img": os.path.join(IMG7_DIR, "ai_lop7_bai4_social_networks.jpg"),
                "question": "👉 Kết nối bạn bè năm châu & Luôn giữ vững tấm khiên bảo vệ an toàn thông tin cá nhân!"
            }),
            ("three_cards", {
                "title": "Nhận diện 2 mặt của Mạng xã hội",
                "cards": [
                    {"header": "🌟 Mặt tích cực", "text": "• Giao lưu bạn bè bốn phương\n• Học tập trực tuyến bổ ích\n• Cập nhật tri thức nhanh chóng"},
                    {"header": "⚠️ Cạm bẫy & Rủi ro", "text": "• Nguy cơ lộ thông tin cá nhân\n• Tiếp xúc tin giả (Fake News)\n• Bắt nạt qua mạng (Cyberbullying)"},
                    {"header": "🛡️ Phòng ngừa văn minh", "text": "• Kiểm chứng thông tin trước khi chia sẻ\n• Tôn trọng, không công kích thô bạo\n• Kịp thời báo cho người lớn khi gặp nguy hiểm"}
                ]
            }),
            ("conclusion_orange", {
                "question": "Tổng kết kiến thức trọng tâm Tuần 04 (Lớp 7)",
                "suggestions": [
                    "👉 Tiết 3: Tổ chức tệp ngăn nắp, sao lưu dữ liệu và diệt virus định kỳ.",
                    "👉 Tiết 4: Khai thác tiện ích mạng xã hội, ứng xử văn minh và bảo vệ danh tính số."
                ],
                "img": os.path.join(IMG7_DIR, "ai_lop7_bai4_social_networks.jpg"),
                "conclusion": "Quản lý dữ liệu máy tính khoa học kết hợp với văn hóa ứng xử số văn minh là nền tảng của công dân số thời đại 4.0!"
            }),
            ("quiz_mascot", {
                "question": "Hành vi nào sau đây là KHÔNG NÊN LÀM trên mạng xã hội?",
                "options": [
                    "A. Chia sẻ bài học hay cho bạn bè",
                    "B. Đăng số CCCD và mật khẩu lên trang cá nhân",
                    "C. Báo cáo bài viết lừa đảo cho quản trị viên",
                    "D. Chúc mừng sinh nhật bạn bè lịch sự"
                ]
            }),
            ("thanks", {
                "title": "BÀI HỌC KẾT THÚC! 🚀",
                "content": "BTVN: Hoàn thành phiếu học tập và thực hiện bảo mật 2 lớp cho các tài khoản trực tuyến!"
            })
        ]
    }
]

BUILDERS_MAP = {
    "cover": build_cover,
    "section_banner": build_section_banner,
    "hero_example": build_hero_example,
    "full_hero": build_full_hero_slide,
    "arrow_badges": build_arrow_badges_diagram,
    "three_cards": build_three_cards_diagram,
    "conclusion_orange": build_conclusion_orange,
    "quiz_mascot": build_quiz_mascot,
    "thanks": build_thanks
}

def generate_master_decks():
    print("=== BẮT ĐẦU TẠO BỘ SLIDE CAO CẤP LỚP 5 & LỚP 7 (TUẦN 04) ===")
    for config in DECKS_MASTER:
        prs = Presentation(TEMPLATE)
        # Clean default template slides
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
            
        blank_layout = prs.slide_layouts[6]
        pal = config["pal"]
        file_path = config["file_path"]
        
        print(f"\n--> Đang tạo: {os.path.basename(file_path)}")
        for b_name, b_data in config["slides"]:
            if b_name in BUILDERS_MAP:
                BUILDERS_MAP[b_name](prs, b_data, pal, blank_layout)
            else:
                print(f"Unknown builder: {b_name}")
                
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        prs.save(file_path)
        print(f"  [HOÀN TẤT] {file_path} ({len(prs.slides)} slides)")

if __name__ == "__main__":
    generate_master_decks()
