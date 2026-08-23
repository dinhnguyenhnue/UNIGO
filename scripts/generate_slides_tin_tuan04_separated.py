# -*- coding: utf-8 -*-
"""
Hệ thống Tạo Slide Bài Giảng Chuẩn UNIGO Tuần 04 (Môn Tin Học)
Quy chuẩn: TÁCH BIỆT TỪNG TIẾT / TỪNG KHBD = 1 FILE SLIDE (.PPTX) ĐỘC LẬP
Áp dụng 20 Dạng Bố Cục Chuẩn (Visual-First & Hero Image)
"""
import sys, io, os, glob
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
from PIL import Image, ImageDraw, ImageFont

TEMPLATE  = r'D:\UNIGO\Hệ thống mẫu văn bản\Mẫu slide UNIGO.pptx'
KHBD_BASE = r'D:\UNIGO\KHBD_Tin_học'

SAFE_TOP    = 1.15
SAFE_BOTTOM = 6.35
SAFE_LEFT   = 0.3
SAFE_RIGHT  = 13.0
SLIDE_W     = 13.33
SLIDE_H     = 7.50

COLOR_PALETTES = [
    # 0: Blue (Tiền TH)
    {"primary": "1B4F9B", "accent": "2D7DD2", "bg": "EBF3FE", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "1A2744", "text_on_card": "0F172A"},
    # 1: Purple (Lớp 1)
    {"primary": "5B21B6", "accent": "7C3AED", "bg": "F3EEFF", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "2E1065", "text_on_card": "0F172A"},
    # 2: Teal (Lớp 2)
    {"primary": "0F766E", "accent": "14B8A6", "bg": "ECFDF5", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "134E4A", "text_on_card": "0F172A"},
    # 3: Orange (Lớp 3)
    {"primary": "C2410C", "accent": "EA580C", "bg": "FFF7ED", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "7C2D12", "text_on_card": "0F172A"},
    # 4: Indigo (Lớp 4)
    {"primary": "3730A3", "accent": "4F46E5", "bg": "EEF2FF", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "1E1B4B", "text_on_card": "0F172A"},
    # 5: Rose (Lớp 5 - Tiết 3)
    {"primary": "BE185D", "accent": "EC4899", "bg": "FDF2F8", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "831843", "text_on_card": "0F172A"},
    # 6: Ruby (Lớp 5 - Tiết 4)
    {"primary": "9F1239", "accent": "F43F5E", "bg": "FFF1F2", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "881337", "text_on_card": "0F172A"},
    # 7: Emerald (Lớp 6 - Tiết 3)
    {"primary": "047857", "accent": "10B981", "bg": "ECFDF5", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "064E3B", "text_on_card": "0F172A"},
    # 8: Forest (Lớp 6 - Tiết 4)
    {"primary": "065F46", "accent": "059669", "bg": "F0FDF4", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "064E3B", "text_on_card": "0F172A"},
    # 9: Sky (Lớp 7 - Tiết 3)
    {"primary": "0369A1", "accent": "0EA5E9", "bg": "F0F9FF", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "0C4A6E", "text_on_card": "0F172A"},
    # 10: Cyan (Lớp 7 - Tiết 4)
    {"primary": "0E7490", "accent": "06B6D4", "bg": "ECFEFF", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "164E63", "text_on_card": "0F172A"},
    # 11: Amber (Lớp 8 - Tiết 3)
    {"primary": "B45309", "accent": "F59E0B", "bg": "FFFBEB", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "78350F", "text_on_card": "0F172A"},
    # 12: Violet (Lớp 8 - Tiết 4)
    {"primary": "6D28D9", "accent": "8B5CF6", "bg": "F5F3FF", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "4C1D95", "text_on_card": "0F172A"},
]

def hex_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_safe_shape(slide, shape_type, left, top, width, height, fill_hex, border_hex=None, send_to_back=False):
    actual_top = max(top, SAFE_TOP)
    actual_bottom = min(top + height, SAFE_BOTTOM)
    actual_height = max(actual_bottom - actual_top, 0.1)

    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(actual_top), Inches(width), Inches(actual_height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(fill_hex)
    if border_hex:
        shape.line.color.rgb = hex_rgb(border_hex)
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()

    if send_to_back:
        sp = shape._element
        spTree = sp.getparent()
        spTree.remove(sp)
        spTree.insert(2, sp)
    return shape

def add_textbox(slide, left, top, width, height, text, size_pt=18, bold=False, color_hex="0F172A", alignment=PP_ALIGN.LEFT, font_name="Arial"):
    actual_top = max(top, SAFE_TOP)
    actual_bottom = min(top + height, SAFE_BOTTOM)
    actual_height = max(actual_bottom - actual_top, 0.2)

    txbox = slide.shapes.add_textbox(Inches(left), Inches(actual_top), Inches(width), Inches(actual_height))
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    p.font.size = Pt(size_pt)
    p.font.bold = bold
    p.font.color.rgb = hex_rgb(color_hex)
    p.font.name = font_name
    return txbox

def add_multiline_textbox(slide, left, top, width, height, lines, size_pt=18, bold=False, color_hex="0F172A", alignment=PP_ALIGN.LEFT):
    actual_top = max(top, SAFE_TOP)
    actual_bottom = min(top + height, SAFE_BOTTOM)
    actual_height = max(actual_bottom - actual_top, 0.2)

    txbox = slide.shapes.add_textbox(Inches(left), Inches(actual_top), Inches(width), Inches(actual_height))
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = alignment
        p.font.size = Pt(size_pt)
        p.font.bold = bold
        p.font.color.rgb = hex_rgb(color_hex)
        p.font.name = "Arial"
        p.space_after = Pt(6)
    return txbox

def add_picture_safe(slide, img_path, left, top, width, height):
    actual_top = max(top, SAFE_TOP)
    actual_bottom = min(top + height, SAFE_BOTTOM)
    actual_height = max(actual_bottom - actual_top, 0.5)
    try:
        return slide.shapes.add_picture(img_path, Inches(left), Inches(actual_top), Inches(width), Inches(actual_height))
    except Exception as e:
        return None

def add_slide_transition(slide, transition_type="fade"):
    transitions = {
        'fade': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>',
        'push': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:push dir="l"/></p:transition>',
        'wipe': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:wipe/></p:transition>',
        'cover': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:cover/></p:transition>',
    }
    xml_str = transitions.get(transition_type, transitions['fade'])
    slide._element.append(etree.fromstring(xml_str))

# --- SLIDE BUILDERS (Applying the 20 visual layouts) ---

def build_cover(prs, data, pal, lesson, layout):
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["primary"], send_to_back=True)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, 1.8, pal["accent"], send_to_back=True)
    add_textbox(slide, 0.8, 2.0, 7.5, 1.5, data["title"], size_pt=34, bold=True, color_hex="FFFFFF")
    add_textbox(slide, 0.8, 3.8, 7.5, 0.8, data.get("subtitle", ""), size_pt=20, color_hex="FFFFFF")
    add_slide_transition(slide, "fade")
    return slide

def build_section_banner(prs, data, pal, lesson, layout):
    # Layout 2: Section Banner
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)
    
    # Search header
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 2.5, SAFE_TOP + 0.3, SLIDE_W - 5.0, 0.7, "#FFFFFF", border_hex="#CBD5E1")
    add_textbox(slide, 2.8, SAFE_TOP + 0.38, SLIDE_W - 5.6, 0.5, f"🔍  TIẾT {data.get('tiet', '1')}. {data['title'].upper()}", size_pt=20, bold=True, color_hex=pal["primary"], alignment=PP_ALIGN.CENTER)
    
    # Large Banner
    banner_y = SAFE_TOP + 1.3
    banner_h = SAFE_BOTTOM - banner_y - 0.4
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.2, banner_y, SLIDE_W - 2.4, banner_h, pal["primary"])
    add_textbox(slide, 1.5, banner_y + 0.8, SLIDE_W - 3.0, 1.8, data["title"], size_pt=36, bold=True, color_hex="FFFFFF", alignment=PP_ALIGN.CENTER)
    if data.get("subtitle"):
        add_textbox(slide, 1.5, banner_y + 2.2, SLIDE_W - 3.0, 0.8, data["subtitle"], size_pt=22, bold=False, color_hex="FFFFFF", alignment=PP_ALIGN.CENTER)
    
    add_slide_transition(slide, "push")
    return slide

def build_hero_example(prs, data, pal, lesson, layout):
    # Layout 5: Hero Image (Left 35% Cloud card + Right 65% Big illustration)
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, pal["primary"])
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, f"  {lesson['lop_label'].upper()} • VÍ DỤ MINH HỌA TRỰC QUAN", size_pt=13, bold=True, color_hex="FFFFFF")

    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, SLIDE_W-1, 0.65, data["title"], size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    content_y = title_y + 0.75
    content_h = SAFE_BOTTOM - content_y - 0.15

    # Left Cloud Card (35%)
    card_w = 4.6
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.5, content_y, card_w, content_h, "#FFFFFF", border_hex=pal["accent"])
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0.5, content_y, 0.12, content_h, pal["accent"])
    lines = data.get("text", "").split("\n")
    add_multiline_textbox(slide, 0.8, content_y + 0.3, card_w - 0.5, content_h - 0.6, lines, size_pt=20, bold=True, color_hex="#0F172A")

    # Right Hero Image (65%)
    img_x = 5.4
    img_w = SLIDE_W - img_x - 0.5
    img_path = data.get("img")
    if img_path and os.path.isfile(img_path):
        add_picture_safe(slide, img_path, img_x, content_y, img_w, content_h)

    add_slide_transition(slide, "wipe")
    return slide

def build_posture_compare(prs, data, pal, lesson, layout):
    # Layout 16 & 17: Posture 3 images a, b, c with ❌ ❌ ✅
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, pal["primary"])
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, f"  {lesson['lop_label'].upper()} • KHÁM PHÁ BÀI HỌC", size_pt=13, bold=True, color_hex="FFFFFF")

    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, SLIDE_W-1, 0.65, data["title"], size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    # 3 Cards for a, b, c
    items = [
        ("Hình a: Ngồi sai (Gù lưng)", "❌ SAI", "#DC2626", data.get("img_a")),
        ("Hình b: Ngồi sai (Cúi sát mắt)", "❌ SAI", "#DC2626", data.get("img_b")),
        ("Hình c: Ngồi đúng (Thẳng lưng)", "✅ ĐÚNG", "#059669", data.get("img_c"))
    ]
    card_w = 3.8
    card_h = 3.6
    for i, (label, status, scol, ipath) in enumerate(items):
        cx = 0.5 + i * (card_w + 0.3)
        cy = title_y + 0.8
        add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, card_w, card_h, "#FFFFFF", border_hex=scol)
        add_safe_shape(slide, MSO_SHAPE.RECTANGLE, cx, cy, card_w, 0.45, scol)
        add_textbox(slide, cx, cy + 0.05, card_w, 0.4, status, size_pt=18, bold=True, color_hex="FFFFFF", alignment=PP_ALIGN.CENTER)
        
        if ipath and os.path.isfile(ipath):
            add_picture_safe(slide, ipath, cx + 0.2, cy + 0.6, card_w - 0.4, 2.0)
        
        add_textbox(slide, cx + 0.1, cy + 2.7, card_w - 0.2, 0.8, label, size_pt=17, bold=True, color_hex="#0F172A", alignment=PP_ALIGN.CENTER)

    add_slide_transition(slide, "wipe")
    return slide

def build_arrow_badges_diagram(prs, data, pal, lesson, layout):
    # Layout 19: Left 5 Arrow Badges + Right Big diagram
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, pal["primary"])
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, f"  {lesson['lop_label'].upper()} • QUY CHUẨN TRỰC QUAN", size_pt=13, bold=True, color_hex="FFFFFF")

    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, SLIDE_W-1, 0.65, data["title"], size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    content_y = title_y + 0.75
    badges = data.get("badges", [])
    b_colors = ["#EA580C", "#2563EB", "#059669", "#D97706", "#DB2777"]
    
    # Left Arrow Badges
    badge_w = 6.2
    badge_h = 0.70
    gap = 0.14
    for i, btext in enumerate(badges[:5]):
        by = content_y + i * (badge_h + gap)
        bcol = b_colors[i % len(b_colors)]
        add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.5, by, badge_w, badge_h, "#FFFFFF", border_hex=bcol)
        add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0.5, by, 0.20, badge_h, bcol)
        add_textbox(slide, 0.9, by + 0.12, badge_w - 0.6, badge_h - 0.24, btext, size_pt=17, bold=True, color_hex="#0F172A")

    # Right Diagram
    img_x = 7.0
    img_w = SLIDE_W - img_x - 0.5
    img_h = SAFE_BOTTOM - content_y - 0.15
    img_path = data.get("img")
    if img_path and os.path.isfile(img_path):
        add_picture_safe(slide, img_path, img_x, content_y, img_w, img_h)

    add_slide_transition(slide, "wipe")
    return slide

def build_conclusion_orange(prs, data, pal, lesson, layout):
    # Layout 14: Speech bubble + Suggestions + Orange Conclusion Box
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, pal["primary"])
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, f"  {lesson['lop_label'].upper()} • BÀI HỌC VÀ KẾT LUẬN", size_pt=13, bold=True, color_hex="FFFFFF")

    # Top Speech Bubble Question
    q_y = SAFE_TOP + 0.65
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, q_y, 8.5, 0.9, "#FFFFFF", border_hex=pal["accent"])
    add_textbox(slide, 0.9, q_y + 0.15, 8.0, 0.6, f"💬 {data.get('question', '')}", size_pt=20, bold=True, color_hex=pal["primary"])
    
    # Suggestion bullets
    s_y = q_y + 1.1
    s_lines = data.get("suggestions", [])
    add_multiline_textbox(slide, 0.8, s_y, 8.0, 1.3, s_lines, size_pt=18, bold=False, color_hex="#0F172A")

    # Right Picture
    img_path = data.get("img")
    if img_path and os.path.isfile(img_path):
        add_picture_safe(slide, img_path, 9.5, q_y, 3.2, 2.3)

    # Orange Conclusion Box at Bottom
    c_y = SAFE_BOTTOM - 1.4
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, c_y, SLIDE_W - 1.2, 1.25, "#FFF7ED", border_hex="#EA580C")
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0.6, c_y, 0.20, 1.25, "#EA580C")
    c_text = f"💡 KẾT LUẬN: {data.get('conclusion', '')}"
    add_textbox(slide, 1.0, c_y + 0.2, SLIDE_W - 1.8, 0.85, c_text, size_pt=20, bold=True, color_hex="#9A3412")

    add_slide_transition(slide, "cover")
    return slide

def build_quiz_mascot(prs, data, pal, lesson, layout):
    # Layout 10: Quiz A B C D grid with Mascot
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, "#D97706")
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, f"  {lesson['lop_label'].upper()} • THỬ THÁCH TRẮC NGHIỆM VUI", size_pt=13, bold=True, color_hex="FFFFFF")

    # Question
    q_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, q_y, SLIDE_W - 1.0, 0.7, f"❓ {data['question']}", size_pt=22, bold=True, color_hex="#1E293B")

    # 4 Options 2x2
    options = data.get("options", [])
    card_w = (SLIDE_W - 1.0 - 0.35) / 2
    card_h = 1.6
    for i, opt in enumerate(options[:4]):
        col, row = i % 2, i // 2
        cx = 0.5 + col * (card_w + 0.35)
        cy = q_y + 0.85 + row * (card_h + 0.25)
        add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, card_w, card_h, "#FFFBEB", border_hex="#F59E0B")
        add_safe_shape(slide, MSO_SHAPE.RECTANGLE, cx, cy, 0.15, card_h, "#F59E0B")
        add_textbox(slide, cx + 0.3, cy + 0.25, card_w - 0.5, card_h - 0.5, opt, size_pt=19, bold=True, color_hex="#78350F")

    add_slide_transition(slide, "push")
    return slide

def build_thanks(prs, data, pal, lesson, layout):
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["primary"], send_to_back=True)
    add_textbox(slide, 1.0, 2.0, SLIDE_W - 2.0, 1.2, data["title"], size_pt=36, bold=True, color_hex="FFFFFF", alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1.0, 3.5, SLIDE_W - 2.0, 1.0, data["content"], size_pt=22, bold=False, color_hex="FFFFFF", alignment=PP_ALIGN.CENTER)
    add_slide_transition(slide, "fade")
    return slide

# --- 12 SEPARATE LESSONS FOR TUẦN 04 ---

SEPARATED_LESSONS_T4 = [
    # 1. Tiền Tiểu học (Tiết 3)
    {
        "folder": "Tiền_tiểu_học", "lop_label": "Tiền Tiểu học", "palette_idx": 0,
        "file_name": "Slide_Tin_hoc_Tiền_tiểu_học_Tiet03_Em_ngoi_may_tinh_an_toan.pptx",
        "builders": [
            ("cover", {"title": "Em ngồi máy tính an toàn 🪑", "subtitle": "Tin học • Tiền Tiểu học • Tiết 3"}),
            ("posture_compare", {"title": "Bạn nào ngồi học đúng tư thế?",
                                 "img_a": r"D:\UNIGO\KHBD_Tin_học\Tiền_tiểu_học\Tuần_04\images\ai_bad_posture.jpg",
                                 "img_b": r"D:\UNIGO\KHBD_Tin_học\Tiền_tiểu_học\Tuần_04\images\learn1_card2.png",
                                 "img_c": r"D:\UNIGO\KHBD_Tin_học\Tiền_tiểu_học\Tuần_04\images\ai_good_posture.jpg"}),
            ("arrow_badges", {"title": "4 Quy tắc ngồi máy tính an toàn",
                              "badges": ["1. Ngồi thẳng lưng, tựa nhẹ vào ghế", "2. Hai bàn chân đặt bằng phẳng trên sàn", "3. Mắt cách màn hình bằng 1 cánh tay", "4. Tuyệt đối không chạm vào ổ cắm điện"],
                              "img": r"D:\UNIGO\KHBD_Tin_học\Tiền_tiểu_học\Tuần_04\images\ai_good_posture.jpg"}),
            ("conclusion_orange", {"question": "Vì sao con phải ngồi thẳng lưng khi học máy tính?",
                                   "suggestions": ["👉 Giúp con không bị gù lưng và vẹo cột sống.", "👉 Giúp đôi mắt của con luôn sáng và khỏe mạnh."],
                                   "img": r"D:\UNIGO\KHBD_Tin_học\Tiền_tiểu_học\Tuần_04\images\ai_good_posture.jpg",
                                   "conclusion": "Ngồi thẳng lưng, mắt cách màn hình 50-70cm và luôn giữ an toàn điện!"}),
            ("quiz_mascot", {"question": "Khoảng cách an toàn từ mắt đến màn hình là bao nhiêu?",
                             "options": ["A. Rất gần (cúi sát mắt)", "B. Bằng một cánh tay (50-70cm)", "C. Nằm ra bàn để xem", "D. Nhắm mắt lại khi dùng máy"]}),
            ("thanks", {"title": "Các con giỏi lắm! 🌟", "content": "BTVN: Hãy nhắc nhở bố mẹ và anh chị ngồi đúng tư thế nhé!"})
        ]
    },
    # 2. Lớp 1 (Tiết 3)
    {
        "folder": "Lớp_1", "lop_label": "Lớp 1", "palette_idx": 1,
        "file_name": "Slide_Tin_hoc_Lớp_1_Tiet03_Tu_the_va_an_toan_khi_dung_may.pptx",
        "builders": [
            ("cover", {"title": "Tư thế & An toàn khi dùng máy 🪑", "subtitle": "Tin học • Lớp 1 • Tiết 3"}),
            ("posture_compare", {"title": "Quan sát và nhận biết tư thế chuẩn",
                                 "img_a": r"D:\UNIGO\KHBD_Tin_học\Lớp_1\Tuần_04\images\ai_bad_posture.jpg",
                                 "img_b": r"D:\UNIGO\KHBD_Tin_học\Lớp_1\Tuần_04\images\learn1_card2.png",
                                 "img_c": r"D:\UNIGO\KHBD_Tin_học\Lớp_1\Tuần_04\images\ai_good_posture.jpg"}),
            ("arrow_badges", {"title": "Quy tắc an toàn điện trong phòng máy",
                              "badges": ["1. Không dùng que nhọn chọc vào ổ điện", "2. Không mang đồ ăn, nước uống gần máy", "3. Không tự ý cắm hoặc rút dây nguồn", "4. Báo ngay thầy cô khi có mùi khét lạ"],
                              "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_1\Tuần_04\images\learn1_card4.png"}),
            ("conclusion_orange", {"question": "Em cần làm gì để bảo vệ mắt khi dùng máy tính?",
                                   "suggestions": ["👉 Giữ khoảng cách mắt 50-70cm so với màn hình.", "👉 Chớp mắt thư giãn và nghỉ ngơi sau 30 phút học."],
                                   "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_1\Tuần_04\images\ai_good_posture.jpg",
                                   "conclusion": "Ngồi đúng tư thế bảo vệ cột sống và luôn tuân thủ nội quy an toàn điện phòng máy!"}),
            ("quiz_mascot", {"question": "Khi thấy máy tính bốc khói hoặc có mùi khét, em phải làm gì?",
                             "options": ["A. Tự ý lấy nước dội vào máy", "B. Báo ngay cho giáo viên quản lý", "C. Tiếp tục ngồi chơi game", "D. Rủ bạn lại gần xem"]}),
            ("thanks", {"title": "Em học rất chăm chỉ! ⭐", "content": "BTVN: Hướng dẫn em nhỏ ở nhà ngồi học đúng tư thế!"})
        ]
    },
    # 3. Lớp 2 (Tiết 3)
    {
        "folder": "Lớp_2", "lop_label": "Lớp 2", "palette_idx": 2,
        "file_name": "Slide_Tin_hoc_Lớp_2_Tiet03_Ban_phim_va_go_cau_ngan.pptx",
        "builders": [
            ("cover", {"title": "Bàn phím và gõ câu ngắn ⌨️", "subtitle": "Tin học • Lớp 2 • Tiết 3"}),
            ("hero_example", {"title": "Các khu vực chính trên bàn phím",
                              "text": "1. Khu vực phím chữ (A - Z)\n2. Hàng phím số (0 - 9)\n3. Phím Cách (Spacebar)\n4. Phím Enter xuống dòng",
                              "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_2\Tuần_04\images\learn1_card1.png"}),
            ("arrow_badges", {"title": "Các phím chức năng quan trọng",
                              "badges": ["Phím Backspace (←): Xóa chữ bên trái", "Phím Caps Lock: Bật/Tắt viết IN HOA", "Phím Shift: Gõ chữ hoa và ký tự trên", "Phím Delete: Xóa chữ bên phải con trỏ"],
                              "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_2\Tuần_04\images\learn2_card1.png"}),
            ("conclusion_orange", {"question": "Làm thế nào để tạo khoảng cách giữa hai từ?",
                                   "suggestions": ["👉 Nhấn phím Cách (Spacebar) dài nhất ở hàng dưới cùng.", "👉 Chỉ nhấn 1 lần giữa hai từ liên tiếp."],
                                   "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_2\Tuần_04\images\learn1_card3.png",
                                   "conclusion": "Bàn phím là công cụ nhập văn bản quan trọng. Nắm vững các phím chức năng để soạn thảo nhanh!"}),
            ("quiz_mascot", {"question": "Phím nào dùng để xóa ký tự bên trái con trỏ soạn thảo?",
                             "options": ["A. Phím Enter", "B. Phím Backspace (←)", "C. Phím Spacebar", "D. Phím Caps Lock"]}),
            ("thanks", {"title": "Tuyệt vời! 🎉", "content": "BTVN: Luyện gõ tên của em và các bạn vào WordPad!"})
        ]
    },
    # 4. Lớp 3 (Tiết 3)
    {
        "folder": "Lớp_3", "lop_label": "Lớp 3", "palette_idx": 3,
        "file_name": "Slide_Tin_hoc_Lớp_3_Tiet03_Bai_3_May_tinh_va_em.pptx",
        "builders": [
            ("cover", {"title": "Bài 3. Máy tính và em 💻", "subtitle": "Tin học • Lớp 3 • Tiết 3"}),
            ("hero_example", {"title": "Lợi ích tuyệt vời của máy tính",
                              "text": "• Học tập trực tuyến, tìm kiếm tri thức\n• Vẽ tranh, nghe nhạc, giải trí lành mạnh\n• Liên lạc với người thân, bạn bè\n• Lưu trữ tài liệu học tập ngăn nắp",
                              "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_3\Tuần_04\images\learn1_card1.png"}),
            ("arrow_badges", {"title": "Tác hại khi sử dụng máy tính sai cách",
                              "badges": ["Ngồi lâu gây mỏi mắt và cận thị", "Nghiện trò chơi điện tử, lười vận động", "Nguy cơ tiếp xúc thông tin không lành mạnh", "Thức khuya ảnh hưởng đến sức khỏe"],
                              "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_3\Tuần_04\images\learn2_card1.png"}),
            ("conclusion_orange", {"question": "Thời gian sử dụng máy tính hợp lý mỗi ngày là bao nhiêu?",
                                   "suggestions": ["👉 Mỗi lần sử dụng không quá 30 - 45 phút.", "👉 Nghỉ ngơi mắt và vận động nhẹ giữa các giờ học."],
                                   "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_3\Tuần_04\images\learn2_card2.png",
                                   "conclusion": "Máy tính là người bạn hữu ích. Hãy sử dụng điều độ, khoa học để bảo vệ sức khỏe bản thân!"}),
            ("quiz_mascot", {"question": "Hành động nào sau đây là NÊN LÀM khi sử dụng máy tính?",
                             "options": ["A. Chơi game liên tục 4 tiếng", "B. Dùng máy tính học Tiếng Anh 30 phút", "C. Vừa ăn bánh vừa gõ phím", "D. Để máy tính gần cốc nước ngọt"]}),
            ("thanks", {"title": "Bài học kết thúc! 🌟", "content": "BTVN: Lập thời gian biểu học máy tính trong tuần của em!"})
        ]
    },
    # 5. Lớp 4 (Tiết 3)
    {
        "folder": "Lớp_4", "lop_label": "Lớp 4", "palette_idx": 4,
        "file_name": "Slide_Tin_hoc_Lớp_4_Tiet03_Bai_3_Thong_tin_tren_trang_Web.pptx",
        "builders": [
            ("cover", {"title": "Thông tin trên trang Web 🌐", "subtitle": "Tin học • Lớp 4 • Tiết 3"}),
            ("hero_example", {"title": "4 Dạng thông tin đa phương tiện",
                              "text": "1. Văn bản (Text): Bài viết, tin tức\n2. Hình ảnh (Images): Tranh vẽ, ảnh chụp\n3. Âm thanh (Audio): Bài hát, giọng đọc\n4. Video: Phim hoạt hình, bài giảng",
                              "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_4\Tuần_04\images\learn1_card1.png"}),
            ("arrow_badges", {"title": "Khám phá Siêu liên kết (Hyperlink)",
                              "badges": ["Đoạn chữ/hình ảnh dẫn sang trang mới", "Con trỏ chuột biến thành HÌNH BÀN TAY 👆", "Nháy chuột vào để mở liên kết", "Thanh địa chỉ URL chỉ rõ vị trí trang web"],
                              "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_4\Tuần_04\images\learn2_card1.png"}),
            ("conclusion_orange", {"question": "Làm sao biết một đoạn chữ là Siêu liên kết?",
                                   "suggestions": ["👉 Đoạn chữ thường có gạch chân hoặc màu xanh nổi bật.", "👉 Khi trỏ chuột vào, con trỏ đổi thành hình bàn tay chỉ trỏ."],
                                   "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_4\Tuần_04\images\learn2_card2.png",
                                   "conclusion": "Trang web chứa đa dạng thông tin. Siêu liên kết giúp kết nối và duyệt web nhanh chóng!"}),
            ("quiz_mascot", {"question": "Khi trỏ chuột vào Siêu liên kết, con trỏ chuột sẽ đổi thành hình gì?",
                             "options": ["A. Hình mũi tên đen", "B. Hình bàn tay chỉ ngón 👆", "C. Hình đồng hồ cát", "D. Hình dấu cộng lớn"]}),
            ("thanks", {"title": "Em giỏi lắm! 🚀", "content": "BTVN: Tìm 3 trang web học tập bổ ích và ghi vào vở nhé!"})
        ]
    },
    # 6. Lớp 5 (Tiết 3: Bài 3. Tìm kiếm thông tin)
    {
        "folder": "Lớp_5", "lop_label": "Lớp 5", "palette_idx": 5,
        "file_name": "Slide_Tin_hoc_Lớp_5_Tiet03_Bai_3_Tim_kiem_thong_tin.pptx",
        "builders": [
            ("cover", {"title": "Tìm kiếm thông tin hiệu quả 🔍", "subtitle": "Tin học • Lớp 5 • Tiết 3"}),
            ("hero_example", {"title": "Kỹ thuật sử dụng Từ khóa tìm kiếm",
                              "text": "• Chọn từ khóa ngắn gọn, chính xác\n• Đặt từ khóa trong dấu ngoặc kép \" \"\n• Tìm chính xác nguyên cụm từ\n• Tiết kiệm thời gian tra cứu dữ liệu",
                              "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_5\Tuần_04\images\learn1_card1.png"}),
            ("arrow_badges", {"title": "Đánh giá & Kiểm chứng thông tin",
                              "badges": ["Ưu tiên website uy tín (.gov.vn, .edu.vn)", "Kiểm tra tác giả và ngày đăng bài", "Đối chiếu nhiều nguồn thông tin khác nhau", "Không chia sẻ thông tin chưa kiểm chứng"],
                              "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_5\Tuần_04\images\learn1_card3.png"}),
            ("conclusion_orange", {"question": "Tại sao nên đặt từ khóa trong dấu ngoặc kép \" \"?",
                                   "suggestions": ["👉 Máy tìm kiếm sẽ tìm chính xác cả cụm từ theo đúng thứ tự.", "👉 Loại bỏ hàng triệu kết quả rời rạc không liên quan."],
                                   "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_5\Tuần_04\images\learn1_card2.png",
                                   "conclusion": "Tìm kiếm thông minh với từ khóa chuẩn xác và luôn kiểm chứng độ tin cậy của thông tin!"}),
            ("quiz_mascot", {"question": "Đuôi tên miền nào sau đây đại diện cho các cơ quan giáo dục?",
                             "options": ["A. .gov.vn", "B. .edu.vn", "C. .com", "D. .net"]}),
            ("thanks", {"title": "Hoàn thành Tiết 3! ⭐", "content": "Chuẩn bị bài tiếp theo: Cây thư mục (Tiết 4) nhé!"})
        ]
    },
    # 7. Lớp 5 (Tiết 4: Bài 4. Cây thư mục)
    {
        "folder": "Lớp_5", "lop_label": "Lớp 5", "palette_idx": 6,
        "file_name": "Slide_Tin_hoc_Lớp_5_Tiet04_Bai_4_Cay_thu_muc.pptx",
        "builders": [
            ("cover", {"title": "Cấu trúc Cây thư mục (Folder Tree) 📁", "subtitle": "Tin học • Lớp 5 • Tiết 4"}),
            ("hero_example", {"title": "Phân cấp Cây thư mục khoa học",
                              "text": "• Thư mục gốc (Root): Ổ đĩa C:, D:, E:\n• Thư mục mẹ và Thư mục con (Subfolder)\n• Tệp tin (File): .docx, .pptx, .png\n• Đường dẫn (Path) chỉ rõ vị trí lưu trữ",
                              "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_5\Tuần_04\images\ai_folder_tree.jpg"}),
            ("arrow_badges", {"title": "Các thành phần trong Cây thư mục",
                              "badges": ["Ổ đĩa gốc (C:, D:) chứa toàn bộ dữ liệu", "Thư mục mẹ HOC_TAP chứa các môn học", "Thư mục con TOAN, VAN, TIN_HOC", "Đường dẫn D:\\HOC_TAP\\TIN_HOC\\Bai4.docx"],
                              "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_5\Tuần_04\images\ai_folder_tree.jpg"}),
            ("conclusion_orange", {"question": "Lợi ích của việc tổ chức tệp theo Cây thư mục là gì?",
                                   "suggestions": ["👉 Giúp dữ liệu ngăn nắp, dễ quản lý.", "👉 Tìm kiếm tài liệu nhanh chóng, không bị thất lạc."],
                                   "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_5\Tuần_04\images\ai_folder_tree.jpg",
                                   "conclusion": "Tổ chức tệp tin theo cấu trúc hình cây giúp máy tính ngăn nắp và nâng cao hiệu suất học tập!"}),
            ("quiz_mascot", {"question": "Trong đường dẫn D:\\ANH\\Gia_dinh.png, tệp tin cần mở là gì?",
                             "options": ["A. Ổ đĩa D:", "B. Thư mục ANH", "C. Tệp Gia_dinh.png", "D. Hệ điều hành"]}),
            ("thanks", {"title": "Xuất sắc! ⭐", "content": "BTVN: Tạo cây thư mục HOC_TAP trên máy tính của em!"})
        ]
    },
    # 8. Lớp 6 (Tiết 3: Bài 3. Thông tin trong máy tính)
    {
        "folder": "Lớp_6", "lop_label": "Lớp 6", "palette_idx": 7,
        "file_name": "Slide_Tin_hoc_Lớp_6_Tiet03_Bai_3_Thong_tin_trong_may_tinh.pptx",
        "builders": [
            ("cover", {"title": "Thông tin trong máy tính 🔢", "subtitle": "Tin học • Lớp 6 • Tiết 3"}),
            ("hero_example", {"title": "Bit — Đơn vị đo thông tin cơ bản",
                              "text": "• Bit (Binary digit): Gồm 2 ký hiệu 0 và 1\n• 8 Bit ghép lại thành 1 Byte dữ liệu\n• Mọi văn bản, ảnh, âm thanh đều số hóa thành Bit\n• Dung lượng bộ nhớ đo khả năng lưu trữ",
                              "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_6\Tuần_04\images\learn1_card1.png"}),
            ("arrow_badges", {"title": "Bảng đơn vị đo dung lượng thông tin",
                              "badges": ["1 Byte (B) = 8 Bit", "1 Kilobyte (KB) = 1024 Byte", "1 Megabyte (MB) = 1024 KB", "1 Gigabyte (GB) = 1024 MB", "1 Terabyte (TB) = 1024 GB"],
                              "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_6\Tuần_04\images\learn1_card2.png"}),
            ("conclusion_orange", {"question": "Máy tính xử lý dữ liệu bằng ký hiệu nào?",
                                   "suggestions": ["👉 Máy tính chỉ hiểu và xử lý dãy bit gồm 2 số 0 và 1.", "👉 Thiết bị lưu trữ có dung lượng càng lớn thì chứa được càng nhiều dữ liệu."],
                                   "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_6\Tuần_04\images\learn1_card3.png",
                                   "conclusion": "Dãy Bit là ngôn ngữ cốt lõi của máy tính điện tử. Nắm vững các đơn vị đo dữ liệu để quản lý dung lượng!"}),
            ("quiz_mascot", {"question": "1 Byte bằng bao nhiêu Bit?",
                             "options": ["A. 2 Bit", "B. 4 Bit", "C. 8 Bit", "D. 1024 Bit"]}),
            ("thanks", {"title": "Hoàn thành Tiết 3! 🚀", "content": "Chuẩn bị bài tiếp theo: Mạng máy tính (Tiết 4) nhé!"})
        ]
    },
    # 9. Lớp 6 (Tiết 4: Bài 4. Mạng máy tính)
    {
        "folder": "Lớp_6", "lop_label": "Lớp 6", "palette_idx": 8,
        "file_name": "Slide_Tin_hoc_Lớp_6_Tiet04_Bai_4_Mang_may_tinh.pptx",
        "builders": [
            ("cover", {"title": "Khám phá Mạng máy tính 🌐", "subtitle": "Tin học • Lớp 6 • Tiết 4"}),
            ("hero_example", {"title": "Khái niệm và Phân loại mạng",
                              "text": "• Mạng LAN: Mạng cục bộ trong trường học, gia đình\n• Mạng Internet: Mạng toàn cầu kết nối triệu máy tính\n• Mạng có dây: Dùng cáp mạng truyền tín hiệu\n• Mạng không dây: Dùng sóng Wi-Fi linh hoạt",
                              "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_6\Tuần_04\images\ai_lan_network.jpg"}),
            ("arrow_badges", {"title": "Các thành phần của mạng máy tính",
                              "badges": ["Thiết bị đầu cuối: PC, Laptop, Smartphone", "Thiết bị kết nối: Switch, Hub, Router, Modem", "Môi trường truyền dẫn: Cáp mạng và sóng Wi-Fi", "Lợi ích: Dùng chung dữ liệu và máy in"],
                              "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_6\Tuần_04\images\ai_lan_network.jpg"}),
            ("conclusion_orange", {"question": "Lợi ích lớn nhất của mạng máy tính là gì?",
                                   "suggestions": ["👉 Chia sẻ thông tin, tài liệu nhanh chóng tức thì.", "👉 Dùng chung các thiết bị phần cứng đắt tiền như máy in."],
                                   "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_6\Tuần_04\images\ai_lan_network.jpg",
                                   "conclusion": "Mạng máy tính kết nối thế giới, giúp việc trao đổi dữ liệu và hợp tác trở nên dễ dàng hơn bao giờ hết!"}),
            ("quiz_mascot", {"question": "Thiết bị nào sau đây dùng để phát sóng mạng không dây?",
                             "options": ["A. Bàn phím", "B. Bộ định tuyến không dây (Wi-Fi Router)", "C. Màn hình máy tính", "D. Chuột quang"]}),
            ("thanks", {"title": "Bài học kết thúc! ⭐", "content": "BTVN: Quan sát các thiết bị mạng trong gia đình em!"})
        ]
    },
    # 10. Lớp 7 (Tiết 3: Bài 3. Quản lý dữ liệu trong máy tính)
    {
        "folder": "Lớp_7", "lop_label": "Lớp 7", "palette_idx": 9,
        "file_name": "Slide_Tin_hoc_Lớp_7_Tiet03_Bai_3_Quan_ly_du_lieu_trong_may_tinh.pptx",
        "builders": [
            ("cover", {"title": "Quản lý dữ liệu trong máy tính 💾", "subtitle": "Tin học • Lớp 7 • Tiết 3"}),
            ("hero_example", {"title": "Sao lưu & Bảo vệ dữ liệu an toàn",
                              "text": "• Sao lưu dữ liệu định kỳ (Backup)\n• Lưu trữ đám mây: Google Drive, OneDrive\n• Nén tệp giảm dung lượng: .zip, .rar\n• Đặt mật khẩu bảo vệ tài liệu quan trọng",
                              "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_7\Tuần_04\images\learn1_card1.png"}),
            ("arrow_badges", {"title": "Quy trình nén và giải nén tệp tin",
                              "badges": ["Chọn tệp cần nén -> Nhấp chuột phải", "Chọn Send to -> Compressed (zipped) folder", "Tệp .zip giúp gửi thư điện tử nhanh chóng", "Nhấp chuột phải -> Extract All để giải nén"],
                              "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_7\Tuần_04\images\learn1_card3.png"}),
            ("conclusion_orange", {"question": "Tại sao phải sao lưu dữ liệu sang nhiều nơi?",
                                   "suggestions": ["👉 Đề phòng máy tính bị hỏng ổ cứng hoặc nhiễm virus.", "👉 Có thể khôi phục lại dữ liệu học tập bất cứ lúc nào."],
                                   "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_7\Tuần_04\images\learn1_card2.png",
                                   "conclusion": "Chủ động sao lưu và bảo mật dữ liệu là thói quen quan trọng của công dân số thông minh!"}),
            ("quiz_mascot", {"question": "Định dạng tệp nào sau đây là tệp nén dữ liệu?",
                             "options": ["A. .docx", "B. .zip", "C. .pptx", "D. .mp3"]}),
            ("thanks", {"title": "Hoàn thành Tiết 3! 🚀", "content": "Chuẩn bị bài tiếp theo: Mạng xã hội (Tiết 4) nhé!"})
        ]
    },
    # 11. Lớp 7 (Tiết 4: Bài 4. Mạng xã hội)
    {
        "folder": "Lớp_7", "lop_label": "Lớp 7", "palette_idx": 10,
        "file_name": "Slide_Tin_hoc_Lớp_7_Tiet04_Bai_4_Mang_xa_hoi.pptx",
        "builders": [
            ("cover", {"title": "Mạng xã hội & Văn hóa số 📱", "subtitle": "Tin học • Lớp 7 • Tiết 4"}),
            ("hero_example", {"title": "Khai thác Mạng xã hội an toàn",
                              "text": "• Kết nối bạn bè, trao đổi kiến thức học tập\n• Cảnh giác trước tin giả và lừa đảo qua mạng\n• Bảo vệ thông tin cá nhân và mật khẩu tài khoản\n• Ứng xử văn minh, tôn trọng người khác",
                              "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_7\Tuần_04\images\learn2_card1.png"}),
            ("arrow_badges", {"title": "Nguyên tắc vàng khi tham gia Mạng xã hội",
                              "badges": ["Không chia sẻ địa chỉ nhà, số CCCD công khai", "Không bấm vào đường link lạ có thưởng", "Không bình luận thô tục hay bắt nạt qua mạng", "Báo người lớn khi gặp phiền toái trên mạng"],
                              "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_7\Tuần_04\images\learn2_card3.png"}),
            ("conclusion_orange", {"question": "Làm gì khi bắt gặp thông tin sai sự thật trên mạng xã hội?",
                                   "suggestions": ["👉 Tuyệt đối không bấm chia sẻ (Share) tiếp.", "👉 Báo cáo (Report) bài viết vi phạm cho nền tảng."],
                                   "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_7\Tuần_04\images\learn2_card4.png",
                                   "conclusion": "Mạng xã hội là con dao hai lưỡi. Hãy là người dùng tỉnh táo, có trách nhiệm và văn minh!"}),
            ("quiz_mascot", {"question": "Thông tin nào sau đây TUYỆT ĐỐI KHÔNG chia sẻ công khai trên mạng?",
                             "options": ["A. Sở thích âm nhạc", "B. Mật khẩu tài khoản và mã OTP", "C. Cảm nhận về một cuốn sách hay", "D. Ảnh phong cảnh trường học"]}),
            ("thanks", {"title": "Bài học kết thúc! 🌟", "content": "BTVN: Rà soát lại quyền riêng tư trên các tài khoản mạng của em!"})
        ]
    },
    # 12. Lớp 8 (Tiết 3: Bài 3. Thực hành khai thác thông tin số)
    {
        "folder": "Lớp_8", "lop_label": "Lớp 8", "palette_idx": 11,
        "file_name": "Slide_Tin_hoc_Lớp_8_Tiet03_Bai_3_Thuc_hanh_khai_thac_thong_tin_so.pptx",
        "builders": [
            ("cover", {"title": "Thực hành khai thác thông tin số 🔍", "subtitle": "Tin học • Lớp 8 • Tiết 3"}),
            ("hero_example", {"title": "Cú pháp tìm kiếm nâng cao",
                              "text": "• filetype:pdf tìm chính xác định dạng tệp\n• site:edu.vn giới hạn trên website giáo dục\n• Dấu ngoặc kép \" \" tìm chính xác cụm từ\n• Trừ từ khóa (-) để loại bỏ nội dung rác",
                              "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_8\Tuần_04\images\learn1_card1.png"}),
            ("arrow_badges", {"title": "Quy tắc trích dẫn nguồn & Bản quyền",
                              "badges": ["Ghi rõ tên tác giả và năm xuất bản", "Dẫn đường link (URL) của nguồn thông tin", "Không sao chép nguyên văn mà không trích dẫn", "Tôn trọng quyền sở hữu trí tuệ của tác giả số"],
                              "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_8\Tuần_04\images\learn1_card3.png"}),
            ("conclusion_orange", {"question": "Tại sao việc trích dẫn nguồn lại bắt buộc trong học tập?",
                                   "suggestions": ["👉 Thể hiện sự tôn trọng công sức nghiên cứu của tác giả.", "👉 Giúp bài thuyết trình có độ tin cậy và học thuật cao."],
                                   "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_8\Tuần_04\images\learn1_card4.png",
                                   "conclusion": "Khai thác thông tin số chuyên nghiệp đi đôi với văn hóa trích dẫn và tôn trọng bản quyền!"}),
            ("quiz_mascot", {"question": "Cú pháp nào tìm các tệp PDF về chủ đề Trí tuệ nhân tạo?",
                             "options": ["A. Trí tuệ nhân tạo .doc", "B. \"Trí tuệ nhân tạo\" filetype:pdf", "C. find PDF Trí tuệ nhân tạo", "D. AI.download.pdf"]}),
            ("thanks", {"title": "Hoàn thành Tiết 3! ⭐", "content": "Chuẩn bị ôn tập Đánh giá định kỳ 1 (Tiết 4) nhé!"})
        ]
    },
    # 13. Lớp 8 (Tiết 4: Ôn tập Đánh giá định kỳ 1)
    {
        "folder": "Lớp_8", "lop_label": "Lớp 8", "palette_idx": 12,
        "file_name": "Slide_Tin_hoc_Lớp_8_Tiet04_On_tap_DGDK1.pptx",
        "builders": [
            ("cover", {"title": "Ôn tập Đánh giá định kỳ 1 📝", "subtitle": "Tin học • Lớp 8 • Tiết 4"}),
            ("hero_example", {"title": "Hệ thống hóa kiến thức trọng tâm",
                              "text": "1. Lịch sử phát triển 5 thế hệ máy tính\n2. Thông tin trong môi trường số & Độ tin cậy\n3. Cú pháp khai thác thông tin nâng cao\n4. Đạo đức và bản quyền công nghệ số",
                              "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_8\Tuần_04\images\learn2_card1.png"}),
            ("quiz_mascot", {"question": "Máy tính thế hệ thứ nhất sử dụng linh kiện điện tử nào?",
                             "options": ["A. Đèn điện tử chân không", "B. Bóng bán dẫn (Transistor)", "C. Mạch tích hợp (IC)", "D. Bộ vi xử lý (VLSI)"]}),
            ("quiz_mascot", {"question": "Hành vi nào sau đây là VI PHẠM bản quyền phần mềm?",
                             "options": ["A. Mua bản quyền sử dụng chính hãng", "B. Sử dụng phần mềm mã nguồn mở", "C. Bẻ khóa (Crack) phần mềm có phí", "D. Tự viết phần mềm phục vụ học tập"]}),
            ("conclusion_orange", {"question": "Bí quyết đạt điểm tối đa trong bài kiểm tra định kỳ là gì?",
                                   "suggestions": ["👉 Ôn lại các mốc lịch sử máy tính và đơn vị đo dữ liệu.", "👉 Nắm vững quy tắc an toàn thông tin và trích dẫn bản quyền."],
                                   "img": r"D:\UNIGO\KHBD_Tin_học\Lớp_8\Tuần_04\images\learn2_card4.png",
                                   "conclusion": "Nắm chắc lý thuyết và vận dụng linh hoạt để tự tin đạt điểm 10 trong bài kiểm tra định kỳ 1!"}),
            ("thanks", {"title": "Chúc các em làm bài đạt điểm tối đa! 🌟", "content": "Tự tin, bình tĩnh và hoàn thành xuất sắc bài kiểm tra nhé!"})
        ]
    }
]

def generate_all_separated_slides():
    print("============================================================")
    print(" TẠO LẠI TOÀN BỘ SLIDE TIN HỌC TUẦN 04 (TÁCH RIÊNG TỪNG KHBD)")
    print("============================================================")

    for item in SEPARATED_LESSONS_T4:
        prs = Presentation(TEMPLATE)
        layout = prs.slide_layouts[6]
        pal = COLOR_PALETTES[item["palette_idx"] % len(COLOR_PALETTES)]
        
        # Build slides based on builders
        for btype, bdata in item["builders"]:
            if btype == "cover":
                build_cover(prs, bdata, pal, item, layout)
            elif btype == "section_banner":
                build_section_banner(prs, bdata, pal, item, layout)
            elif btype == "hero_example":
                build_hero_example(prs, bdata, pal, item, layout)
            elif btype == "posture_compare":
                build_posture_compare(prs, bdata, pal, item, layout)
            elif btype == "arrow_badges":
                build_arrow_badges_diagram(prs, bdata, pal, item, layout)
            elif btype == "conclusion_orange":
                build_conclusion_orange(prs, bdata, pal, item, layout)
            elif btype == "quiz_mascot":
                build_quiz_mascot(prs, bdata, pal, item, layout)
            elif btype == "thanks":
                build_thanks(prs, bdata, pal, item, layout)

        out_dir = os.path.join(KHBD_BASE, item["folder"], "Tuần_04")
        os.makedirs(out_dir, exist_ok=True)
        out_path = os.path.join(out_dir, item["file_name"])
        prs.save(out_path)
        print(f"  [+] Đã tạo: {item['folder']} -> {item['file_name']} ({len(prs.slides)} slides)")

    print("\n[OK] Hoàn tất tạo toàn bộ Slide Tin học Tuần 04 tách riêng từng KHBD!")

if __name__ == "__main__":
    generate_all_separated_slides()
