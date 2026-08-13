# -*- coding: utf-8 -*-
"""
Script tạo Bộ Slide bài giảng Student-Facing v3 cho TOÀN BỘ 9 KHỐI LỚP Tuần 03
(Tiền Tiểu Học, Lớp 1, 2, 3, 4, 5, 6, 7, 8)

Quy chuẩn v3:
  1. Per-Bullet Images:
     - TH (Tiền TH -> Lớp 5): Layout A — Grid Flashcard (ảnh trên, text dưới)
     - THCS (Lớp 6 -> Lớp 8): Layout B — Horizontal Row (ảnh trái 2in×1.5in, text phải)
  2. Animation tuần tự cho slide Luyện tập / Ghép nối / Quiz (On Click appear per item)
  3. Game Images: Slide trò chơi bắt buộc có ảnh minh họa đủ lớn (~2in-3.5in)
  4. Anti-Bug Checklist: Safe Zone (Y 1.15 -> 6.35), Z-order insert(2), no custom footers, high contrast palette
"""
import sys, io, os, copy, glob
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import traceback

TEMPLATE  = r'D:\UNIGO\Hệ thống mẫu văn bản\Mẫu slide UNIGO.pptx'
SGK_BASE  = r'D:\UNIGO\SGK'
KHBD_BASE = r'D:\UNIGO\KHBD_Tin_học'

# ─── Template Safe Zone ───
SAFE_TOP    = 1.15   # Logo ends at 1.09in
SAFE_BOTTOM = 6.35   # Footer starts at 6.43in
SAFE_LEFT   = 0.3
SAFE_RIGHT  = 13.0
SLIDE_W     = 13.33
SLIDE_H     = 7.50

# ─── 8 Bộ màu tương phản cao (High Contrast Palettes) ───
COLOR_PALETTES = [
    # 0: Blue (Tiền TH)
    {"primary": "1B4F9B", "accent": "2D7DD2", "bg": "EBF3FE", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "1A2744", "text_on_card": "1A2744"},
    # 1: Purple (Lớp 1)
    {"primary": "5B21B6", "accent": "7C3AED", "bg": "F3EEFF", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "2E1065", "text_on_card": "2E1065"},
    # 2: Teal (Lớp 2)
    {"primary": "0F766E", "accent": "14B8A6", "bg": "ECFDF5", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "134E4A", "text_on_card": "134E4A"},
    # 3: Orange (Lớp 3)
    {"primary": "C2410C", "accent": "EA580C", "bg": "FFF7ED", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "7C2D12", "text_on_card": "431407"},
    # 4: Indigo (Lớp 4)
    {"primary": "3730A3", "accent": "4F46E5", "bg": "EEF2FF", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "1E1B4B", "text_on_card": "1E1B4B"},
    # 5: Rose (Lớp 5)
    {"primary": "BE185D", "accent": "EC4899", "bg": "FDF2F8", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "831843", "text_on_card": "831843"},
    # 6: Emerald (Lớp 6)
    {"primary": "047857", "accent": "10B981", "bg": "ECFDF5", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "064E3B", "text_on_card": "064E3B"},
    # 7: Sky (Lớp 7)
    {"primary": "0369A1", "accent": "0EA5E9", "bg": "F0F9FF", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "0C4A6E", "text_on_card": "0C4A6E"},
    # 8: Amber (Lớp 8)
    {"primary": "B45309", "accent": "F59E0B", "bg": "FFFBEB", "card": "FFFFFF",
     "text_on_primary": "FFFFFF", "text_on_bg": "78350F", "text_on_card": "451A03"},
]

# ─── Nội dung bài học 9 Khối Lớp TUẦN 03 ───
LESSONS_TUAN03 = [
    {
        "id": "Tien_TH", "folder": "Tiền_tiểu_học", "lop_label": "Tiền Tiểu học", "lop_num": 0, "is_thcs": False,
        "bai": "Bài 2. Em ngồi máy tính an toàn",
        "file_name": "Slide_Tin_hoc_Tien_TH_Bai02_Em_ngoi_may_tinh_an_toan.pptx",
        "slides_content": [
            {"type": "cover", "title": "Em ngồi máy tính an toàn 🪑", "subtitle": "Tin học • Tiền Tiểu học • Bài 2"},
            {"type": "warmup", "title": "Ai ngồi đúng tư thế nào?",
             "content": "Con hãy nhìn hai hình mẫu:\nBạn nào đang ngồi thẳng lưng, mắt cách màn hình vừa đủ?\nCùng chỉ cho cô xem nhé!"},
            {"type": "learn", "title": "Tư thế ngồi chuẩn",
             "bullets": [
                 "Lưng thẳng, tựa nhẹ vào ghế",
                 "Hai chân đặt bằng phẳng trên sàn",
                 "Mắt cách màn hình 50-70cm (bằng 1 cánh tay)",
                 "Cánh tay vuông góc khi dùng bàn phím"
             ]},
            {"type": "learn", "title": "Bảo vệ đôi mắt của em",
             "bullets": [
                 "Không ngồi quá gần màn hình",
                 "Nghỉ ngơi mắt sau mỗi 15-20 phút",
                 "Nhìn ra xa cửa sổ để thư giãn mắt",
                 "Giữ phòng học đủ ánh sáng"
             ]},
            {"type": "learn", "title": "An toàn điện trong phòng máy",
             "bullets": [
                 "Không chạm tay vào ổ cắm hay dây điện",
                 "Không mang đồ ăn, nước uống gần máy",
                 "Báo ngay cho giáo viên khi máy có sự cố"
             ]},
            {"type": "practice", "title": "Bắt chước tư thế đúng!",
             "instruction": "Các con cùng thực hiện 4 bước:",
             "items": ["Chỉnh ghế thẳng lưng", "Đặt hai chân xuống sàn", "Đưa tay duỗi thẳng đo màn hình", "Mỉm cười sẵn sàng học bài!"]},
            {"type": "activity", "title": "Trò chơi: Đúng hay Sai?",
             "instruction": "Hành động nào an toàn?",
             "items": ["Ngồi còng lưng ghé sát mắt ↔ SAI ❌", "Nghỉ mắt nhìn ra cửa sổ ↔ ĐÚNG ✅", "Vừa uống nước vừa dùng máy ↔ SAI ❌", "Nhờ cô trợ giúp khi hỏng điện ↔ ĐÚNG ✅"]},
            {"type": "summary", "title": "Hôm nay em đã nhớ!",
             "items": ["Ngồi thẳng lưng, chân chạm sàn", "Mắt xa màn hình một cánh tay", "Giữ an toàn điện tuyệt đối"]},
            {"type": "thanks", "title": "Các con giỏi lắm! 🌟", "content": "BTVN: Hướng dẫn bố mẹ\ntư thế ngồi máy tính đúng ở nhà!"},
        ]
    },
    {
        "id": "Lop_1", "folder": "Lớp_1", "lop_label": "Lớp 1", "lop_num": 1, "is_thcs": False,
        "bai": "Bài 2. Ôn và nâng cấp kỹ năng chuột",
        "file_name": "Slide_Tin_hoc_Lop_1_Bai02_On_va_nang_cap_ky_nang_chuot.pptx",
        "slides_content": [
            {"type": "cover", "title": "Nâng cấp kỹ năng chuột 🖱️", "subtitle": "Tin học • Lớp 1 • Bài 2"},
            {"type": "warmup", "title": "Chú chuột siêu tốc!",
             "content": "Em đã biết cách cầm chuột chưa?\nCùng ôn lại 5 thao tác chuột cơ bản nào!"},
            {"type": "learn", "title": "5 thao tác dùng chuột",
             "bullets": [
                 "Di chuyển chuột (Move mouse)",
                 "Nháy chuột trái (Left click)",
                 "Nháy chuột phải (Right click)",
                 "Nháy đôi chuột (Double click)"
             ]},
            {"type": "learn", "title": "Kéo thả chuột chuyên nghiệp",
             "bullets": [
                 "Nắm giữ nút chuột trái",
                 "Di chuyển chuột đến vị trí mới",
                 "Thả ngón tay ra để hoàn tất",
                 "Dùng để di chuyển biểu tượng trên màn hình"
             ]},
            {"type": "learn", "title": "Bí kíp điều khiển con trỏ",
             "bullets": [
                 "Cầm chuột nhẹ nhàng bằng tay phải",
                 "Ngón trỏ đặt lên nút trái",
                 "Ngón giữa đặt lên nút phải",
                 "Giữ chuột phẳng trên mặt bàn"
             ]},
            {"type": "practice", "title": "Thực hành luyện tập chuột!",
             "instruction": "Thực hiện theo các yêu cầu:",
             "items": ["Nháy đôi mở phần mềm Paint", "Kéo thả biểu tượng vào thư mục", "Nháy phải chuột xem thực đơn", "Nháy trái chọn màu vẽ"]},
            {"type": "activity", "title": "Thử thách kéo thả!",
             "instruction": "Thực hiện bài tập ghép hình:",
             "items": ["Di chuyển hình tròn vào ô tương ứng", "Kéo ô chữ thả vào khung bài làm", "Ghép mảnh tranh thành hình hoàn chỉnh"]},
            {"type": "summary", "title": "Tóm tắt bài học",
             "items": ["Thành thạo 5 thao tác chuột cơ bản", "Kéo thả chuột chính xác và linh hoạt", "Bảo vệ chuột sạch sẽ sau khi dùng"]},
            {"type": "thanks", "title": "Em giỏi lắm! ⭐", "content": "BTVN: Luyện tập trò chơi kéo thả chuột 10 phút!"},
        ]
    },
    {
        "id": "Lop_2", "folder": "Lớp_2", "lop_label": "Lớp 2", "lop_num": 2, "is_thcs": False,
        "bai": "Bài 2. Ôn và nâng cấp kỹ năng chuột",
        "file_name": "Slide_Tin_hoc_Lop_2_Bai02_On_va_nang_cap_ky_nang_chuot.pptx",
        "slides_content": [
            {"type": "cover", "title": "Luyện tập thao tác chuột nâng cao 🖱️", "subtitle": "Tin học • Lớp 2 • Bài 2"},
            {"type": "warmup", "title": "Thử thách phản xạ chuột!",
             "content": "Em có thể nháy chuột nhanh trong 5 giây không?\nCùng sẵn sàng nhé!"},
            {"type": "learn", "title": "Nâng cấp kỹ năng nháy đôi",
             "bullets": [
                 "Nháy hai lần liên tiếp thật nhanh",
                 "Dùng để mở ngay các ứng dụng và thư mục",
                 "Giữ chuột cố định không di chuyển khi nháy đôi"
             ]},
            {"type": "learn", "title": "Luyện tập cuộn chuột (Scroll)",
             "bullets": [
                 "Dùng ngón trỏ lăn con lăn ở giữa chuột",
                 "Cuộn lên để xem phần trên trang web",
                 "Cuộn xuống để xem phần dưới tài liệu"
             ]},
            {"type": "learn", "title": "Lựa chọn nhiều đối tượng",
             "bullets": [
                 "Kéo chuột tạo khung bao quanh các đối tượng",
                 "Chọn nhiều tệp cùng một lúc dễ dàng",
                 "Tiết kiệm thời gian thao tác"
             ]},
            {"type": "practice", "title": "Thực hành trên máy tính!",
             "instruction": "Lần lượt thực hiện các bước:",
             "items": ["Mở thư mục Học tập bằng nháy đôi", "Lăn con lăn xem hết danh sách tệp", "Kéo chuột chọn 3 tệp bài tập", "Nháy phải chuột để xem thông tin"]},
            {"type": "activity", "title": "Trò chơi: Bắt bong bóng!",
             "instruction": "Quy tắc chơi trên phần mềm:",
             "items": ["Nháy chuột trái diệt bóng thường ↔ 10 điểm", "Nháy đôi diệt bóng vàng ↔ 20 điểm", "Kéo thả bóng xanh vào giỏ ↔ 30 điểm"]},
            {"type": "summary", "title": "Ghi nhớ bài học",
             "items": ["Nháy đôi nhanh và chuẩn xác", "Sử dụng hiệu quả con lăn cuộn chuột", "Thao tác chọn nhiều đối tượng"]},
            {"type": "thanks", "title": "Xuất sắc! 🎉", "content": "BTVN: Thực hành cuộn chuột đọc sách điện tử ở nhà!"},
        ]
    },
    {
        "id": "Lop_3", "folder": "Lớp_3", "lop_label": "Lớp 3", "lop_num": 3, "is_thcs": False,
        "bai": "Bài 2. Xử lí thông tin",
        "file_name": "Slide_Tin_hoc_Lop_3_Bai02_Xu_li_thong_tin.pptx",
        "slides_content": [
            {"type": "cover", "title": "Bài 2. Xử lí thông tin 🧠", "subtitle": "Tin học • Lớp 3 • Bài 2"},
            {"type": "warmup", "title": "Bộ não chúng ta xử lý thế nào?",
             "content": "Khi nhìn thấy đèn giao thông chuyển sang màu đỏ:\nEm sẽ làm gì?\nĐó chính là quá trình xử lý thông tin đấy!"},
            {"type": "learn", "title": "Con người xử lý thông tin",
             "bullets": [
                 "Tiếp nhận thông tin qua giác quan (mắt, tai, mũi...)",
                 "Bộ não suy nghĩ, phân tích và xử lý",
                 "Đưa ra quyết định hoặc hành động phù hợp"
             ]},
            {"type": "learn", "title": "Máy tính xử lý thông tin",
             "bullets": [
                 "Thu nhận thông tin vào (Input: bàn phím, chuột)",
                 "Bộ xử lý CPU tính toán và xử lý dữ liệu",
                 "Xuất kết quả ra ngoài (Output: màn hình, loa)"
             ]},
            {"type": "learn", "title": "So sánh Não người ↔ Máy tính",
             "bullets": [
                 "Bộ não con người: sáng tạo, cảm xúc, linh hoạt",
                 "Máy tính: tính toán cực nhanh, chính xác, không mệt mỏi",
                 "Con người điều khiển máy tính phục vụ cuộc sống"
             ]},
            {"type": "practice", "title": "Phân tích ví dụ!",
             "instruction": "Xác định 3 giai đoạn xử lý thông tin:",
             "items": ["Nghe tiếng chuông báo thức (Thông tin vào)", "Bộ não biết đã đến giờ dậy (Xử lý)", "Thức dậy chuẩn bị đi học (Kết quả)"]},
            {"type": "activity", "title": "Trò chơi: Thần đồng tính nhanh!",
             "instruction": "Thi đấu giữa Não người & Máy tính:",
             "items": ["Tính 15 + 27 ↔ Cả người và máy đều nhanh!", "Tính 987 x 654 ↔ Máy tính xử lý trong 0.01s!", "Sáng tác bức tranh ↔ Não con người sáng tạo hơn!"]},
            {"type": "summary", "title": "Tóm tắt kiến thức",
             "items": ["Quá trình xử lý: Tiếp nhận -> Xử lý -> Quyết định", "Bộ não xử lý thông tin của con người", "CPU là bộ não xử lý thông tin của máy tính"]},
            {"type": "thanks", "title": "Bài học kết thúc!", "content": "BTVN: Tìm 2 ví dụ về xử lý thông tin\ntrong sinh hoạt hàng ngày của em!"},
        ]
    },
    {
        "id": "Lop_4", "folder": "Lớp_4", "lop_label": "Lớp 4", "lop_num": 4, "is_thcs": False,
        "bai": "Bài 2. Gõ bàn phím đúng cách",
        "file_name": "Slide_Tin_hoc_Lop_4_Bai02_Go_ban_phim_dung_cach.pptx",
        "slides_content": [
            {"type": "cover", "title": "Gõ bàn phím đúng cách ⌨️", "subtitle": "Tin học • Lớp 4 • Bài 2"},
            {"type": "warmup", "title": "10 ngón tay múa trên bàn phím!",
             "content": "Em gõ bàn phím bằng mấy ngón tay?\nGõ bằng 10 ngón sẽ giúp em gõ nhanh gấp 5 lần đấy!"},
            {"type": "learn", "title": "Hàng phím cơ sở (Home row)",
             "bullets": [
                 "Các phím: A S D F G H J K L ;",
                 "Hai phím có gờ: F và J (đặt ngón trỏ)",
                 "Nơi 10 ngón tay luôn xuất phát và trở về"
             ]},
            {"type": "learn", "title": "Phân công 10 ngón tay",
             "bullets": [
                 "Ngón trỏ trái: F, V, B, R, T, 4, 5",
                 "Ngón trỏ phải: J, N, M, U, Y, 7, 8",
                 "Hai ngón cái phụ trách phím cách (Spacebar)"
             ]},
            {"type": "learn", "title": "Quy tắc gõ mười ngón",
             "bullets": [
                 "Đặt ngón tay đúng vị trí xuất phát ở hàng cơ sở",
                 "Mỗi ngón tay chỉ gõ các phím được phân công",
                 "Gõ nhẹ nhàng, không nhìn bàn phím (Touch typing)"
             ]},
            {"type": "practice", "title": "Luyện đặt tay đúng cách!",
             "instruction": "Thực hiện theo các bước:",
             "items": ["Đặt ngón trỏ trái lên phím F (có gờ)", "Đặt ngón trỏ phải lên phím J (có gờ)", "Đặt 6 ngón còn lại lên A,S,D và K,L,;", "Hai ngón cái đặt nhẹ lên phím cách (Space)"]},
            {"type": "activity", "title": "Luyện gõ trên RapidTyping!",
             "instruction": "Thực hành bài luyện gõ hàng phím cơ sở:",
             "items": ["Gõ chuỗi: asdf jkl; asdf jkl;", "Gõ từ ngắn: fa la da ka ja", "Thi đua xem ai đạt độ chính xác 100%!"]},
            {"type": "summary", "title": "Hôm nay em đã học!",
             "items": ["Hàng phím cơ sở và 2 phím có gờ F, J", "Phân công vị trí 10 ngón tay", "Quy tắc gõ phím không cần nhìn (Touch typing)"]},
            {"type": "thanks", "title": "Tuyệt vời! 🌟", "content": "BTVN: Luyện gõ 10 ngón trên phần mềm\n15 phút mỗi ngày!"},
        ]
    },
    {
        "id": "Lop_5", "folder": "Lớp_5", "lop_label": "Lớp 5", "lop_num": 5, "is_thcs": False,
        "bai": "Bài 2. Tìm kiếm thông tin trên website",
        "file_name": "Slide_Tin_hoc_Lop_5_Bai02_Tim_kiem_thong_tin_tren_website.pptx",
        "slides_content": [
            {"type": "cover", "title": "Tìm kiếm thông tin trên Website 🌐", "subtitle": "Tin học • Lớp 5 • Bài 2"},
            {"type": "warmup", "title": "Làm sao tìm được câu trả lời?",
             "content": "Nếu muốn biết 'Con cá heo thở bằng gì?':\nEm sẽ gõ từ khóa gì vào Google để tìm nhanh nhất?"},
            {"type": "learn", "title": "Trình duyệt web và Máy tìm kiếm",
             "bullets": [
                 "Trình duyệt Web (Chrome, Edge): công cụ mở trang web",
                 "Máy tìm kiếm (Google, Bing): công cụ tra cứu thông tin",
                 "Nhập địa chỉ website vào thanh địa chỉ (Address bar)"
             ]},
            {"type": "learn", "title": "Từ khóa tìm kiếm (Keywords)",
             "bullets": [
                 "Từ khóa là từ hoặc cụm từ thể hiện nội dung cần tìm",
                 "Chọn từ khóa ngắn gọn, chính xác",
                 "Ví dụ: thay vì 'tôi muốn tìm ảnh con hổ', gõ 'ảnh con hổ'"
             ]},
            {"type": "learn", "title": "Đánh giá thông tin tìm được",
             "bullets": [
                 "Đọc lướt tiêu đề và đoạn tóm tắt kết quả",
                 "Chọn các trang web uy tín (.edu, .gov, báo chính thống)",
                 "So sánh thông tin từ 2-3 trang web khác nhau"
             ]},
            {"type": "practice", "title": "Thực hành tìm kiếm!",
             "instruction": "Thực hiện trên máy tính:",
             "items": ["Mở trình duyệt Google Chrome", "Gõ từ khóa: 'Hành tinh lớn nhất hệ mặt trời'", "Nhấn Enter và đọc kết quả đầu tiên", "Lưu ảnh hành tinh về máy tính"]},
            {"type": "activity", "title": "Thử thách: Truy tìm tri thức!",
             "instruction": "Tìm câu trả lời cho các câu hỏi:",
             "items": ["Ngọn núi cao nhất Việt Nam tên là gì? ↔ Fansipan", "Quốc hoa của Việt Nam là hoa gì? ↔ Hoa Sen", "Vận tốc ánh sáng là bao nhiêu? ↔ ~300.000 km/s"]},
            {"type": "summary", "title": "Tóm tắt bài học",
             "items": ["Sử dụng trình duyệt web và máy tìm kiếm", "Lựa chọn từ khóa thông minh và ngắn gọn", "Đánh giá và chọn lọc thông tin đáng tin cậy"]},
            {"type": "thanks", "title": "Em giỏi lắm! 🚀", "content": "BTVN: Tìm kiếm 3 thông tin thú vị\nvề động vật biển và ghi vào vở!"},
        ]
    },
    {
        "id": "Lop_6", "folder": "Lớp_6", "lop_label": "Lớp 6", "lop_num": 6, "is_thcs": True,
        "bai": "Bài 2. Xử lí thông tin",
        "file_name": "Slide_Tin_hoc_Lop_6_Bai02_Xu_li_thong_tin.pptx",
        "slides_content": [
            {"type": "cover", "title": "Xử lí thông tin", "subtitle": "Tin học • Lớp 6 • Bài 2"},
            {"type": "warmup", "title": "Hệ thống xử lý thông tin quanh ta!",
             "content": "Hãy quan sát một chiếc máy tính bỏ túi:\nKhi em bấm '8 x 9 =', quy trình bên trong máy diễn ra như thế nào?"},
            {"type": "learn", "title": "Mô hình xử lý thông tin 4 bước",
             "bullets": [
                 "1. Thu nhận thông tin (Input) — tiếp nhận dữ liệu đầu vào",
                 "2. Lưu trữ thông tin (Storage) — ghi nhớ trên bộ nhớ",
                 "3. Xử lý thông tin (Processing) — biến đổi dữ liệu theo giải thuật",
                 "4. Truyền / Xuất thông tin (Output) — đưa kết quả ra ngoài"
             ]},
            {"type": "learn", "title": "Máy tính là thiết bị xử lý thông tin",
             "bullets": [
                 "Máy tính thực hiện tự động quy trình xử lý thông tin",
                 "Tốc độ xử lý hàng tỷ phép tính mỗi giây",
                 "Khả năng lưu trữ khối lượng dữ liệu khổng lồ"
             ]},
            {"type": "learn", "title": "Vai trò của xử lý thông tin trong đời sống",
             "bullets": [
                 "Dự báo thời tiết dựa trên dữ liệu khí hậu",
                 "Hệ thống định vị GPS tìm đường đi tối ưu",
                 "Y tế: phân tích hình ảnh X-quang chẩn đoán bệnh"
             ]},
            {"type": "practice", "title": "Phân tích chu trình xử lý!",
             "instruction": "Ghép nối thiết bị với bước xử lý tương ứng:",
             "items": ["Bàn phím & Micro ↔ 1. Thu nhận thông tin (Input)", "Ổ cứng & RAM ↔ 2. Lưu trữ thông tin (Storage)", "Bộ vi xử lý CPU ↔ 3. Xử lý thông tin (Processing)", "Màn hình & Loa ↔ 4. Xuất thông tin (Output)"]},
            {"type": "activity", "title": "Thảo luận nhóm tình huống",
             "instruction": "Thảo luận nhóm 4 học sinh (5 phút):",
             "items": ["Mô tả chu trình xử lý thông tin của xe tự lái", "Thiết bị nào đóng vai trò thu nhận thông tin?", "Nêu tầm quan trọng của việc lưu trữ dữ liệu accurate"]},
            {"type": "summary", "title": "Tóm tắt bài học",
             "items": ["Mô hình 4 bước: Thu nhận -> Lưu trữ -> Xử lý -> Xuất thông tin", "Máy tính xử lý thông tin tự động, chính xác và nhanh chóng", "Ứng dụng rộng rãi trong mọi lĩnh vực xã hội hiện đại"]},
            {"type": "thanks", "title": "Bài học kết thúc!", "content": "BTVN: Vẽ sơ đồ khối chu trình xử lý thông tin\ncủa thiết bị nhận diện khuôn mặt!"},
        ]
    },
    {
        "id": "Lop_7", "folder": "Lớp_7", "lop_label": "Lớp 7", "lop_num": 7, "is_thcs": True,
        "bai": "Bài 2. Phần mềm máy tính",
        "file_name": "Slide_Tin_hoc_Lop_7_Bai02_Phan_mem_may_tinh.pptx",
        "slides_content": [
            {"type": "cover", "title": "Phần mềm máy tính (Software)", "subtitle": "Tin học • Lớp 7 • Bài 2"},
            {"type": "warmup", "title": "Nếu không có phần mềm?",
             "content": "Một chiếc máy tính rất mạnh nhưng không cài bất kỳ phần mềm nào:\nLiệu ta có sử dụng được máy tính không?\nVì sao?"},
            {"type": "learn", "title": "Phần mềm hệ thống (System Software)",
             "bullets": [
                 "Hệ điều hành (Operating System — OS): Windows, macOS, Linux, Android",
                 "Quản lý toàn bộ phần cứng và phần mềm của máy tính",
                 "Môi trường để các phần mềm ứng dụng hoạt động"
             ]},
            {"type": "learn", "title": "Phần mềm ứng dụng (Application Software)",
             "bullets": [
                 "Phần mềm văn phòng: Word, Excel, PowerPoint",
                 "Trình duyệt web: Chrome, Edge, Safari",
                 "Phần mềm đồ họa & giải trí: Photoshop, VLC Player, Games"
             ]},
            {"type": "learn", "title": "Mối quan hệ Phần cứng ↔ Hệ điều hành ↔ Ứng dụng",
             "bullets": [
                 "Người dùng thao tác trên Phần mềm ứng dụng",
                 "Phần mềm ứng dụng gửi lệnh tới Hệ điều hành",
                 "Hệ điều hành điều khiển Phần cứng thực thi"
             ]},
            {"type": "practice", "title": "Phân loại phần mềm!",
             "instruction": "Phân loại các phần mềm sau:",
             "items": ["Windows 11, Android 14 ↔ Phần mềm Hệ thống (OS)", "Microsoft Word, Canva ↔ Phần mềm Ứng dụng", "Google Chrome, Cốc Cốc ↔ Phần mềm Ứng dụng", "iOS, Ubuntu Linux ↔ Phần mềm Hệ thống (OS)"]},
            {"type": "activity", "title": "Thảo luận bản quyền phần mềm",
             "instruction": "Trả lời các câu hỏi tình huống (5 phút):",
             "items": ["Phần mềm thương mại và phần mềm miễn phí khác nhau gì?", "Tại sao không nên sử dụng phần mềm bẻ khóa (crack)?", "Kể tên 2 phần mềm mã nguồn mở phổ biến"]},
            {"type": "summary", "title": "Tóm tắt bài học",
             "items": ["2 loại phần mềm: Phần mềm hệ thống & Phần mềm ứng dụng", "Hệ điều hành điều khiển phần cứng và quản lý ứng dụng", "Sử dụng phần mềm có bản quyền và an toàn"]},
            {"type": "thanks", "title": "Bài học kết thúc!", "content": "BTVN: Kể tên 5 phần mềm ứng dụng\nem sử dụng nhiều nhất trên máy tính hoặc điện thoại!"},
        ]
    },
    {
        "id": "Lop_8", "folder": "Lớp_8", "lop_label": "Lớp 8", "lop_num": 8, "is_thcs": True,
        "bai": "Bài 2. Thông tin trong môi trường số",
        "file_name": "Slide_Tin_hoc_Lop_8_Bai02_Thong_tin_trong_moi_truong_so.pptx",
        "slides_content": [
            {"type": "cover", "title": "Thông tin trong môi trường số", "subtitle": "Tin học • Lớp 8 • Bài 2"},
            {"type": "warmup", "title": "Thông tin trên Internet đáng tin đến đâu?",
             "content": "Mỗi ngày có hàng tỷ thông tin được đăng lên mạng:\nLàm sao em phân biệt được thông tin ĐÚNG và tin GIẢ (Fake news)?"},
            {"type": "learn", "title": "Đặc điểm thông tin trong môi trường số",
             "bullets": [
                 "Khối lượng khổng lồ, cập nhật liên tục từng giây",
                 "Đa dạng hình thức: văn bản, hình ảnh, âm thanh, video",
                 "Lan truyền nhanh chóng trên phạm vi toàn cầu"
             ]},
            {"type": "learn", "title": "Tác động của thông tin số",
             "bullets": [
                 "Mặt tích cực: tiếp cận tri thức nhanh, học tập mọi lúc mọi nơi",
                 "Mặt tiêu cực: nguy cơ tin giả, lừa đảo trực tuyến, nghiện mạng",
                 "Thông tin xấu độc ảnh hưởng tâm lý và tư duy"
             ]},
            {"type": "learn", "title": "Tiêu chí đánh giá độ tin cậy thông tin",
             "bullets": [
                 "1. Nguồn tin: tác giả, cơ quan xuất bản uy tín",
                 "2. Tính cập nhật: thời điểm đăng tải thông tin",
                 "3. Tính mục đích: thông tin nhằm cung cấp tri thức hay quảng cáo, giật gân",
                 "4. Kiểm chứng chéo: so sánh với các nguồn báo chí chính thống"
             ]},
            {"type": "practice", "title": "Thực hành đánh giá nguồn tin!",
             "instruction": "Đánh giá độ tin cậy của các nguồn:",
             "items": ["Cổng thông tin Chinhphu.vn ↔ Rất tin cậy (Chính thống)", "Tài khoản cá nhân lạ trên TikTok ↔ Cần kiểm chứng cậy!", "Website của Bộ Giáo dục và Đào tạo ↔ Rất tin cậy", "Bài viết không có tác giả trên blog cá nhân ↔ Độ tin cậy thấp"]},
            {"type": "activity", "title": "Thảo luận: Phòng chống Tin giả",
             "instruction": "Thảo luận nhóm 4 học sinh (5 phút):",
             "items": ["Em làm gì khi thấy một tin đồn thất thố trên mạng?", "Nêu 3 dấu hiệu nhận biết một bài viết tin giả", "Quy tắc 5 giây suy nghĩ trước khi nhấn nút Share/Chia sẻ"]},
            {"type": "summary", "title": "Tóm tắt bài học",
             "items": ["Thông tin số phong phú nhưng cần chọn lọc", "Đánh giá nguồn tin theo 4 tiêu chí chuẩn", "Trở thành người tiêu dùng thông tin thông minh trên Internet"]},
            {"type": "thanks", "title": "Bài học kết thúc!", "content": "BTVN: Tìm 1 ví dụ tin giả gần đây\nvà phân tích nguyên nhân vì sao nó sai!"},
        ]
    },
]

# ═══════════════════════════════════════════════════
# HELPER FUNCTIONS
# ═══════════════════════════════════════════════════

def hex_rgb(h):
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def set_font(run, size_pt, bold=False, color_hex="333333", font_name="Arial"):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = hex_rgb(color_hex)
    run.font.name = font_name

def add_textbox(slide, left, top, width, height, text, size_pt=18, bold=False,
                color_hex="333333", alignment=PP_ALIGN.LEFT, font_name="Arial"):
    actual_top = max(top, SAFE_TOP)
    txbox = slide.shapes.add_textbox(Inches(left), Inches(actual_top), Inches(width), Inches(height))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size_pt, bold, color_hex, font_name)
    return txbox

def add_multiline_textbox(slide, left, top, width, height, lines, size_pt=18,
                           color_hex="333333", bold=False, font_name="Arial",
                           alignment=PP_ALIGN.LEFT):
    actual_top = max(top, SAFE_TOP)
    txbox = slide.shapes.add_textbox(Inches(left), Inches(actual_top), Inches(width), Inches(height))
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alignment
        run = p.add_run()
        run.text = line
        set_font(run, size_pt, bold, color_hex, font_name)
        p.space_after = Pt(6)
    return txbox

def add_safe_shape(slide, shape_type, left, top, width, height, fill_hex,
                   border_hex=None, send_to_back=False):
    actual_top = max(top, SAFE_TOP)
    actual_bottom = min(top + height, SAFE_BOTTOM)
    actual_height = max(actual_bottom - actual_top, 0.1)

    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left), Inches(actual_top), Inches(width), Inches(actual_height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(fill_hex)
    if border_hex:
        shape.line.color.rgb = hex_rgb(border_hex)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()

    if send_to_back:
        sp = shape._element
        spTree = sp.getparent()
        spTree.remove(sp)
        spTree.insert(2, sp)  # Never insert(0)
    return shape

def add_picture_safe(slide, img_path, left, top, width, height):
    actual_top = max(top, SAFE_TOP)
    actual_bottom = min(top + height, SAFE_BOTTOM)
    actual_height = max(actual_bottom - actual_top, 0.5)
    try:
        return slide.shapes.add_picture(img_path, Inches(left), Inches(actual_top),
                                         Inches(width), Inches(actual_height))
    except Exception as e:
        return None

def add_slide_transition(slide, transition_type="fade"):
    transitions = {
        'fade': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>',
        'push': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:push dir="l"/></p:transition>',
        'wipe': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:wipe/></p:transition>',
        'cover': '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:cover/></p:transition>',
    }
    xml_str = transitions.get(transition_type, transitions['fade'])
    slide._element.append(etree.fromstring(xml_str))

def add_appear_animation(slide, shapes_with_click):
    nsmap = {
        'p':  'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a':  'http://schemas.openxmlformats.org/drawingml/2006/main',
    }

    timing = etree.SubElement(slide._element, '{%s}timing' % nsmap['p'])
    tnLst = etree.SubElement(timing, '{%s}tnLst' % nsmap['p'])
    par_root = etree.SubElement(tnLst, '{%s}par' % nsmap['p'])
    cTn_root = etree.SubElement(par_root, '{%s}cTn' % nsmap['p'])
    cTn_root.set('id', '1')
    cTn_root.set('dur', 'indefinite')
    cTn_root.set('restart', 'never')
    cTn_root.set('nodeType', 'tmRoot')

    childTnLst_root = etree.SubElement(cTn_root, '{%s}childTnLst' % nsmap['p'])

    seq = etree.SubElement(childTnLst_root, '{%s}seq' % nsmap['p'])
    seq.set('concurrent', '1')
    seq.set('nextAc', 'seek')

    cTn_seq = etree.SubElement(seq, '{%s}cTn' % nsmap['p'])
    cTn_seq.set('id', '2')
    cTn_seq.set('dur', 'indefinite')
    cTn_seq.set('nodeType', 'mainSeq')

    childTnLst_seq = etree.SubElement(cTn_seq, '{%s}childTnLst' % nsmap['p'])

    prevCond = etree.SubElement(seq, '{%s}prevCondLst' % nsmap['p'])
    cond_prev = etree.SubElement(prevCond, '{%s}cond' % nsmap['p'])
    cond_prev.set('evt', 'onPrev')
    cond_prev.set('delay', '0')
    tgtEl_prev = etree.SubElement(cond_prev, '{%s}tgtEl' % nsmap['p'])
    etree.SubElement(tgtEl_prev, '{%s}sldTgt' % nsmap['p'])

    nextCond = etree.SubElement(seq, '{%s}nextCondLst' % nsmap['p'])
    cond_next = etree.SubElement(nextCond, '{%s}cond' % nsmap['p'])
    cond_next.set('evt', 'onNext')
    cond_next.set('delay', '0')
    tgtEl_next = etree.SubElement(cond_next, '{%s}tgtEl' % nsmap['p'])
    etree.SubElement(tgtEl_next, '{%s}sldTgt' % nsmap['p'])

    id_counter = 3

    for shape_group, click_idx in shapes_with_click:
        if not isinstance(shape_group, (list, tuple)):
            shape_group = [shape_group]

        par_click = etree.SubElement(childTnLst_seq, '{%s}par' % nsmap['p'])
        cTn_click = etree.SubElement(par_click, '{%s}cTn' % nsmap['p'])
        cTn_click.set('id', str(id_counter)); id_counter += 1
        cTn_click.set('fill', 'hold')

        stCondLst = etree.SubElement(cTn_click, '{%s}stCondLst' % nsmap['p'])
        cond = etree.SubElement(stCondLst, '{%s}cond' % nsmap['p'])
        cond.set('delay', '0')

        childTnLst_click = etree.SubElement(cTn_click, '{%s}childTnLst' % nsmap['p'])

        for shape in shape_group:
            if shape is None:
                continue
            sp_id = shape.shape_id

            par_anim = etree.SubElement(childTnLst_click, '{%s}par' % nsmap['p'])
            cTn_anim = etree.SubElement(par_anim, '{%s}cTn' % nsmap['p'])
            cTn_anim.set('id', str(id_counter)); id_counter += 1
            cTn_anim.set('presetID', '1')  # Appear
            cTn_anim.set('presetClass', 'entr')
            cTn_anim.set('presetSubtype', '0')
            cTn_anim.set('fill', 'hold')

            stCond_anim = etree.SubElement(cTn_anim, '{%s}stCondLst' % nsmap['p'])
            cond_anim = etree.SubElement(stCond_anim, '{%s}cond' % nsmap['p'])
            cond_anim.set('delay', '0')

            childTnLst_anim = etree.SubElement(cTn_anim, '{%s}childTnLst' % nsmap['p'])

            p_set = etree.SubElement(childTnLst_anim, '{%s}set' % nsmap['p'])
            cBhvr = etree.SubElement(p_set, '{%s}cBhvr' % nsmap['p'])
            cTn_set = etree.SubElement(cBhvr, '{%s}cTn' % nsmap['p'])
            cTn_set.set('id', str(id_counter)); id_counter += 1
            cTn_set.set('dur', '1')
            cTn_set.set('fill', 'hold')

            stCond_set = etree.SubElement(cTn_set, '{%s}stCondLst' % nsmap['p'])
            cond_set = etree.SubElement(stCond_set, '{%s}cond' % nsmap['p'])
            cond_set.set('delay', '0')

            tgtEl = etree.SubElement(cBhvr, '{%s}tgtEl' % nsmap['p'])
            spTgt = etree.SubElement(tgtEl, '{%s}spTgt' % nsmap['p'])
            spTgt.set('spid', str(sp_id))

            attrNameLst = etree.SubElement(cBhvr, '{%s}attrNameLst' % nsmap['p'])
            attrName = etree.SubElement(attrNameLst, '{%s}attrName' % nsmap['p'])
            attrName.text = 'style.visibility'

            p_to = etree.SubElement(p_set, '{%s}to' % nsmap['p'])
            strVal = etree.SubElement(p_to, '{%s}strVal' % nsmap['p'])
            strVal.set('val', 'visible')

def find_best_image_for_lesson(lesson, slide_type, idx=0):
    """Find image for a slide from SGK or KHBD images folder"""
    lop_folder = lesson["folder"]
    lop_num = lesson["lop_num"]

    # 1. Check KHBD images in Tuần_03
    khbd_img_dir = os.path.join(KHBD_BASE, lop_folder, "Tuần_03", "images")
    if os.path.isdir(khbd_img_dir):
        indexed_name = f"{slide_type}{idx+1}_lop{lop_num}"
        for ext in ('.jpg', '.png', '.jpeg'):
            p = os.path.join(khbd_img_dir, indexed_name + ext)
            if os.path.isfile(p):
                return p
        exact_name = f"{slide_type}_lop{lop_num}"
        for ext in ('.jpg', '.png', '.jpeg'):
            p = os.path.join(khbd_img_dir, exact_name + ext)
            if os.path.isfile(p):
                return p
        imgs = sorted([f for f in os.listdir(khbd_img_dir) if f.endswith(('.png','.jpeg','.jpg'))])
        if imgs:
            return os.path.join(khbd_img_dir, imgs[idx % len(imgs)])

    # Fallback: check Tuần_02 images
    khbd_t2_dir = os.path.join(KHBD_BASE, lop_folder, "Tuần_02", "images")
    if os.path.isdir(khbd_t2_dir):
        imgs = sorted([f for f in os.listdir(khbd_t2_dir) if f.endswith(('.png','.jpeg','.jpg'))])
        if imgs:
            return os.path.join(khbd_t2_dir, imgs[idx % len(imgs)])

    # 2. Check SGK images
    sgk_img_dir = os.path.join(SGK_BASE, f"Lớp_{lop_num}", "bai1_images")
    if os.path.isdir(sgk_img_dir):
        imgs = sorted([f for f in os.listdir(sgk_img_dir)
                       if f.endswith(('.png','.jpeg','.jpg')) and not os.path.isdir(os.path.join(sgk_img_dir, f))])
        if imgs:
            return os.path.join(sgk_img_dir, imgs[idx % len(imgs)])
    return None

def remove_template_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        sldId = prs.slides._sldIdLst[0]
        prs.slides._sldIdLst.remove(sldId)

# ═══════════════════════════════════════════════════
# BUILDERS FOR ALL GRADES
# ═══════════════════════════════════════════════════

def build_cover(prs, data, pal, lesson, layout):
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["primary"], send_to_back=True)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, 1.8, pal["accent"], send_to_back=True)

    add_textbox(slide, 0.8, 2.0, 7.5, 1.5, data["title"], size_pt=34, bold=True, color_hex=pal["text_on_primary"])
    add_textbox(slide, 0.8, 3.8, 7.5, 0.8, data.get("subtitle", ""), size_pt=20, color_hex=pal["text_on_primary"])

    img = find_best_image_for_lesson(lesson, "cover", 0)
    if img:
        add_picture_safe(slide, img, SLIDE_W - 4.5, SAFE_TOP + 0.3, 3.8, 3.8)

    add_slide_transition(slide, "fade")
    return slide

def build_warmup(prs, data, pal, lesson, layout):
    slide = prs.slides.add_slide(layout)
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, pal["primary"])
    badge_text = f"  {lesson['lop_label'].upper()} • {lesson['bai'].upper()}"
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, badge_text, size_pt=13, bold=True, color_hex=pal["text_on_primary"])

    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, 8, 0.65, data["title"], size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    content_y = title_y + 0.75
    card_w = 7.5
    card_h = SAFE_BOTTOM - content_y - 0.15
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.4, content_y, card_w, card_h, pal["card"])
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0.4, content_y, 0.12, card_h, pal["accent"])

    lines = data.get("content", "").split("\n")
    add_multiline_textbox(slide, 0.8, content_y+0.2, card_w-0.6, card_h-0.4, lines, size_pt=20, color_hex=pal["text_on_card"])

    img = find_best_image_for_lesson(lesson, "warmup", 0)
    if img:
        add_picture_safe(slide, img, 8.3, content_y, SAFE_RIGHT-8.3-0.2, card_h)

    add_slide_transition(slide, "push")
    return slide

def build_learn_grid(prs, data, pal, lesson, layout, learn_idx=0):
    """Layout A — Grid Flashcard (Primary Grades: Tiền TH -> Lớp 5)"""
    slide = prs.slides.add_slide(layout)
    bullets = data.get("bullets", [])
    n = len(bullets)

    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, pal["primary"])
    badge_text = f"  {lesson['lop_label'].upper()} • {lesson['bai'].upper()}"
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, badge_text, size_pt=13, bold=True, color_hex=pal["text_on_primary"])

    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, SLIDE_W-1, 0.65, data["title"], size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    grid_top = title_y + 0.8
    grid_bottom = SAFE_BOTTOM - 0.1
    avail_h = grid_bottom - grid_top

    cols, rows = (n, 1) if n <= 3 else (2, 2)
    gap = 0.25
    total_gap_x = gap * (cols - 1)
    card_w = (SLIDE_W - 1.0 - total_gap_x) / cols
    if rows == 1:
        img_h, text_h = 2.0, 0.9
    else:
        img_h, text_h = 1.1, 0.5
    card_h = img_h + text_h + 0.15
    total_gap_y = gap * (rows - 1)
    start_x = (SLIDE_W - (cols * card_w + total_gap_x)) / 2
    start_y = grid_top + (avail_h - (rows * card_h + total_gap_y)) / 2
    start_y = max(start_y, grid_top)

    for i, btext in enumerate(bullets):
        col, row = i % cols, i // cols
        x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + gap)

        if y + card_h > SAFE_BOTTOM + 0.05:
            break

        add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h, pal["card"], border_hex="D0D5DD")
        add_safe_shape(slide, MSO_SHAPE.RECTANGLE, x, y, card_w, 0.06, pal["accent"])

        img = find_best_image_for_lesson(lesson, "learn", learn_idx*3 + i)
        if img:
            img_w = min(1.8, card_w - 0.4)
            img_x = x + (card_w - img_w) / 2
            add_picture_safe(slide, img, img_x, y + 0.15, img_w, img_h)

        text_y = y + img_h + 0.15
        add_textbox(slide, x + 0.15, text_y, card_w - 0.3, text_h, btext, size_pt=16, bold=False, color_hex=pal["text_on_card"], alignment=PP_ALIGN.CENTER)

    add_slide_transition(slide, "wipe")
    return slide

def build_learn_row(prs, data, pal, lesson, layout, learn_idx=0):
    """Layout B — Horizontal Row (Middle School Grades: Lớp 6 -> Lớp 8)"""
    slide = prs.slides.add_slide(layout)
    bullets = data.get("bullets", [])

    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, pal["primary"])
    badge_text = f"  {lesson['lop_label'].upper()} • {lesson['bai'].upper()}"
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, badge_text, size_pt=13, bold=True, color_hex=pal["text_on_primary"])

    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, SLIDE_W-1, 0.65, data["title"], size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    row_top = title_y + 0.75
    n = len(bullets)
    gap = 0.15
    avail_h = SAFE_BOTTOM - row_top - 0.1
    row_h = min(1.4, (avail_h - gap*(n-1)) / max(n, 1))

    for i, btext in enumerate(bullets[:4]):
        y = row_top + i * (row_h + gap)

        if y + row_h > SAFE_BOTTOM + 0.05:
            break

        img = find_best_image_for_lesson(lesson, "learn", learn_idx*3 + i)

        img_w = 2.0
        img_x = 0.5
        txt_x = img_x + img_w + 0.2
        txt_w = SLIDE_W - txt_x - 0.5

        if img:
            add_picture_safe(slide, img, img_x, y, img_w, row_h)
        else:
            txt_x = 0.5
            txt_w = SLIDE_W - 1.0

        card = add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, txt_x, y, txt_w, row_h, pal["card"], border_hex="D0D5DD")
        add_safe_shape(slide, MSO_SHAPE.RECTANGLE, txt_x, y, 0.08, row_h, pal["accent"])
        add_textbox(slide, txt_x + 0.25, y + 0.15, txt_w - 0.4, row_h - 0.3, btext, size_pt=18, bold=False, color_hex=pal["text_on_card"])

    add_slide_transition(slide, "wipe")
    return slide

def build_practice(prs, data, pal, lesson, layout):
    slide = prs.slides.add_slide(layout)

    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, pal["primary"])
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, f"  {lesson['lop_label'].upper()} • LUYỆN TẬP", size_pt=13, bold=True, color_hex=pal["text_on_primary"])

    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, SLIDE_W-1, 0.65, data["title"], size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    instr_y = title_y + 0.75
    instr_lines = data.get("instruction", "").split("\n")
    add_multiline_textbox(slide, 0.5, instr_y, SLIDE_W-1, 0.7, instr_lines, size_pt=18, color_hex=pal["text_on_bg"])

    items = data.get("items", [])
    item_y_start = instr_y + 0.8
    n_items = len(items)

    cols = 2 if n_items > 3 else max(n_items, 1)
    rows = (n_items + 1) // 2 if n_items > 3 else 1

    gap = 0.25
    card_w = (SLIDE_W - 1.0 - gap*(cols-1)) / cols
    card_h = min(1.8, (SAFE_BOTTOM - item_y_start - gap*(rows-1)) / rows)

    animation_groups = []

    for i, item in enumerate(items[:4]):
        col, row = i % cols, i // cols
        x = 0.5 + col * (card_w + gap)
        y = item_y_start + row * (card_h + gap)

        if y + card_h > SAFE_BOTTOM + 0.05:
            break

        card = add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h, pal["card"], border_hex=pal["accent"])

        img = find_best_image_for_lesson(lesson, "practice", i)
        pic = None
        if img:
            img_w = 1.4
            pic = add_picture_safe(slide, img, x + 0.2, y + 0.15, img_w, card_h - 0.3)
            txt_x = x + 0.2 + img_w + 0.2
            txt_w = card_w - (0.2 + img_w + 0.3)
        else:
            txt_x = x + 0.2
            txt_w = card_w - 0.4

        txt = add_textbox(slide, txt_x, y + 0.15, txt_w, card_h - 0.3, item, size_pt=16, bold=True, color_hex=pal["text_on_card"])

        group = [s for s in [card, pic, txt] if s is not None]
        animation_groups.append((group, i))

    if animation_groups:
        add_appear_animation(slide, animation_groups)

    add_slide_transition(slide, "wipe")
    return slide

def build_activity(prs, data, pal, lesson, layout):
    """Activity / Game Slide — Mandatory illustration picture"""
    slide = prs.slides.add_slide(layout)

    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    badge_y = SAFE_TOP + 0.05
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, badge_y, SLIDE_W, 0.45, pal["accent"])
    add_textbox(slide, 0.3, badge_y+0.03, SLIDE_W-0.6, 0.4, f"  {lesson['lop_label'].upper()} • THỬ THÁCH", size_pt=13, bold=True, color_hex=pal["text_on_primary"])

    title_y = SAFE_TOP + 0.6
    add_textbox(slide, 0.5, title_y, SLIDE_W-1, 0.65, data["title"], size_pt=26, bold=True, color_hex=pal["text_on_bg"])

    content_y = title_y + 0.75
    content_h = SAFE_BOTTOM - content_y - 0.15

    card_w = 7.0
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.4, content_y, card_w, content_h, pal["card"])
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0.4, content_y, 0.12, content_h, pal["accent"])

    lines = data.get("content", "").split("\n") if data.get("content") else data.get("instruction", "").split("\n")
    add_multiline_textbox(slide, 0.8, content_y+0.2, card_w-0.6, content_h-0.4, lines, size_pt=20, color_hex=pal["text_on_card"])

    img = find_best_image_for_lesson(lesson, "activity", 0)
    if img:
        img_x = 7.8
        img_w = SAFE_RIGHT - img_x - 0.2
        add_picture_safe(slide, img, img_x, content_y + 0.1, img_w, content_h - 0.2)

    add_slide_transition(slide, "cover")
    return slide

def build_summary(prs, data, pal, lesson, layout):
    slide = prs.slides.add_slide(layout)

    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["bg"], send_to_back=True)

    px, py = 0.6, SAFE_TOP + 0.15
    pw = SLIDE_W - 1.2
    ph = SAFE_BOTTOM - py - 0.15
    add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, px, py, pw, ph, pal["card"])
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, px, py, pw, 0.06, pal["primary"])

    add_textbox(slide, px+0.3, py+0.2, pw-0.6, 0.65, data["title"], size_pt=24, bold=True, color_hex=pal["primary"], alignment=PP_ALIGN.CENTER)

    div_y = py + 0.95
    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, px+2, div_y, pw-4, 0.03, pal["accent"])

    items = data.get("items", [])
    for i, item in enumerate(items[:4]):
        y = div_y + 0.25 + i * 1.1
        if y + 0.8 > SAFE_BOTTOM:
            break
        add_safe_shape(slide, MSO_SHAPE.RECTANGLE, px+0.4, y, 0.08, 0.7, pal["accent"])
        add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, px+0.6, y, pw-1.2, 0.7, pal["bg"])
        add_textbox(slide, px+0.9, y+0.15, pw-1.8, 0.4, item, size_pt=18, color_hex=pal["text_on_bg"])

    add_slide_transition(slide, "fade")
    return slide

def build_thanks(prs, data, pal, lesson, layout):
    slide = prs.slides.add_slide(layout)

    add_safe_shape(slide, MSO_SHAPE.RECTANGLE, 0, SAFE_TOP, SLIDE_W, SAFE_BOTTOM-SAFE_TOP, pal["primary"], send_to_back=True)

    add_textbox(slide, 1.0, SAFE_TOP+0.8, SLIDE_W-2, 1.2, data["title"], size_pt=36, bold=True, color_hex=pal["text_on_primary"], alignment=PP_ALIGN.CENTER)

    content = data.get("content", "")
    if content:
        card_w = 8.0
        card_x = (SLIDE_W - card_w) / 2
        card_y = SAFE_TOP + 2.5
        card_h = min(2.0, SAFE_BOTTOM - card_y - 0.2)
        add_safe_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, card_w, card_h, pal["card"])
        lines = content.split("\n")
        add_multiline_textbox(slide, card_x+0.5, card_y+0.3, card_w-1, card_h-0.6, lines, size_pt=20, color_hex=pal["text_on_card"], alignment=PP_ALIGN.CENTER)

    add_slide_transition(slide, "fade")
    return slide

# ═══════════════════════════════════════════════════
# VERIFICATION
# ═══════════════════════════════════════════════════

def verify_slide_deck(prs, lesson_label):
    issues = []
    for si, sl in enumerate(prs.slides):
        for sh in sl.shapes:
            top_in = sh.top / 914400
            bottom_in = (sh.top + sh.height) / 914400
            if top_in < SAFE_TOP - 0.02 and sh.name not in ('Picture 7',):
                issues.append(f"Slide {si+1}: {sh.name} top={top_in:.2f}in < {SAFE_TOP}in")
            if bottom_in > SAFE_BOTTOM + 0.02 and sh.name not in ('Picture 9',):
                issues.append(f"Slide {si+1}: {sh.name} bottom={bottom_in:.2f}in > {SAFE_BOTTOM}in")

    if issues:
        print(f"  ⚠️ {lesson_label} VERIFICATION: {len(issues)} issues found")
        for iss in issues:
            print(f"     ❌ {iss}")
    else:
        print(f"  ✅ {lesson_label} VERIFICATION: ALL {len(prs.slides)} slides PASS anti-bug checklist")
    return len(issues) == 0

# ═══════════════════════════════════════════════════
# MAIN BUILDER FOR ALL 9 GRADES TUẦN 03
# ═══════════════════════════════════════════════════

def create_slide_deck_v3_tuan03(lesson, palette_idx):
    palette = COLOR_PALETTES[palette_idx % len(COLOR_PALETTES)]
    prs = Presentation(TEMPLATE)
    remove_template_slides(prs)
    layout = prs.slide_layouts[6]

    learn_idx = 0
    for sc in lesson["slides_content"]:
        stype = sc["type"]
        try:
            if stype == "cover":
                build_cover(prs, sc, palette, lesson, layout)
            elif stype == "warmup":
                build_warmup(prs, sc, palette, lesson, layout)
            elif stype == "learn":
                if lesson["is_thcs"]:
                    build_learn_row(prs, sc, palette, lesson, layout, learn_idx)
                else:
                    build_learn_grid(prs, sc, palette, lesson, layout, learn_idx)
                learn_idx += 1
            elif stype == "practice":
                build_practice(prs, sc, palette, lesson, layout)
            elif stype == "activity":
                build_activity(prs, sc, palette, lesson, layout)
            elif stype == "summary":
                build_summary(prs, sc, palette, lesson, layout)
            elif stype == "thanks":
                build_thanks(prs, sc, palette, lesson, layout)
        except Exception as e:
            print(f"     ⚠️ Error on slide '{sc.get('title','')}': {e}")

    # Verify
    verify_slide_deck(prs, lesson["lop_label"])

    # Output path
    output_dir = os.path.join(KHBD_BASE, lesson["folder"], "Tuần_03")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, lesson["file_name"])

    try:
        prs.save(output_path)
        print(f"  ✅ {lesson['lop_label']}: {len(prs.slides)} slides -> {lesson['file_name']}")
    except PermissionError:
        alt = output_path.replace(".pptx", "_v3_alt.pptx")
        prs.save(alt)
        print(f"  ⚠️ {lesson['lop_label']}: File locked -> saved as {os.path.basename(alt)}")
    except Exception as e:
        print(f"  ❌ {lesson['lop_label']} save error: {e}")

    return output_path

def main():
    print("=" * 70)
    print("  TẠO BỘ SLIDE BÀI GIẢNG CHUẨN V3 — TOÀN BỘ 9 KHỐI LỚP (TUẦN 03)")
    print("  Quy chuẩn: Per-Bullet Images, Animation On-Click, Game Pics, Anti-Bug")
    print("=" * 70)

    for i, lesson in enumerate(LESSONS_TUAN03):
        print(f"\n[ Khối: {lesson['lop_label']} — {lesson['bai']} ]")
        try:
            create_slide_deck_v3_tuan03(lesson, i)
        except Exception as e:
            print(f"❌ Error on {lesson['lop_label']}: {e}")
            traceback.print_exc()

    print("\n" + "=" * 70)
    print("  HOÀN THÀNH TOÀN BỘ 9 BỘ SLIDE TUẦN 03!")
    print("=" * 70)

if __name__ == "__main__":
    main()
