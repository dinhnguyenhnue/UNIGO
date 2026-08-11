# -*- coding: utf-8 -*-
import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation

f = r'D:\UNIGO\KHBD_Tin_học\Lớp_1\Tuần_02\Slide_Tin_hoc_Lop_1_Bai01_Chiec_may_tinh_cua_em.pptx'
prs = Presentation(f)

for si, sl in enumerate(prs.slides):
    print(f'=== Slide {si} ({len(sl.shapes)} shapes) ===')
    for idx, sh in enumerate(sl.shapes):
        t = sh.top / 914400
        b = (sh.top + sh.height) / 914400
        l = sh.left / 914400
        w = sh.width / 914400
        text_preview = ''
        if hasattr(sh, 'text') and sh.text:
            text_preview = f' | "{sh.text[:60]}"'
        fill_info = ''
        try:
            if sh.fill and sh.fill.fore_color:
                fill_info = f' fill=#{sh.fill.fore_color.rgb}'
        except:
            pass
        print(f'  [{idx}] {sh.name}: Y={t:.2f}-{b:.2f} X={l:.2f} W={w:.2f}{fill_info}{text_preview}')
    print()
